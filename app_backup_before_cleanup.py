# Grade Level (lower is better)
import os
from dotenv import load_dotenv
from flask import Flask, request, send_file, render_template_string, redirect, url_for, jsonify
from werkzeug.utils import secure_filename
import uuid

# Load environment variables
load_dotenv()
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import anthropic
import re
import io
import base64
import matplotlib.pyplot as plt
import numpy as np
import threading
import hashlib
import time
from tenacity import retry, stop_after_attempt, wait_exponential
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import json
from api_utils import ApiUtils, API_CHECK_SUCCESS_TEMPLATE, API_CHECK_ERROR_TEMPLATE
from migrate_pdf_functions import PDFMigrationHelper
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.colors import HexColor
from reportlab.lib.enums import TA_JUSTIFY
from io import BytesIO
from urllib.parse import urlparse



# Global dictionaries to store status information
global_analysis_status = {}
processing_tasks = {}
api_call_counter = {}
results = []
current_batch = []
current_batch_tokens = 0

# Define the upload folder path
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')

# Set up Flask
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER  # Use the absolute path
app.config['OUTPUT_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')
app.config['REPORT_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'reports')
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB max upload

# Create directories if they don't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
os.makedirs(app.config['REPORT_FOLDER'], exist_ok=True)


# Set up API client (not exposed in UI)
api_key = os.getenv('ANTHROPIC_API_KEY')
if not api_key:
    raise ValueError("ANTHROPIC_API_KEY not found in environment variables")

api_utils = ApiUtils(api_key)
client = anthropic.Anthropic(api_key=api_utils)

# Initialize PDF Migration Helper for new service-based PDF processing
pdf_migration_helper = PDFMigrationHelper({
    'output_folder': app.config['OUTPUT_FOLDER'],
    'upload_folder': app.config['UPLOAD_FOLDER'],
    'anthropic_api_key': str(api_utils)
})

# Initialize enhanced services
from services import (
    PDFService, PowerPointService, ConversionService, 
    EducationalContentService
)

service_config = {
    'output_folder': app.config['OUTPUT_FOLDER'],
    'upload_folder': app.config['UPLOAD_FOLDER'],
    'anthropic_api_key': str(api_utils)
}

pdf_service = PDFService(service_config)
pptx_service = PowerPointService(service_config)
conversion_service = ConversionService(service_config)
educational_service = EducationalContentService(service_config)

# Define colors for different learning profiles
PROFILE_COLORS = {
    "dyslexia": RGBColor(0, 102, 204),  # Blue
    "adhd": RGBColor(46, 139, 87),      # Sea Green
    "esl": RGBColor(148, 0, 211)        # Purple
}

# Readability thresholds for recommendations
READABILITY_THRESHOLDS = {
    "dyslexia": {
        "flesch_reading_ease": 70,  # Higher is easier
        "flesch_kincaid_grade": 6,  # Lower grade level is easier
        "smog_index": 8,            # Lower is easier
        "sentence_length": 12,      # Average words per sentence
        "complex_word_percent": 15  # Percentage of complex words
    },
    "adhd": {
        "flesch_reading_ease": 60,
        "flesch_kincaid_grade": 8,
        "smog_index": 10,
        "sentence_length": 15,
        "complex_word_percent": 20
    },
    "esl": {
        "flesch_reading_ease": 80,
        "flesch_kincaid_grade": 5,
        "smog_index": 7,
        "sentence_length": 10,
        "complex_word_percent": 10
    }
}

# Cache for storing adaptation results
class AdaptationCache:
    """Enhanced cache for storing text adaptation results with better hit rate"""
    def __init__(self, max_size=2000):  # Increased max size
        self.cache = {}
        self.max_size = max_size
        self.access_times = {}
        self.hit_count = 0
        self.miss_count = 0
    
    def get_key(self, text, profile):
        """Generate a cache key based on text content and profile"""
        # Process text to improve cache hit rate
        text_to_hash = self._normalize_text(text)
        text_hash = hashlib.md5(text_to_hash.encode()).hexdigest()
        # Make sure profile is included in the key and is case-insensitive
        profile_key = profile.lower() if profile else "default"
        return f"{text_hash}-{profile_key}"
    
    def _normalize_text(self, text):
        """Normalize text to improve cache hit rates"""
        if not text:
            return ""
            
        # Truncate very long texts for hashing
        text_to_process = text[:2000] if len(text) > 2000 else text
        
        # Normalize whitespace
        normalized = ' '.join(text_to_process.split())
        
        # Convert to lowercase for case-insensitive matching
        normalized = normalized.lower()
        
        # Optionally, you could remove punctuation as well if that helps cache hits
        # normalized = re.sub(r'[^\w\s]', '', normalized)
        
        return normalized
    
    # Cache management methods
    def get(self, text, profile):
        """Get adaptation from cache if available with metrics tracking"""
        key = self.get_key(text, profile)
        if key in self.cache:
            # Update access time
            self.access_times[key] = time.time()
            self.hit_count += 1
            
            # Log cache hit rate periodically
            if (self.hit_count + self.miss_count) % 100 == 0:
                hit_rate = self.hit_count / (self.hit_count + self.miss_count) * 100
                print(f"Cache hit rate: {hit_rate:.1f}% ({self.hit_count} hits, {self.miss_count} misses)")
                
            return self.cache[key]
        
        self.miss_count += 1
        return None
    
    def set(self, text, profile, adapted_text):
        """Store adaptation in cache"""
        key = self.get_key(text, profile)
        
        # Check if we need to evict entries because cache is full
        if len(self.cache) >= self.max_size:
            # Evict 10% of least recently used items at once for better performance
            items_to_evict = max(1, int(self.max_size * 0.1))
            oldest_keys = sorted(self.access_times, key=self.access_times.get)[:items_to_evict]
            
            for old_key in oldest_keys:
                if old_key in self.cache:
                    del self.cache[old_key]
                    del self.access_times[old_key]
        
        # Store new item
        self.cache[key] = adapted_text
        self.access_times[key] = time.time()

# Initialize the cache globally
adaptation_cache = AdaptationCache()

# HTML Templates
from html_templates import (DOWNLOAD_TEMPLATE_UNIVERSAL, INDEX_TEMPLATE, SCAFFOLDING_TEMPLATE, FRAMEWORK_TEMPLATE,
                           ASSESSMENT_TEMPLATE_SIMPLIFIED, ASSESSMENT_TEMPLATE, 
                           PROCESSING_TEMPLATE_WITH_PROGRESS, DOWNLOAD_TEMPLATE_WITH_TRANSLATION, ERROR_TEMPLATE)

@app.route('/')
def index():
    """Render the home page with a form"""
    return render_template_string(INDEX_TEMPLATE)

@app.route('/analysis_status/<file_id>', methods=['GET'])
def check_analysis_status(file_id):
    """Check the status of structural analysis"""
    if file_id in global_analysis_status:
        return jsonify(global_analysis_status[file_id])
    else:
        return jsonify({'status': 'not_found'})

# Clean Upload Route - Replace the entire broken upload section with this

@app.route('/upload', methods=['POST'])
def upload():
    """Handle file upload and direct to appropriate process - supports both PPTX and PDF"""
    try:
        # Get the uploaded file (check common field names)
        file = None
        if 'pptx' in request.files and request.files['pptx'].filename:
            file = request.files['pptx']
        elif 'pdf' in request.files and request.files['pdf'].filename:
            file = request.files['pdf']
        elif 'file' in request.files and request.files['file'].filename:
            file = request.files['file']
        
        # Validate file presence
        if not file or file.filename == '':
            return render_template_string(ERROR_TEMPLATE, 
                message="No file selected. Please choose a PowerPoint (.pptx) or PDF file."), 400
        
        # Get form parameters
        profile = request.form.get('profile')
        action = request.form.get('action', 'assess')
        target_language = request.form.get('target_language', '')
        export_format = request.form.get('export_format', 'pdf')
        
        # Debug form parameters
        print(f"DEBUG: Form parameters - profile: {profile}, action: {action}, target_language: '{target_language}', export_format: {export_format}")
        
        # Validate required parameters
        if not profile:
            return render_template_string(ERROR_TEMPLATE, 
                message="Please select a learning profile."), 400
        
        # Secure filename and validate extension
        filename = secure_filename(file.filename)
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext not in ['.pptx', '.pdf']:
            return render_template_string(ERROR_TEMPLATE, 
                message=f"Invalid file type '{file_ext}'. Please upload a PowerPoint (.pptx) or PDF file."), 400
        
        # Generate unique file ID and save file
        file_id = str(uuid.uuid4())
        processing_tasks[file_id] = {
            'status': 'upload', 
            'filename': filename, 
            'profile': profile,
            'file_type': file_ext
        }
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        file.save(file_path)
        
        # Route based on action
        if action == 'assess':
            # Assessment route
            if target_language:
                return redirect(url_for('analyze_scaffolding', file_id=file_id, profile=profile, target_language=target_language))
            else:
                return redirect(url_for('analyze_scaffolding', file_id=file_id, profile=profile))
        
        else:  # action == 'adapt'
            # Start adaptation process
            processing_tasks[file_id].update({
                'status': 'processing', 
                'message': 'Starting adaptation...',
                'export_format': export_format,
                'progress': {'total': 0, 'processed': 0, 'percentage': 0}
            })
            
            # Choose processing method based on file type
            if file_ext == '.pdf':
                # Use new service-based PDF processing
                print(f"Processing PDF with service-based system: {filename}")
                thread_target = process_pdf_with_services
                thread_args = (file_path, file_id, filename, profile, export_format, target_language)
            else:  # .pptx
                print(f"Processing PowerPoint: {filename}")
                thread_target = process_presentation_efficiently
                thread_args = (file_path, file_id, filename, profile, target_language)
            
            # Start background processing
            thread = threading.Thread(
                target=thread_target,
                args=thread_args,
                name=f"process-{file_id}"
            )
            thread.daemon = True
            thread.start()
            
            # Profile display names
            profile_names = {
                "dyslexia": "Dyslexia Support",
                "adhd": "ADHD Support", 
                "esl": "English Language Learners"
            }
            
            # Return processing page
            return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                file_id=file_id, 
                filename=f"adapted_{filename}",
                profile=profile,
                profile_name=profile_names.get(profile, profile)
            )
            
    except Exception as e:
        print(f"Error in upload route: {str(e)}")
        return render_template_string(ERROR_TEMPLATE, 
            message=f"Upload failed: {str(e)}"), 500

# Add this helper function to gracefully handle missing PDF libraries
# API Call with Retry Logic
@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
def call_claude_api_with_retry(prompt, model="claude-3-5-sonnet-20240620", max_tokens=1024):
    """Call Claude API with retry using the ApiUtils module"""
    return api_utils.call_with_retry(prompt, model, max_tokens)

def debug_processing_status(file_id, message):
    """Print debug info about processing status"""
    print(f"DEBUG [{file_id}]: {message}")
    if file_id in processing_tasks:
        status = processing_tasks[file_id].get('status', 'unknown')
        progress = processing_tasks[file_id].get('progress', {})
        print(f"  Status: {status}, Progress: {progress}")

def analyze_instructional_framework(pptx_path):
    """Use Claude to analyze the instructional framework of a presentation"""
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Extract all slide content with metadata
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = ""
        title = ""
        
        # Get slide title
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title.text
            
        # Extract all text from the slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text:
                slide_content += shape.text + "\n"
        
        slides_data.append({
            "slide_number": i+1,
            "title": title,
            "content": slide_content
        })
    
    # Prepare prompt for Claude
    slide_descriptions = "\n\n".join([
        f"SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
        for s in slides_data
    ])
    
    # Use a simplified prompt to avoid potential issues with JSON example
    prompt = f"""
    You are an expert in instructional design. Analyze this PowerPoint presentation and identify its instructional framework and teaching pattern.
    
    SLIDES:
    {slide_descriptions}
    
    Focus on identifying these specific instructional patterns:
    
    1. Recall/I do/We do/You do/Review framework
    2. 5E Instructional Model (Engage, Explore, Explain, Elaborate, Evaluate)
    3. Other common frameworks (specify which one)
    
    Provide your analysis as JSON with this structure:
    {{
      "framework": {{
        "identified_framework": "Name of framework",
        "framework_alignment_score": 0-100,
        "strengths": ["strength 1", "strength 2"],
        "weaknesses": ["weakness 1", "weakness 2"],
        "missing_phases": ["phase 1", "phase 2"],
        "recommendations": "brief recommendation",
        "balance_analysis": "brief analysis of instructional balance"
      }},
      "slides": [
        {{
          "slide_number": 1,
          "framework_phase": "phase name",
          "effectiveness": "strong/adequate/weak",
          "elements": ["element 1", "element 2"]
        }}
      ]
    }}
    
    Respond only with the JSON, no other text.
    """
    
    try:
        # Call Claude API
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        
        # Extract and parse JSON response
        content = response.content[0].text.strip()
        
        # Try to find JSON in the response
        try:
            # First attempt: try to parse the entire response as JSON
            framework_data = json.loads(content)
            return framework_data
        except json.JSONDecodeError:
            # Second attempt: try to find JSON block within the response
            json_match = re.search(r'(\{.*\})', content, re.DOTALL)
            if json_match:
                try:
                    framework_data = json.loads(json_match.group(0))
                    return framework_data
                except json.JSONDecodeError:
                    print("Error parsing JSON from LLM response")
                    return {"error": "Could not parse JSON response", "raw_response": content[:200] + "..."}
            else:
                return {"error": "Could not extract JSON from response", "raw_response": content[:200] + "..."}
        
    except Exception as e:
        print(f"Error in framework analysis: {str(e)}")
        return {"error": f"Analysis failed: {str(e)}"}

# Function removed - now using pdf_service.extract_content_from_pdf()
# def extract_content_from_pdf(file_path):
    """Extract content from PDF using available libraries with enhanced visual preservation"""
    print(f"Extracting content from PDF: {file_path}")
    
    content_data = {
        'pages': [],
        'metadata': {},
        'images': []
    }
    
    # Try PyMuPDF first for comprehensive extraction (text + images + formatting)
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(file_path)
        content_data['metadata'] = {
            'total_pages': len(doc),
            'title': doc.metadata.get('title', '') if doc.metadata else '',
        }
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_content = {
                'page_number': page_num + 1,
                'text': '',
                'tables': [],
                'images': [],
                'text_blocks': []
            }
            
            # Extract text with formatting information
            text_dict = page.get_text("dict")
            full_text = ""
            text_blocks = []
            
            for block in text_dict["blocks"]:
                if "lines" in block:
                    block_text = ""
                    for line in block["lines"]:
                        line_text = ""
                        for span in line["spans"]:
                            span_text = span["text"]
                            line_text += span_text
                            
                            # Store formatting info
                            text_blocks.append({
                                'text': span_text,
                                'font': span.get('font', 'Arial'),
                                'size': span.get('size', 12),
                                'color': span.get('color', 0),  # RGB color as int
                                'bbox': span.get('bbox', [0,0,0,0])  # bounding box
                            })
                        
                        block_text += line_text + "\n"
                    full_text += block_text + "\n\n"
            
            page_content['text'] = full_text.strip()
            page_content['text_blocks'] = text_blocks
            
            # Extract images
            image_list = page.get_images()
            for img_index, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image temporarily
                    image_filename = f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}"
                    image_path = os.path.join(app.config['OUTPUT_FOLDER'], f"temp_{image_filename}")
                    
                    with open(image_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    page_content['images'].append({
                        'filename': image_filename,
                        'path': image_path,
                        'width': base_image.get('width', 0),
                        'height': base_image.get('height', 0)
                    })
                    
                except Exception as img_error:
                    print(f"Error extracting image {img_index}: {img_error}")
            
            content_data['pages'].append(page_content)
        
        doc.close()
        print(f"✓ Extracted content from {len(content_data['pages'])} pages using PyMuPDF with images")
        return content_data
        
    except Exception as e:
        print(f"PyMuPDF extraction failed: {e}, falling back to pdfplumber")
    
    # Fallback to pdfplumber (original method)
    if pdf_capabilities.get('pdfplumber', False):
        try:
            import pdfplumber
            
            with pdfplumber.open(file_path) as pdf:
                content_data['metadata'] = {
                    'total_pages': len(pdf.pages),
                    'title': pdf.metadata.get('Title', '') if pdf.metadata else '',
                }
                
                for page_num, page in enumerate(pdf.pages):
                    page_content = {
                        'page_number': page_num + 1,
                        'text': '',
                        'tables': [],
                        'images': [],
                        'text_blocks': []
                    }
                    
                    # Extract text
                    if page.extract_text():
                        page_content['text'] = page.extract_text()
                    
                    # Extract tables
                    tables = page.extract_tables()
                    if tables:
                        page_content['tables'] = tables
                    
                    content_data['pages'].append(page_content)
            
            print(f"✓ Extracted content from {len(content_data['pages'])} pages using pdfplumber")
            return content_data
            
        except Exception as e:
            print(f"pdfplumber extraction failed: {e}")
    
    # Fallback to PyPDF2
    if pdf_capabilities.get('pypdf2', False):
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                content_data['metadata'] = {
                    'total_pages': len(pdf_reader.pages),
                    'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                }
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_content = {
                        'page_number': page_num + 1,
                        'text': '',
                        'tables': []
                    }
                    
                    try:
                        page_content['text'] = page.extract_text()
                    except Exception as e:
                        print(f"Error extracting text from page {page_num + 1}: {e}")
                        page_content['text'] = ""
                    
                    content_data['pages'].append(page_content)
            
            print(f"✓ Extracted content from {len(content_data['pages'])} pages using PyPDF2")
            return content_data
            
        except Exception as e:
            print(f"PyPDF2 extraction failed: {e}")
    
    # If all else fails, return empty structure
    print("✗ No PDF extraction method available")
    return content_data

def create_adapted_pdf(adapted_content, output_path, profile):
    """Create a PDF from adapted content with enhanced visual preservation"""
    try:
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT, TA_JUSTIFY
        
        # Create PDF document
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Profile-specific colors
        profile_colors = {
            "dyslexia": colors.HexColor("#0066cc"),
            "adhd": colors.HexColor("#2e8b57"),
            "esl": colors.HexColor("#9400d3")
        }
        profile_color = profile_colors.get(profile, colors.black)
        
        # Get hex string for HTML formatting
        profile_color_hex = {
            "dyslexia": "#0066cc",
            "adhd": "#2e8b57",
            "esl": "#9400d3"
        }.get(profile, "#000000")
        
        # Custom styles for different profiles with enhanced visual design
        if profile == "dyslexia":
            # Dyslexia-friendly styling with background
            content_style = ParagraphStyle(
                'DyslexiaContent',
                parent=styles['Normal'],
                fontSize=14,  # Larger font
                leading=20,   # More line spacing
                fontName='Helvetica',  # Sans-serif font
                leftIndent=20,
                rightIndent=20,
                spaceAfter=12,
                backColor=colors.HexColor("#FFFFF0"),  # Light yellow background
                borderColor=colors.HexColor("#E0E0E0"),
                borderWidth=0.5,
                borderPadding=10,
                alignment=TA_LEFT
            )
            # Special style for bullet points
            bullet_style = ParagraphStyle(
                'DyslexiaBullet',
                parent=content_style,
                leftIndent=35,
                bulletIndent=20
            )
        elif profile == "adhd":
            # ADHD-friendly styling with visual breaks
            content_style = ParagraphStyle(
                'ADHDContent',
                parent=styles['Normal'],
                fontSize=12,
                leading=16,
                spaceAfter=10,
                bulletFontName='Helvetica-Bold',
                borderColor=colors.HexColor("#D0D0D0"),
                borderWidth=0.5,
                borderPadding=8,
                leftIndent=15,
                rightIndent=15
            )
            bullet_style = ParagraphStyle(
                'ADHDBullet',
                parent=content_style,
                leftIndent=30,
                bulletIndent=15
            )
        else:  # ESL
            content_style = ParagraphStyle(
                'ESLContent',
                parent=styles['Normal'],
                fontSize=12,
                leading=18,
                spaceAfter=10,
                alignment=TA_JUSTIFY,
                leftIndent=10,
                rightIndent=10
            )
            bullet_style = ParagraphStyle(
                'ESLBullet',
                parent=content_style,
                leftIndent=25,
                bulletIndent=10
            )
        
        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=20,
            textColor=profile_color
        )
        
        # Add title page
        title = adapted_content.get('metadata', {}).get('title', 'Adapted Document')
        story.append(Paragraph(f"<b>{title}</b>", title_style))
        story.append(Spacer(1, 12))
        story.append(Paragraph(f"Adapted for {profile.title()} Support", content_style))
        story.append(Spacer(1, inch))
        
        # Process each page
        for page in adapted_content.get('pages', []):
            # Page header
            page_title = f"Page {page.get('page_number', 'Unknown')}"
            story.append(Paragraph(page_title, styles['Heading2']))
            story.append(Spacer(1, 12))
            
            # Page content
            if page.get('text'):
                # Split into paragraphs
                paragraphs = page['text'].split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        text_to_process = para_text.strip()
                        
                        # Check if it's a bullet point or list item
                        is_bullet = text_to_process.startswith('•') or text_to_process.startswith('-')
                        is_numbered = any(text_to_process.startswith(f"{i}.") or text_to_process.startswith(f"{i})") for i in range(1, 10))
                        
                        if is_bullet or is_numbered:
                            # Remove bullet/number and clean up
                            if is_bullet:
                                text_to_process = text_to_process[1:].strip()
                            else:
                                # Remove number and period/parenthesis
                                text_to_process = text_to_process.split('.', 1)[1].strip() if '.' in text_to_process[:3] else text_to_process.split(')', 1)[1].strip()
                            
                            # Apply special formatting for list items
                            text_escaped = text_to_process.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            
                            # Use bullet style for list items
                            if 'bullet_style' in locals():
                                formatted_text = f'• <font color="{profile_color_hex}"><b>{text_escaped.split()[0]}</b></font> {" ".join(text_escaped.split()[1:])}'
                                story.append(Paragraph(formatted_text, bullet_style))
                            else:
                                formatted_text = f'• {text_escaped}'
                                story.append(Paragraph(formatted_text, content_style))
                            story.append(Spacer(1, 4))
                        else:
                            # Regular paragraph formatting
                            # Break long sentences for better readability
                            if profile == "dyslexia" and len(text_to_process) > 100:
                                # Split at punctuation marks
                                sentences = text_to_process.replace('. ', '.<br/><br/>').replace('? ', '?<br/><br/>').replace('! ', '!<br/><br/>')
                                text_escaped = sentences.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            else:
                                text_escaped = text_to_process.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                            
                            # Apply color to first word
                            words = text_escaped.split(' ', 1)
                            if words and words[0]:
                                first_word_colored = f'<font color="{profile_color_hex}"><b>{words[0]}</b></font>'
                                rest_text = words[1] if len(words) > 1 else ''
                                formatted_text = f'{first_word_colored} {rest_text}'
                            else:
                                formatted_text = text_escaped
                            
                            try:
                                story.append(Paragraph(formatted_text, content_style))
                                story.append(Spacer(1, 8))
                            except Exception as para_error:
                                # If paragraph fails, add as plain text
                                print(f"Paragraph formatting error: {para_error}")
                                plain_text = para_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                                story.append(Paragraph(plain_text, content_style))
                                story.append(Spacer(1, 8))
            
            # Handle images if present
            if page.get('images'):
                for img_index, image_data in enumerate(page.get('images', [])):
                    try:
                        image_path = image_data['path']
                        if os.path.exists(image_path):
                            from reportlab.platypus import Image as RLImage
                            
                            # Calculate appropriate size based on profile
                            if profile == "dyslexia":
                                max_width = 5 * inch  # Slightly smaller for clarity
                                max_height = 3.5 * inch
                            elif profile == "adhd":
                                max_width = 4.5 * inch  # Smaller to reduce distraction
                                max_height = 3 * inch
                            else:  # ESL
                                max_width = 5.5 * inch
                                max_height = 4 * inch
                            
                            # Get original dimensions
                            original_width = image_data.get('width', 400)
                            original_height = image_data.get('height', 300)
                            
                            # Calculate proportional size
                            width_ratio = max_width / original_width if original_width > max_width else 1
                            height_ratio = max_height / original_height if original_height > max_height else 1
                            ratio = min(width_ratio, height_ratio)
                            
                            img_width = original_width * ratio
                            img_height = original_height * ratio
                            
                            # Add spacing before image
                            story.append(Spacer(1, 0.2 * inch))
                            
                            # Add image with proper alignment
                            img = RLImage(image_path, width=img_width, height=img_height, kind='proportional')
                            story.append(img)
                            
                            # Add image caption/description
                            img_caption_style = ParagraphStyle(
                                'ImageCaption',
                                parent=styles['Normal'],
                                fontSize=10,
                                fontName='Helvetica-Oblique',
                                textColor=colors.HexColor("#666666"),
                                alignment=TA_LEFT if profile == "dyslexia" else 1,  # Center for others
                                spaceAfter=12
                            )
                            
                            caption_text = f"<i>[Image {img_index + 1} from page {page.get('page_number', 'Unknown')}]</i>"
                            if profile == "dyslexia":
                                caption_text += "<br/><i>Description: Add image description for accessibility</i>"
                            
                            story.append(Paragraph(caption_text, img_caption_style))
                            story.append(Spacer(1, 0.3 * inch))
                            
                    except Exception as img_error:
                        print(f"Error adding image to PDF: {img_error}")
            
            # Handle tables if present
            if page.get('tables'):
                for table_data in page['tables']:
                    if table_data and len(table_data) > 0:
                        # Create table with profile-specific styling
                        t = Table(table_data)
                        
                        # Profile-specific table styles
                        if profile == "dyslexia":
                            table_style = TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#B0C4DE")),  # Light steel blue headers
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),  # Left align for dyslexia
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, -1), 12),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 15),
                                ('TOPPADDING', (0, 0), (-1, -1), 10),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor("#FFFFF0")),  # Light yellow
                                ('GRID', (0, 0), (-1, -1), 1.5, colors.HexColor("#4169E1")),  # Royal blue grid
                                ('LINEBELOW', (0, 0), (-1, 0), 2, colors.HexColor("#4169E1")),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor("#FFFFF0"), colors.HexColor("#F0F8FF")])
                            ])
                        elif profile == "adhd":
                            table_style = TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#90EE90")),  # Light green headers
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, -1), 11),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 10),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor("#228B22")),  # Forest green
                                ('LINEBELOW', (0, 0), (-1, 0), 2, colors.HexColor("#228B22")),
                                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE')
                            ])
                        else:  # ESL
                            table_style = TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#DDA0DD")),  # Plum headers
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, -1), 11),
                                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor("#F8F8FF")),  # Ghost white
                                ('GRID', (0, 0), (-1, -1), 1, colors.HexColor("#9370DB")),  # Medium purple
                                ('LINEBELOW', (0, 0), (-1, 0), 2, colors.HexColor("#9370DB"))
                            ])
                        
                        t.setStyle(table_style)
                        story.append(t)
                        story.append(Spacer(1, 12))
            
            # Add page break between pages (except for the last page)
            if page != adapted_content['pages'][-1]:
                story.append(Spacer(1, 0.5 * inch))
        
        # Build PDF
        doc.build(story)
        
        # Clean up temporary image files
        cleanup_temp_images(adapted_content)
        
        print(f"✓ Created adapted PDF: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error creating adapted PDF: {e}")
        import traceback
        traceback.print_exc()
        return False

def cleanup_temp_images(adapted_content):
    """Clean up temporary image files created during PDF extraction"""
    try:
        for page in adapted_content.get('pages', []):
            for image_data in page.get('images', []):
                image_path = image_data.get('path')
                if image_path and os.path.exists(image_path) and 'temp_' in image_path:
                    os.remove(image_path)
                    print(f"Cleaned up temporary image: {image_path}")
    except Exception as e:
        print(f"Error cleaning up temporary images: {e}")

def diagnose_pdf_content(pdf_path):
    """Diagnose PDF content to determine best preservation approach"""
    try:
        import fitz
        
        print(f"\n🔍 Analyzing PDF: {pdf_path}")
        print("=" * 50)
        
        doc = fitz.open(pdf_path)
        pdf_info = {
            'total_pages': len(doc),
            'has_images': False,
            'has_vectors': False,
            'has_backgrounds': False,
            'has_forms': False,
            'total_images': 0,
            'total_vectors': 0
        }
        
        for page_num, page in enumerate(doc):
            # Check for images
            images = page.get_images()
            if images:
                pdf_info['has_images'] = True
                pdf_info['total_images'] += len(images)
                print(f"Page {page_num + 1}: {len(images)} images found")
            
            # Check for vector drawings
            drawings = page.get_drawings()
            if drawings:
                pdf_info['has_vectors'] = True
                pdf_info['total_vectors'] += len(drawings)
                
                # Check for background fills
                for drawing in drawings:
                    if drawing["type"] == "f" and drawing["rect"].width > page.rect.width * 0.8:
                        pdf_info['has_backgrounds'] = True
                        break
            
            # Check for forms/annotations
            annots = list(page.annots())
            if annots:
                pdf_info['has_forms'] = True
        
        doc.close()
        
        # Print summary
        print(f"\n📊 PDF Summary:")
        print(f"Total pages: {pdf_info['total_pages']}")
        print(f"Contains images: {'Yes' if pdf_info['has_images'] else 'No'} ({pdf_info['total_images']} total)")
        print(f"Contains vector graphics: {'Yes' if pdf_info['has_vectors'] else 'No'} ({pdf_info['total_vectors']} total)")
        print(f"Has background fills: {'Yes' if pdf_info['has_backgrounds'] else 'No'}")
        print(f"Has forms/annotations: {'Yes' if pdf_info['has_forms'] else 'No'}")
        print("=" * 50)
        
        return pdf_info
        
    except Exception as e:
        print(f"Error diagnosing PDF: {e}")
        return None

def estimate_text_width(text, font_size):
    """Estimate text width based on character count and font size"""
    # Approximate character widths for common fonts
    char_widths = {
        ' ': 0.25,
        'i': 0.3, 'l': 0.3, 'I': 0.3, '1': 0.5,
        'f': 0.4, 'j': 0.4, 't': 0.4, 'r': 0.4,
        'm': 0.8, 'w': 0.8, 'M': 0.8, 'W': 0.9,
        'default': 0.5  # Average character width
    }
    
    width = 0
    for char in text:
        width += char_widths.get(char, char_widths['default'])
    
    return width * font_size

def create_visual_preserved_pdf(input_path, output_path, profile, adapted_content=None):
    """
    Preserve ALL visual elements while updating text content
    Enhanced implementation with better text mapping and preservation
    """
    try:
        import fitz  # PyMuPDF
        
        print(f"Creating visually-preserved adapted PDF with updated text for {profile} profile")
        
        # Open the original PDF
        doc = fitz.open(input_path)
        
        # Profile-specific settings
        profile_settings = {
            "dyslexia": {
                "fill_color": (1, 1, 0.941),     # Light yellow
                "fill_opacity": 0.08,            # 8% opacity
                "font": "helv",                  # Sans-serif
                "font_size_min": 12,             # Minimum font size
                "highlight_color": (0, 0.4, 0.8), # Blue for first word
            },
            "adhd": {
                "fill_color": (0.941, 1, 0.941), # Light green
                "fill_opacity": 0.06,            # 6% opacity
                "font": "helv",
                "font_size_min": 11,
                "highlight_color": (0.18, 0.545, 0.341), # Green for first word
            },
            "esl": {
                "fill_color": (0.98, 0.941, 1),  # Light purple
                "fill_opacity": 0.05,            # 5% opacity
                "font": "helv",
                "font_size_min": 11,
                "highlight_color": (0.58, 0, 0.827), # Purple for first word
            }
        }
        
        settings = profile_settings.get(profile, profile_settings["dyslexia"])
        
        # Process each page
        for page_num, page in enumerate(doc):
            print(f"Processing page {page_num + 1}/{len(doc)}")
            
            # Get page dimensions for later use
            page_rect = page.rect
            
            # Method: Smart text replacement with visual preservation
            print(f"DEBUG: adapted_content available: {bool(adapted_content)}")
            if adapted_content:
                print(f"DEBUG: adapted_content has {len(adapted_content.get('pages', []))} pages")
                
            if adapted_content and page_num < len(adapted_content.get('pages', [])):
                adapted_page = adapted_content['pages'][page_num]
                adapted_text = adapted_page.get('text', '')
                print(f"DEBUG: Page {page_num + 1} adapted text length: {len(adapted_text) if adapted_text else 0}")
                print(f"DEBUG: Page {page_num + 1} adapted text preview: {adapted_text[:200] if adapted_text else 'No text'}")
                
                if adapted_text.strip():
                    print(f"DEBUG: Using adapted text for page {page_num + 1}")
                    text_to_use = adapted_text
                else:
                    print(f"DEBUG: No adapted text available, using original text for page {page_num + 1}")
                    # Get original text from the page
                    original_text = page.get_text()
                    text_to_use = original_text
                
                if text_to_use and text_to_use.strip():
                    # Get text with detailed formatting info
                    text_dict = page.get_text("dict", sort=True)
                    
                    # Build text blocks with full formatting details
                    text_blocks = []
                    for block in text_dict["blocks"]:
                        if block["type"] == 0:  # Text block
                            block_info = {
                                'bbox': block['bbox'],
                                'lines': []
                            }
                            
                            for line in block["lines"]:
                                line_info = {
                                    'bbox': line['bbox'],
                                    'spans': [],
                                    'text': ''
                                }
                                
                                for span in line["spans"]:
                                    span_text = span["text"]
                                    line_info['text'] += span_text
                                    line_info['spans'].append({
                                        'text': span_text,
                                        'font': span.get('font', settings['font']),
                                        'size': span.get('size', 12),
                                        'flags': span.get('flags', 0),
                                        'color': span.get('color', 0),
                                        'bbox': span['bbox']
                                    })
                                
                                if line_info['text'].strip():
                                    block_info['lines'].append(line_info)
                            
                            if block_info['lines']:
                                text_blocks.append(block_info)
                    
                    # Process text into paragraphs
                    text_paragraphs = [p.strip() for p in text_to_use.split('\n\n') if p.strip()]
                    text_lines = []
                    for para in text_paragraphs:
                        para_lines = [line.strip() for line in para.split('\n') if line.strip()]
                        text_lines.extend(para_lines)
                    
                    # Advanced text mapping: match paragraphs and lines
                    print(f"Mapping {len(text_blocks)} text blocks to {len(text_lines)} text lines")
                    
                    # Clear all text areas first (preserve images and graphics)
                    for block in text_blocks:
                        for line in block['lines']:
                            # Create a slightly expanded rectangle to ensure complete coverage
                            line_rect = fitz.Rect(line['bbox'])
                            line_rect.y0 -= 2  # Expand up
                            line_rect.y1 += 2  # Expand down
                            line_rect.x0 -= 2  # Expand left
                            line_rect.x1 += 2  # Expand right
                            
                            # DO NOT fill with white - this covers backgrounds!
                            # The issue is that white rectangles cover any background images/colors
                            # page.draw_rect(line_rect, color=None, fill=(1, 1, 1), fill_opacity=1)
                    
                    # Now add text with proper formatting
                    text_line_idx = 0
                    
                    for block_idx, block in enumerate(text_blocks):
                        if text_line_idx >= len(text_lines):
                            break
                        
                        for line_idx, line in enumerate(block['lines']):
                            if text_line_idx >= len(text_lines):
                                break
                            
                            # Get the text line to display
                            display_line = text_lines[text_line_idx]
                            text_line_idx += 1
                            
                            # Use original line's formatting as template
                            if line['spans']:
                                template_span = line['spans'][0]
                                font_name = template_span['font']
                                font_size = max(template_span['size'], settings['font_size_min'])
                                font_color = template_span['color']
                                
                                # Convert color from int to RGB tuple
                                if isinstance(font_color, int):
                                    r = (font_color >> 16) & 0xFF
                                    g = (font_color >> 8) & 0xFF
                                    b = font_color & 0xFF
                                    font_color = (r/255.0, g/255.0, b/255.0)
                                else:
                                    font_color = (0, 0, 0)
                                
                                # Calculate text position
                                x_pos = line['bbox'][0]
                                y_pos = line['bbox'][3] - 2  # Baseline position
                                
                                # Handle first word highlighting for content
                                words = display_line.split()
                                if words and line_idx == 0 and block_idx == 0:  # First line of first block
                                    # Insert first word with profile color
                                    first_word = words[0]
                                    page.insert_text(
                                        fitz.Point(x_pos, y_pos),
                                        first_word,
                                        fontname=font_name,
                                        fontsize=font_size,
                                        color=settings['highlight_color']
                                    )
                                    
                                    # Estimate first word width
                                    text_width = estimate_text_width(first_word, font_size)
                                    x_pos += text_width + font_size * 0.25  # Add space
                                    
                                    # Insert rest of line
                                    if len(words) > 1:
                                        rest_text = ' '.join(words[1:])
                                        page.insert_text(
                                            fitz.Point(x_pos, y_pos),
                                            rest_text,
                                            fontname=font_name,
                                            fontsize=font_size,
                                            color=font_color
                                        )
                                else:
                                    # Insert entire line with original color
                                    page.insert_text(
                                        fitz.Point(x_pos, y_pos),
                                        display_line,
                                        fontname=font_name,
                                        fontsize=font_size,
                                        color=font_color
                                    )
            
            # Add subtle accessibility tint overlay
            page.draw_rect(
                page_rect,
                color=None,
                fill=settings['fill_color'],
                fill_opacity=settings['fill_opacity'],
                overlay=True
            )
            
            # Optional: Add subtle reading guides for dyslexia
            if profile == "dyslexia":
                # Add very faint horizontal lines every 40 points
                line_spacing = 40
                for y in range(int(page_rect.height / line_spacing)):
                    y_pos = y * line_spacing
                    page.draw_line(
                        fitz.Point(0, y_pos),
                        fitz.Point(page_rect.width, y_pos),
                        color=(0.97, 0.97, 0.97),
                        width=0.5
                    )
        
        # Save with optimal settings
        doc.save(output_path, garbage=0, deflate=True, clean=False)
        doc.close()
        
        print(f"✓ Created visually-preserved adapted PDF with updated text: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error with enhanced text replacement: {e}")
        import traceback
        traceback.print_exc()
        # Fall back to overlay method
        return create_visual_preserved_with_text_overlay(input_path, output_path, profile, adapted_content)

def create_visual_preserved_pdf_simple(input_path, output_path, profile):
    """
    Simplest fallback using pdf2image - just adds a tint overlay
    Preserves EVERYTHING in the original PDF
    """
    try:
        from pdf2image import convert_from_path
        from PIL import Image
        
        print(f"Using pdf2image method for visual preservation")
        
        # Convert PDF to images at high quality
        images = convert_from_path(input_path, dpi=200)
        
        if not images:
            raise Exception("No pages found in PDF")
        
        # Profile-specific tints (very subtle)
        profile_tints = {
            "dyslexia": (255, 255, 235, 30),    # Light yellow, 30/255 opacity
            "adhd": (240, 255, 240, 25),        # Light green, 25/255 opacity  
            "esl": (250, 245, 255, 20)          # Light purple, 20/255 opacity
        }
        
        tint = profile_tints.get(profile, profile_tints["dyslexia"])
        processed_images = []
        
        for i, img in enumerate(images):
            print(f"Processing page {i + 1}/{len(images)}")
            
            # Just add a subtle tint overlay - preserves everything
            overlay = Image.new('RGBA', img.size, tint)
            img_tinted = Image.alpha_composite(img.convert('RGBA'), overlay)
            
            # Optional: Slight contrast boost for readability
            if profile == "dyslexia" or profile == "adhd":
                from PIL import ImageEnhance
                enhancer = ImageEnhance.Contrast(img_tinted)
                img_tinted = enhancer.enhance(1.1)  # 10% contrast boost
            
            # Convert back to RGB for PDF
            processed_images.append(img_tinted.convert('RGB'))
        
        # Save as PDF
        print("Saving as PDF...")
        processed_images[0].save(
            output_path,
            save_all=True,
            append_images=processed_images[1:] if len(processed_images) > 1 else [],
            format='PDF',
            resolution=200.0
        )
        
        print(f"✓ Created visually-preserved PDF: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error in pdf2image method: {e}")
        return False

def create_visual_preserved_with_text_overlay(input_path, output_path, profile, adapted_content):
    """
    Alternative method: Use redaction to remove text while preserving backgrounds
    Then add adapted text on top
    """
    try:
        import fitz
        
        print(f"Creating visual PDF with text overlay for {profile} - using redaction method")
        
        # Profile-specific settings
        profile_settings = {
            "dyslexia": {
                "tint": (1, 1, 0.94),
                "opacity": 0.08,
                "highlight_color": (0, 0.4, 0.8),
            },
            "adhd": {
                "tint": (0.94, 1, 0.94),
                "opacity": 0.06,
                "highlight_color": (0.18, 0.545, 0.341),
            },
            "esl": {
                "tint": (0.98, 0.94, 1),
                "opacity": 0.05,
                "highlight_color": (0.58, 0, 0.827),
            }
        }
        
        settings = profile_settings.get(profile, profile_settings["dyslexia"])
        
        # Open original PDF
        doc = fitz.open(input_path)
        
        for page_num, page in enumerate(doc):
            print(f"Processing page {page_num + 1} with redaction method")
            
            # Get text areas for redaction
            text_dict = page.get_text("dict", sort=True)
            
            # Redact all text areas - this removes text while preserving backgrounds
            for block in text_dict["blocks"]:
                if block["type"] == 0:  # Text block
                    # Add redaction annotation
                    rect = fitz.Rect(block['bbox'])
                    page.add_redact_annot(rect)
            
            # Apply redactions - this removes text but preserves backgrounds
            page.apply_redactions()
            
            # Add accessibility tint overlay
            page.draw_rect(
                page.rect,
                color=None,
                fill=settings['tint'],
                fill_opacity=settings['opacity'],
                overlay=False
            )
            
            # Add adapted text as visible layer
            if adapted_content and page_num < len(adapted_content.get('pages', [])):
                adapted_page = adapted_content['pages'][page_num]
                adapted_text = adapted_page.get('text', '')
                
                if adapted_text.strip():
                    # Use default text positioning
                    margin = 50
                    y_position = 80
                    font_size = 12
                    max_width = page.rect.width - (2 * margin)
                    
                    # Insert the adapted text as a text block
                    text_rect = fitz.Rect(margin, y_position - 10, margin + max_width, page.rect.height - 50)
                    
                    # Use insert_textbox for better text flow
                    overflow = page.insert_textbox(
                        text_rect,
                        adapted_text,
                        fontsize=font_size,
                        fontname="helv",
                        color=(0, 0, 0),
                        align=fitz.TEXT_ALIGN_LEFT
                    )
                    
                    if overflow:
                        print(f"Warning: Text overflow on page {page_num + 1}, some content may be cut off")
            
        
        # Save the new PDF with optimization
        doc.save(output_path, garbage=0, deflate=True)
        doc.close()
        
        print(f"✓ Created visual PDF with enhanced text overlay: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error creating text overlay PDF: {e}")
        import traceback
        traceback.print_exc()
        # Final fallback - just use the simple visual preservation
        return create_visual_preserved_pdf_simple(input_path, output_path, profile)

# New service-based PDF processing function
def process_pdf_with_services(file_path, file_id, filename, profile, export_format='pdf', target_language=None):
    """Service-based PDF processing function using the new architecture"""
    print(f"Starting service-based PDF processing for {filename} with profile: {profile}")
    
    try:
        processing_tasks[file_id]['message'] = 'Initializing PDF processing services...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 5, 'percentage': 5}
        
        # Use the migration helper to process the PDF
        processing_tasks[file_id]['message'] = 'Processing PDF with visual preservation...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 20, 'percentage': 20}
        
        # Process with services - this handles adaptation and visual preservation
        result = pdf_migration_helper.process_with_pdf_template_system(
            file_path, 
            profile, 
            direct_adapt=True,
            preserve_visuals=True,
            translate=(target_language and target_language.strip() != ''),
            target_language=target_language
        )
        
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 60, 'percentage': 60}
        
        if not result['success']:
            raise Exception(f"PDF processing failed: {result.get('error', 'Unknown error')}")
        
        # Update processing status
        processing_tasks[file_id]['message'] = 'PDF adaptation completed successfully!'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 90, 'percentage': 90}
        
        # Store the output path for downloads
        output_path = result['output_path']
        processing_tasks[file_id]['output_path'] = output_path
        processing_tasks[file_id]['adapted_content'] = result.get('adapted_content', {})
        
        # Handle translation if requested and available
        if target_language and target_language.strip():
            processing_tasks[file_id]['message'] = f'Creating translation in {target_language}...'
            processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 95, 'percentage': 95}
            
            # Check if translation was successful
            if 'translated_output_path' in result:
                processing_tasks[file_id]['translated_output_path'] = result['translated_output_path']
                processing_tasks[file_id]['translated_language'] = result['translated_language']
                processing_tasks[file_id]['translation_languages'] = [target_language]
                print(f"✓ Translation completed: {result['translated_output_path']}")
            else:
                print(f"⚠ Translation was requested but failed for {target_language}")
                processing_tasks[file_id]['translation_error'] = f'Translation to {target_language} failed'
        
        # Mark as completed
        processing_tasks[file_id]['status'] = 'completed'
        processing_tasks[file_id]['message'] = 'Processing completed successfully!'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 100, 'percentage': 100}
        
        print(f"✓ Service-based PDF processing completed successfully for {filename}")
        return True
        
    except Exception as e:
        print(f"✗ Error in service-based PDF processing: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Update task with error
        processing_tasks[file_id]['status'] = 'error'
        processing_tasks[file_id]['message'] = f'Error: {str(e)}'
        processing_tasks[file_id]['error'] = str(e)
        
        return False

# Original PDF processing function (kept for fallback)
def process_with_pdf_template_system(file_path, file_id, filename, profile, export_format='pdf', target_language=None):
    """Main PDF processing function - properly handles export format"""
    print(f"Starting PDF template processing for {filename} with export format: {export_format}")
    
    try:
        # Check if we have basic PDF capabilities
        if not pdf_capabilities.get('pdfplumber', False) and not pdf_capabilities.get('pypdf2', False):
            raise Exception("No PDF processing libraries available. Please install pdfplumber or PyPDF2.")
        
        processing_tasks[file_id]['message'] = 'Extracting content from PDF...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 10, 'percentage': 10}
        
        # Step 1: Extract content from PDF
        pdf_content = pdf_service.extract_content_from_pdf(file_path)
        
        if not pdf_content['pages']:
            raise Exception("No content could be extracted from the PDF")
        
        processing_tasks[file_id]['message'] = 'Adapting PDF content...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 30, 'percentage': 30}
        
        # Step 2: Adapt content (with optional translation)
        print(f"\n=== ADAPTATION DEBUG ===")
        print(f"Profile: {profile}")
        print(f"Target language: {target_language}")
        print(f"Force adaptation: True")
        
        result = adapt_pdf_content(pdf_content, profile, target_language, force_adaptation=True)
        if isinstance(result, tuple):
            adapted_content, translated_content = result
            print(f"Got tuple result - has translation: {translated_content is not None}")
            if translated_content:
                print(f"Translation has {len(translated_content.get('pages', []))} pages")
        else:
            adapted_content = result
            translated_content = None
            print(f"Got single result - no translation")
        print(f"=======================\n")
        
        processing_tasks[file_id]['message'] = f'Creating adapted {export_format.upper()}...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 70, 'percentage': 70}
        
        # Step 3: Create output based on requested format
        success = False
        output_path = None
        translated_path = None
        
        if export_format.lower() == 'pdf':
            # Create PDF output
            base_name = os.path.splitext(filename)[0]
            if base_name.startswith('adapted_'):
                output_filename = f"{base_name}.pdf"
            else:
                output_filename = f"adapted_{base_name}.pdf"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
            print(f"Creating PDF output at: {output_path}")
            print(f"Output filename: {output_filename}")
            print(f"File ID: {file_id}")
            
            # Check if visual preservation is requested (can be added as a form parameter)
            preserve_visuals = processing_tasks.get(file_id, {}).get('preserve_visuals', True)
            
            if preserve_visuals:
                try:
                    # Diagnose PDF content first
                    diagnose_pdf_content(file_path)
                    
                    # Use the visual-preserving method that maintains ALL original graphics
                    print("Using visual-preserving PDF adaptation method")
                    success = create_visual_preserved_pdf(file_path, output_path, profile, adapted_content)
                    
                    # If PyMuPDF method fails, try visual preservation with text overlay
                    if not success:
                        print("PyMuPDF method failed, trying visual preservation with text overlay")
                        success = create_visual_preserved_with_text_overlay(file_path, output_path, profile, adapted_content)
                    
                    # If visual preservation fails completely, fall back to standard method
                    if not success:
                        print("All visual preservation methods failed, falling back to standard method")
                        success = pdf_service.create_adapted_pdf(adapted_content, output_path, profile)
                        
                except Exception as e:
                    print(f"Error in visual preservation pipeline: {e}")
                    import traceback
                    traceback.print_exc()
                    # Fall back to standard method
                    print("Falling back to standard PDF creation due to visual preservation error")
                    success = pdf_service.create_adapted_pdf(adapted_content, output_path, profile)
            else:
                # Use standard text-focused adaptation
                success = pdf_service.create_adapted_pdf(adapted_content, output_path, profile)
            
            # Create translated version if available
            if success and translated_content and target_language:
                print(f"\n=== CREATING TRANSLATED PDF ===")
                print(f"Target language: {target_language}")
                print(f"Original filename: {filename}")
                
                # Clean the base filename to avoid double prefixes
                base_for_translation = os.path.splitext(filename)[0]
                if base_for_translation.startswith('adapted_'):
                    base_for_translation = base_for_translation[8:]  # Remove 'adapted_' prefix
                
                translated_filename = f"translated_{target_language}_{base_for_translation}.pdf"
                translated_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{translated_filename}")
                print(f"Translation output path: {translated_path}")
                
                translation_success = False
                if preserve_visuals:
                    translation_success = create_visual_preserved_pdf(file_path, translated_path, profile, translated_content)
                    if not translation_success:
                        print("Visual preservation failed for translation, trying text overlay method")
                        translation_success = create_visual_preserved_with_text_overlay(file_path, translated_path, profile, translated_content)
                    if not translation_success:
                        print("All visual methods failed for translation, using standard method")
                        translation_success = pdf_service.create_adapted_pdf(translated_content, translated_path, profile)
                else:
                    translation_success = pdf_service.create_adapted_pdf(translated_content, translated_path, profile)
                
                # Verify the translated file was created
                if translation_success and os.path.exists(translated_path):
                    print(f"✓ Successfully created translated PDF at: {translated_path}")
                    print(f"✓ File size: {os.path.getsize(translated_path)} bytes")
                else:
                    print(f"✗ Failed to create translated PDF")
                    translated_path = None  # Clear the path if creation failed
                print("==============================\n")
            
        elif export_format.lower() == 'pptx':
            # Create PowerPoint output
            base_name = os.path.splitext(filename)[0]
            if base_name.startswith('adapted_'):
                output_filename = f"{base_name}.pptx"
            else:
                output_filename = f"adapted_{base_name}.pptx"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
            success = create_adapted_presentation_from_pdf(adapted_content, output_path, profile)
        
        else:
            raise Exception(f"Unsupported export format: {export_format}")
        
        if not success:
            raise Exception(f"Failed to create adapted {export_format}")
        
        # Verify the file was actually created
        if output_path and not os.path.exists(output_path):
            raise Exception(f"File creation reported success but file not found at: {output_path}")
        
        print(f"✓ Successfully created file at: {output_path}")
        print(f"✓ File size: {os.path.getsize(output_path) if output_path and os.path.exists(output_path) else 'Unknown'} bytes")
        
        # Update processing status
        status_data = {
            'status': 'complete',
            'message': f'PDF content successfully adapted to {export_format.upper()}',
            'adapted_path': output_path,
            'export_format': export_format,
            'profile': profile,
            'progress': {
                'total': 100,
                'processed': 100,
                'percentage': 100
            }
        }
        
        # Add translation information if available
        if translated_path and target_language:
            print(f"\n=== SETTING TRANSLATION INFO ===")
            print(f"translated_path: {translated_path}")
            print(f"target_language: {target_language}")
            print(f"File exists: {os.path.exists(translated_path)}")
            
            status_data['has_translation'] = True
            status_data['translated_path'] = translated_path
            status_data['target_language'] = target_language
            status_data['translated_language'] = target_language.title()
            # Extract just the filename without the file_id prefix for download links
            full_filename = os.path.basename(translated_path)
            if full_filename.startswith(f"{file_id}_"):
                status_data['translated_filename'] = full_filename[len(f"{file_id}_"):]
            else:
                status_data['translated_filename'] = full_filename
            status_data['message'] += f' and translated to {target_language}'
            
            print(f"translated_filename set to: {status_data['translated_filename']}")
            print(f"has_translation set to: {status_data['has_translation']}")
            print(f"================================\n")
        
        processing_tasks[file_id] = status_data
        
        print(f"✓ PDF processing complete for {file_id}")
        return output_path
        
    except Exception as e:
        print(f"✗ Error in PDF processing: {str(e)}")
        processing_tasks[file_id] = {
            'status': 'error',
            'message': f"PDF processing failed: {str(e)}"
        }
        return None

# Fix the upload route's PDF handling
def get_fixed_upload_route():
    """Fixed upload route with proper PDF handling"""
    return '''
@app.route('/upload', methods=['POST'])
def upload():
    """Handle file upload and direct to appropriate process - supports both PPTX and PDF"""
    try:
        # Get the uploaded file
        file = None
        if 'pptx' in request.files and request.files['pptx'].filename:
            file = request.files['pptx']
        elif 'pdf' in request.files and request.files['pdf'].filename:
            file = request.files['pdf']
        elif 'file' in request.files and request.files['file'].filename:
            file = request.files['file']
        
        # Validate file presence
        if not file or file.filename == '':
            return render_template_string(ERROR_TEMPLATE, 
                message="No file selected. Please choose a PowerPoint (.pptx) or PDF file."), 400
        
        # Get form parameters
        profile = request.form.get('profile')
        action = request.form.get('action', 'assess')
        target_language = request.form.get('target_language', '')
        export_format = request.form.get('export_format', 'pdf')  # Default to PDF
        
        # Debug form parameters for second route
        print(f"DEBUG: Form parameters (route 2) - profile: {profile}, action: {action}, target_language: '{target_language}', export_format: {export_format}")
        
        # Validate required parameters
        if not profile:
            return render_template_string(ERROR_TEMPLATE, 
                message="Please select a learning profile."), 400
        
        # Secure filename and validate extension
        filename = secure_filename(file.filename)
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext not in ['.pptx', '.pdf']:
            return render_template_string(ERROR_TEMPLATE, 
                message=f"Invalid file type '{file_ext}'. Please upload a PowerPoint (.pptx) or PDF file."), 400
        
        # Generate unique file ID and save file
        file_id = str(uuid.uuid4())
        processing_tasks[file_id] = {
            'status': 'upload', 
            'filename': filename, 
            'profile': profile,
            'file_type': file_ext,
            'export_format': export_format  # Store export format
        }
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        file.save(file_path)
        
        # Route based on action
        if action == 'assess':
            # Assessment route - PDFs get converted to PPTX for assessment
            if file_ext == '.pdf':
                # For PDF assessment, we need to convert to PPTX first
                return render_template_string(ERROR_TEMPLATE,
                    message="PDF assessment is not yet implemented. Please use Direct Adaptation instead."), 400
            else:
                if target_language:
                    return redirect(url_for('analyze_scaffolding', file_id=file_id, profile=profile, target_language=target_language))
                else:
                    return redirect(url_for('analyze_scaffolding', file_id=file_id, profile=profile))
        
        else:  # action == 'adapt'
            # Start adaptation process
            processing_tasks[file_id].update({
                'status': 'processing', 
                'message': 'Starting adaptation...',
                'export_format': export_format,
                'progress': {'total': 0, 'processed': 0, 'percentage': 0}
            })
            
            # Choose processing method based on file type
            if file_ext == '.pdf':
                # Use new service-based PDF processing
                print(f"Processing PDF with service-based system (route 2): {filename}")
                thread_target = process_pdf_with_services
                thread_args = (file_path, file_id, filename, profile, export_format, target_language)
            else:  # .pptx
                print(f"Processing PowerPoint: {filename}")
                thread_target = process_presentation_efficiently
                thread_args = (file_path, file_id, filename, profile, target_language)
            
            # Start background processing
            thread = threading.Thread(
                target=thread_target,
                args=thread_args,
                name=f"process-{file_id}"
            )
            thread.daemon = True
            thread.start()
            
            # Profile display names
            profile_names = {
                "dyslexia": "Dyslexia Support",
                "adhd": "ADHD Support", 
                "esl": "English Language Learners"
            }
            
            # Return processing page
            return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                file_id=file_id, 
                filename=f"adapted_{filename}",
                profile=profile,
                profile_name=profile_names.get(profile, profile)
            )
            
    except Exception as e:
        print(f"Error in upload route: {str(e)}")
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, 
            message=f"Upload failed: {str(e)}"), 500
'''

def convert_pptx_to_pdf_template(pptx_path, pdf_path):
    """Convert PPTX to PDF while preserving layout and converting videos to links"""
    try:
        # First extract content from PPTX
        prs = Presentation(pptx_path)
        
        # Create PDF document
        doc = SimpleDocTemplate(pdf_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles for different content types
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=20,
            textColor=colors.darkblue
        )
        
        content_style = ParagraphStyle(
            'CustomContent',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=12,
            leftIndent=20
        )
        
        for slide_index, slide in enumerate(prs.slides):
            # Add slide number
            slide_title = Paragraph(f"<b>Slide {slide_index + 1}</b>", title_style)
            story.append(slide_title)
            story.append(Spacer(1, 12))
            
            # Process slide content
            for shape in slide.shapes:
                # Handle text content
                if hasattr(shape, "text_frame") and shape.text:
                    text_para = Paragraph(shape.text, content_style)
                    story.append(text_para)
                    story.append(Spacer(1, 6))
                
                # Handle images
                elif hasattr(shape, "image"):
                    try:
                        # Save image temporarily
                        image_stream = BytesIO(shape.image.blob)
                        temp_img_path = os.path.join(app.config['UPLOAD_FOLDER'], f'temp_slide_{slide_index}_img.png')
                        
                        with open(temp_img_path, 'wb') as f:
                            f.write(image_stream.getvalue())
                        
                        # Add image to PDF
                        img = RLImage(temp_img_path, width=4*inch, height=3*inch)
                        story.append(img)
                        story.append(Spacer(1, 12))
                        
                        # Clean up temp file
                        os.remove(temp_img_path)
                    except Exception as e:
                        print(f"Error processing image: {e}")
                
                # Handle media (videos) - convert to links
                elif hasattr(shape, "media_format"):
                    media_name = getattr(shape, 'media_name', 'Video Content')
                    link_text = f"[VIDEO LINK: {media_name}]"
                    link_para = Paragraph(f"<i>{link_text}</i>", content_style)
                    story.append(link_para)
                    story.append(Spacer(1, 6))
                
                # Handle tables
                elif hasattr(shape, "table"):
                    table_text = "TABLE CONTENT:\n"
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        table_text += row_text + "\n"
                    
                    table_para = Paragraph(table_text.replace('\n', '<br/>'), content_style)
                    story.append(table_para)
                    story.append(Spacer(1, 12))
            
            # Add page break between slides
            story.append(Spacer(1, 20))
        
        # Build PDF
        doc.build(story)
        return True
        
    except Exception as e:
        print(f"Error converting PPTX to PDF template: {e}")
        return False

def adapt_pdf_content(pdf_content, profile, target_language=None, force_adaptation=True):
    """Adapt PDF content for specific learning profile with optional translation"""
    print(f"Adapting PDF content for {profile}")
    if target_language:
        print(f"Translation to {target_language} will be created")
    if force_adaptation:
        print("DEBUG: Force adaptation is enabled - will adapt all text content")
    
    # Collect text elements that need adaptation
    text_elements = []
    element_references = []
    
    for page_index, page in enumerate(pdf_content['pages']):
        page_text = page.get('text', '')
        print(f"DEBUG: Page {page_index + 1} has {len(page_text)} characters")
        print(f"DEBUG: Page {page_index + 1} text preview: {page_text[:100] if page_text else 'No text'}")
        
        if page_text:
            if force_adaptation:
                needs_adapt = len(page_text.strip()) >= 15  # Simple length check when forcing
                print(f"DEBUG: Force adaptation - Page {page_index + 1} needs adaptation: {needs_adapt}")
            else:
                needs_adapt = needs_adaptation(page_text, profile)
                print(f"DEBUG: Page {page_index + 1} needs adaptation: {needs_adapt}")
            
            if needs_adapt:
                text_elements.append(page_text)
                element_references.append({
                    'type': 'page_text',
                    'page_index': page_index
                })
            else:
                print(f"DEBUG: Page {page_index + 1} skipped - doesn't need adaptation")
        else:
            print(f"DEBUG: Page {page_index + 1} skipped - no text")
    
    print(f"Found {len(text_elements)} text elements requiring adaptation")
    
    # Adapt text content in batches
    adapted_content = {'pages': [], 'metadata': pdf_content['metadata'].copy()}
    translated_content = None
    
    if text_elements:
        print(f"DEBUG: Processing {len(text_elements)} text elements")
        adapted_texts = process_text_batch(text_elements, profile)
        print(f"DEBUG: Got {len(adapted_texts)} adapted texts")
        
        # Create translations if requested (for any profile with target_language)
        translated_texts = []
        translated_content = None
        if target_language and target_language.strip():
            print(f"\n=== TRANSLATION CREATION ===")
            print(f"Profile: {profile}")
            print(f"Target language: {target_language}")
            print(f"Creating translations to {target_language}")
            
            for i, text in enumerate(text_elements):
                print(f"DEBUG: Translating text {i+1}/{len(text_elements)}")
                translated_text = translate_text(text, target_language)
                translated_texts.append(translated_text)
            
            print(f"DEBUG: Created {len(translated_texts)} translations")
            # Create translated content structure
            translated_content = {'pages': [], 'metadata': pdf_content['metadata'].copy()}
            print("===========================\n")
        
        # Apply adaptations back to content
        for i, adapted_text in enumerate(adapted_texts):
            ref = element_references[i]
            page_index = ref['page_index']
            
            # Ensure we have enough pages in adapted_content
            while len(adapted_content['pages']) <= page_index:
                adapted_content['pages'].append({
                    'page_number': len(adapted_content['pages']) + 1,
                    'text': '',
                    'tables': []
                })
            
            adapted_content['pages'][page_index]['text'] = adapted_text
            
            # Handle translated content
            if translated_content and i < len(translated_texts):
                while len(translated_content['pages']) <= page_index:
                    translated_content['pages'].append({
                        'page_number': len(translated_content['pages']) + 1,
                        'text': '',
                        'tables': []
                    })
                translated_content['pages'][page_index]['text'] = translated_texts[i]
    else:
        print("DEBUG: No text elements found requiring adaptation - using original content")
        # If no adaptation needed, just return the original content
        adapted_content = pdf_content.copy()
    
    # Copy over non-adapted pages for both adapted and translated content
    for page_index, original_page in enumerate(pdf_content['pages']):
        if page_index >= len(adapted_content['pages']):
            adapted_content['pages'].append(original_page.copy())
        else:
            # Copy tables and other non-text elements
            adapted_content['pages'][page_index]['tables'] = original_page.get('tables', [])
            adapted_content['pages'][page_index]['page_number'] = original_page.get('page_number', page_index + 1)
        
        # Handle translated content as well
        if translated_content:
            if page_index >= len(translated_content['pages']):
                translated_content['pages'].append(original_page.copy())
            else:
                # Copy tables and other non-text elements
                translated_content['pages'][page_index]['tables'] = original_page.get('tables', [])
                translated_content['pages'][page_index]['page_number'] = original_page.get('page_number', page_index + 1)
    
    print(f"Adaptation complete: {len(adapted_content['pages'])} pages")
    
    # Debug: Check what content we actually have
    for i, page in enumerate(adapted_content['pages']):
        text_len = len(page.get('text', ''))
        print(f"DEBUG: Adapted page {i+1} has {text_len} characters of text")
        if text_len > 0:
            print(f"DEBUG: Page {i+1} text preview: {page.get('text', '')[:100]}")
    
    if translated_content:
        print(f"Translation complete: {len(translated_content['pages'])} pages")
        for i, page in enumerate(translated_content['pages']):
            text_len = len(page.get('text', ''))
            print(f"DEBUG: Translated page {i+1} has {text_len} characters of text")
        return adapted_content, translated_content
    return adapted_content

def create_adapted_presentation_from_pdf(adapted_content, output_path, profile):
    """Create PowerPoint presentation from adapted PDF content"""
    print(f"Creating PowerPoint from adapted PDF content")
    
    try:
        prs = Presentation()
        
        # Profile-specific colors
        profile_colors = {
            "dyslexia": RGBColor(0, 102, 204),  # Blue
            "adhd": RGBColor(46, 139, 87),      # Green
            "esl": RGBColor(148, 0, 211)        # Purple
        }
        profile_color = profile_colors.get(profile, RGBColor(0, 0, 0))
        
        # Add title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        
        title_text = adapted_content['metadata'].get('title', 'Adapted Content')
        if title_slide.shapes.title:
            title_slide.shapes.title.text = f"Adapted: {title_text}"
        
        # Add content slides
        content_layout = prs.slide_layouts[1]  # Title and Content
        
        for page in adapted_content['pages']:
            if page['text'].strip():
                slide = prs.slides.add_slide(content_layout)
                
                # Set slide title
                if slide.shapes.title:
                    slide.shapes.title.text = f"Page {page['page_number']}"
                
                # Add adapted content
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Content placeholder
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    apply_text_to_text_frame(content_placeholder.text_frame, page['text'], profile)
        
        # Save presentation
        prs.save(output_path)
        print(f"✓ Saved adapted presentation to: {output_path}")
        return True
        
    except Exception as e:
        print(f"Error creating presentation: {e}")
        return False
    
# Duplicate function removed - using the one defined earlier



def create_pptx_from_pdf_content(adapted_content, output_path, profile):
    """Create PowerPoint from adapted PDF content"""
    try:
        prs = Presentation()
        
        for page in adapted_content['pages']:
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = f"Page {page['page_number']}"
                # Apply profile formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, 
                                            slide.shapes.title.text, profile)
            
            # Add content
            if page['text']:
                # Find content placeholder
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Content placeholder
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    apply_text_to_text_frame(content_placeholder.text_frame, 
                                           page['text'], profile)
            
            # Note: Images would need additional handling to be added to PPTX
            # This is a simplified version focusing on text content
        
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"Error creating PPTX from PDF content: {e}")
        return False

def process_pdf_presentation(file_path, file_id, filename, profile, export_format='pdf'):
    """Main function to process PDF presentations"""
    try:
        # Extract content from PDF
        print(f"Extracting content from PDF: {filename}")
        pdf_content = pdf_service.extract_content_from_pdf(file_path)
        
        # Create template PDF (save original structure)
        template_filename = f"template_{filename}"
        template_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{template_filename}")
        
        # For PDFs, the template is essentially a copy of the original
        import shutil
        shutil.copy2(file_path, template_path)
        
        # Adapt content
        print(f"Adapting PDF content for {profile}")
        adapted_content = adapt_pdf_content(pdf_content, profile)
        
        # Save adapted version based on export format
        if export_format.lower() == 'pdf':
            output_filename = f"adapted_{filename}"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
            
            success = pdf_service.create_adapted_pdf(adapted_content, output_path, profile)
            if not success:
                raise Exception("Failed to create adapted PDF")
        
        elif export_format.lower() == 'pptx':
            output_filename = f"adapted_{os.path.splitext(filename)[0]}.pptx"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
            
            success = create_pptx_from_pdf_content(adapted_content, output_path, profile)
            if not success:
                raise Exception("Failed to create PowerPoint from PDF")
        
        else:
            raise Exception(f"Unsupported export format: {export_format}")
        
        # Update processing status
        processing_tasks[file_id] = {
            'status': 'complete',
            'message': 'PDF processing completed successfully',
            'template_path': template_path,
            'adapted_path': output_path,
            'export_format': export_format,
            'progress': {
                'total': 100,
                'processed': 100,
                'percentage': 100
            }
        }
        
        print(f"PDF processing complete for {file_id}")
        return output_path
        
    except Exception as e:
        print(f"Error processing PDF: {str(e)}")
        processing_tasks[file_id] = {
            'status': 'error',
            'message': str(e)
        }
        return None

@app.route('/convert_to_pdf/<file_id>')
def convert_to_pdf(file_id):
    """Convert an uploaded PPTX to PDF template"""
    try:
        if file_id not in processing_tasks:
            return jsonify({'error': 'File not found'}), 404
        
        filename = processing_tasks[file_id].get('filename', '')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found on server'}), 404
        
        # Convert to PDF
        pdf_filename = f"{os.path.splitext(filename)[0]}_template.pdf"
        pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{pdf_filename}")
        
        success = convert_pptx_to_pdf_template(file_path, pdf_path)
        
        if success:
            processing_tasks[file_id]['pdf_template'] = pdf_path
            return jsonify({
                'status': 'success',
                'pdf_path': pdf_path,
                'message': 'PDF template created successfully'
            })
        else:
            return jsonify({'error': 'Failed to convert to PDF'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/check_api')
def check_api():
    """API connection check endpoint"""
    result = api_utils.check_connection()
    
    if result["status"] == "connected":
        # Return success page
        return render_template_string(API_CHECK_SUCCESS_TEMPLATE, result=result)
    else:
        # Return error page
        return render_template_string(API_CHECK_ERROR_TEMPLATE, result=result)

# Replace the download route in app.py with this fixed version

@app.route('/download/<file_id>/<filename>')
def download(file_id, filename):
    """Universal download page for both PDF and PPTX adaptations"""
    try:
        # Get the task info
        task_info = processing_tasks.get(file_id, {})
        profile = task_info.get('profile', "dyslexia")
        export_format = task_info.get('export_format', 'pdf')
        
        # Clean the filename - remove any double 'adapted_' prefixes
        if filename.startswith('adapted_adapted_'):
            filename = filename.replace('adapted_adapted_', 'adapted_', 1)
            print(f"Fixed double prefix in download route, cleaned filename: {filename}")
        
        # Determine original file format from the filename
        original_ext = os.path.splitext(filename)[1].lower()
        original_format = 'pdf' if original_ext == '.pdf' else 'pptx'
        
        # Build expected filenames for different formats
        base_name = os.path.splitext(filename)[0]
        
        # Ensure we have exactly one 'adapted_' prefix
        if not base_name.startswith('adapted_'):
            base_name = f"adapted_{base_name}"
        
        pdf_filename = f"{base_name}.pdf"
        pptx_filename = f"{base_name}.pptx"
        
        # Check what files actually exist
        pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{pdf_filename}")
        pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{pptx_filename}")
        
        pdf_available = os.path.exists(pdf_path)
        pptx_available = os.path.exists(pptx_path)
        
        print(f"Download route - PDF available: {pdf_available}, PPTX available: {pptx_available}")
        print(f"Export format requested: {export_format}")
        
        # Check for template
        template_available = 'template_path' in task_info and os.path.exists(task_info.get('template_path', ''))
        
        # Map profile to display name (moved up for use in filesystem check)
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Translation info
        has_translation = task_info.get('has_translation', False)
        translated_filename = task_info.get('translated_filename', '')
        translated_language = task_info.get('translated_language', 'Translated')
        
        # If translation info not in task_info, check filesystem
        # Check for any profile since task_info might be empty
        if not has_translation:
            # Look for translated files in the output directory
            output_files = os.listdir(app.config['OUTPUT_FOLDER'])
            for output_file in output_files:
                if output_file.startswith(f"{file_id}_translated_"):
                    has_translation = True
                    # Extract the filename without file_id prefix
                    translated_filename = output_file[len(f"{file_id}_"):]
                    # Extract language from filename (e.g., translated_polish_...)
                    parts = translated_filename.split('_')
                    if len(parts) >= 2 and parts[0] == 'translated':
                        translated_language = parts[1].title()
                    # Since we found a translation, assume ESL profile if not set
                    if not task_info:
                        profile = 'esl'
                    print(f"Found translated file on filesystem: {output_file}")
                    print(f"Extracted language: {translated_language}")
                    break
        
        # Debug logging for translation
        # Force check filesystem for translations if none found in task
        if not has_translation:
            # More aggressive filesystem check
            all_output_files = []
            try:
                all_output_files = os.listdir(app.config['OUTPUT_FOLDER'])
                for output_file in all_output_files:
                    # Check for any translated file with this file_id
                    if output_file.startswith(f"{file_id}_") and "translated_" in output_file:
                        has_translation = True
                        translated_filename = output_file[len(f"{file_id}_"):]
                        # Extract language from filename
                        if "_translated_" in output_file:
                            lang_part = output_file.split("_translated_")[1].split("_")[0]
                            translated_language = lang_part.title()
                        print(f"✅ Found translation file: {output_file}")
                        break
            except Exception as e:
                print(f"Error checking for translations: {e}")
        
        print(f"Final translation status: has_translation={has_translation}, file={translated_filename}")
        
        # Get profile display name
        profile_name = profile_names.get(profile, profile)
        
        # Determine which filename to use for download
        if export_format == 'pdf' and pdf_available:
            download_filename = pdf_filename
        elif export_format == 'pptx' and pptx_available:
            download_filename = pptx_filename
        else:
            # Fallback to whatever is available
            if pdf_available:
                download_filename = pdf_filename
            elif pptx_available:
                download_filename = pptx_filename
            else:
                download_filename = filename
        
        # Use the universal template
        template = DOWNLOAD_TEMPLATE_UNIVERSAL
        
        return render_template_string(
            template,
            file_id=file_id,
            filename=download_filename,  # This should be the clean filename without file_id
            profile=profile,
            profile_name=profile_name,
            original_format=original_format,
            export_format=export_format,
            pdf_available=pdf_available,
            pptx_available=pptx_available,
            pdf_filename=pdf_filename,    # Add these for template use
            pptx_filename=pptx_filename,  # Add these for template use
            template_available=template_available,
            has_translation=has_translation,
            translated_filename=translated_filename,
            translated_language=translated_language,
            processing_time=None,
            file_count=1 + (1 if has_translation else 0)
        )
    
    except Exception as e:
        print(f"Error in download route: {e}")
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, 
                                     message=f"Error loading download page: {str(e)}"), 500

# Also add the missing download_file route if it doesn't exist
@app.route('/download_file/<file_id>/<filename>')
def download_file(file_id, filename):
    """Actually download the file with proper filename handling"""
    try:
        print(f"Download request: file_id={file_id}, filename={filename}")
        
        # Get task info for additional context
        task_info = processing_tasks.get(file_id, {})
        print(f"Task info for {file_id}: {task_info}")
        
        # Clean the filename - remove any double 'adapted_' prefixes
        if filename.startswith('adapted_adapted_'):
            filename = filename.replace('adapted_adapted_', 'adapted_', 1)
            print(f"Fixed double prefix, cleaned filename: {filename}")
        
        # Ensure filename has exactly one 'adapted_' prefix for lookup
        base_name = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1]
        
        if not base_name.startswith('adapted_'):
            lookup_filename = f"adapted_{filename}"
        else:
            lookup_filename = filename
        
        # Construct the file path with file_id prefix
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{lookup_filename}")
        print(f"Looking for file at: {file_path}")
        
        # Check if file exists
        if os.path.exists(file_path):
            # Use the clean filename (without file_id) for download
            return send_file(file_path, as_attachment=True, download_name=lookup_filename)
        
        # If not found with standard naming, try alternatives
        if True:  # This maintains the original indentation
            print(f"File not found at primary path: {file_path}")
            
            # Try without the file_id prefix
            alt_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
            print(f"Trying alternative path: {alt_path}")
            
            if os.path.exists(alt_path):
                file_path = alt_path
                print(f"Found file at alternative path: {alt_path}")
            else:
                # Check if the task has a specific adapted_path
                adapted_path = task_info.get('adapted_path')
                if adapted_path and os.path.exists(adapted_path):
                    file_path = adapted_path
                    print(f"Found file at adapted_path: {adapted_path}")
                else:
                    # List all files in output folder for debugging
                    output_files = os.listdir(app.config['OUTPUT_FOLDER'])
                    matching_files = [f for f in output_files if file_id in f or filename.replace('adapted_', '') in f]
                    
                    print(f"Files in output folder containing file_id or base filename: {matching_files}")
                    
                    # Try to find a matching file
                    if matching_files:
                        # Use the first matching file
                        found_file = matching_files[0]
                        file_path = os.path.join(app.config['OUTPUT_FOLDER'], found_file)
                        print(f"Using closest match: {found_file}")
                    else:
                        error_msg = f"File not found: {filename}. Checked paths: {file_path}, {alt_path}"
                        if adapted_path:
                            error_msg += f", {adapted_path}"
                        return render_template_string(ERROR_TEMPLATE, message=error_msg), 404
        
        print(f"Sending file: {file_path}")
        # Send the file with clean filename
        clean_filename = os.path.basename(file_path).replace(f"{file_id}_", "", 1)
        return send_file(file_path, as_attachment=True, download_name=clean_filename)
        
    except Exception as e:
        print(f"Error downloading file: {e}")
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE,
            message=f"Error downloading file: {str(e)}"), 500

@app.route('/debug/files/<file_id>')
def debug_files(file_id):
    """Debug route to check file status and locations"""
    try:
        # Get task info
        task_info = processing_tasks.get(file_id, {})
        
        # List all files in output folder
        output_files = []
        try:
            output_files = os.listdir(app.config['OUTPUT_FOLDER'])
        except:
            output_files = ["Error reading output folder"]
        
        # Find files matching this file_id
        matching_files = [f for f in output_files if file_id in f]
        
        # Debug info
        debug_info = {
            'file_id': file_id,
            'task_info': task_info,
            'output_folder': app.config['OUTPUT_FOLDER'],
            'all_output_files_count': len(output_files),
            'matching_files': matching_files,
            'processing_tasks_keys': list(processing_tasks.keys())
        }
        
        return jsonify(debug_info)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/analyze/framework/<file_id>')
def analyze_framework(file_id):
    """Analyze a presentation's instructional framework"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing_framework'
        processing_tasks[file_id]['message'] = 'Analyzing instructional framework...'
        
        # Analyze framework
        framework_data = analyze_instructional_framework(file_path)
        
        # Store the framework data
        processing_tasks[file_id]['framework'] = framework_data
        processing_tasks[file_id]['status'] = 'complete'
        
        # Create a template to display the results
        return render_template_string(
            FRAMEWORK_TEMPLATE,
            framework=framework_data,
            file_id=file_id
        )
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing framework: {str(e)}"
        }), 500

def extract_learning_scaffolding_with_llm(pptx_path):
    """Use Claude to identify learning scaffolding elements in a presentation"""
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Extract all slide content with metadata
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = ""
        title = ""
        
        # Get slide title
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title.text
            
        # Extract all text from the slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text:
                slide_content += shape.text + "\n"
        
        slides_data.append({
            "slide_number": i+1,
            "title": title,
            "content": slide_content
        })
    
    # Prepare prompt for Claude
    slide_descriptions = "\n\n".join([
        f"SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
        for s in slides_data
    ])
    
    prompt = f"""
    Analyze the following PowerPoint presentation slides and identify the learning scaffolding elements.
    
    SLIDES:
    {slide_descriptions}
    
    For each slide, determine its instructional purpose (e.g., objectives, key concept, example, practice, assessment, review).
    
    Then extract the following scaffolding elements from the entire presentation:
    1. Learning Objectives
    2. Key Concepts and their definitions
    3. Examples
    4. Practice Activities
    5. Assessment Items
    6. Review Elements
    
    Format your response as JSON with the following structure:
    {{
      "slides": [
        {{
          "slide_number": 1,
          "purpose": "title",
          "elements": []
        }},
        // other slides...
      ],
      "scaffolding": {{
        "learning_objectives": ["objective 1", "objective 2"],
        "key_concepts": [
          {{ "term": "concept 1", "definition": "definition 1" }}
        ],
        "examples": [
          {{ "slide_number": 5, "content": "example content" }}
        ],
        "practice_activities": [],
        "assessment_items": [],
        "review_elements": []
      }}
    }}
    """
    
    try:
        # Call Claude API
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        
        # Extract and parse JSON response
        content = response.content[0].text.strip()
        json_match = re.search(r'(\{.*\})', content, re.DOTALL)
        
        if json_match:
            try:
                scaffolding_data = json.loads(json_match.group(0))
                return scaffolding_data
            except json.JSONDecodeError:
                print("Error parsing JSON from LLM response")
        
        # Fallback to simplified extraction if JSON parsing fails
        return {
            "slides": [{"slide_number": s["slide_number"], "purpose": "content"} for s in slides_data],
            "scaffolding": {
                "learning_objectives": [],
                "key_concepts": [],
                "examples": [],
                "practice_activities": [],
                "assessment_items": [],
                "review_elements": []
            }
        }
    except Exception as e:
        print(f"Error in scaffolding extraction: {str(e)}")
        # Return empty scaffolding structure as fallback
        return {
            "slides": [],
            "scaffolding": {
                "learning_objectives": [],
                "key_concepts": [],
                "examples": [],
                "practice_activities": [],
                "assessment_items": [],
                "review_elements": []
            }
        }

@app.route('/analyze/scaffolding/<file_id>')
def analyze_scaffolding(file_id):
    """Analyze a presentation for learning scaffolding elements"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        # Debug output
        print(f"Looking for file at: {file_path}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Check file type
        file_type = processing_tasks[file_id].get('file_type', '').lower()
        
        if file_type == '.pdf':
            # PDF files cannot be analyzed for scaffolding with current implementation
            return jsonify({
                "status": "error",
                "message": "Scaffolding analysis is currently only available for PowerPoint presentations (.pptx files). Please upload a PowerPoint file or use the adaptation feature for PDFs."
            }), 400
        
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing'
        processing_tasks[file_id]['message'] = 'Extracting learning scaffolding elements...'
        
        # Extract scaffolding elements
        scaffolding_data = extract_learning_scaffolding_with_llm(file_path)
        
        # Store the scaffolding data in the processing task
        processing_tasks[file_id]['scaffolding'] = scaffolding_data
        processing_tasks[file_id]['status'] = 'complete'
        processing_tasks[file_id]['message'] = 'Scaffolding analysis complete'
        
        # Redirect to the scaffolding results page
        return redirect(url_for('view_presentation_scaffolding', file_id=file_id))        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        if file_id in processing_tasks:
            processing_tasks[file_id]['status'] = 'error'
            processing_tasks[file_id]['message'] = f"Error analyzing scaffolding: {str(e)}"
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing scaffolding: {str(e)}"
        }), 500

@app.route('/view_scaffolding/<file_id>')
def view_presentation_scaffolding(file_id):
    """Show the scaffolding analysis results"""
    if file_id not in processing_tasks:
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Presentation not found. Please upload again.")
                                     
    if 'scaffolding' not in processing_tasks[file_id]:
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Scaffolding analysis not found. Please analyze the presentation first.")
    
    scaffolding_data = processing_tasks[file_id]['scaffolding']
    
    return render_template_string(
        SCAFFOLDING_TEMPLATE,
        scaffolding=scaffolding_data,
        slides=scaffolding_data.get('slides', [])
    )

def extract_content_from_pptx(file_path):
    """Extract all text content from a PowerPoint file"""
    try:
        # Load the presentation
        prs = Presentation(file_path)
        
        # Extract text from all slides
        all_content = ""
        
        for slide_index, slide in enumerate(prs.slides):
            # Add slide number
            all_content += f"Slide {slide_index + 1}:\n"
            
            # Process every shape that might contain text
            for shape in slide.shapes:
                # Process text content in text frames
                if hasattr(shape, "text_frame") and shape.text:
                    all_content += shape.text + "\n\n"
                
                # Check for tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_content = []
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text:
                                row_content.append(cell.text)
                        all_content += " | ".join(row_content) + "\n"
                    all_content += "\n"
                
                # Group shapes
                if hasattr(shape, "shapes"):
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text_frame") and subshape.text:
                            all_content += subshape.text + "\n\n"
            
            all_content += "\n---\n\n"
        
        return all_content
        
    except Exception as e:
        print(f"Error extracting content from PPTX: {str(e)}")
        return "Error extracting content from presentation"

def update_task_status(file_id, status, message=None, progress=None, **kwargs):
    """Helper function to update task status consistently"""
    if file_id not in processing_tasks:
        processing_tasks[file_id] = {}
    
    processing_tasks[file_id]['status'] = status
    
    if message is not None:
        processing_tasks[file_id]['message'] = message
    
    if progress is not None:
        processing_tasks[file_id]['progress'] = progress
    
    # Add any additional fields
    for key, value in kwargs.items():
        processing_tasks[file_id][key] = value
    
    # Debug logging
    print(f"Status update for {file_id}: {status} - {message}")

def update_task_progress(file_id, processed, total, message=None):
    """Helper function to update task progress"""
    percentage = min(100, int((processed / max(1, total)) * 100))
    
    progress = {
        'total': total,
        'processed': processed,
        'percentage': percentage
    }
    
    if message:
        update_task_status(file_id, 'processing', message, progress)
    else:
        if file_id in processing_tasks:
            processing_tasks[file_id]['progress'] = progress


def process_structure_analysis(file_id, file_path, profile):
    """Process the structural analysis in a background thread"""
    try:
        # Store analysis status
        analysis_status = {
            'status': 'processing',
            'scaffolding': None,
            'framework': None,
            'error': None
        }
        
        # Store in global dictionary or database
        global_analysis_status[file_id] = analysis_status
        
        # Extract learning scaffolding with LLM
        scaffolding_analysis = extract_learning_scaffolding_with_llm(file_path)
        analysis_status['scaffolding'] = scaffolding_analysis
        
        # Identify instructional framework
        framework_analysis = analyze_instructional_framework(file_path)
        analysis_status['framework'] = framework_analysis
        
        # Mark as complete
        analysis_status['status'] = 'complete'
        
    except Exception as e:
        print(f"Error in structure analysis: {str(e)}")
        analysis_status['status'] = 'error'
        analysis_status['error'] = str(e)

## Text Processing Functions
def process_text_batch(texts, profile, max_batch_size=5, max_tokens_per_batch=4000):
    """Process multiple text elements in efficient batches"""
    global api_call_counter
    
    # Get current file_id from thread name if possible
    file_id = None
    current_thread = threading.current_thread()
    for fid in processing_tasks.keys():
        if fid in current_thread.name:
            file_id = fid
            break
    
    results = []
    current_batch = []
    current_batch_tokens = 0
    
    # Pre-calculate estimated number of batches
    total_batches = 0
    temp_batch = []
    temp_tokens = 0
    
    def estimate_tokens(text):
        """Roughly estimate token count based on character count"""
        return len(text) // 4  # Approximation: ~4 chars per token
    
    # Estimate number of batches
    for text in texts:
        text_tokens = estimate_tokens(text)
        if len(temp_batch) >= max_batch_size or (temp_tokens + text_tokens > max_tokens_per_batch and temp_batch):
            total_batches += 1
            temp_batch = []
            temp_tokens = 0
        temp_batch.append(text)
        temp_tokens += text_tokens
    
    # Add the last batch if non-empty
    if temp_batch:
        total_batches += 1
    
    # Initialize counter for this file_id
    if file_id:
        api_call_counter[file_id] = {
            'total': total_batches,
            'current': 0
        }
        # Set initial progress in processing_tasks
        if file_id in processing_tasks:
            processing_tasks[file_id]['progress'] = {
                'total': total_batches,
                'processed': 0,
                'percentage': 0
            }
            processing_tasks[file_id]['message'] = f"Starting to process {total_batches} batches"
            print(f"Estimated {total_batches} batches to process for {len(texts)} text elements")
    
    # Group texts into efficient batches
    for text in texts:
        text_tokens = estimate_tokens(text)
        
        # If this text would make the batch too large, process the current batch first
        if len(current_batch) >= max_batch_size or (current_batch_tokens + text_tokens > max_tokens_per_batch and current_batch):
            # Process the current batch
            batch_results = process_single_batch(current_batch, profile)
            results.extend(batch_results)
            current_batch = []
            current_batch_tokens = 0
            
            # Small delay between batches to avoid rate limits
            time.sleep(0.5)
        
        # Add this text to the current batch
        current_batch.append(text)
        current_batch_tokens += text_tokens
    
    # Process any remaining texts in the final batch
    if current_batch:
        batch_results = process_single_batch(current_batch, profile)
        results.extend(batch_results)
    
    # Mark as 100% complete when done
    if file_id and file_id in processing_tasks:
        processing_tasks[file_id]['progress'] = {
            'total': total_batches,
            'processed': total_batches,
            'percentage': 100
        }
        processing_tasks[file_id]['message'] = f"Completed processing {total_batches} batches"
    
    return results

# Helper function to get the current file_id from the thread context
def get_current_file_id():
    """Get the file_id for the current processing thread"""
    # This is a placeholder - you'll need to implement a way to track which file_id 
    # is being processed in the current thread
    current_thread = threading.current_thread()
    thread_name = current_thread.name
    
    # Check if thread name contains a file_id (if you named your threads with file_ids)
    for file_id in processing_tasks.keys():
        if file_id in thread_name:
            return file_id
    
    # If we can't determine, return None
    return None

@app.route('/status/<file_id>')
def status(file_id):
    """Check the status of a processing task"""
    try:
        if file_id not in processing_tasks:
            return jsonify({
                'status': 'not_found',
                'message': 'Task not found'
            }), 404
        
        task = processing_tasks[file_id]
        
        # Return the task status as JSON
        response_data = {
            'status': task.get('status', 'unknown'),
            'message': task.get('message', ''),
            'progress': task.get('progress', {
                'total': 0,
                'processed': 0,
                'percentage': 0
            })
        }
        
        # Add additional fields if they exist
        if 'error' in task:
            response_data['error'] = task['error']
        
        if 'has_translation' in task:
            response_data['has_translation'] = task['has_translation']
            
        if 'export_format' in task:
            response_data['export_format'] = task['export_format']
        
        # Debug logging
        print(f"Status check for {file_id}: {response_data}")
        
        return jsonify(response_data)
        
    except Exception as e:
        print(f"Error in status route: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': f'Error checking status: {str(e)}'
        }), 500

# Helper function to process a single batch of texts
def process_single_batch(texts, profile):
    """Process a single batch of texts using one API call"""
    global api_call_counter
    
    # Get current file_id from thread name if possible
    file_id = None
    current_thread = threading.current_thread()
    for fid in processing_tasks.keys():
        if fid in current_thread.name:
            file_id = fid
            break
    
    # For very small batches, process individually
    if len(texts) <= 1:
        # Update counter if we have file_id
        if file_id and file_id in api_call_counter:
            api_call_counter[file_id]['current'] += 1
            # Update processing task with progress info
            if file_id in processing_tasks:
                progress = api_call_counter[file_id]
                percentage = min(95, int((progress['current'] / max(1, progress['total'])) * 100))
                processing_tasks[file_id]['progress'] = {
                    'total': progress['total'],
                    'processed': progress['current'],
                    'percentage': percentage
                }
                processing_tasks[file_id]['message'] = f"Processing: {progress['current']}/{progress['total']} batches ({percentage}%)"
        
        return [adapt_text_with_matcha(texts[0], profile)]
    
    # Create a single prompt with multiple texts
    combined_prompt = f"Adapt the following texts for {profile} users. Format your response using exactly '### TEXT N ###' before each adapted text (where N is the text number).\n\n"
    
    for i, text in enumerate(texts):
        combined_prompt += f"### TEXT {i+1} ###\n{text}\n\n"
    
    try:
        # Call API with the combined prompt
        response = call_claude_api_with_retry(combined_prompt)
        
        # Update counter if we have file_id
        if file_id and file_id in api_call_counter:
            api_call_counter[file_id]['current'] += 1
            # Update processing task with progress info
            if file_id in processing_tasks:
                progress = api_call_counter[file_id]
                percentage = min(95, int((progress['current'] / max(1, progress['total'])) * 100))
                processing_tasks[file_id]['progress'] = {
                    'total': progress['total'],
                    'processed': progress['current'],
                    'percentage': percentage
                }
                processing_tasks[file_id]['message'] = f"Processing: {progress['current']}/{progress['total']} batches ({percentage}%)"
                print(f"Progress update: {percentage}% ({progress['current']}/{progress['total']} batches)")
        
        # Parse the response to extract individual adapted texts
        adapted_texts = []
        content = response.content[0].text
        
        # Better parsing approach
        pattern = r'###\s*TEXT\s*(\d+)\s*###\s*(.*?)(?=###\s*TEXT\s*\d+\s*###|$)'
        matches = re.findall(pattern, content, re.DOTALL)
        
        # Create a dictionary to maintain order
        text_dict = {}
        for match in matches:
            text_num = int(match[0])
            if 1 <= text_num <= len(texts):  # Ensure valid text number
                text_dict[text_num] = match[1].strip()
        
        # If we don't have matches for all texts, fall back to individual processing
        if len(text_dict) != len(texts):
            print(f"Warning: Batch processing response parsing failed. Falling back to individual processing.")
            return [adapt_text_with_matcha(text, profile) for text in texts]
            
        # Convert dictionary to ordered list
        for i in range(1, len(texts) + 1):
            if i in text_dict:
                adapted_texts.append(text_dict[i])
            else:
                # If missing a specific text, adapt it individually
                adapted_texts.append(adapt_text_with_matcha(texts[i-1], profile))
        
        return adapted_texts
        
    except Exception as e:
        print(f"Error in batch processing: {str(e)}")
        # Fall back to individual processing
        return [adapt_text_with_matcha(text, profile) for text in texts]

def initialize_pdf_capabilities():
    """Properly initialize PDF processing capabilities"""
    capabilities = {
        'pymupdf': False,
        'pdfplumber': False,
        'reportlab': False,
        'pypdf2': False
    }
    
    # Test PyMuPDF
    try:
        import fitz
        capabilities['pymupdf'] = True
        print("✓ PyMuPDF (fitz) available")
    except ImportError:
        print("✗ PyMuPDF not available")
    
    # Test pdfplumber
    try:
        import pdfplumber
        capabilities['pdfplumber'] = True
        print("✓ pdfplumber available")
    except ImportError:
        print("✗ pdfplumber not available")
    
    # Test reportlab
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph
        capabilities['reportlab'] = True
        print("✓ reportlab available")
    except ImportError:
        print("✗ reportlab not available")
    
    # Test PyPDF2
    try:
        import PyPDF2
        capabilities['pypdf2'] = True
        print("✓ PyPDF2 available")
    except ImportError:
        print("✗ PyPDF2 not available")
    
    return capabilities

# Initialize PDF capabilities at startup
pdf_capabilities = initialize_pdf_capabilities()

def needs_adaptation(text, profile):
    """Determine if text needs API-based adaptation"""
    # Skip very short text or placeholder text
    if len(text.strip()) < 15 or "Click to add text" in text:
        print(f"DEBUG: Skipping text - too short or placeholder: '{text[:50]}'")
        return False
    
    # Calculate metrics
    metrics = calculate_simple_readability(text)
    print(f"DEBUG: Text metrics - Flesch: {metrics['flesch_reading_ease']:.1f}, Grade: {metrics['flesch_kincaid_grade']:.1f}, Sentence length: {metrics['sentence_length']:.1f}")
    
    # Check against thresholds for the profile
    threshold = READABILITY_THRESHOLDS.get(profile, {})
    print(f"DEBUG: Thresholds for {profile} - Flesch: {threshold.get('flesch_reading_ease', 70)}, Grade: {threshold.get('flesch_kincaid_grade', 8)}, Sentence: {threshold.get('sentence_length', 15)}")
    
    # For debugging: temporarily be more aggressive about adaptation
    # Reduce the thresholds to make adaptation more likely
    flesch_threshold = threshold.get('flesch_reading_ease', 70) - 20  # Lower threshold
    grade_threshold = threshold.get('flesch_kincaid_grade', 8) + 3    # Higher threshold
    sentence_threshold = threshold.get('sentence_length', 15) + 5    # Higher threshold
    
    # If already meeting key thresholds, no need to adapt
    if (metrics['flesch_reading_ease'] >= flesch_threshold and
        metrics['flesch_kincaid_grade'] <= grade_threshold and
        metrics['sentence_length'] <= sentence_threshold):
        print(f"DEBUG: Text meets adjusted thresholds - skipping adaptation")
        return False
    
    print(f"DEBUG: Text needs adaptation")
    return True

def calculate_simple_readability(text):
    """Calculate simplified readability metrics to determine if adaptation is needed"""
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = max(1, len(sentences))  # Avoid division by zero
    
    # Calculate average sentence length
    avg_sentence_length = word_count / sentence_count
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter
        vowels = "aeiouy"
        word = word.lower()
        count = 0
        if word and word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word and word.endswith("e"):
            count -= 1
        return max(1, count)  # Always at least 1 syllable
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = (len(complex_words) / max(1, word_count)) * 100
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / max(1, word_count))
        flesch = max(0, min(100, flesch))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / max(1, word_count)) - 15.59
        fk_grade = max(0, fk_grade)
    else:
        fk_grade = 0
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

def process_presentation_efficiently(file_path, file_id, filename, profile, target_language=None):
    print(f"DEBUG: Starting process_presentation_efficiently with profile: {profile}")
    try:
        # Normalize profile to prevent case issues
        profile = profile.lower() if profile else "dyslexia"

        # Initialize processing task status
        if file_id not in processing_tasks:
            processing_tasks[file_id] = {'status': 'processing', 'message': 'Initializing processing...'}

        # Check API connection first
        processing_tasks[file_id]['message'] = 'Checking API connection...'
        api_status = api_utils.check_connection()

        if api_status["status"] != "connected":
            # API connection failed - update status and exit
            processing_tasks[file_id] = {
                'status': 'error', 
                'message': f"Claude API connection error: {api_status['message']}"
            }
            return None
        # Load the presentation
        prs = Presentation(file_path)

        # Create a copy for the translated version if a target language is specified
        translated_prs = Presentation(file_path) if target_language else None
        # Collect text elements that need processing
        text_elements = []
        element_references = []

        print(f"Analyzing presentation with {len(prs.slides)} slides...")

        # First pass: collect all text
        for slide_index, slide in enumerate(prs.slides):
            slide_elements = collect_text_elements_from_slide(slide, slide_index)

            for element in slide_elements:
                text = element['text']
                # Only add if it needs adaptation
                if needs_adaptation(text, profile):
                    text_elements.append(text)
                    element_references.append(element['reference'])

        print(f"Found {len(text_elements)} text elements requiring adaptation")

        # Process collected text in batches
        if text_elements:
            print(f"Processing {len(text_elements)} text elements in batches")
            adapted_texts = process_text_batch(text_elements, profile)

            # Create translations if target language is specified
            translated_texts = []
            if target_language:
                print(f"Translating {len(text_elements)} text elements to {target_language}")
                for original_text in text_elements:
                    translated_text = translate_text(original_text, target_language)
                    translated_texts.append(translated_text)

            # Apply adapted texts back to their elements in the English version
            for i, adapted_text in enumerate(adapted_texts):
                reference = element_references[i]
                apply_adapted_text(reference, adapted_text, profile)

                # Apply translations to the translated presentation
                if translated_prs and target_language and i < len(translated_texts):
                    slide_num = reference.get('slide_num', 0)
                    shape_index = reference.get('shape_index', 0)

                    if slide_num < len(translated_prs.slides):
                        trans_slide = translated_prs.slides[slide_num]
                        shapes = list(trans_slide.shapes)
                        if shape_index < len(shapes):
                            trans_shape = shapes[shape_index]

                            if reference['type'] == 'shape' and hasattr(trans_shape, "text_frame"):
                                apply_text_to_text_frame(trans_shape.text_frame, translated_texts[i], profile)
                            elif reference['type'] == 'cell':
                                # Handle table cells
                                row_index = reference.get('row_index', 0)
                                cell_index = reference.get('cell_index', 0)
                                if hasattr(trans_shape, "table") and row_index < len(trans_shape.table.rows):
                                    row = trans_shape.table.rows[row_index]
                                    if cell_index < len(row.cells):
                                        cell = row.cells[cell_index]
                                        if hasattr(cell, "text_frame"):
                                            apply_text_to_text_frame(cell.text_frame, translated_texts[i], profile)
                            elif reference['type'] == 'subshape':
                                # Handle grouped shapes
                                subshape_index = reference.get('subshape_index', 0)
                                if hasattr(trans_shape, "shapes") and subshape_index < len(trans_shape.shapes):
                                    subshape = trans_shape.shapes[subshape_index]
                                    if hasattr(subshape, "text_frame"):
                                        apply_text_to_text_frame(subshape.text_frame, translated_texts[i], profile)

             # Save the adapted presentation
        output_filename = f"adapted_{filename}"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
        prs.save(output_path)

        # Save the translated version if created
        translated_path = None
        translated_filename = ""
        if translated_prs and target_language:
            # Code for saving the translated presentation
            translated_filename = f"translated_{target_language}_{filename}"
            translated_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{translated_filename}")
            translated_prs.save(translated_path)

        # Update processing status to complete
        processing_tasks[file_id] = {
            'status': 'complete',
            'message': 'Adaptation completed successfully',
            'profile': profile,
            'has_translation': bool(translated_path),
            'translated_filename': translated_filename,
            'translated_language': target_language.title() if target_language else '',
            'target_language': target_language if target_language else '',
            'progress': {
                'total': 100,
                'processed': 100,
                'percentage': 100
            }
        }

        print(f"Processing complete for {file_id}. Status set to complete.")
        return output_path

    except Exception as e:
        print(f"Error during processing: {str(e)}")
        # Update task status on error
        processing_tasks[file_id] = {
            'status': 'error', 
            'message': str(e)
        }
# Helper function to collect text elements from a slide
def collect_text_elements_from_slide(slide, slide_index=0):
    """Collect text elements from a slide and their references"""
    elements = []
    
    # Process shapes with text
    for shape_index, shape in enumerate(slide.shapes):
        if hasattr(shape, "text_frame") and shape.text:
            elements.append({
                'text': shape.text,
                'reference': {
                    'type': 'shape', 
                    'object': shape,
                    'slide_num': slide_index,
                    'shape_index': shape_index
                }
            })
        
        # Tables
        if hasattr(shape, "table"):
            for row_index, row in enumerate(shape.table.rows):
                for cell_index, cell in enumerate(row.cells):
                    if hasattr(cell, "text_frame") and cell.text:
                        elements.append({
                            'text': cell.text,
                            'reference': {
                                'type': 'cell', 
                                'object': cell,
                                'slide_num': slide_index,
                                'shape_index': shape_index,
                                'row_index': row_index,
                                'cell_index': cell_index
                            }
                        })
        
        # Group shapes
        if hasattr(shape, "shapes"):
            for subshape_index, subshape in enumerate(shape.shapes):
                if hasattr(subshape, "text_frame") and subshape.text:
                    elements.append({
                        'text': subshape.text,
                        'reference': {
                            'type': 'subshape', 
                            'object': subshape,
                            'slide_num': slide_index,
                            'shape_index': shape_index,
                            'subshape_index': subshape_index
                        }
                    })
    
    return elements

def apply_dyslexia_formatting(text_frame):
    """Apply special formatting for dyslexia support"""
    # Only apply this formatting if we're actually handling dyslexia content
    # This is a safeguard in case this function is called inappropriately
    try:
        # Adjust text alignment to left
        for paragraph in text_frame.paragraphs:
            if hasattr(paragraph, "alignment"):
                paragraph.alignment = PP_ALIGN.LEFT
                
            # Increase line spacing if not already set
            # Check if line_spacing is None before comparing
            if hasattr(paragraph, "line_spacing") and paragraph.line_spacing is not None and paragraph.line_spacing < 1.2:
                paragraph.line_spacing = 1.5
            elif hasattr(paragraph, "line_spacing") and paragraph.line_spacing is None:
                # If it's None, just set it without comparison
                paragraph.line_spacing = 1.5
                
            # Apply sans-serif font to all runs if not already a sans-serif font
            for run in paragraph.runs:
                if hasattr(run, "font") and hasattr(run.font, "name"):
                    current_font = run.font.name.lower() if hasattr(run.font.name, "lower") else ""
                    # Only change if not already using a dyslexia-friendly font
                    if not any(font in current_font for font in ["arial", "verdana", "tahoma", "calibri", "helvetica", "open sans"]):
                        run.font.name = "Arial"
                    
                    # Increase font size slightly for better readability if it's too small
                    if hasattr(run.font, "size") and hasattr(run.font.size, "pt"):
                        current_size = run.font.size.pt
                        # Only increase if not already large
                        if current_size < 12:
                            run.font.size = Pt(max(12, current_size * 1.2))
    except Exception as e:
        # Log error but continue processing
        print(f"Error applying dyslexia formatting: {e}")
        pass

# Helper function to apply adapted text back to elements
def apply_adapted_text(reference, adapted_text, profile):
    """Apply adapted text to the appropriate element with formatting"""
    element_type = reference['type']
    obj = reference['object']
    
    if element_type == 'shape' or element_type == 'subshape' or element_type == 'cell':
        if hasattr(obj, "text_frame"):
            apply_text_to_text_frame(obj.text_frame, adapted_text, profile)
            
            # For dyslexia, apply additional formatting
            if profile == "dyslexia" and hasattr(obj, "text_frame"):
                apply_dyslexia_formatting(obj.text_frame)

# Function to apply text to a text frame with color formatting
def apply_text_to_text_frame(text_frame, adapted_text, profile):
    """Apply adapted text to a text frame with proper formatting"""
    # Skip if text frame is empty
    if not hasattr(text_frame, "paragraphs"):
        return
    
    # Get the profile color - only apply if it matches the passed profile
    actual_profile = profile.lower() if profile else "default"
    profile_color = PROFILE_COLORS.get(actual_profile)
    
    try:
        # Store original formatting settings we want to preserve
        original_formatting = []
        for p in text_frame.paragraphs:
            p_format = {
                "level": p.level if hasattr(p, "level") else 0,
                "alignment": p.alignment if hasattr(p, "alignment") else None,
                "font_name": None,
                "font_size": None,
                "font_bold": None,
                "font_italic": None
            }
            
            # Get formatting from first run if available
            if p.runs and hasattr(p.runs[0], "font"):
                font = p.runs[0].font
                p_format["font_name"] = font.name if hasattr(font, "name") else None
                p_format["font_size"] = font.size if hasattr(font, "size") else None
                p_format["font_bold"] = font.bold if hasattr(font, "bold") else None
                p_format["font_italic"] = font.italic if hasattr(font, "italic") else None
            
            original_formatting.append(p_format)
        
        # Clear all paragraphs except the first one
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            tr = p._p
            tr.getparent().remove(tr)
        
        # Clear the first paragraph content but preserve its format
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            while len(p.runs) > 0:
                try:
                    p._p.remove(p.runs[0]._r)
                except:
                    break
        else:
            p = text_frame.add_paragraph()
        
        # Split the adapted text into paragraphs
        text_paragraphs = adapted_text.split('\n')
        
        # Format the first paragraph with colored first word
        if text_paragraphs and text_paragraphs[0].strip():
            words = text_paragraphs[0].split()
            if words:
                # First word with color
                first_word = words[0]
                first_run = p.add_run()
                first_run.text = first_word
                # Apply profile color to first word
                if profile_color and hasattr(first_run, "font") and hasattr(first_run.font, "color"):
                    first_run.font.color.rgb = profile_color
                
                # Apply original formatting to first run
                if original_formatting and hasattr(first_run, "font"):
                    font = first_run.font
                    if original_formatting[0]["font_name"]:
                        font.name = original_formatting[0]["font_name"]
                    if original_formatting[0]["font_size"]:
                        font.size = original_formatting[0]["font_size"]
                    if original_formatting[0]["font_bold"]:
                        font.bold = original_formatting[0]["font_bold"]
                    if original_formatting[0]["font_italic"]:
                        font.italic = original_formatting[0]["font_italic"]
                
                # Rest of first paragraph
                rest_of_text = ' '.join(words[1:]) if len(words) > 1 else ''
                if rest_of_text:
                    rest_run = p.add_run()
                    rest_run.text = ' ' + rest_of_text
                    # Apply original formatting to rest run
                    if original_formatting and hasattr(rest_run, "font"):
                        font = rest_run.font
                        if original_formatting[0]["font_name"]:
                            font.name = original_formatting[0]["font_name"]
                        if original_formatting[0]["font_size"]:
                            font.size = original_formatting[0]["font_size"]
                        if original_formatting[0]["font_bold"]:
                            font.bold = original_formatting[0]["font_bold"]
                        if original_formatting[0]["font_italic"]:
                            font.italic = original_formatting[0]["font_italic"]
            
            # Apply original paragraph formatting
            if original_formatting:
                if original_formatting[0]["alignment"]:
                    p.alignment = original_formatting[0]["alignment"]
                if original_formatting[0]["level"]:
                    p.level = original_formatting[0]["level"]
        
        # Add remaining paragraphs
        for i, p_text in enumerate(text_paragraphs[1:], 1):
            if p_text.strip():
                new_p = text_frame.add_paragraph()
                run = new_p.add_run()
                run.text = p_text
                
                # Apply formatting if available
                fmt_index = min(i, len(original_formatting) - 1) if original_formatting else -1
                
                if fmt_index >= 0:
                    # Apply paragraph formatting
                    if original_formatting[fmt_index]["alignment"]:
                        new_p.alignment = original_formatting[fmt_index]["alignment"]
                    if original_formatting[fmt_index]["level"]:
                        new_p.level = original_formatting[fmt_index]["level"]
                    
                    # Apply run formatting
                    if hasattr(run, "font"):
                        font = run.font
                        if original_formatting[fmt_index]["font_name"]:
                            font.name = original_formatting[fmt_index]["font_name"]
                        if original_formatting[fmt_index]["font_size"]:
                            font.size = original_formatting[fmt_index]["font_size"]
                        if original_formatting[fmt_index]["font_bold"]:
                            font.bold = original_formatting[fmt_index]["font_bold"]
                        if original_formatting[fmt_index]["font_italic"]:
                            font.italic = original_formatting[fmt_index]["font_italic"]
    
    except Exception as e:
        print(f"Error formatting text frame: {str(e)}")
        # Fallback: just set the text directly if there's an error
        try:
            text_frame.text = adapted_text
        except:
            pass


@app.route('/error')
def error():
    """Show error page"""
    message = request.args.get('message', 'An unknown error occurred')
    return render_template_string(ERROR_TEMPLATE, message=message)

def analyze_pptx(file_path):
    """Analyze PowerPoint content for readability"""
    # Load the presentation
    prs = Presentation(file_path)
    
    # Extract all text from the presentation
    all_text = ""
    slide_texts = []  # Text per slide for analysis
    
    for slide_index, slide in enumerate(prs.slides):
        slide_text = ""
        
        # Process every shape that might contain text
        for shape in slide.shapes:
            # Process text content in text frames
            if hasattr(shape, "text_frame") and shape.text:
                all_text += shape.text + "\n\n"
                slide_text += shape.text + "\n\n"
            
            # Check for tables
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, "text_frame") and cell.text:
                            all_text += cell.text + "\n"
                            slide_text += cell.text + "\n"
            
            # Group shapes
            if hasattr(shape, "shapes"):
                for subshape in shape.shapes:
                    if hasattr(subshape, "text_frame") and subshape.text:
                        all_text += subshape.text + "\n\n"
                        slide_text += subshape.text + "\n\n"
        
        slide_texts.append({
            "slide_number": slide_index + 1,
            "text": slide_text
        })
    
    # Calculate readability metrics
    metrics = calculate_readability_metrics(all_text)
    
    # Get complex words
    complex_words = find_complex_words(all_text)
    
    return all_text, slide_texts, metrics, complex_words

def calculate_readability_metrics(text):
    """Calculate readability metrics for the text"""
    # Simplified calculation of readability metrics
    # In a production environment, use a proper textstat library
    
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = len(sentences)
    
    # Calculate average sentence length
    avg_sentence_length = round(word_count / max(1, sentence_count), 1)
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter - not perfect but adequate for demo
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = round((len(complex_words) / max(1, word_count)) * 100, 1)
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / word_count)
        flesch = max(0, min(100, round(flesch)))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / word_count) - 15.59
        fk_grade = max(0, round(fk_grade, 1))
    else:
        fk_grade = 0
    
    # Calculate SMOG Index (simplified)
    if sentence_count >= 30:
        smog = 1.043 * ((len(complex_words) * (30 / sentence_count)) ** 0.5) + 3.1291
    else:
        smog = 1.043 * ((len(complex_words) * (30 / max(1, sentence_count))) ** 0.5) + 3.1291
    smog = max(0, round(smog, 1))
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "smog_index": smog,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

def generate_lesson_plan(topic, grade_level, slide_count, profile, extra_notes, subject_area):
    """Generate a lesson plan structure using Claude with UK curriculum"""
    # Updated UK education level descriptions
    level_descriptions = {
        "ks3-year7-8": "Key Stage 3 (Year 7-8, ages 11-13)",
        "ks3-year9": "Key Stage 3 (Year 9, age 14)",
        "ks4-gcse": "Key Stage 4 (GCSE, Year 10-11, ages 15-16)",
        "ks5-alevel": "Key Stage 5 (A-Level, Year 12-13, ages 17-18)",
        "btec": "BTEC qualification level"
    }
    
    level_desc = level_descriptions.get(grade_level, "secondary school")
    
    adaptation_needs = {
        "dyslexia": "simplified vocabulary, short sentences, clear structure, minimal text per slide",
        "adhd": "visually engaging, chunked information, clear headings, bulleted lists, minimal distractions",
        "esl": "basic vocabulary with more complex terms in parentheses, visual support, simple sentence structure"
    }
    
    adaptation_desc = adaptation_needs.get(profile, "")
    
    # Subject-specific guidance - abbreviated for code clarity
    subject_guidance = {
        "english": "Include key literary terms, text analysis techniques, and writing structures appropriate for this level.",
        "maths": "Include clear mathematical notations, step-by-step explanations, and visual representations of concepts.",
        "science": "Include scientific diagrams, experimental procedures, and key terminology with explanations."
        # Other subjects removed for brevity
    }
    
    subject_specific_guidance = subject_guidance.get(subject_area, "")
    
    prompt = f"""
    Create a detailed lesson plan for a {slide_count}-slide PowerPoint presentation about "{topic}" for {level_desc} students following the UK National Curriculum for {subject_area}.
    
    This presentation needs to be adapted for students with {profile} and should include: {adaptation_desc}.
    
    Additional subject-specific guidance: {subject_specific_guidance}
    
    Additional instructions from the teacher: {extra_notes}
    
    For each slide, provide:
    1. Slide type (title slide, content slide, image slide, etc.)
    2. Title of the slide
    3. Key content points (in bullet form)
    4. A brief image description that would support the content (for slides that should have images)
    
    Format your response as a structured JSON array where each element is a slide with the following properties:
    - layout_index: (0 for title slide, 1 for content slide with title)
    - title: The slide title
    - content: The main content (bulleted list where appropriate)
    - image_prompt: A description for image generation (if applicable)
    
    Keep the language appropriate for {profile} needs and align with UK curriculum standards.
    """
    
    try:
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        content = response.content[0].text.strip()
        
        # Extract JSON from the response
        json_match = re.search(r'\[\s*{.*}\s*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            lesson_plan = json.loads(json_str)
            return lesson_plan
        else:
            # Fallback handling when JSON structure isn't found
            slides = []
            sections = content.split('\n\n')
            for i, section in enumerate(sections):
                if i == 0:
                    # First section is the title slide
                    slides.append({
                        "layout_index": 0,
                        "title": topic,
                        "content": section
                    })
                else:
                    # Extract title and content
                    lines = section.split('\n')
                    if lines:
                        title = lines[0].strip('#- ').strip()
                        content = '\n'.join(lines[1:]).strip()
                        slides.append({
                            "layout_index": 1,
                            "title": title,
                            "content": content,
                            "image_prompt": f"Educational illustration about {title} for {level_desc} students studying {subject_area}"
                        })
            
            return slides
    
    except Exception as e:
        print(f"Error generating lesson plan: {str(e)}")
        # Create a basic fallback plan
        return [
            {"layout_index": 0, "title": topic, "content": f"A presentation about {topic} for {level_desc} students studying {subject_area}"},
            {"layout_index": 1, "title": "Introduction", "content": f"• Introduction to {topic}\n• Key concepts\n• Learning objectives", "image_prompt": f"Introduction to {topic} for UK {subject_area} curriculum"},
            {"layout_index": 1, "title": "Main Concepts", "content": f"• Important aspects of {topic}\n• Key information", "image_prompt": f"Main concepts of {topic} for UK {subject_area} curriculum"}
        ]

def generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area):
    """Generate an enriched lesson plan based on the adapted content for UK curriculum"""
    # Implementation similar to generate_lesson_plan but enhanced with the adapted content
    # Shortened for brevity as the logic is similar to generate_lesson_plan
    prompt = f"""
    Create an enriched, enhanced lesson plan for a PowerPoint presentation about "{topic}" based on existing content.
    [Content details omitted for brevity]
    """
    
    try:
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        content = response.content[0].text.strip()
        
        # Extract and process JSON response
        # [Processing logic omitted as it's similar to generate_lesson_plan]
        # ...
        
    except Exception as e:
        print(f"Error generating enriched lesson plan: {str(e)}")
        # Return fallback plan
        # ...

def find_complex_words(text):
    """Find and count complex words (3+ syllables)"""
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    complex_words = {}
    
    # Count syllables function (simplified)
    def count_syllables(word):
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    for word in words:
        syllables = count_syllables(word)
        if syllables >= 3:
            if word in complex_words:
                complex_words[word] += 1
            else:
                complex_words[word] = 1
    
    # Sort by frequency (most frequent first)
    sorted_words = sorted(complex_words.items(), key=lambda x: x[1], reverse=True)
    
    # Return the top 10 most frequent complex words
    return dict(sorted_words[:10])

def calculate_adaptation_score(metrics, thresholds):
    """Calculate an overall adaptation score based on metrics and thresholds"""
    # Initialize score components
    flesch_score = 0
    grade_score = 0
    smog_score = 0
    sentence_score = 0
    complex_score = 0
    
    # Flesch Reading Ease (higher is better)
    flesch_threshold = thresholds['flesch_reading_ease']
    if metrics['flesch_reading_ease'] >= flesch_threshold:
        flesch_score = 100
    elif metrics['flesch_reading_ease'] >= flesch_threshold - 20:
        # Partial score for being close
        flesch_score = 50 + ((metrics['flesch_reading_ease'] - (flesch_threshold - 20)) / 20) * 50
    else:
        flesch_score = max(0, (metrics['flesch_reading_ease'] / (flesch_threshold - 20)) * 50)
    
    # Grade Level (lower is better)
    grade_threshold = thresholds['flesch_kincaid_grade']
    if metrics['flesch_kincaid_grade'] <= grade_threshold:
        grade_score = 100
    elif metrics['flesch_kincaid_grade'] <= grade_threshold + 4:
        # Partial score for being close
        grade_score = 100 - ((metrics['flesch_kincaid_grade'] - grade_threshold) / 4) * 50
    else:
        grade_score = max(0, 50 - ((metrics['flesch_kincaid_grade'] - (grade_threshold + 4)) / 4) * 50)
    
    # SMOG Index (lower is better)
    smog_threshold = thresholds['smog_index']
    if metrics['smog_index'] <= smog_threshold:
        smog_score = 100
    elif metrics['smog_index'] <= smog_threshold + 3:
        # Partial score for being close
        smog_score = 100 - ((metrics['smog_index'] - smog_threshold) / 3) * 50
    else:
        smog_score = max(0, 50 - ((metrics['smog_index'] - (smog_threshold + 3)) / 3) * 50)
    
    # Sentence Length (lower is better)
    sentence_threshold = thresholds['sentence_length']
    if metrics['sentence_length'] <= sentence_threshold:
        sentence_score = 100
    elif metrics['sentence_length'] <= sentence_threshold + 8:
        # Partial score for being close
        sentence_score = 100 - ((metrics['sentence_length'] - sentence_threshold) / 8) * 50
    else:
        sentence_score = max(0, 50 - ((metrics['sentence_length'] - (sentence_threshold + 8)) / 5) * 50)
    
    # Complex Word Percentage (lower is better)
    complex_threshold = thresholds['complex_word_percent']
    if metrics['complex_word_percent'] <= complex_threshold:
        complex_score = 100
    elif metrics['complex_word_percent'] <= complex_threshold + 15:
        # Partial score for being close
        complex_score = 100 - ((metrics['complex_word_percent'] - complex_threshold) / 15) * 50
    else:
        complex_score = max(0, 50 - ((metrics['complex_word_percent'] - (complex_threshold + 15)) / 10) * 50)
    
    # Weighted average (weights can be adjusted based on importance)
    weights = {
        'flesch': 0.3,
        'grade': 0.2,
        'smog': 0.1,
        'sentence': 0.2,
        'complex': 0.2
    }
    
    final_score = (
        flesch_score * weights['flesch'] +
        grade_score * weights['grade'] +
        smog_score * weights['smog'] +
        sentence_score * weights['sentence'] +
        complex_score * weights['complex']
    )
    
    # Round to whole number
    return round(final_score)

def generate_recommendation(score, profile, metrics, thresholds):
    """Generate a recommendation based on the score and metrics"""
    profile_names = {
        "dyslexia": "readers with dyslexia",
        "adhd": "readers with ADHD",
        "esl": "English language learners"
    }
    
    audience = profile_names.get(profile, profile)
    
    if score >= 85:
        return f"Your content is already well-suited for {audience}. Adaptation will provide minimal benefits but may still enhance readability."
    
    elif score >= 70:
        # Medium need for adaptation
        issues = []
        
        if metrics['flesch_reading_ease'] < thresholds['flesch_reading_ease']:
            issues.append("general readability")
            
        if metrics['flesch_kincaid_grade'] > thresholds['flesch_kincaid_grade']:
            issues.append("reading level")
            
        if metrics['sentence_length'] > thresholds['sentence_length']:
            issues.append("sentence length")
            
        if metrics['complex_word_percent'] > thresholds['complex_word_percent']:
            issues.append("vocabulary complexity")
            
        if issues:
            issue_text = ", ".join(issues[:-1])
            if len(issues) > 1:
                issue_text += f", and {issues[-1]}"
            else:
                issue_text = issues[0]
                
            return f"Your content would benefit from moderate adaptation for {audience}, particularly improving {issue_text}."
        else:
            return f"Your content would benefit from moderate adaptation to better suit {audience}."
    
    else:
        # High need for adaptation
        key_issues = []
        
        if metrics['flesch_reading_ease'] < thresholds['flesch_reading_ease'] - 15:
            key_issues.append("significantly improve general readability")
            
        if metrics['flesch_kincaid_grade'] > thresholds['flesch_kincaid_grade'] + 3:
            key_issues.append(f"reduce the reading level (currently at grade {metrics['flesch_kincaid_grade']})")
            
        if metrics['sentence_length'] > thresholds['sentence_length'] + 5:
            key_issues.append(f"shorten sentences (currently averaging {metrics['sentence_length']} words)")
            
        if metrics['complex_word_percent'] > thresholds['complex_word_percent'] + 10:
            key_issues.append(f"simplify vocabulary (currently {metrics['complex_word_percent']}% complex words)")
            
        if key_issues:
            issue_text = ", ".join(key_issues[:-1])
            if len(key_issues) > 1:
                issue_text += f", and {key_issues[-1]}"
            else:
                issue_text = key_issues[0]
                
            return f"Your content requires significant adaptation for {audience}. We recommend you: {issue_text}."
        else:
            return f"Your content requires significant adaptation to better suit the needs of {audience}."

def generate_complexity_chart(slide_texts, profile):
    """Generate a chart showing content complexity across slides"""
    # Create data for visualization
    slide_numbers = []
    complexity_scores = []
    threshold_values = []
    
    # Get threshold for the selected profile
    threshold = READABILITY_THRESHOLDS.get(profile, {}).get('flesch_kincaid_grade', 8)
    
    # Calculate complexity for each slide
    for slide in slide_texts:
        # Skip slides with very little text
        if len(slide['text'].strip()) < 20:
            continue
            
        metrics = calculate_simple_readability(slide['text'])
        slide_numbers.append(slide['slide_number'])
        complexity_scores.append(metrics['flesch_kincaid_grade'])
        threshold_values.append(threshold)
    
    # If no valid slides, return placeholder
    if not slide_numbers:
        return generate_placeholder_chart("No substantial text content found in slides")
    
    try:
        # Create the chart
        plt.figure(figsize=(8, 4))
        plt.bar(slide_numbers, complexity_scores, color='#3498db', alpha=0.7)
        plt.axhline(y=threshold, color='#e74c3c', linestyle='-', label=f'Target ({threshold})')
        
        plt.xlabel('Slide Number')
        plt.ylabel('Reading Grade Level')
        plt.title('Content Complexity by Slide')
        plt.legend()
        plt.grid(True, alpha=0.3)
        
        # Add value labels on bars
        for i, v in enumerate(complexity_scores):
            plt.text(slide_numbers[i], v + 0.3, str(round(v, 1)), ha='center')
        
        # Ensure y-axis starts at 0
        plt.ylim(bottom=0)
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
    
    except Exception as e:
        print(f"Error generating complexity chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_complex_words_chart(complex_words):
    """Generate a bar chart showing most frequent complex words"""
    try:
        # If no complex words found
        if not complex_words:
            return generate_placeholder_chart("No complex words found")
            
        # Sort and limit to top 10 words
        words = list(complex_words.keys())[:10]
        frequencies = list(complex_words.values())[:10]
        
        # Create the chart
        plt.figure(figsize=(8, 4))
        bars = plt.barh(words, frequencies, color='#2ecc71', alpha=0.7)
        
        plt.xlabel('Frequency')
        plt.ylabel('Words')
        plt.title('Most Frequent Complex Words')
        plt.grid(True, axis='x', alpha=0.3)
        
        # Add value labels
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 0.3, bar.get_y() + bar.get_height()/2, str(int(width)), va='center')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating complex words chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_placeholder_chart(message="No data available"):
    """Generate a placeholder chart with a message"""
    try:
        plt.figure(figsize=(8, 4))
        
        # Create an empty chart with a message
        plt.text(0.5, 0.5, message, ha='center', va='center', fontsize=14)
        plt.axis('off')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating placeholder chart: {e}")
        # Last resort - return empty image data
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="

def create_efficient_prompt(text, profile):
    print(f"DEBUG: Creating prompt for profile: {profile}")
    # First line contains essential instructions, rest is just the content
    if profile == "dyslexia":
        instructions = "Adapt for dyslexia: short sentences (max 15 words), simple words, active voice, bullet points where appropriate. Keep meaning intact."
    elif profile == "adhd":
        instructions = "Adapt for ADHD: clear structure, short chunks, bullet points, highlight key info, remove unnecessary details. Keep meaning intact."
    else:  # ESL
        instructions = "Adapt for English learners: simpler words (original in parentheses), short sentences, explain idioms, consistent terms. Keep meaning intact."
    
    print(f"DEBUG: Using instructions for {profile}: {instructions}")
    
    # Minimal format with clear sections
    prompt = f"{instructions}\n\nOriginal:\n{text}\n\nAdapted:"
    print(f"DEBUG: Generated prompt length: {len(prompt)} characters")
    return prompt

def adapt_text_with_matcha(text, profile):
    """Adapt text using Claude for the specified learning profile"""
    # Skip processing very short text
    if len(text.strip()) < 15:
        return text
    
    # Normalize the profile to prevent case issues
    profile = profile.lower() if profile else "dyslexia"
    
    # Check cache first
    cached_result = adaptation_cache.get(text, profile)
    if cached_result:
        return cached_result
    
    try:
        # Create efficient prompt with normalized profile
        prompt = create_efficient_prompt(text, profile)
        
        # Call Claude API
        response = call_claude_api_with_retry(prompt)
        
        # Extract adapted text from response
        adapted_text = response.content[0].text.strip()
        
        # Store in cache for future use
        adaptation_cache.set(text, profile, adapted_text)
        
        return adapted_text
        
    except Exception as e:
        print(f"Error in text adaptation: {e}")
        # Return original text if adaptation fails
        return text

def translate_text(text, target_language):
    """
    Translate text to the specified target language using Claude
    """
    print(f"DEBUG: translate_text called with target_language: '{target_language}', text length: {len(text) if text else 0}")
    
    if not text or not text.strip():
        print("DEBUG: Empty text provided to translate_text, returning as-is")
        return text
        
    try:
        # Special handling instructions for specific languages
        special_instructions = ""
        
        if target_language.lower() == "polish":
            special_instructions = """
            When translating to Polish:
            - Maintain proper Polish diacritical marks (ą, ć, ę, ł, ń, ó, ś, ź, ż)
            - Pay attention to grammatical cases and gender agreement
            - For technical terms, consider including the English original in parentheses first time
            """
        elif target_language.lower() == "ukrainian":
            special_instructions = """
            When translating to Ukrainian:
            - Use modern Ukrainian vocabulary rather than Russified terms where possible
            - Properly handle Ukrainian specific characters (є, і, ї, ґ)
            - For technical/scientific terms, consider providing the English original in parentheses on first occurrence
            """
        
        # Create a prompt for the translation with language-specific instructions
        prompt = f"""Translate the following text to {target_language}. 
Maintain the original formatting, including line breaks and bullet points.
Ensure the translation sounds natural and conversational.
{special_instructions}

Text to translate:
"{text}"

Translation:"""

        # Call Claude API with the translation prompt
        response = call_claude_api_with_retry(prompt)
        
        # Extract the translated content from the response
        translated_text = response.content[0].text.strip()
        
        # Log success for monitoring
        print(f"Successfully translated text of length {len(text)} to {target_language}")
        
        return translated_text
        
    except Exception as e:
        print(f"Error translating text: {str(e)}")
        # Return original text if translation fails
        return text

# Error handling and user feedback routes
@app.route('/adapt/<file_id>/<profile>')
def adapt(file_id, profile):
    print(f"DEBUG: Adaptation requested with profile: {profile}")
    """Adapt the PowerPoint based on assessment with API check"""
    try:
        # Check API connection first
        api_status = api_utils.check_connection()
        if api_status["status"] != "connected":
            # API connection failed - show error
            return render_template_string(ERROR_TEMPLATE, 
                                         message=f"Claude API connection error: {api_status['message']}"), 400
        
        # Validate if file_id exists in our system
        if file_id not in processing_tasks:
            return render_template_string(ERROR_TEMPLATE, 
                                         message="Presentation not found. Please upload again."), 404
        
        # Set initial processing status with empty progress
        processing_tasks[file_id] = {
            'status': 'processing', 
            'message': 'Starting adaptation...',
            'progress': {
                'total': 0,
                'processed': 0,
                'percentage': 0
            }
        }
        
        # Get the original filename from the processing tasks
        filename = processing_tasks[file_id].get('filename', 'presentation.pptx')
        
        # IMPORTANT FIX: Use the correct file path pattern
        # Look for files in the format "file_id_filename.pptx"
        file_path = None
        for file in os.listdir(app.config['UPLOAD_FOLDER']):
            if file.startswith(f"{file_id}_"):
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], file)
                break
        
        # If file not found, check alternate patterns
        if not file_path:
            # Try with exact file_id
            potential_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}.pptx")
            if os.path.exists(potential_path):
                file_path = potential_path
        
        # Double check the file exists
        if not file_path or not os.path.exists(file_path):
            # Debug info
            print(f"File not found. Looking for file with ID: {file_id}")
            print(f"Files in upload folder: {os.listdir(app.config['UPLOAD_FOLDER'])}")
            
            processing_tasks[file_id] = {'status': 'error', 'message': 'Presentation file not found on server'}
            return redirect(url_for('error', message="File not found on server"))
        
        # Validate profile type
        valid_profiles = ['dyslexia', 'esl', 'adhd', 'visual', 'cognitive']
        if profile not in valid_profiles:
            processing_tasks[file_id] = {'status': 'error', 'message': f'Invalid profile: {profile}'}
            return redirect(url_for('error', message=f"Invalid profile: {profile}"))
        
        # Start processing in a background thread to avoid blocking
        def process_task():
            try:
                process_presentation_efficiently(file_path, file_id, filename, profile)
            except Exception as e:
                import traceback
                traceback.print_exc()
                processing_tasks[file_id] = {'status': 'error', 'message': str(e)}
        
        # Start the processing thread with a name containing the file_id
        processing_thread = threading.Thread(
            target=process_task,
            name=f"process-{file_id}"  # This helps identify which file_id is being processed
        )
        processing_thread.daemon = True
        processing_thread.start()
        
        # Map profile to display name for the processing template
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Determine which template to use (with or without progress)
        template = PROCESSING_TEMPLATE_WITH_PROGRESS if 'PROCESSING_TEMPLATE_WITH_PROGRESS' in globals() else PROCESSING_TEMPLATE_WITH_PROGRESS
        
        # Return the processing template with progress tracking
        return render_template_string(
            template,
            file_id=file_id,
            filename=f"adapted_{filename}",
            profile=profile,
            profile_name=profile_names.get(profile, profile)
        )
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        error_message = f"Error starting adaptation: {str(e)}"
        return redirect(url_for('error', message=error_message))

# Route to handle the status check for long-running tasks    
@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Handle generation of a new presentation"""
    try:
        # Get parameters from form
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Generate a unique ID for this presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        filename = f"{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_new_presentation, 
            args=(file_id, filename, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                                   file_id=file_id, 
                                   filename=filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

@app.route('/enrich_and_generate', methods=['POST'])
def enrich_and_generate():
    """Handle generating an enhanced presentation based on an existing adapted file"""
    try:
        # Get parameters from form
        original_file_id = request.form.get('original_file_id')
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Find the original adapted file
        adapted_filename = None
        adapted_file_path = None
        
        for filename in os.listdir(app.config['OUTPUT_FOLDER']):
            if filename.startswith(f"{original_file_id}_"):
                adapted_filename = filename.replace(f"{original_file_id}_", "")
                adapted_file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
                break
        
        if not adapted_file_path:
            return render_template_string(ERROR_TEMPLATE, message="Original adapted file not found"), 404
        
        # Generate a unique ID for this new presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename for the new presentation
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        new_filename = f"enriched-{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_enriched_presentation, 
            args=(file_id, new_filename, adapted_file_path, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                                   file_id=file_id, 
                                   filename=new_filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

def generate_new_presentation(file_id, filename, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate a new presentation based on the specified parameters"""
    try:
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Step 1: Plan the content structure using Claude
        lesson_plan = educational_service.generate_lesson_plan(
            topic, grade_level, slide_count,
            learning_objectives=[f"Understand {topic}", f"Apply {topic} concepts"],
            profile=profile
        )
        
        # Step 2: Create slides based on the plan
        for slide_data in lesson_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content
                    adapted_content = adapt_text_with_matcha(content, profile)
                    text_frame = content_placeholder.text_frame
                    apply_text_to_text_frame(text_frame, adapted_content, profile)
        
        # Save the presentation
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

def convert_pptx_to_pdf_windows(pptx_path, pdf_path):
    """
    Convert PowerPoint to PDF using COM objects (Windows only)
    """
    try:
        import comtypes.client  # Windows-only library for COM automation
    except ImportError:
        print("comtypes not available (Windows only)")
        return False

    
    try:
        # Constants for PowerPoint conversion
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = True
        
        # Open the presentation
        deck = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
        
        # Save as PDF
        deck.SaveAs(pdf_path, 32)  # 32 is the constant for PDF format
        
        # Close the presentation and quit PowerPoint
        deck.Close()
        powerpoint.Quit()
        
        return True
    except Exception as e:
        print(f"Error converting to PDF using Windows COM: {str(e)}")
        return False

def convert_pptx_to_pdf_libreoffice(pptx_path, pdf_path):
    """
    Convert PowerPoint to PDF using LibreOffice/OpenOffice (Linux/Mac)
    Requires soffice or unoconv to be installed
    """
    import subprocess
    import platform
    
    try:
        # Check if using Mac
        if platform.system() == 'Darwin':  # macOS
            # On macOS, the path to soffice might be different
            soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            result = subprocess.run([
                soffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(pdf_path),
                pptx_path
            ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        else:  # Linux
            # Try using unoconv first
            try:
                result = subprocess.run([
                    'unoconv',
                    '-f', 'pdf',
                    '-o', pdf_path,
                    pptx_path
                ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            except:
                # Fall back to soffice if unoconv fails
                result = subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(pdf_path),
                    pptx_path
                ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Check if the PDF file was created
        if os.path.exists(pdf_path):
            return True
        else:
            # LibreOffice might have used a different filename
            base_filename = os.path.splitext(os.path.basename(pptx_path))[0]
            potential_pdf_path = os.path.join(os.path.dirname(pdf_path), f"{base_filename}.pdf")
            
            if os.path.exists(potential_pdf_path):
                # Rename to the expected path
                os.rename(potential_pdf_path, pdf_path)
                return True
        
        print(f"Error converting to PDF using LibreOffice: {result.stderr.decode()}")
        return False
    except Exception as e:
        print(f"Error converting to PDF: {str(e)}")
        return False

def convert_pptx_to_pdf_fallback(pptx_path, pdf_path):
    """
    Fallback conversion method using pdf2image and python-pptx
    Not as good quality but works cross-platform without dependencies
    """
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        import tempfile
        import io
        
        # Create a temporary directory to store images
        with tempfile.TemporaryDirectory() as temp_dir:
            # Load the presentation
            prs = Presentation(pptx_path)
            
            # Create a PDF 
            c = canvas.Canvas(pdf_path, pagesize=letter)
            
            # Get dimensions
            width, height = letter
            
            # For each slide, export as image and add to PDF
            for i, slide in enumerate(prs.slides):
                # Export slide as image (this part is tricky without COM/LibreOffice)
                # Using a placeholder approach here - in practice, you'd need a more robust solution
                
                # Create a temporary file for the slide
                slide_pptx = os.path.join(temp_dir, f"slide_{i}.pptx")
                slide_pdf = os.path.join(temp_dir, f"slide_{i}.pdf")
                
                # Create a new presentation with just this slide
                temp_prs = Presentation()
                # Copy the slide layout
                slide_layout = temp_prs.slide_layouts[0]  
                new_slide = temp_prs.slides.add_slide(slide_layout)
                
                # Copy content (simplified, actual implementation would be more complex)
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text:
                        for shape_placeholder in new_slide.placeholders:
                            if hasattr(shape_placeholder, "text_frame"):
                                shape_placeholder.text = shape.text
                                break
                
                # Save the temporary presentation
                temp_prs.save(slide_pptx)
                
                # Try to convert this slide to PDF using any available method
                if not convert_pptx_to_pdf_windows(slide_pptx, slide_pdf) and \
                   not convert_pptx_to_pdf_libreoffice(slide_pptx, slide_pdf):
                    # If both methods fail, we'll just create a simple page
                    c.drawString(100, height - 100, f"Slide {i+1}")
                    if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                        c.drawString(100, height - 150, f"Title: {slide.shapes.title.text}")
                    c.showPage()
                    continue
                    
                # If we succeeded in creating a PDF for this slide
                if os.path.exists(slide_pdf):
                    # Convert PDF to images
                    images = convert_from_path(slide_pdf)
                    
                    # Add each image to our PDF
                    for img in images:
                        # Convert PIL image to reportlab-compatible format
                        img_bytes = io.BytesIO()
                        img.save(img_bytes, format='JPEG')
                        img_bytes.seek(0)
                        
                        # Add to PDF
                        c.drawImage(img_bytes, 0, 0, width, height)
                        c.showPage()
            
            c.save()
            return True
            
    except Exception as e:
        print(f"Error in fallback PDF conversion: {str(e)}")
        return False

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """
    Main function to convert PowerPoint to PDF, trying different methods
    """
    import platform
    
    # Determine system type
    system = platform.system()
    
    # Try the appropriate method first
    if system == 'Windows':
        if convert_pptx_to_pdf_windows(pptx_path, pdf_path):
            return True
    else:  # Linux or Mac
        if convert_pptx_to_pdf_libreoffice(pptx_path, pdf_path):
            return True
    
    # If the primary method fails, try the fallback
    return convert_pptx_to_pdf_fallback(pptx_path, pdf_path)

# 3. Add PDF conversion capability to your processing flow

def process_presentation_with_pdf(file_path, file_id, filename, profile, target_language=None):
    """Modified presentation processing function that also generates PDF"""
    # First do the normal PowerPoint adaptation
    output_path = process_presentation_efficiently(file_path, file_id, filename, profile, target_language)
    
    # Also generate PDF version
    if output_path:
        pdf_filename = f"adapted_{os.path.splitext(filename)[0]}.pdf"
        pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{pdf_filename}")
        
        try:
            # Convert the adapted PowerPoint to PDF
            pdf_success = conversion_service.convert_pptx_to_pdf(output_path, pdf_path)
            
            # Update processing task with PDF info
            if pdf_success and os.path.exists(pdf_path):
                processing_tasks[file_id]['has_pdf'] = True
                processing_tasks[file_id]['pdf_filename'] = pdf_filename
            else:
                processing_tasks[file_id]['has_pdf'] = False
                
        except Exception as e:
            print(f"Error generating PDF: {str(e)}")
            processing_tasks[file_id]['has_pdf'] = False
    
    return output_path

def generate_enriched_presentation(file_id, filename, adapted_file_path, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate an enriched presentation based on an adapted file"""
    try:
        # Extract content from the adapted file
        adapted_content = pptx_service.extract_content_from_pptx(adapted_file_path)
        
        # Generate an enriched presentation plan
        enriched_plan = generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Create slides based on the plan
        for slide_data in enriched_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content - already adapted content
                    text_frame = content_placeholder.text_frame
                    apply_text_to_text_frame(text_frame, content, profile)
        
        # Save the presentation
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

# HTML templates module
# Create a separate file html_templates.py with all the HTML templates

if __name__ == "__main__":
    # Run the Flask app
    app.run(debug=True, host='0.0.0.0', port=5000)