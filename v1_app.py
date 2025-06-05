import os
from flask import Flask, request, jsonify, send_from_directory, send_file, make_response
from flask_cors import CORS
from werkzeug.utils import secure_filename
import sqlite3
import json
import uuid
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import PyPDF2
import anthropic
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
api_key = os.getenv("ANTHROPIC_API_KEY")
if not api_key:
    raise ValueError("ANTHROPIC_API_KEY not found in environment variables")

def test_claude_models():
    """Test available Claude models and return one that works"""
    if not api_key:
        print("No API key available, cannot test models")
        return None
        
    # Models to try in order of preference
    models_to_try = [
        "claude-3-opus-20240229",
        "claude-3-sonnet-20240229", 
        "claude-3-haiku-20240307",
        "claude-3.5-sonnet-20240620",
        "claude-3-opus",
        "claude-3-sonnet",
        "claude-3-haiku",
        "claude-3.5-sonnet"
    ]
    
    client = anthropic.Anthropic(api_key=api_key)
    
    for model in models_to_try:
        try:
            print(f"Testing model: {model}")
            response = client.messages.create(
                model=model,
                max_tokens=10,
                messages=[{"role": "user", "content": "Hello"}]
            )
            print(f"Model {model} works!")
            return model
        except Exception as e:
            print(f"Model {model} failed: {str(e)}")
    
    print("No models worked, defaulting to claude-3-opus-20240229")
    return "claude-3-opus-20240229"  # Default to this even if it failed

# Call this function during initialization
WORKING_MODEL = test_claude_models()
print(f"Using Claude model: {WORKING_MODEL}")

if api_key:
    api_key = api_key.strip()  # Remove any whitespace
    print(f"API Key after stripping whitespace: Length={len(api_key)}")
    
    # Sometimes .env files can add quotes to the key, remove them
    if api_key.startswith('"') and api_key.endswith('"'):
        api_key = api_key[1:-1]
        print("Removed quotes from API key")
    
    if api_key.startswith("'") and api_key.endswith("'"):
        api_key = api_key[1:-1]
        print("Removed single quotes from API key")
        
    print(f"Final API Key format check: {api_key.startswith('sk-ant-')}")

# Create the Claude client with the cleaned API key
try:
    claude = anthropic.Anthropic(api_key=api_key)
    print("Claude client initialized successfully")
except Exception as e:
    print(f"Error initializing Claude client: {str(e)}")
    claude = None  # Set to None so we can check later

print(f"API Key loaded: {api_key[:8] if api_key and len(api_key) > 8 else 'None'}...")
if not api_key or not api_key.startswith('sk-ant-'):
    print("WARNING: Claude API key is missing or invalid!")
    print("Please set a valid Claude API key in your .env file or environment variables")

print("==== API KEY DEBUGGING ====")
print(f"API Key loaded: {api_key[:8] if api_key and len(api_key) > 8 else 'None'}...")
print(f"API Key length: {len(api_key) if api_key else 0}")
print(f"API Key format check: {api_key.startswith('sk-ant-') if api_key else False}")

try:
    print("Attempting to connect to Claude API...")
    client = anthropic.Anthropic(api_key=api_key)
    response = client.messages.create(
        model="claude-3-opus-20240229",
        max_tokens=100,
        messages=[{"role": "user", "content": "Hello, Claude! This is a test message."}]
    )
    print("Success! Claude responded with:", response.content[0].text)
except Exception as e:
    print("Error connecting to Claude API:", e)
    print("Error type:", type(e).__name__)
    print("Full error details:", repr(e))

print("==== END DEBUGGING ====")

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = 'static/uploads'
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB

claude = anthropic.Anthropic(api_key=api_key)

def get_db_connection():
    conn = sqlite3.connect('adaptive_learning.db')
    conn.row_factory = sqlite3.Row
    return conn

def update_db_schema():
    print("Checking database schema for updates...")
    try:
        conn = get_db_connection()
        # Check if column exists
        cursor = conn.execute("PRAGMA table_info(adaptations)")
        columns = [info[1] for info in cursor.fetchall()]
        
        if 'user_prompt' not in columns:
            print("Adding user_prompt column to adaptations table...")
            conn.execute("ALTER TABLE adaptations ADD COLUMN user_prompt TEXT;")
            conn.commit()
            print("Database schema updated successfully.")
        else:
            print("user_prompt column already exists.")
        conn.close()
    except Exception as e:
        print(f"Error updating database schema: {str(e)}")

def init_db():
    conn = get_db_connection()
    
    # Create documents table
    conn.execute('''
    CREATE TABLE IF NOT EXISTS documents (
        id TEXT PRIMARY KEY,
        file_name TEXT NOT NULL,
        file_type TEXT NOT NULL,
        content JSON NOT NULL,
        upload_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    ''')
    
    # Create profiles table
    conn.execute('''
    CREATE TABLE IF NOT EXISTS profiles (
        id TEXT PRIMARY KEY,
        name TEXT NOT NULL,
        description TEXT NOT NULL,
        adaptation_instructions TEXT NOT NULL
    )
    ''')

    # Create adaptations table
    conn.execute('''
    CREATE TABLE IF NOT EXISTS adaptations (
        id TEXT PRIMARY KEY,
        document_id TEXT NOT NULL,
        profile_id TEXT NOT NULL,
        adapted_content JSON NOT NULL,
        user_prompt TEXT,
        creation_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY (document_id) REFERENCES documents (id),
        FOREIGN KEY (profile_id) REFERENCES profiles (id)
    )
    ''')

    # Initialize default profiles
    profiles = [
        {
            'id': 'dyslexia',
            'name': 'Dyslexia Support',
            'description': 'Adapts content for students with dyslexia',
            'adaptation_instructions': '1. Use shorter sentences (max 12 words)\n2. Simplify vocabulary while preserving meaning\n3. Use bullet points instead of paragraphs\n4. Add more white space between text\n5. Bold key terms'
        },
        {
            'id': 'adhd',
            'name': 'ADHD Support',
            'description': 'Adapts content for students with attention difficulties',
            'adaptation_instructions': '1. Break content into smaller chunks\n2. Use more headings and subheadings\n3. Add emphasis to important points\n4. Use numbered steps for processes\n5. Include engagement questions'
        },
        {
            'id': 'esl',
            'name': 'English Language Learner',
            'description': 'Adapts content for students learning English',
            'adaptation_instructions': '1. Use simpler vocabulary\n2. Add brief definitions for difficult terms\n3. Use shorter sentences\n4. Provide clear structure\n5. Include more examples'
        }
    ]

    for profile in profiles:
        conn.execute(
            'INSERT OR IGNORE INTO profiles (id, name, description, adaptation_instructions) VALUES (?, ?, ?, ?)',
            (profile['id'], profile['name'], profile['description'], profile['adaptation_instructions'])
        )

    conn.commit()
    conn.close()

def update_dyslexia_profile():
    """
    Updates the dyslexia profile in the database to include Tom Gates styling
    """
    conn = get_db_connection()
    
    # Enhanced adaptation instructions with Tom Gates styling
    new_instructions = """1. Use shorter sentences (max 12 words)
2. Simplify vocabulary while preserving meaning
3. Use bullet points instead of paragraphs
4. Add more white space between text sections
5. Bold key terms
6. Apply Tom Gates-style formatting:
   - Vary font sizes for emphasis (larger for important points)
   - Mix different font styles (use Comic Sans or similar dyslexia-friendly fonts)
   - Create visual emphasis with CAPITALS, bold, and *asterisks* for key words
   - Add decorative elements around important information like boxes or clouds
   - Break content into small, visually distinct chunks
7. Create a fun, engaging visual hierarchy that aids reading and comprehension"""
    
    # Update the dyslexia profile
    conn.execute(
        'UPDATE profiles SET adaptation_instructions = ? WHERE id = ?',
        (new_instructions, 'dyslexia')
    )
    
    conn.commit()
    conn.close()
    
    print("Dyslexia profile updated with Tom Gates styling instructions")
    return True

def replace_text_with_tom_gates_styling(shape, text):
    """
    Replace text in a shape with Tom Gates-style formatting applied.
    Parses text for markdown-like formatting indicators and applies appropriate formatting.
    """
    if not hasattr(shape, "text_frame"):
        print(f"Shape has no text_frame attribute")
        return
        
    text_frame = shape.text_frame
    
    # Clear existing paragraphs
    for i in range(len(text_frame.paragraphs) - 1, -1, -1):
        p = text_frame.paragraphs[i]
        try:
            p._element.getparent().remove(p._element)
        except:
            pass
    
    # Split text into paragraphs
    paragraphs = text.split('\n')
    
    # Process each paragraph
    for para_text in paragraphs:
        if not para_text.strip():
            continue
            
        p = text_frame.add_paragraph()
        
        # Check for overall paragraph styling
        font_size = None
        if '[LARGE FONT]' in para_text:
            font_size = Pt(18)  # Larger font
            para_text = para_text.replace('[LARGE FONT]', '')
        elif '[small font]' in para_text:
            font_size = Pt(10)  # Smaller font
            para_text = para_text.replace('[small font]', '')
        
        # Default font size if not specified
        if font_size is None:
            font_size = Pt(14)  # Regular size
        
        # Split for inline styling with special markers
        # This is a simplified parsing - in a real implementation, you'd want more robust parsing
        parts = []
        current_part = ""
        bold_active = False
        italic_active = False
        
        # Very simple parsing of formatting markers
        i = 0
        while i < len(para_text):
            if para_text[i:i+2] == '**' and not italic_active:
                if bold_active:
                    parts.append((current_part, {'bold': True}))
                    current_part = ""
                    bold_active = False
                else:
                    if current_part:
                        parts.append((current_part, {}))
                        current_part = ""
                    bold_active = True
                i += 2
            elif para_text[i:i+1] == '*' and not bold_active:
                if italic_active:
                    parts.append((current_part, {'italic': True}))
                    current_part = ""
                    italic_active = False
                else:
                    if current_part:
                        parts.append((current_part, {}))
                        current_part = ""
                    italic_active = True
                i += 1
            else:
                current_part += para_text[i]
                i += 1
        
        # Add any remaining text
        if current_part:
            if bold_active:
                parts.append((current_part, {'bold': True}))
            elif italic_active:
                parts.append((current_part, {'italic': True}))
            else:
                parts.append((current_part, {}))
        
        # If we didn't parse any special formatting, add the whole text at once
        if not parts:
            parts = [(para_text, {})]
        
        # Create runs for each part with appropriate styling
        for text_part, style in parts:
            # Check for ALL CAPS styling
            all_caps = False
            if text_part.isupper() and len(text_part) > 3:  # Only apply to longer text
                all_caps = True
            
            run = p.add_run()
            run.text = text_part
            
            # Apply font size
            run.font.size = font_size
            
            # Apply parsed styles
            if style.get('bold') or all_caps:
                run.font.bold = True
            if style.get('italic'):
                run.font.italic = True
            
            # Apply font variations for visual interest
            # Use a small set of dyslexia-friendly fonts
            if text_part.isupper() and len(text_part) > 2:
                run.font.name = "Comic Sans MS"
            elif style.get('bold'):
                run.font.name = "Arial"
            elif style.get('italic'):
                run.font.name = "Verdana"
            else:
                # Default font - use a dyslexia-friendly font
                run.font.name = "Trebuchet MS"
    
    return True

# Functions for TOM GATES export using direct creation approach
def create_tom_gates_pptx(adapted_content, output_path):
    """
    Creates a PowerPoint with full Tom Gates styling.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import random
    
    print(f"Creating Tom Gates PowerPoint at {output_path}")
    
    # Create new presentation
    prs = Presentation()
    
    # Get blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    
    # Process each slide
    for i, slide_data in enumerate(adapted_content.get('slides', [])):
        print(f"Processing slide {i+1}")
        
        # Create new slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        if slide_data.get('title'):
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(1.2))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = slide_data.get('title')
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.bold = True
            title_para.font.size = Pt(32)
            title_para.font.color.rgb = RGBColor(0, 0, 128)  # Dark blue
            
            # Use Comic Sans if available
            try:
                title_para.font.name = "Comic Sans MS"
            except:
                pass
        
        # Set up for content items
        top_position = Inches(2.0)
        
        # Process text elements
        for element in slide_data.get('text_elements', []):
            # Skip empty elements
            if not element.get('text'):
                continue
                
            # Get text content
            content = element.get('text')
            
            # Create textbox for this element
            content_box = slide.shapes.add_textbox(
                Inches(1), top_position, Inches(8), Inches(0.8))
            content_frame = content_box.text_frame
            content_para = content_frame.add_paragraph()
            
            # Check if this is a question/answer item
            if ":" in content or "?" in content:
                # Split into parts for different formatting
                if ":" in content:
                    parts = content.split(":", 1)
                    question = parts[0] + ":"
                    answer = parts[1].strip() if len(parts) > 1 else ""
                else:
                    parts = content.split("?", 1)
                    question = parts[0] + "?"
                    answer = parts[1].strip() if len(parts) > 1 else ""
                
                # Add question part
                q_run = content_para.add_run()
                q_run.text = question + " "
                q_run.font.bold = True
                q_run.font.size = Pt(18)
                
                # Add answer part in red
                if answer:
                    a_run = content_para.add_run()
                    a_run.text = answer
                    a_run.font.color.rgb = RGBColor(255, 0, 0)  # Red
                    a_run.font.size = Pt(18)
                
                # Use Comic Sans if available
                try:
                    content_para.font.name = "Comic Sans MS"
                except:
                    pass
                    
                # Add yellow highlighting for answers
                if "=" in content or "mass" in content.lower():
                    # Add yellow highlight shape behind the text
                    highlight = slide.shapes.add_shape(
                        1,  # Rectangle shape
                        content_box.left - Inches(0.1),
                        content_box.top - Inches(0.05),
                        content_box.width + Inches(0.2),
                        content_box.height + Inches(0.1)
                    )
                    highlight.fill.solid()
                    highlight.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow
                    
                    # Move highlight to back
                    highlight.z_order = 0
            else:
                # Regular content
                content_para.text = content
                content_para.font.size = Pt(18)
                
                # Use Comic Sans if available
                try:
                    content_para.font.name = "Comic Sans MS"
                except:
                    pass
            
            # Move down for next item with slight variation for visual interest
            spacing = random.uniform(0.8, 1.1)
            top_position += Inches(spacing)
        
        # Add decorative element for important slides
        if any(keyword in slide_data.get('title', '').lower() 
               for keyword in ['mass', 'weight', 'product']):
            # Create special emphasis text at bottom
            emphasis_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(9), Inches(1))
            emphasis_frame = emphasis_box.text_frame
            emphasis_para = emphasis_frame.add_paragraph()
            emphasis_para.text = "IF PRODUCTS WEIGH 88G REACTANTS MUST WEIGH... 88G!"
            emphasis_para.alignment = PP_ALIGN.CENTER
            emphasis_para.font.bold = True
            emphasis_para.font.size = Pt(28)
            
            # Use Comic Sans if available
            try:
                emphasis_para.font.name = "Comic Sans MS"
            except:
                pass
    
    # Save the presentation
    prs.save(output_path)
    print(f"Tom Gates PowerPoint created at {output_path}")
    return output_path

def create_tom_gates_pdf(adapted_content, output_path):
    """
    Creates a PDF with Tom Gates styling.
    """
    try:
        # Import required libraries
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.colors import black, white, red, blue, yellow
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
    except ImportError:
        print("Installing reportlab...")
        import subprocess
        subprocess.check_call(["pip", "install", "reportlab"])
        
        # Import after installation
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.colors import black, white, red, blue, yellow
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        
    print(f"Creating Tom Gates PDF at {output_path}")
    
    # Create PDF document
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    elements = []
    
    # Create styles
    styles = getSampleStyleSheet()
    
    # Create custom Tom Gates styles
    styles.add(ParagraphStyle(
        name='TomGatesTitle',
        parent=styles['Heading1'],
        fontSize=24,
        alignment=TA_CENTER,
        textColor=blue,
        fontName='Helvetica-Bold',  # Fallback if Comic Sans not available
        spaceAfter=20
    ))
    
    styles.add(ParagraphStyle(
        name='TomGatesQuestion',
        parent=styles['Normal'],
        fontSize=16,
        textColor=black,
        fontName='Helvetica',
        spaceAfter=5
    ))
    
    styles.add(ParagraphStyle(
        name='TomGatesAnswer',
        parent=styles['Normal'],
        fontSize=16,
        textColor=red,
        fontName='Helvetica',
        spaceAfter=15,
        backColor=yellow
    ))
    
    styles.add(ParagraphStyle(
        name='TomGatesEmphasis',
        parent=styles['Normal'],
        fontSize=20,
        textColor=black,
        fontName='Helvetica-Bold',
        alignment=TA_CENTER,
        spaceAfter=20
    ))
    
    # Process content based on type
    if 'pages' in adapted_content:
        # PDF document format
        for page_idx, page in enumerate(adapted_content.get('pages', [])):
            # Add page break between pages except first
            if page_idx > 0:
                elements.append(PageBreak())
                
            # Add page title
            elements.append(Paragraph(f"Page {page_idx + 1}", styles['TomGatesTitle']))
            elements.append(Spacer(1, 12))
            
            # Process paragraphs
            for para_text in page.get('paragraphs', []):
                if not para_text.strip():
                    continue
                    
                # Handle question/answer format
                if ":" in para_text:
                    parts = para_text.split(":", 1)
                    elements.append(Paragraph(parts[0] + ":", styles['TomGatesQuestion']))
                    if len(parts) > 1:
                        elements.append(Paragraph(parts[1].strip(), styles['TomGatesAnswer']))
                else:
                    # Regular paragraph
                    elements.append(Paragraph(para_text, styles['TomGatesQuestion']))
                    
                elements.append(Spacer(1, 8))
                
    elif 'slides' in adapted_content:
        # Presentation format converted to PDF
        for slide_idx, slide in enumerate(adapted_content.get('slides', [])):
            # Add page break between slides except first
            if slide_idx > 0:
                elements.append(PageBreak())
                
            # Add slide title
            elements.append(Paragraph(slide.get('title', f'Slide {slide_idx + 1}'), 
                                     styles['TomGatesTitle']))
            elements.append(Spacer(1, 12))
            
            # Process text elements
            for element in slide.get('text_elements', []):
                text = element.get('text', '')
                if not text.strip():
                    continue
                    
                # Handle question/answer format
                if ":" in text:
                    parts = text.split(":", 1)
                    elements.append(Paragraph(parts[0] + ":", styles['TomGatesQuestion']))
                    if len(parts) > 1:
                        elements.append(Paragraph(parts[1].strip(), styles['TomGatesAnswer']))
                elif "?" in text:
                    parts = text.split("?", 1)
                    elements.append(Paragraph(parts[0] + "?", styles['TomGatesQuestion']))
                    if len(parts) > 1:
                        elements.append(Paragraph(parts[1].strip(), styles['TomGatesAnswer']))
                else:
                    # Regular paragraph
                    elements.append(Paragraph(text, styles['TomGatesQuestion']))
                    
                elements.append(Spacer(1, 8))
                
            # Add emphasis for important slides
            if any(keyword in slide.get('title', '').lower() 
                   for keyword in ['mass', 'weight', 'product']):
                elements.append(Spacer(1, 20))
                elements.append(Paragraph(
                    "IF PRODUCTS WEIGH 88G REACTANTS MUST WEIGH... 88G!", 
                    styles['TomGatesEmphasis']
                ))
                
    # Build PDF
    doc.build(elements)
    print(f"Tom Gates PDF created at {output_path}")
    return output_path

# This new route doesn't conflict with existing routes
@app.route('/enable-tom-gates-style', methods=['GET'])
def enable_tom_gates_style():
    """
    Updates the app to incorporate Tom Gates styling for dyslexia support.
    This is an admin route to apply the changes without modifying the core files.
    """
    try:
        # Update the dyslexia profile with new instructions
        success = update_dyslexia_profile()
        
        # Return success message
        return """
        <html>
        <head>
            <title>Tom Gates Style Update</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
                .container { max-width: 800px; margin: 0 auto; }
                h1 { color: #4CAF50; }
                .success { background-color: #dff0d8; padding: 15px; border-radius: 5px; margin: 20px 0; }
                pre { background-color: #f5f5f5; padding: 15px; border-radius: 5px; overflow-x: auto; }
                .btn { background-color: #4CAF50; color: white; padding: 10px 15px; text-decoration: none; display: inline-block; border-radius: 4px; margin-top: 20px; }
            </style>
        </head>
        <body>
            <div class="container">
                <h1>Tom Gates Style Update</h1>
                <div class="success">
                    <p><strong>Success!</strong> The dyslexia profile has been updated with Tom Gates styling instructions.</p>
                </div>
                <h2>How to Use Tom Gates Styling:</h2>
                <ol>
                    <li>Go back to the home page</li>
                    <li>Select a document to adapt</li>
                    <li>Choose the "Dyslexia Support" profile</li>
                    <li>Add custom instructions if needed (e.g., "Please use Tom Gates style with varied fonts and emphasis")</li>
                    <li>Click "Generate Adaptation"</li>
                </ol>
                <p><strong>Note:</strong> To export presentations with Tom Gates styling, you'll need to modify the export function.</p>
                <a href="/" class="btn">Return to Home</a>
            </div>
        </body>
        </html>
        """
        
    except Exception as e:
        return f"Error updating Tom Gates style: {str(e)}", 500

# IMPORTANT: Add this function to modify your existing create_adaptation_prompt function
def add_tom_gates_prompt(prompt, profile_id):
    """
    Add Tom Gates styling instructions to an existing prompt for dyslexia profile
    """
    if profile_id != 'dyslexia':
        return prompt
        
    tom_gates_instructions = """
    ## Tom Gates Style Guidelines
    The Tom Gates books by Liz Pichon use a distinctive visual style that's highly effective for dyslexic readers:
    
    1. VARY the presentation of text with different sizes, styles, and emphasis
    2. Use a mix of FONTS and text sizes to create visual interest and highlight important information
    3. Add *decorative elements* around key terms or concepts
    4. Break content into SMALL, visually distinct chunks
    5. Create visual hierarchy with text formatting rather than relying only on plain text
    6. Use informal, conversational language that feels approachable
    
    When adapting this content, please imagine how it would appear in a Tom Gates book. Indicate formatting changes using markdown:
    - **Bold text** for emphasis
    - ALL CAPS for important terms
    - *Asterisks* for decorative elements
    - Indicate [LARGE FONT] or [small font] where appropriate
    - Suggest simple doodles or visual elements with descriptions like [box around text] or [squiggly underline]
    """
    
    # Find the position to insert the Tom Gates instructions
    output_format_pos = prompt.find("## Output Format")
    if output_format_pos > 0:
        # Insert before Output Format section
        return prompt[:output_format_pos] + tom_gates_instructions + "\n\n" + prompt[output_format_pos:]
    else:
        # Append to the end if Output Format section not found
        return prompt + "\n\n" + tom_gates_instructions

init_db()
update_db_schema()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def process_pptx(file_path):
    presentation = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(presentation.slides):
        slide_content = {
            'slide_num': slide_num,
            'title': '',
            'text_elements': [],
            'has_images': any(shape.shape_type == 13 for shape in slide.shapes)
        }

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if not slide_content['title'] and (shape.name.startswith('Title') or slide_num == 0):
                    slide_content['title'] = shape.text
                else:
                    text_level = 0
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        text_level = shape.text_frame.paragraphs[0].level

                    slide_content['text_elements'].append({
                        'text': shape.text,
                        'level': text_level,
                        'shape_id': shape.shape_id if hasattr(shape, 'shape_id') else None
                    })
        slides.append(slide_content)

    return {'slides': slides}

def process_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        pages = []

        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            text = page.extract_text()
            paragraphs = [p for p in text.split('\n\n') if p.strip()]
            pages.append({
                'page_num': page_num,
                'paragraphs': paragraphs
            })

    return {'pages': pages}

def create_adaptation_prompt(document_data, profile_data, user_prompt=None):
    """
    Creates a prompt for Claude to adapt educational content for a specific learning profile.
    Enhanced with clearer instructions for content replacement.
    """
    profile_name = profile_data['name']
    instructions = profile_data['adaptation_instructions']

    prompt = f"""
    # Content Adaptation Task

    I need you to adapt educational content for a student with {profile_name}.

    ## Adaptation Instructions
    {instructions}
    """
    
    # Add Tom Gates style instructions for dyslexia profile
    if profile_data['id'] == 'dyslexia':
        prompt += """
    ## Tom Gates Style Guidelines
    The Tom Gates books by Liz Pichon use a distinctive visual style that's highly effective for dyslexic readers:
    
    1. VARY the presentation of text with different sizes, styles, and emphasis
    2. Use a mix of FONTS and text sizes to create visual interest and highlight important information
    3. Add *decorative elements* around key terms or concepts
    4. Break content into SMALL, visually distinct chunks
    5. Create visual hierarchy with text formatting rather than relying only on plain text
    6. Use informal, conversational language that feels approachable
    
    When adapting this content, please imagine how it would appear in a Tom Gates book. Indicate formatting changes using markdown:
    - **Bold text** for emphasis
    - ALL CAPS for important terms
    - *Asterisks* for decorative elements
    - Indicate [LARGE FONT] or [small font] where appropriate
    - Suggest simple doodles or visual elements with descriptions like [box around text] or [squiggly underline]
    """
    
    # Add user's custom instructions if provided
    if user_prompt and user_prompt.strip():
        prompt += f"""
    ## Additional Customization Instructions
    {user_prompt}
        """

    prompt += """
    ## Original Content
    """

    if 'slides' in document_data:
        slides = document_data['slides']
        for slide in slides:
            prompt += f"\n\nSLIDE {slide['slide_num'] + 1}:\n"
            prompt += f"TITLE: {slide.get('title', 'Untitled')}\n"
            prompt += "CONTENT:\n"
            for element in slide.get('text_elements', []):
                indent = "  " * element.get('level', 0)
                prompt += f"{indent}- {element['text']}\n"

    elif 'pages' in document_data:
        pages = document_data['pages']
        for page in pages:
            prompt += f"\n\nPAGE {page['page_num'] + 1}:\n"
            for paragraph in page.get('paragraphs', []):
                prompt += f"{paragraph}\n\n"

    prompt += """
    ## IMPORTANT OUTPUT INSTRUCTIONS
    
    1. You MUST maintain the EXACT same slide numbers and structure as the original
    2. Do NOT skip or combine any slides
    3. Each slide must have a title and all bullet points from the original
    4. Your adapted content should be a direct replacement that could be copied over the original
    
    ## Output Format

    For presentations, return content in this format precisely:
    SLIDE [number]:
    TITLE: [adapted title]
    CONTENT:
    - [adapted bullet point with Tom Gates styling if applicable]
    - [adapted bullet point with styling]

    For documents, return:
    PAGE [number]:
    [adapted paragraph 1 with Tom Gates styling if applicable]

    [adapted paragraph 2 with styling]

    Keep your response focused only on the adapted content. Don't explain your changes.
    Do not include any introductory or explanatory text about the adaptation process.
    """

    return prompt

def process_claude_response(response_text, original_document):
    """
    Process Claude's response and convert it to the appropriate format
    Improved with better error handling and debugging
    """
    print(f"Processing Claude response - length: {len(response_text)} characters")
    print(f"First 100 chars of response: {response_text[:100]}")
    
    try:
        if 'slides' in original_document:
            print("Processing presentation response")
            slides = []
            current_slide = None
            content_mode = False

            for line in response_text.split('\n'):
                line = line.strip()
                if not line:
                    continue

                if line.startswith('SLIDE '):
                    if current_slide is not None:
                        slides.append(current_slide)
                        print(f"Added slide {current_slide['slide_num']} with {len(current_slide['text_elements'])} elements")

                    try:
                        slide_num_text = line.replace('SLIDE ', '').replace(':', '')
                        print(f"Parsing slide number from: '{slide_num_text}'")
                        slide_num = int(slide_num_text) - 1
                        current_slide = {
                            'slide_num': slide_num,
                            'title': '',
                            'text_elements': []
                        }
                        print(f"Created new slide with number {slide_num}")
                    except ValueError as e:
                        print(f"Could not parse slide number from '{line}': {str(e)}")
                        current_slide = {
                            'slide_num': len(slides),
                            'title': '',
                            'text_elements': []
                        }
                        print(f"Using default slide number {len(slides)}")
                    content_mode = False

                elif current_slide is not None and line.startswith('TITLE:'):
                    current_slide['title'] = line.replace('TITLE:', '').strip()
                    print(f"Set slide title: '{current_slide['title']}'")
                    content_mode = False

                elif current_slide is not None and line.startswith('CONTENT:'):
                    print("Entering content mode")
                    content_mode = True

                elif current_slide is not None and content_mode:
                    original_line = line
                    line = line.lstrip()
                    leading_spaces = len(original_line) - len(line)
                    level = leading_spaces // 2

                    if line.startswith('- '):
                        text = line[2:].strip()
                    elif line.startswith('• '):
                        text = line[1:].strip()
                    else:
                        text = line

                    if text:  
                        current_slide['text_elements'].append({
                            'text': text,
                            'level': level
                        })
                        print(f"Added text element: '{text[:30]}...' at level {level}")

            if current_slide is not None:
                slides.append(current_slide)
                print(f"Added final slide {current_slide['slide_num']} with {len(current_slide['text_elements'])} elements")

            print(f"Processed {len(slides)} slides in total")
            return {'slides': slides}

        elif 'pages' in original_document:
            print("Processing document response")
            pages = []
            current_page = None
            current_paragraphs = []

            for line in response_text.split('\n'):
                if line.strip().startswith('PAGE '):
                    if current_page is not None:
                        current_page['paragraphs'] = current_paragraphs
                        pages.append(current_page)
                        print(f"Added page {current_page['page_num']} with {len(current_paragraphs)} paragraphs")
                        current_paragraphs = []

                    try:
                        page_num_text = line.strip().replace('PAGE ', '').replace(':', '')
                        print(f"Parsing page number from: '{page_num_text}'")
                        page_num = int(page_num_text) - 1
                        current_page = {'page_num': page_num}
                        print(f"Created new page with number {page_num}")
                    except ValueError as e:
                        print(f"Could not parse page number from '{line}': {str(e)}")
                        current_page = {'page_num': len(pages)}
                        print(f"Using default page number {len(pages)}")

                elif current_page is not None:
                    if line.strip():
                        current_paragraphs.append(line.strip())
                        if len(current_paragraphs) % 5 == 0:  # Log every 5 paragraphs
                            print(f"Added {len(current_paragraphs)} paragraphs to current page")

            if current_page is not None:
                current_page['paragraphs'] = current_paragraphs
                pages.append(current_page)
                print(f"Added final page {current_page['page_num']} with {len(current_paragraphs)} paragraphs")

            print(f"Processed {len(pages)} pages in total")
            return {'pages': pages}

        # Fallback if we can't determine the document type
        print("Could not determine document type, returning raw text")
        return {'text': response_text}
        
    except Exception as e:
        print(f"Error processing Claude response: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Return a minimal structure to avoid breaking the application
        if 'slides' in original_document:
            return {'slides': [{'slide_num': 0, 'title': 'Error Processing Response', 'text_elements': [{'text': f'Error: {str(e)}', 'level': 0}]}]}
        elif 'pages' in original_document:
            return {'pages': [{'page_num': 0, 'paragraphs': [f'Error processing response: {str(e)}', 'Please try again or contact support.']}]}
        else:
            return {'text': f'Error processing response: {str(e)}'}

@app.route('/', methods=['GET'])
def home():
    """
    Render the main page with customized interface for adaptive learning
    """
    html = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Adaptive Learning Platform</title>
        <style>
            body {
                font-family: Arial, sans-serif;
                margin: 0;
                padding: 20px;
                background-color: #f5f5f5;
            }
            .container {
                max-width: 1000px;
                margin: 0 auto;
                background-color: white;
                padding: 20px;
                border-radius: 8px;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }
            h1, h2, h3 {
                color: #333;
            }
            .form-group {
                margin-bottom: 20px;
            }
            label {
                display: block;
                margin-bottom: 5px;
                font-weight: bold;
            }
            select, textarea, input[type="file"] {
                width: 100%;
                padding: 10px;
                border: 1px solid #ddd;
                border-radius: 4px;
                box-sizing: border-box;
            }
            textarea {
                min-height: 150px;
                resize: vertical;
            }
            .btn {
                background-color: #4CAF50;
                color: white;
                border: none;
                padding: 10px 20px;
                border-radius: 4px;
                cursor: pointer;
                font-size: 16px;
            }
            .btn:hover {
                background-color: #45a049;
            }
            .card {
                border: 1px solid #ddd;
                border-radius: 4px;
                padding: 15px;
                margin-bottom: 15px;
            }
            .info-text {
                color: #666;
                font-style: italic;
                margin-top: 5px;
            }
            .adaptations-list {
                margin-top: 20px;
            }
            .adaptation-item {
                padding: 10px;
                border-bottom: 1px solid #eee;
            }
            .adaptation-item:hover {
                background-color: #f9f9f9;
            }
            .message {
                padding: 10px;
                margin: 10px 0;
                border-radius: 5px;
            }
            .success {
                background-color: #dff0d8;
                color: #3c763d;
                border: 1px solid #d6e9c6;
            }
            .error {
                background-color: #f2dede;
                color: #a94442;
                border: 1px solid #ebccd1;
            }
            #debug-info {
                margin-top: 20px;
                padding: 10px;
                background-color: #f0f0f0;
                border: 1px solid #ddd;
                border-radius: 4px;
                display: none;
            }
            .action-link {
                margin-right: 8px;
                text-decoration: none;
            }
            .view-link {
                color: #3498db;
            }
            .export-link {
                color: #27ae60;
            }
            .delete-link {
                color: #e74c3c;
                cursor: pointer;
            }
            
            /* Modal styles */
            .modal {
                display: none;
                position: fixed;
                z-index: 1000;
                left: 0;
                top: 0;
                width: 100%;
                height: 100%;
                background-color: rgba(0,0,0,0.5);
            }
            
            .modal-content {
                background-color: #fefefe;
                margin: 15% auto;
                padding: 20px;
                border-radius: 5px;
                width: 400px;
                box-shadow: 0 4px 8px rgba(0,0,0,0.1);
            }
            
            .modal-title {
                margin-top: 0;
                color: #333;
            }
            
            .modal-actions {
                text-align: right;
                margin-top: 20px;
            }
            
            .modal-btn {
                padding: 8px 16px;
                margin-left: 10px;
                border: none;
                border-radius: 4px;
                cursor: pointer;
            }
            
            .modal-cancel {
                background-color: #ccc;
            }
            
            .modal-confirm {
                background-color: #e74c3c;
                color: white;
            }
        </style>
    </head>
    <body>
        <div class="container">
            <h1>Adaptive Learning Platform</h1>
            
            <div id="messageArea"></div>
            
            <div class="card">
                <h2>Document Adaptation</h2>
                
                <div class="form-group">
                    <label for="document">Select Document:</label>
                    <select id="document" name="document">
                        <option value="">-- Select a document --</option>
                        <!-- Documents will be loaded here -->
                    </select>
                </div>
                
                <div class="form-group">
                    <label for="profile">Select Learning Profile:</label>
                    <select id="profile" name="profile">
                        <option value="">-- Select a profile --</option>
                        <!-- Profiles will be loaded here -->
                    </select>
                </div>
                
                <div class="form-group">
                    <label for="customPrompt">Custom Adaptation Instructions (Optional):</label>
                    <textarea id="customPrompt" name="customPrompt" placeholder="Enter any additional instructions to help Claude adapt the content more effectively..."></textarea>
                    <p class="info-text">Example: "Please use more visual metaphors," "Incorporate greater repetition of key terms," "Add questions at the end of each section to check understanding"</p>
                </div>
                
                <button id="adaptBtn" class="btn">Generate Adaptation</button>
            </div>
            
            <!-- File upload form -->
            <div class="card">
                <h2>Upload New Document</h2>
                
                <div class="form-group">
                    <label for="fileUpload">Select PDF or PowerPoint file:</label>
                    <input type="file" id="fileUpload" name="fileUpload" accept=".pdf,.pptx">
                </div>
                
                <button id="uploadBtn" class="btn">Upload Document</button>
            </div>
            
            <!-- Recent adaptations -->
            <div class="card">
                <h2>Recent Adaptations</h2>
                <div id="adaptationsList" class="adaptations-list">
                    <!-- Adaptations will be loaded here -->
                    <p>Loading recent adaptations...</p>
                </div>
            </div>
            
            <div id="debug-info"></div>
        </div>
        
        <!-- Delete confirmation modal -->
        <div id="deleteModal" class="modal">
            <div class="modal-content">
                <h3 class="modal-title">Confirm Deletion</h3>
                <p>Are you sure you want to delete this adaptation? This action cannot be undone.</p>
                <div class="modal-actions">
                    <button id="cancelDelete" class="modal-btn modal-cancel">Cancel</button>
                    <button id="confirmDelete" class="modal-btn modal-confirm">Delete</button>
                </div>
            </div>
        </div>

        <script>
            // Global variable to store the adaptation ID to delete
            let adaptationToDelete = null;
            
            // Debug function to log information
            function debugLog(message) {
                console.log(message);
                const debugArea = document.getElementById('debug-info');
                debugArea.style.display = 'block';
                debugArea.innerHTML += '<p>' + message + '</p>';
            }
            
            // Load documents and profiles on page load
            document.addEventListener('DOMContentLoaded', function() {
                debugLog('Page loaded, initializing...');
                
                try {
                    fetchDocuments();
                    fetchProfiles();
                    
                    // Error handling for recent adaptations
                    try {
                        fetchRecentAdaptations();
                    } catch (e) {
                        debugLog('Error loading recent adaptations: ' + e.message);
                        document.getElementById('adaptationsList').innerHTML = '<p>Could not load recent adaptations.</p>';
                    }
                    
                    // Set up event listeners
                    const adaptBtn = document.getElementById('adaptBtn');
                    if (adaptBtn) {
                        debugLog('Found adapt button, adding event listener');
                        adaptBtn.addEventListener('click', function() {
                            debugLog('Adapt button clicked');
                            try {
                                createAdaptation();
                            } catch (e) {
                                debugLog('Error in createAdaptation: ' + e.message);
                                showMessage('Error: ' + e.message, 'error');
                            }
                        });
                    } else {
                        debugLog('ERROR: Adapt button not found!');
                    }
                    
                    const uploadBtn = document.getElementById('uploadBtn');
                    if (uploadBtn) {
                        uploadBtn.addEventListener('click', function() {
                            try {
                                uploadDocument();
                            } catch (e) {
                                debugLog('Error in uploadDocument: ' + e.message);
                                showMessage('Error: ' + e.message, 'error');
                            }
                        });
                    }
                    
                    // Set up modal event listeners
                    document.getElementById('cancelDelete').addEventListener('click', hideDeleteModal);
                    document.getElementById('confirmDelete').addEventListener('click', performDelete);
                    
                } catch (e) {
                    debugLog('Initialization error: ' + e.message);
                }
            });
            
            function showMessage(message, type) {
                debugLog('Showing message: ' + message + ' (type: ' + type + ')');
                const messageArea = document.getElementById('messageArea');
                messageArea.innerHTML = `<div class="message ${type}">${message}</div>`;
                
                // Auto-hide success messages after 5 seconds
                if (type === 'success') {
                    setTimeout(() => {
                        messageArea.innerHTML = '';
                    }, 5000);
                }
            }
            
            // Fetch available documents
            function fetchDocuments() {
                debugLog('Fetching documents...');
                fetch('/documents')
                    .then(response => response.json())
                    .then(data => {
                        debugLog('Received ' + data.length + ' documents');
                        const select = document.getElementById('document');
                        select.innerHTML = '<option value="">-- Select a document --</option>';
                        
                        data.forEach(doc => {
                            const option = document.createElement('option');
                            option.value = doc.id;
                            option.textContent = doc.file_name;
                            select.appendChild(option);
                        });
                    })
                    .catch(error => {
                        debugLog('Error fetching documents: ' + error.message);
                        showMessage('Error loading documents. Please refresh the page.', 'error');
                    });
            }
            
            // Fetch available profiles
            function fetchProfiles() {
                debugLog('Fetching profiles...');
                fetch('/profiles')
                    .then(response => response.json())
                    .then(data => {
                        debugLog('Received ' + data.length + ' profiles');
                        const select = document.getElementById('profile');
                        select.innerHTML = '<option value="">-- Select a profile --</option>';
                        
                        data.forEach(profile => {
                            const option = document.createElement('option');
                            option.value = profile.id;
                            option.textContent = profile.name;
                            select.appendChild(option);
                        });
                    })
                    .catch(error => {
                        debugLog('Error fetching profiles: ' + error.message);
                        showMessage('Error loading profiles. Please refresh the page.', 'error');
                    });
            }
            
            // Fetch recent adaptations
            function fetchRecentAdaptations() {
                debugLog('Fetching recent adaptations...');
                fetch('/recent-adaptations')
                    .then(response => {
                        if (!response.ok) {
                            throw new Error('Server returned ' + response.status);
                        }
                        return response.json();
                    })
                    .then(data => {
                        debugLog('Received ' + data.length + ' adaptations');
                        const list = document.getElementById('adaptationsList');
                        
                        if (data.length === 0) {
                            list.innerHTML = '<p>No adaptations yet. Create one by selecting a document and profile above.</p>';
                            return;
                        }
                        
                        list.innerHTML = '';
                        data.forEach(adaptation => {
                            const item = document.createElement('div');
                            item.className = 'adaptation-item';
                            item.innerHTML = `
                                <strong>${adaptation.document_name}</strong> - 
                                <em>${adaptation.profile_name}</em>
                                <span style="color: #888; font-size: 0.9em;"> (${new Date(adaptation.creation_time).toLocaleString()})</span>
                                <br>
                                <a href="/adaptations/${adaptation.id}/view" class="action-link view-link">View</a> | 
                                <a href="/simple-export/${adaptation.id}" class="action-link export-link">Export</a> | 
                                <a class="action-link delete-link" data-id="${adaptation.id}">Delete</a>
                            `;
                            list.appendChild(item);
                        });
                        
                        // Add event listeners to delete links
                        document.querySelectorAll('.delete-link').forEach(link => {
                            link.addEventListener('click', function() {
                                const adaptationId = this.getAttribute('data-id');
                                showDeleteModal(adaptationId);
                            });
                        });
                    })
                    .catch(error => {
                        debugLog('Error fetching recent adaptations: ' + error.message);
                        document.getElementById('adaptationsList').innerHTML = 
                            '<p>Error loading recent adaptations: ' + error.message + '</p>';
                    });
            }
            
            // Show delete confirmation modal
            function showDeleteModal(adaptationId) {
                adaptationToDelete = adaptationId;
                document.getElementById('deleteModal').style.display = 'block';
            }
            
            // Hide delete confirmation modal
            function hideDeleteModal() {
                document.getElementById('deleteModal').style.display = 'none';
                adaptationToDelete = null;
            }
            
            // Perform the deletion
            function performDelete() {
                if (!adaptationToDelete) {
                    hideDeleteModal();
                    return;
                }
                
                debugLog('Deleting adaptation: ' + adaptationToDelete);
                
                fetch(`/adaptations/${adaptationToDelete}/delete`, {
                    method: 'POST'
                })
                .then(response => response.json())
                .then(data => {
                    hideDeleteModal();
                    
                    if (data.success) {
                        showMessage('Adaptation deleted successfully', 'success');
                        // Refresh the adaptations list
                        fetchRecentAdaptations();
                    } else {
                        showMessage('Error: ' + (data.error || 'Could not delete adaptation'), 'error');
                    }
                })
                .catch(error => {
                    hideDeleteModal();
                    debugLog('Error deleting adaptation: ' + error.message);
                    showMessage('Error deleting adaptation: ' + error.message, 'error');
                });
            }
            
            // Create adaptation with custom prompt
            function createAdaptation() {
                debugLog('Creating adaptation...');
                const documentId = document.getElementById('document').value;
                const profileId = document.getElementById('profile').value;
                const userPrompt = document.getElementById('customPrompt').value;
                
                debugLog('Selected document: ' + documentId);
                debugLog('Selected profile: ' + profileId);
                debugLog('Custom prompt length: ' + userPrompt.length + ' characters');
                
                if (!documentId || !profileId) {
                    debugLog('Missing required fields');
                    showMessage('Please select both a document and a profile', 'error');
                    return;
                }
                
                showMessage('Generating adaptation... This might take a minute.', 'success');
                
                const data = {
                    document_id: documentId,
                    profile_id: profileId,
                    user_prompt: userPrompt
                };
                
                debugLog('Sending request to /adapt endpoint...');
                
                fetch('/adapt', {
                    method: 'POST',
                    headers: {
                        'Content-Type': 'application/json'
                    },
                    body: JSON.stringify(data)
                })
                .then(response => {
                    debugLog('Received response with status: ' + response.status);
                    if (!response.ok) {
                        return response.json().then(errData => {
                            throw new Error(errData.error || 'Server error: ' + response.status);
                        });
                    }
                    return response.json();
                })
                .then(data => {
                    debugLog('Adaptation created with ID: ' + data.adaptation_id);
                    if (data.success) {
                        showMessage('Adaptation generated successfully!', 'success');
                        setTimeout(() => {
                            debugLog('Redirecting to view page...');
                            window.location.href = `/adaptations/${data.adaptation_id}/view`;
                        }, 1000);
                    } else {
                        debugLog('Error in success response: ' + JSON.stringify(data));
                        showMessage(`Error: ${data.error || 'Unknown error'}`, 'error');
                    }
                })
                .catch(error => {
                    debugLog('Fetch error: ' + error.message);
                    showMessage('Error creating adaptation: ' + error.message, 'error');
                });
            }
            
            // Upload new document
            function uploadDocument() {
                debugLog('Uploading document...');
                const fileInput = document.getElementById('fileUpload');
                if (!fileInput.files.length) {
                    showMessage('Please select a file to upload', 'error');
                    return;
                }
                
                const file = fileInput.files[0];
                debugLog('Selected file: ' + file.name + ' (' + file.type + ')');
                
                const formData = new FormData();
                formData.append('file', file);
                
                showMessage('Uploading document...', 'success');
                
                fetch('/upload', {
                    method: 'POST',
                    body: formData
                })
                .then(response => response.json())
                .then(data => {
                    if (data.success) {
                        debugLog('Upload successful: ' + data.document_id);
                        showMessage('Document uploaded successfully!', 'success');
                        
                        // Reset the file input
                        fileInput.value = '';
                        
                        // Refresh the document list
                        fetchDocuments();
                    } else {
                        debugLog('Upload error: ' + data.error);
                        showMessage(`Error: ${data.error || 'Unknown error'}`, 'error');
                    }
                })
                .catch(error => {
                    debugLog('Upload fetch error: ' + error.message);
                    showMessage('Error uploading document. Please try again.', 'error');
                });
            }
        </script>
    </body>
    </html>
    """
    return html

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and allowed_file(file.filename):
        document_id = str(uuid.uuid4())
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Create upload directory if it doesn't exist
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        
        file.save(file_path)

        try:
            file_type = filename.rsplit('.', 1)[1].lower()

            if file_type == 'pptx':
                content = process_pptx(file_path)
                doc_type = 'presentation'
            elif file_type == 'pdf':
                content = process_pdf(file_path)
                doc_type = 'document'

            conn = get_db_connection()
            conn.execute(
                'INSERT INTO documents (id, file_name, file_type, content) VALUES (?, ?, ?, ?)',
                (document_id, filename, doc_type, json.dumps(content))
            )
            conn.commit()
            conn.close()

            return jsonify({
                'success': True,
                'document_id': document_id,
                'file_name': filename,
                'document_type': doc_type
            }), 200

        except Exception as e:
            return jsonify({'error': str(e)}), 500

    return jsonify({'error': 'File type not allowed'}), 400

@app.route('/documents', methods=['GET'])
def get_documents():
    conn = get_db_connection()
    documents = conn.execute('SELECT id, file_name, file_type, upload_time FROM documents ORDER BY upload_time DESC').fetchall()
    conn.close()
    result = []
    for doc in documents:
        result.append({
            'id': doc['id'],
            'file_name': doc['file_name'],
            'document_type': doc['file_type'],
            'upload_time': doc['upload_time']
        })

    return jsonify(result)

@app.route('/profiles', methods=['GET'])
def get_profiles():
    conn = get_db_connection()
    profiles = conn.execute('SELECT id, name, description FROM profiles').fetchall()
    conn.close()

    result = []
    for profile in profiles:
        result.append({
            'id': profile['id'],
            'name': profile['name'],
            'description': profile['description']
        })

    return jsonify(result)

@app.route('/recent-adaptations', methods=['GET'])
def get_recent_adaptations():
    """
    Returns recent adaptations for display on the home page
    """
    try:
        conn = get_db_connection()
        
        adaptations = conn.execute('''
            SELECT a.id, a.document_id, a.profile_id, a.creation_time,
                   d.file_name as document_name,
                   p.name as profile_name
            FROM adaptations a
            JOIN documents d ON a.document_id = d.id
            JOIN profiles p ON a.profile_id = p.id
            ORDER BY a.creation_time DESC
            LIMIT 10
        ''').fetchall()
        
        conn.close()
        
        result = []
        for adaptation in adaptations:
            result.append({
                'id': adaptation['id'],
                'document_id': adaptation['document_id'],
                'profile_id': adaptation['profile_id'],
                'document_name': adaptation['document_name'],
                'profile_name': adaptation['profile_name'],
                'creation_time': adaptation['creation_time']
            })
        
        return jsonify(result)
    except Exception as e:
        print(f"Error fetching recent adaptations: {str(e)}")
        return jsonify([])

@app.route('/adaptations/<adaptation_id>/delete', methods=['POST'])
def delete_adaptation(adaptation_id):
    """
    Delete an adaptation by ID
    """
    try:
        print(f"Attempting to delete adaptation {adaptation_id}")
        conn = get_db_connection()
        
        # First check if the adaptation exists
        adaptation = conn.execute('SELECT id FROM adaptations WHERE id = ?', (adaptation_id,)).fetchone()
        
        if not adaptation:
            conn.close()
            print(f"Adaptation {adaptation_id} not found")
            return jsonify({'success': False, 'error': 'Adaptation not found'}), 404
        
        # Delete the adaptation
        conn.execute('DELETE FROM adaptations WHERE id = ?', (adaptation_id,))
        conn.commit()
        conn.close()
        print(f"Successfully deleted adaptation {adaptation_id}")
        
        return jsonify({
            'success': True,
            'message': 'Adaptation deleted successfully'
        })
        
    except Exception as e:
        print(f"Error deleting adaptation: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/adapt', methods=['POST'])
def adapt_content():
    """
    Create an adaptation of a document using Claude API
    Using the model determined to be working during initialization
    """
    try:
        data = request.json
        print(f"Received adaptation request: {data}")

        if not data:
            return jsonify({'success': False, 'error': 'No data provided in request'}), 400

        if 'document_id' not in data:
            return jsonify({'success': False, 'error': 'Missing document_id parameter'}), 400
            
        if 'profile_id' not in data:
            return jsonify({'success': False, 'error': 'Missing profile_id parameter'}), 400

        document_id = data['document_id']
        profile_id = data['profile_id']
        user_prompt = data.get('user_prompt', '')

        print(f"Processing adaptation - document_id: {document_id}, profile_id: {profile_id}")
        print(f"User prompt length: {len(user_prompt)} characters")

        # Database connection
        try:
            conn = get_db_connection()
            print("Database connection established")
        except Exception as db_err:
            print(f"Database connection error: {str(db_err)}")
            return jsonify({'success': False, 'error': f'Database connection failed: {str(db_err)}'}), 500
        
        # Check if adaptation already exists (skip for custom prompts)
        if not user_prompt:
            try:
                existing = conn.execute(
                    'SELECT id FROM adaptations WHERE document_id = ? AND profile_id = ? AND (user_prompt IS NULL OR user_prompt = "")',
                    (document_id, profile_id)
                ).fetchone()

                if existing:
                    print(f"Existing adaptation found: {existing['id']}")
                    return jsonify({
                        'success': True,
                        'adaptation_id': existing['id'],
                        'message': 'Adaptation already exists'
                    }), 200
            except Exception as query_err:
                print(f"Error checking for existing adaptation: {str(query_err)}")
                return jsonify({'success': False, 'error': f'Database query error: {str(query_err)}'}), 500

        # Fetch document
        try:
            document = conn.execute(
                'SELECT content, file_type, file_name FROM documents WHERE id = ?',
                (document_id,)
            ).fetchone()

            if not document:
                print(f"Document not found: {document_id}")
                return jsonify({'success': False, 'error': 'Document not found'}), 404

            document_content = json.loads(document['content'])
            file_name = document['file_name']
            print(f"Document found: {file_name}")
        except json.JSONDecodeError as json_err:
            print(f"Error parsing document content: {str(json_err)}")
            return jsonify({'success': False, 'error': f'Invalid document content format: {str(json_err)}'}), 500
        except Exception as doc_err:
            print(f"Error fetching document: {str(doc_err)}")
            return jsonify({'success': False, 'error': f'Error retrieving document: {str(doc_err)}'}), 500

        # Fetch profile
        try:
            profile = conn.execute(
                'SELECT * FROM profiles WHERE id = ?',
                (profile_id,)
            ).fetchone()

            if not profile:
                print(f"Profile not found: {profile_id}")
                return jsonify({'success': False, 'error': 'Profile not found'}), 404
                
            print(f"Profile found: {profile['name']}")
        except Exception as profile_err:
            print(f"Error fetching profile: {str(profile_err)}")
            return jsonify({'success': False, 'error': f'Error retrieving profile: {str(profile_err)}'}), 500

        # Create adaptation prompt
        try:
            prompt = create_adaptation_prompt(document_content, dict(profile), user_prompt)
            prompt_length = len(prompt)
            print(f"Created adaptation prompt (length: {prompt_length} characters)")
        except Exception as prompt_err:
            print(f"Error creating prompt: {str(prompt_err)}")
            return jsonify({'success': False, 'error': f'Error creating adaptation prompt: {str(prompt_err)}'}), 500
        
        # Call Claude API
        try:
            print("Sending request to Claude API...")
            
            # Use the model that was determined to work during initialization
            if not WORKING_MODEL:
                return jsonify({'success': False, 'error': 'No working Claude model found during initialization'}), 500
                
            print(f"Using Claude model: {WORKING_MODEL}")
            
            # Use the model determined during initialization
            response = claude.messages.create(
                model=WORKING_MODEL,
                max_tokens=4000,
                temperature=0.2,
                system="You are an educational content adaptation expert specializing in creating materials for students with different learning needs.",
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            print("Claude API response received successfully")
            
        except Exception as claude_err:
            print(f"Error calling Claude API: {str(claude_err)}")
            error_message = str(claude_err)
            
            # Try a fallback model if the model error occurs
            if "404" in error_message and "model" in error_message:
                try:
                    print("Attempting with fallback model: claude-3-opus-20240229")
                    response = claude.messages.create(
                        model="claude-3-opus-20240229",
                        max_tokens=4000,
                        temperature=0.2,
                        system="You are an educational content adaptation expert specializing in creating materials for students with different learning needs.",
                        messages=[
                            {"role": "user", "content": prompt}
                        ]
                    )
                    print("Fallback model worked successfully")
                except Exception as fallback_err:
                    print(f"Fallback model also failed: {str(fallback_err)}")
                    return jsonify({
                        'success': False, 
                        'error': f'All Claude models failed. Error: {error_message}. Please check your Anthropic API key and model availability.'
                    }), 500
            elif "401" in error_message:
                return jsonify({'success': False, 'error': 'Claude API authentication failed - invalid API key'}), 500
            elif "429" in error_message:
                return jsonify({'success': False, 'error': 'Claude API rate limit exceeded - please try again later'}), 500
            elif "500" in error_message:
                return jsonify({'success': False, 'error': 'Claude API server error - please try again later'}), 500
            else:
                return jsonify({'success': False, 'error': f'Claude API error: {error_message}'}), 500

        # Process Claude's response
        try:
            adapted_content = process_claude_response(response.content[0].text, document_content)
            print("Claude response processed successfully")
        except Exception as process_err:
            print(f"Error processing Claude response: {str(process_err)}")
            return jsonify({'success': False, 'error': f'Error processing adaptation response: {str(process_err)}'}), 500

        # Save adaptation to database
        try:
            adaptation_id = str(uuid.uuid4())
            conn.execute(
                'INSERT INTO adaptations (id, document_id, profile_id, adapted_content, user_prompt) VALUES (?, ?, ?, ?, ?)',
                (adaptation_id, document_id, profile_id, json.dumps(adapted_content), user_prompt)
            )
            conn.commit()
            print(f"Adaptation saved to database with ID: {adaptation_id}")
        except Exception as save_err:
            print(f"Error saving adaptation: {str(save_err)}")
            return jsonify({'success': False, 'error': f'Error saving adaptation: {str(save_err)}'}), 500
        finally:
            conn.close()

        return jsonify({
            'success': True,
            'adaptation_id': adaptation_id
        }), 200

    except Exception as e:
        print(f"Unexpected error during adaptation: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': str(e)}), 500
    
@app.route('/adaptations/<adaptation_id>', methods=['GET'])
def get_adaptation(adaptation_id):
    conn = get_db_connection()

    adaptation = conn.execute('''
        SELECT a.id, a.document_id, a.profile_id, a.adapted_content, a.user_prompt, a.creation_time,
               d.content as original_content, d.file_type as document_type, d.file_name,
               p.name as profile_name, p.description as profile_description
        FROM adaptations a
        JOIN documents d ON a.document_id = d.id
        JOIN profiles p ON a.profile_id = p.id
        WHERE a.id = ?
    ''', (adaptation_id,)).fetchone()

    conn.close()

    if not adaptation:
        return jsonify({'error': 'Adaptation not found'}), 404

    return jsonify({
        'id': adaptation['id'],
        'document_id': adaptation['document_id'],
        'file_name': adaptation['file_name'],
        'profile_id': adaptation['profile_id'],
        'profile_name': adaptation['profile_name'],
        'profile_description': adaptation['profile_description'],
        'document_type': adaptation['document_type'],
        'user_prompt': adaptation['user_prompt'] if adaptation['user_prompt'] else '',
        'creation_time': adaptation['creation_time'],
        'original_content': json.loads(adaptation['original_content']),
        'adapted_content': json.loads(adaptation['adapted_content'])
    })

# New direct download export function that works
@app.route('/simple-export/<adaptation_id>', methods=['GET'])
def simple_export(adaptation_id):
    """
    Enhanced export that creates a styled document with Tom Gates formatting
    and serves it directly as a download.
    """
    try:
        conn = get_db_connection()
        
        # Get adaptation details
        adaptation = conn.execute('''
            SELECT a.adapted_content, a.profile_id, 
                   d.file_name, d.file_type
            FROM adaptations a
            JOIN documents d ON a.document_id = d.id
            WHERE a.id = ?
        ''', (adaptation_id,)).fetchone()
        
        conn.close()
        
        if not adaptation:
            return "Adaptation not found", 404
            
        # Extract key info
        file_name = adaptation['file_name']
        document_type = adaptation['file_type']
        base_name = os.path.splitext(file_name)[0]
        is_dyslexia = adaptation['profile_id'] == 'dyslexia'
        adapted_content = json.loads(adaptation['adapted_content'])
        
        # Create temp directory for output
        export_folder = os.path.join(app.config['UPLOAD_FOLDER'], 'exports')
        os.makedirs(export_folder, exist_ok=True)
        
        # Format output filename
        if document_type == 'presentation':
            output_filename = f"tomgates_{base_name}_{adaptation_id[:8]}.pptx"
            output_path = os.path.join(export_folder, output_filename)
            
            # Create a Tom Gates styled presentation
            create_tom_gates_pptx(adapted_content, output_path)
            
            # Serve the file
            return send_file(output_path, as_attachment=True)
            
        elif document_type == 'document':
            output_filename = f"tomgates_{base_name}_{adaptation_id[:8]}.pdf"
            output_path = os.path.join(export_folder, output_filename)
            
            # Create a Tom Gates styled PDF
            create_tom_gates_pdf(adapted_content, output_path)
            
            # Serve the file
            return send_file(output_path, as_attachment=True)
        
        else:
            return "Unsupported document type", 400
            
    except Exception as e:
        print(f"Enhanced export error: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Return a simple error file if something goes wrong
        response = make_response(f"Export Error: {str(e)}")
        response.headers["Content-Disposition"] = "attachment; filename=export_error.txt"
        response.headers["Content-Type"] = "text/plain"
        return response

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    export_folder = os.path.join(app.config['UPLOAD_FOLDER'], 'exports')
    return send_from_directory(export_folder, filename, as_attachment=True)



@app.route('/adaptations/<adaptation_id>/view', methods=['GET'])
def view_adaptation(adaptation_id):
    """
    Displays a side-by-side view of original and adapted content
    with support for markdown formatting
    """
    conn = get_db_connection()

    adaptation = conn.execute('''
        SELECT a.id, a.document_id, a.profile_id, a.adapted_content, a.user_prompt,
               d.content as original_content, d.file_type as document_type, d.file_name,
               p.name as profile_name, p.description as profile_description
        FROM adaptations a
        JOIN documents d ON a.document_id = d.id
        JOIN profiles p ON a.profile_id = p.id
        WHERE a.id = ?
    ''', (adaptation_id,)).fetchone()

    conn.close()

    if not adaptation:
        return "Adaptation not found", 404

    adapted_content = json.loads(adaptation['adapted_content'])
    original_content = json.loads(adaptation['original_content'])
    document_type = adaptation['document_type']
    profile_name = adaptation['profile_name']
    file_name = adaptation['file_name']
    user_prompt = adaptation['user_prompt'] if adaptation['user_prompt'] else "No custom instructions provided."

    if document_type == 'presentation':
        export_text = "Export to PowerPoint"
    else:
        export_text = "Export to PDF"

    # Convert Python dictionaries to JSON strings for use in JavaScript
    original_content_json = json.dumps(original_content)
    adapted_content_json = json.dumps(adapted_content)

    # Using triple quotes to avoid issues with JavaScript comments
    html = f"""
    <html>
        <head>
        <title>Adaptation Viewer - {file_name}</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 0; padding: 20px; }}
            h1, h2 {{ color: #333; }}
            .container {{ display: flex; }}
            .column {{ flex: 1; padding: 10px; margin: 5px; border: 1px solid #ddd; border-radius: 5px; overflow-y: auto; max-height: 600px; }}
            .slide-title {{ font-size: 18px; font-weight: bold; margin-top: 15px; }}
            .bullet {{ margin-left: 20px; }}
            .navigation {{ margin: 15px 0; }}
            .top-bar {{ display: flex; justify-content: space-between; align-items: center; margin-bottom: 20px; }}
            button {{ padding: 5px 10px; margin-right: 5px; }}
            .export-btn {{ background-color: #4CAF50; color: white; border: none; border-radius: 4px; cursor: pointer; padding: 8px 15px; }}
            .home-btn {{ background-color: #3498db; color: white; border: none; border-radius: 4px; cursor: pointer; padding: 8px 15px; }}
            .delete-btn {{ background-color: #e74c3c; color: white; border: none; border-radius: 4px; cursor: pointer; padding: 8px 15px; }}
            .message {{ padding: 10px; margin: 10px 0; border-radius: 5px; display: none; }}
            .success {{ background-color: #dff0d8; color: #3c763d; border: 1px solid #d6e9c6; }}
            .error {{ background-color: #f2dede; color: #a94442; border: 1px solid #ebccd1; }}
            .info-box {{ background-color: #f8f9fa; border: 1px solid #ddd; border-radius: 5px; padding: 10px; margin-bottom: 20px; }}
            .custom-prompt {{ background-color: #e6f7ff; border: 1px solid #91d5ff; padding: 10px; margin: 10px 0; border-radius: 4px; }}

            /* Modal styles */
            .modal {{
                display: none;
                position: fixed;
                z-index: 1000;
                left: 0;
                top: 0;
                width: 100%;
                height: 100%;
                background-color: rgba(0,0,0,0.5);
            }}

            .modal-content {{
                background-color: #fefefe;
                margin: 15% auto;
                padding: 20px;
                border-radius: 5px;
                width: 400px;
                box-shadow: 0 4px 8px rgba(0,0,0,0.1);
            }}

            .modal-title {{
                margin-top: 0;
                color: #333;
            }}

            .modal-actions {{
                text-align: right;
                margin-top: 20px;
            }}

            .modal-btn {{
                padding: 8px 16px;
                margin-left: 10px;
                border: none;
                border-radius: 4px;
                cursor: pointer;
            }}

            .modal-cancel {{
                background-color: #ccc;
            }}

            .modal-confirm {{
                background-color: #e74c3c;
                color: white;
            }}

            /* Added styling for Tom Gates formatted content */
            .large-font {{ font-size: 18px; }}
            .small-font {{ font-size: 10px; }}
            .tom-gates-content {{ line-height: 1.6; }}
            .tom-gates-title {{ font-family: "Comic Sans MS", cursive, sans-serif; }}
            .decoration-box {{ border: 2px solid #333; padding: 5px; display: inline-block; }}
            .squiggly {{ text-decoration: underline wavy; }}
        </style>
    </head>
    <body>
        <div class="top-bar">
            <h1>Adaptation Viewer</h1>
            <div>
                <button id="exportBtn" class="export-btn">{export_text}</button>
                <button id="deleteBtn" class="delete-btn">Delete</button>
                <button id="homeBtn" class="home-btn">Home</button>
            </div>
        </div>

        <div id="messageArea" class="message"></div>

        <div class="info-box">
            <h2>{file_name}</h2>
            <p><strong>Profile:</strong> {profile_name}</p>
            <div class="custom-prompt">
                <strong>Custom Instructions:</strong>
                <p>{user_prompt}</p>
            </div>
        </div>

        <div class="container">
            <div class="column">
                <h2>Original Content</h2>
                <div id="originalContent">
                    <!-- Content will be loaded here -->
                </div>
            </div>
            <div class="column">
                <h2>Adapted Content</h2>
                <div id="adaptedContent" class="tom-gates-content">
                    <!-- Content will be loaded here -->
                </div>
            </div>
        </div>

        <!-- Delete confirmation modal -->
        <div id="deleteModal" class="modal">
            <div class="modal-content">
                <h3 class="modal-title">Confirm Deletion</h3>
                <p>Are you sure you want to delete this adaptation? This action cannot be undone.</p>
                <div class="modal-actions">
                    <button id="cancelDelete" class="modal-btn modal-cancel">Cancel</button>
                    <button id="confirmDelete" class="modal-btn modal-confirm">Delete</button>
                </div>
            </div>
        </div>

        <script>
            // Format presentation slides
            function renderPresentationView() {{
                const originalContent = {original_content_json};
                const adaptedContent = {adapted_content_json};

                let originalHtml = '';
                let adaptedHtml = '';

                // Original content
                originalContent.slides.forEach(slide => {{
                    originalHtml += `<div class="slide">`;
                    originalHtml += `<div class="slide-title">Slide ${{slide.slide_num + 1}}: ${{slide.title || 'Untitled'}}</div>`;

                    slide.text_elements.forEach(element => {{
                        const indent = '&nbsp;'.repeat(element.level * 4);
                        originalHtml += `<div class="bullet">${{indent}}• ${{element.text}}</div>`;
                    }});

                    originalHtml += `</div><hr>`;
                }});

                // Adapted content
                adaptedContent.slides.forEach(slide => {{
                    adaptedHtml += `<div class="slide">`;
                    adaptedHtml += `<div class="slide-title tom-gates-title">Slide ${{slide.slide_num + 1}}: ${{slide.title || 'Untitled'}}</div>`;

                    slide.text_elements.forEach(element => {{
                        let text = element.text;

                        // Apply Tom Gates style formatting if present
                        text = text.replace(/\\*\\*(.*?)\\*\\*/g, '<strong>$1</strong>');
                        text = text.replace(/\\*(.*?)\\*/g, '<em>$1</em>');
                        text = text.replace(/\\[LARGE FONT\\](.*?)(?=\\[|$)/g, '<span class="large-font">$1</span>');
                        text = text.replace(/\\[small font\\](.*?)(?=\\[|$)/g, '<span class="small-font">$1</span>');

                        // Convert all caps to styled text (if 4+ characters)
                        text = text.replace(/([A-Z]{{4,}})/g, '<strong>$1</strong>');

                        const indent = '&nbsp;'.repeat(element.level * 4);
                        adaptedHtml += `<div class="bullet">${{indent}}• ${{text}}</div>`;
                    }});

                    adaptedHtml += `</div><hr>`;
                }});

                document.getElementById('originalContent').innerHTML = originalHtml;
                document.getElementById('adaptedContent').innerHTML = adaptedHtml;
            }}

            // Format document pages
            function renderDocumentView() {{
                const originalContent = {original_content_json};
                const adaptedContent = {adapted_content_json};

                let originalHtml = '';
                let adaptedHtml = '';

                // Original content
                originalContent.pages.forEach(page => {{
                    originalHtml += `<div class="page">`;
                    originalHtml += `<h3>Page ${{page.page_num + 1}}</h3>`;

                    page.paragraphs.forEach(paragraph => {{
                        originalHtml += `<p>${{paragraph}}</p>`;
                    }});

                    originalHtml += `</div><hr>`;
                }});

                // Adapted content
                adaptedContent.pages.forEach(page => {{
                    adaptedHtml += `<div class="page">`;
                    adaptedHtml += `<h3 class="tom-gates-title">Page ${{page.page_num + 1}}</h3>`;

                    page.paragraphs.forEach(paragraph => {{
                        let text = paragraph;

                        // Apply Tom Gates style formatting if present
                        text = text.replace(/\\*\\*(.*?)\\*\\*/g, '<strong>$1</strong>');
                        text = text.replace(/\\*(.*?)\\*/g, '<em>$1</em>');
                        text = text.replace(/\\[LARGE FONT\\](.*?)(?=\\[|$)/g, '<span class="large-font">$1</span>');
                        text = text.replace(/\\[small font\\](.*?)(?=\\[|$)/g, '<span class="small-font">$1</span>');

                        // Convert all caps to styled text (if 4+ characters)
                        text = text.replace(/([A-Z]{{4,}})/g, '<strong>$1</strong>');

                        adaptedHtml += `<p>${{text}}</p>`;
                    }});

                    adaptedHtml += `</div><hr>`;
                }});

                document.getElementById('originalContent').innerHTML = originalHtml;
                document.getElementById('adaptedContent').innerHTML = adaptedHtml;
            }}

            // Initialize the view based on document type
            function initializeView() {{
                const documentType = "{document_type}";

                if (documentType === 'presentation') {{
                    renderPresentationView();
                }} else {{
                    renderDocumentView();
                }}

                // Set up event listeners
                document.getElementById('exportBtn').addEventListener('click', exportContent);
                document.getElementById('homeBtn').addEventListener('click', goHome);
                document.getElementById('deleteBtn').addEventListener('click', showDeleteModal);
                document.getElementById('cancelDelete').addEventListener('click', hideDeleteModal);
                document.getElementById('confirmDelete').addEventListener('click', deleteAdaptation);
            }}

            // Go to home page
            function goHome() {{
                window.location.href = '/';
            }}

            // Show delete confirmation modal
            function showDeleteModal() {{
                document.getElementById('deleteModal').style.display = 'block';
            }}

            // Hide delete confirmation modal
            function hideDeleteModal() {{
                document.getElementById('deleteModal').style.display = 'none';
            }}

            // Delete adaptation
            function deleteAdaptation() {{
                fetch('/adaptations/{adaptation_id}/delete', {{
                    method: 'POST'
                }})
                .then(response => response.json())
                .then(data => {{
                    hideDeleteModal();

                    if (data.success) {{
                        // Show success message
                        const messageArea = document.getElementById('messageArea');
                        messageArea.className = 'message success';
                        messageArea.style.display = 'block';
                        messageArea.innerText = 'Adaptation deleted successfully';

                        // Redirect to home after short delay
                        setTimeout(() => {{
                            window.location.href = '/';
                        }}, 1500);
                    }} else {{
                        // Show error message
                        const messageArea = document.getElementById('messageArea');
                        messageArea.className = 'message error';
                        messageArea.style.display = 'block';
                        messageArea.innerText = 'Error: ' + (data.error || 'Could not delete adaptation');
                    }}
                }})
                .catch(error => {{
                    hideDeleteModal();

                    // Show error message
                    const messageArea = document.getElementById('messageArea');
                    messageArea.className = 'message error';
                    messageArea.style.display = 'block';
                    messageArea.innerText = 'Error: ' + error.message;
                }});
            }}

            // Initialize the page
            document.addEventListener('DOMContentLoaded', initializeView);
        </script>
    </body>
    </html>
    """

    return html
# Error handlers
@app.errorhandler(404)
def page_not_found(e):
    return """
    <html>
    <head>
        <title>Page Not Found</title>
        <style>
            body { font-family: Arial, sans-serif; margin: 0; padding: 20px; text-align: center; }
            .container { max-width: 600px; margin: 50px auto; }
            h1 { color: #e74c3c; }
            .btn { background-color: #3498db; color: white; padding: 10px 15px; text-decoration: none; display: inline-block; border-radius: 4px; margin-top: 20px; }
        </style>
    </head>
    <body>
        <div class="container">
            <h1>404 - Page Not Found</h1>
            <p>Sorry, the page you were looking for doesn't exist.</p>
            <a href="/" class="btn">Go to Home Page</a>
        </div>
    </body>
    </html>
    """, 404

@app.errorhandler(500)
def server_error(e):
    return """
    <html>
    <head>
        <title>Server Error</title>
        <style>
            body { font-family: Arial, sans-serif; margin: 0; padding: 20px; text-align: center; }
            .container { max-width: 600px; margin: 50px auto; }
            h1 { color: #e74c3c; }
            .error-details { background-color: #f9f9f9; padding: 15px; text-align: left; border-radius: 5px; }
            .btn { background-color: #3498db; color: white; padding: 10px 15px; text-decoration: none; display: inline-block; border-radius: 4px; margin-top: 20px; }
        </style>
    </head>
    <body>
        <div class="container">
            <h1>500 - Server Error</h1>
            <p>Something went wrong on our end. Please try again later.</p>
            <div class="error-details">
                <p><strong>Error details:</strong> {}</p>
            </div>
            <a href="/" class="btn">Go to Home Page</a>
        </div>
    </body>
    </html>
    """.format(str(e)), 500

# Main application entry point
if __name__ == '__main__':
    # Create static/uploads directory if it doesn't exist
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    
    # Create static/uploads/exports directory for exports
    os.makedirs(os.path.join(app.config['UPLOAD_FOLDER'], 'exports'), exist_ok=True)
    
    # Initialize database
    print("Initializing database...")
    init_db()
    
    # Test Claude API connectivity
    print("Testing Claude API connection...")
    try:
        if claude:
            print("Claude API initialized")
            if WORKING_MODEL:
                print(f"Using Claude model: {WORKING_MODEL}")
            else:
                print("No working Claude model found")
    except Exception as e:
        print(f"Claude API error: {str(e)}")
    
    # Start the Flask application
    port = int(os.environ.get("PORT", 5000))
    debug_mode = os.environ.get("FLASK_DEBUG", "true").lower() == "true"
    
    print(f"Starting application on port {port}, debug mode: {debug_mode}")
    app.run(host='0.0.0.0', port=port, debug=debug_mode)