"""
Educational Content Service

Generates educational content like lesson plans, assessments, etc.
"""
import os
import re
from typing import Dict, Any, List, Optional
from .base_service import BaseService
import anthropic


class EducationalContentService(BaseService):
    """Service for generating educational content"""
    
    def _initialize(self):
        """Initialize educational content service resources"""
        self.api_key = self.config.get('anthropic_api_key')
        if self.api_key:
            self.client = anthropic.Anthropic(api_key=self.api_key)
        else:
            self.client = None
    
    def generate_lesson_plan(self, topic: str, grade_level: str, duration: int,
                           learning_objectives: List[str], profile: str = 'default') -> Dict[str, Any]:
        """
        Generate a complete lesson plan
        
        Args:
            topic: Lesson topic
            grade_level: Target grade level
            duration: Duration in minutes
            learning_objectives: List of learning objectives
            profile: Learning profile for adaptations
            
        Returns:
            Dict containing lesson plan components
        """
        if not self.client:
            raise Exception("No API client available for content generation")
        
        # Create prompt
        prompt = f"""Create a detailed lesson plan for teaching '{topic}' to {grade_level} students.

Duration: {duration} minutes
Learning Profile: {profile}

Learning Objectives:
{chr(10).join(f'- {obj}' for obj in learning_objectives)}

Please include:
1. Lesson Overview
2. Materials Needed
3. Introduction/Hook (5 minutes)
4. Main Content Delivery (specific to {profile} learners)
5. Guided Practice Activities
6. Independent Practice
7. Assessment Strategies
8. Differentiation Strategies
9. Closure/Summary
10. Extension Activities

Format the response as a structured lesson plan."""

        try:
            response = self.client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=3000,
                temperature=0.7,
                messages=[{"role": "user", "content": prompt}]
            )
            
            lesson_plan_text = response.content[0].text
            
            # Parse the response into structured format
            lesson_plan = self._parse_lesson_plan(lesson_plan_text)
            lesson_plan['topic'] = topic
            lesson_plan['grade_level'] = grade_level
            lesson_plan['duration'] = duration
            lesson_plan['learning_objectives'] = learning_objectives
            lesson_plan['profile'] = profile
            
            return lesson_plan
            
        except Exception as e:
            self.logger.error(f"Error generating lesson plan: {str(e)}")
            raise
    
    def _parse_lesson_plan(self, text: str) -> Dict[str, Any]:
        """Parse lesson plan text into structured format"""
        sections = {
            'overview': '',
            'materials': [],
            'introduction': '',
            'main_content': '',
            'guided_practice': '',
            'independent_practice': '',
            'assessment': '',
            'differentiation': '',
            'closure': '',
            'extensions': ''
        }
        
        # Simple parsing - look for section headers
        current_section = None
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for section headers
            lower_line = line.lower()
            if 'overview' in lower_line:
                current_section = 'overview'
            elif 'materials' in lower_line:
                current_section = 'materials'
            elif 'introduction' in lower_line or 'hook' in lower_line:
                current_section = 'introduction'
            elif 'main content' in lower_line:
                current_section = 'main_content'
            elif 'guided practice' in lower_line:
                current_section = 'guided_practice'
            elif 'independent practice' in lower_line:
                current_section = 'independent_practice'
            elif 'assessment' in lower_line:
                current_section = 'assessment'
            elif 'differentiation' in lower_line:
                current_section = 'differentiation'
            elif 'closure' in lower_line or 'summary' in lower_line:
                current_section = 'closure'
            elif 'extension' in lower_line:
                current_section = 'extensions'
            elif current_section:
                # Add content to current section
                if current_section == 'materials':
                    if line.startswith('-') or line.startswith('•'):
                        sections['materials'].append(line[1:].strip())
                else:
                    sections[current_section] += line + '\n'
        
        # Clean up sections
        for key in sections:
            if isinstance(sections[key], str):
                sections[key] = sections[key].strip()
        
        return sections
    
    def generate_assessment(self, topic: str, grade_level: str, 
                          assessment_type: str, num_questions: int,
                          learning_objectives: List[str]) -> Dict[str, Any]:
        """
        Generate assessment questions
        
        Args:
            topic: Topic to assess
            grade_level: Target grade level
            assessment_type: Type of assessment (quiz, test, formative, summative)
            num_questions: Number of questions to generate
            learning_objectives: Learning objectives to assess
            
        Returns:
            Dict containing assessment questions and answer key
        """
        if not self.client:
            raise Exception("No API client available for content generation")
        
        prompt = f"""Create a {assessment_type} assessment for '{topic}' at {grade_level} level.

Number of questions: {num_questions}
Learning Objectives to assess:
{chr(10).join(f'- {obj}' for obj in learning_objectives)}

Please include:
1. A variety of question types (multiple choice, short answer, etc.)
2. Questions that assess different cognitive levels (remember, understand, apply, analyze)
3. Clear instructions for each question
4. Point values for each question
5. An answer key with explanations

Format as:
Question 1: [question text] (X points)
Type: [question type]
Options: (if applicable)
Answer: [correct answer]
Explanation: [why this is correct]
"""

        try:
            response = self.client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=3000,
                temperature=0.5,
                messages=[{"role": "user", "content": prompt}]
            )
            
            assessment_text = response.content[0].text
            
            # Parse into structured format
            assessment = self._parse_assessment(assessment_text)
            assessment['topic'] = topic
            assessment['grade_level'] = grade_level
            assessment['type'] = assessment_type
            assessment['learning_objectives'] = learning_objectives
            
            return assessment
            
        except Exception as e:
            self.logger.error(f"Error generating assessment: {str(e)}")
            raise
    
    def _parse_assessment(self, text: str) -> Dict[str, Any]:
        """Parse assessment text into structured format"""
        questions = []
        current_question = {}
        
        lines = text.split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for question start
            if re.match(r'^Question \d+:', line):
                if current_question:
                    questions.append(current_question)
                current_question = {
                    'question': line.split(':', 1)[1].strip(),
                    'points': 1,  # Default
                    'type': 'short_answer',  # Default
                    'options': [],
                    'answer': '',
                    'explanation': ''
                }
                # Extract points if present
                points_match = re.search(r'\((\d+) points?\)', line)
                if points_match:
                    current_question['points'] = int(points_match.group(1))
            elif line.startswith('Type:'):
                current_question['type'] = line.split(':', 1)[1].strip()
            elif line.startswith('Options:') or re.match(r'^[a-d]\)', line):
                if re.match(r'^[a-d]\)', line):
                    current_question['options'].append(line)
            elif line.startswith('Answer:'):
                current_question['answer'] = line.split(':', 1)[1].strip()
            elif line.startswith('Explanation:'):
                current_question['explanation'] = line.split(':', 1)[1].strip()
        
        # Add last question
        if current_question:
            questions.append(current_question)
        
        # Calculate total points
        total_points = sum(q['points'] for q in questions)
        
        return {
            'questions': questions,
            'total_points': total_points,
            'num_questions': len(questions)
        }
    
    def generate_enrichment_content(self, topic: str, content_type: str,
                                  grade_level: str, duration: int = 30) -> Dict[str, Any]:
        """
        Generate enrichment content
        
        Args:
            topic: Topic for enrichment
            content_type: Type of enrichment (project, research, creative, etc.)
            grade_level: Target grade level
            duration: Duration in minutes
            
        Returns:
            Dict containing enrichment activity details
        """
        if not self.client:
            raise Exception("No API client available for content generation")
        
        prompt = f"""Create an enrichment {content_type} activity for '{topic}' at {grade_level} level.

Duration: {duration} minutes

Please include:
1. Activity Title
2. Overview/Description
3. Learning Goals
4. Materials Needed
5. Step-by-step Instructions
6. Extension Ideas
7. Assessment Rubric
8. Real-world Connections

Make it engaging and challenging for advanced learners."""

        try:
            response = self.client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=2000,
                temperature=0.8,
                messages=[{"role": "user", "content": prompt}]
            )
            
            enrichment_text = response.content[0].text
            
            # Structure the response
            enrichment = {
                'topic': topic,
                'type': content_type,
                'grade_level': grade_level,
                'duration': duration,
                'content': enrichment_text,
                'structured': self._parse_enrichment(enrichment_text)
            }
            
            return enrichment
            
        except Exception as e:
            self.logger.error(f"Error generating enrichment content: {str(e)}")
            raise
    
    def _parse_enrichment(self, text: str) -> Dict[str, Any]:
        """Parse enrichment text into structured format"""
        sections = {
            'title': '',
            'overview': '',
            'goals': [],
            'materials': [],
            'instructions': [],
            'extensions': [],
            'rubric': '',
            'connections': ''
        }
        
        # Simple parsing logic
        lines = text.split('\n')
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            lower_line = line.lower()
            if 'title' in lower_line and len(line) < 100:
                current_section = 'title'
                sections['title'] = line.split(':', 1)[-1].strip()
            elif 'overview' in lower_line or 'description' in lower_line:
                current_section = 'overview'
            elif 'goal' in lower_line:
                current_section = 'goals'
            elif 'material' in lower_line:
                current_section = 'materials'
            elif 'instruction' in lower_line or 'step' in lower_line:
                current_section = 'instructions'
            elif 'extension' in lower_line:
                current_section = 'extensions'
            elif 'rubric' in lower_line:
                current_section = 'rubric'
            elif 'connection' in lower_line or 'real-world' in lower_line:
                current_section = 'connections'
            elif current_section:
                if current_section in ['goals', 'materials', 'instructions', 'extensions']:
                    if line.startswith('-') or line.startswith('•') or re.match(r'^\d+\.', line):
                        sections[current_section].append(line.lstrip('-•0123456789. '))
                else:
                    sections[current_section] += line + '\n'
        
        # Clean up string sections
        for key in ['overview', 'rubric', 'connections']:
            sections[key] = sections[key].strip()
        
        return sections
    
    def analyze_instructional_framework(self, file_path: str) -> Dict[str, Any]:
        """
        Analyze the instructional framework of a PowerPoint presentation or PDF
        
        Args:
            file_path: Path to PowerPoint or PDF file
            
        Returns:
            Dict containing framework analysis results
        """
        if not self.client:
            raise Exception("No API client available for content analysis")
        
        # Extract slide/page content
        slides_data = self._extract_slide_content(file_path)
        
        # Prepare prompt for Claude
        content_descriptions = "\n\n".join([
            f"PAGE/SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
            for s in slides_data
        ])
        
        prompt = f"""You are an expert in instructional design. Analyze this educational content and identify its instructional framework and teaching pattern.

CONTENT:
{content_descriptions}

Focus on identifying these specific instructional patterns:

1. Recall/I do/We do/You do/Review framework
2. 5E Instructional Model (Engage, Explore, Explain, Elaborate, Evaluate)
3. Other common frameworks (specify which one)

Provide your analysis as JSON with this structure:
{{
  "framework": {{
    "identified_framework": "Name of framework",
    "framework_alignment_score": 0-100,
    "strengths": ["strength 1", "strength 2"],
    "weaknesses": ["weakness 1", "weakness 2"],
    "missing_phases": ["phase 1", "phase 2"],
    "recommendations": "brief recommendation",
    "balance_analysis": "brief analysis of instructional balance"
  }},
  "slides": [
    {{
      "slide_number": 1,
      "framework_phase": "phase name",
      "effectiveness": "strong/adequate/weak",
      "elements": ["element 1", "element 2"]
    }}
  ]
}}

Respond only with the JSON, no other text."""

        try:
            response = self.client.messages.create(
                model="claude-3-5-sonnet-20240620",
                max_tokens=4000,
                temperature=0.3,
                messages=[{"role": "user", "content": prompt}]
            )
            
            content = response.content[0].text.strip()
            
            # Parse JSON response
            import json
            try:
                framework_data = json.loads(content)
                return framework_data
            except json.JSONDecodeError:
                # Try to find JSON block within the response
                json_match = re.search(r'(\{.*\})', content, re.DOTALL)
                if json_match:
                    try:
                        framework_data = json.loads(json_match.group(0))
                        return framework_data
                    except json.JSONDecodeError:
                        self.logger.error("Error parsing JSON from LLM response")
                        return {"error": "Could not parse JSON response", "raw_response": content[:200] + "..."}
                else:
                    return {"error": "Could not extract JSON from response", "raw_response": content[:200] + "..."}
                    
        except Exception as e:
            self.logger.error(f"Error in framework analysis: {str(e)}")
            return {"error": f"Analysis failed: {str(e)}"}
    
    def extract_learning_scaffolding(self, file_path: str, profile: str = None) -> Dict[str, Any]:
        """
        Extract learning scaffolding elements from a PowerPoint presentation or PDF
        
        Args:
            file_path: Path to PowerPoint or PDF file
            profile: Learning profile for assessment (dyslexia, adhd, esl)
            
        Returns:
            Dict containing scaffolding analysis results
        """
        if not self.client:
            raise Exception("No API client available for content analysis")
        
        # Extract slide/page content
        slides_data = self._extract_slide_content(file_path)
        print(f"DEBUG - Extracted {len(slides_data)} slides/pages from {file_path}")
        
        # Check if we have content
        if not slides_data:
            return {
                "slides": [],
                "scaffolding_elements": {
                    "learning_objectives": [],
                    "key_concepts": [],
                    "examples": [],
                    "practice_activities": [],
                    "assessment_items": [],
                    "review_elements": []
                },
                "scaffolding_score": 0,
                "scaffolding_analysis": "No content could be extracted from the file",
                "debug_info": "No slides/pages extracted"
            }
        
        # Limit content for better analysis - use first 15 pages/slides for comprehensive analysis
        # This prevents token limit issues and provides more focused results
        analysis_slides = slides_data[:15]
        self.logger.info(f"Analyzing first {len(analysis_slides)} of {len(slides_data)} pages/slides for scaffolding")
        print(f"DEBUG - Using {len(analysis_slides)} slides for analysis")
        
        # Prepare prompt for Claude
        content_descriptions = "\n\n".join([
            f"PAGE/SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
            for s in analysis_slides
        ])
        
        print(f"DEBUG - Content sent to Claude: {content_descriptions[:500]}...")  # Temporary debug output
        
        prompt = f"""Analyze the following educational content and identify the learning scaffolding elements.

CONTENT:
{content_descriptions}

For each page/slide, determine its instructional purpose (e.g., objectives, key concept, example, practice, assessment, review).

Then extract the following scaffolding elements from the entire content:
1. Learning Objectives
2. Key Concepts and their definitions
3. Examples
4. Practice Activities
5. Assessment Items
6. Review Elements

Format your response as JSON with the following structure:
{{
  "slides": [
    {{
      "slide_number": 1,
      "purpose": "title",
      "elements": []
    }}
  ],
  "scaffolding_elements": {{
    "learning_objectives": ["objective 1", "objective 2"],
    "key_concepts": [
      {{
        "concept": "concept name",
        "definition": "definition text"
      }}
    ],
    "examples": [
      {{
        "slide_number": 5,
        "content": "description of example"
      }}
    ],
    "practice_activities": [
      {{
        "slide_number": 8,
        "content": "description of practice activity"
      }}
    ],
    "assessment_items": [
      {{
        "slide_number": 10,
        "content": "description of assessment"
      }}
    ],
    "review_elements": [
      {{
        "slide_number": 12,
        "content": "description of review element"
      }}
    ]
  }},
  "scaffolding_score": 85,
  "scaffolding_analysis": "Analysis of how well the presentation supports learning..."
}}

Respond only with the JSON, no other text."""

        try:
            response = self.client.messages.create(
                model="claude-3-5-sonnet-20240620",
                max_tokens=4000,
                temperature=0.3,
                messages=[{"role": "user", "content": prompt}]
            )
            
            content = response.content[0].text.strip()
            
            # Debug: Log the AI response
            self.logger.info(f"AI response length: {len(content)} characters")
            self.logger.info(f"AI response preview: {content[:200]}...")
            print(f"DEBUG - Full AI response: {content}")  # Temporary debug output
            
            # Parse JSON response
            import json
            try:
                scaffolding_data = json.loads(content)
                # Debug: Log what we found
                elements = scaffolding_data.get('scaffolding_elements', {})
                self.logger.info(f"Found {len(elements.get('assessment_items', []))} assessment items")
                self.logger.info(f"Found {len(elements.get('key_concepts', []))} key concepts")
                print(f"DEBUG - Parsed scaffolding elements: {elements}")  # Temporary debug output
                
                # Add profile assessment if profile is provided
                if profile:
                    profile_assessment = self._assess_content_for_profile(slides_data, profile)
                    scaffolding_data['profile_assessment'] = profile_assessment
                
                # Add debug info to the result
                scaffolding_data['debug_info'] = {
                    'slides_extracted': len(slides_data),
                    'slides_analyzed': len(analysis_slides),
                    'ai_response_length': len(content),
                    'elements_found': {k: len(v) if isinstance(v, list) else str(v) for k, v in elements.items()},
                    'sample_elements': {
                        'learning_objectives_sample': elements.get('learning_objectives', [])[:2],
                        'key_concepts_sample': elements.get('key_concepts', [])[:2],
                        'assessment_items_sample': elements.get('assessment_items', [])[:2]
                    },
                    'profile_used': profile
                }
                return scaffolding_data
            except json.JSONDecodeError:
                # Try to find JSON block within the response
                json_match = re.search(r'(\{.*\})', content, re.DOTALL)
                if json_match:
                    try:
                        scaffolding_data = json.loads(json_match.group(0))
                        return scaffolding_data
                    except json.JSONDecodeError:
                        self.logger.error("Error parsing JSON from LLM response")
                        return {"error": "Could not parse JSON response", "raw_response": content[:200] + "..."}
                else:
                    return {"error": "Could not extract JSON from response", "raw_response": content[:200] + "..."}
                    
        except Exception as e:
            self.logger.error(f"Error in scaffolding analysis: {str(e)}")
            return {"error": f"Analysis failed: {str(e)}"}
    
    def _extract_slide_content(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract content from PowerPoint slides or PDF pages
        
        Args:
            file_path: Path to PowerPoint or PDF file
            
        Returns:
            List of slide/page content dictionaries
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            return self._extract_pdf_content(file_path)
        elif file_ext == '.pptx':
            return self._extract_pptx_content(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    def _extract_pptx_content(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Extract content from PowerPoint slides
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            List of slide content dictionaries
        """
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = ""
            title = ""
            
            # Get slide title
            if hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide.shapes.title.text
                
            # Extract all text from the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text:
                    slide_content += shape.text + "\n"
            
            slides_data.append({
                "slide_number": i+1,
                "title": title,
                "content": slide_content
            })
        
        return slides_data
    
    def _extract_pdf_content(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extract content from PDF pages
        
        Args:
            pdf_path: Path to PDF file
            
        Returns:
            List of page content dictionaries
        """
        try:
            # Use the PDF service to extract content
            from .pdf_service import PDFService
            pdf_service = PDFService(self.config)
            
            pdf_content = pdf_service.extract_content_from_pdf(pdf_path)
            print(f"DEBUG - PDF service returned {len(pdf_content.get('pages', []))} pages")
            
            # Convert PDF pages to slide-like format for consistency
            slides_data = []
            for page in pdf_content.get('pages', []):
                # Try to identify a title from the first line of text
                text = page.get('text', '').strip()
                print(f"DEBUG - Page {page.get('page_number', '?')} text length: {len(text)}")
                lines = text.split('\n') if text else []
                
                # Use first line as title if it's short enough, otherwise use page number
                title = ""
                content = text
                
                if lines and len(lines[0]) < 100:
                    title = lines[0].strip()
                    content = '\n'.join(lines[1:]).strip() if len(lines) > 1 else ""
                
                if not title:
                    title = f"Page {page.get('page_number', '?')}"
                
                slides_data.append({
                    "slide_number": page.get('page_number', len(slides_data) + 1),
                    "title": title,
                    "content": content
                })
            
            return slides_data
            
        except Exception as e:
            self.logger.error(f"Error extracting PDF content: {str(e)}")
            # Fallback: return empty structure
            return [{
                "slide_number": 1,
                "title": "PDF Content",
                "content": "Unable to extract content from PDF"
            }]
    
    def _assess_content_for_profile(self, slides_data: List[Dict[str, Any]], profile: str) -> Dict[str, Any]:
        """
        Assess content suitability for a specific learning profile
        
        Args:
            slides_data: List of slide/page content
            profile: Learning profile (dyslexia, adhd, esl)
            
        Returns:
            Dict containing profile assessment results
        """
        # Use the assessments service for profile evaluation
        from .assessments_service import AssessmentsService
        assessment_service = AssessmentsService(self.config)
        
        # Convert slides data to the format expected by assessments service
        content = {
            'pages': [{
                'page_number': slide.get('slide_number', i+1),
                'text': f"{slide.get('title', '')}\n{slide.get('content', '')}"
            } for i, slide in enumerate(slides_data)]
        }
        
        # Perform profile assessment
        assessment_result = assessment_service.assess_content(content, profile)
        
        return {
            'profile': profile,
            'suitability_score': assessment_result.get('profile_suitability', {}).get(profile, {}).get('overall_score', 0),
            'issues': assessment_result.get('profile_suitability', {}).get(profile, {}).get('issues', []),
            'strengths': assessment_result.get('profile_suitability', {}).get(profile, {}).get('strengths', []),
            'recommendation': assessment_result.get('profile_suitability', {}).get(profile, {}).get('recommendation', ''),
            'readability_metrics': assessment_result.get('readability_metrics', {}),
            'recommendations': assessment_result.get('recommendations', [])
        }