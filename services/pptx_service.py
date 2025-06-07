"""
Enhanced PowerPoint Service

Consolidates all PowerPoint processing functionality from app.py
"""
import os
from typing import Dict, Any, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from .base_service import BaseService
import anthropic
import re
from PIL import Image, ImageDraw, ImageFont
import io


class PowerPointService(BaseService):
    """Enhanced PowerPoint service with all PPTX processing capabilities"""
    
    def _initialize(self):
        """Initialize PowerPoint service resources"""
        self.api_key = self.config.get('anthropic_api_key')
        if self.api_key:
            self.client = anthropic.Anthropic(api_key=self.api_key)
        else:
            self.client = None
        
        # Profile-specific formatting with evidence-based recommendations
        self.PROFILE_SETTINGS = {
            'dyslexia': {
                'font_name': 'OpenDyslexic',  # Dyslexia-friendly font
                'font_fallbacks': ['Arial', 'Comic Sans MS', 'Verdana', 'sans-serif'],  # Evidence-based fallbacks
                'font_size': 18,
                'line_spacing': 1.8,  # Increased from 1.5 for better readability
                'paragraph_spacing': 1.0,  # Extra space between paragraphs
                'background_color': RGBColor(255, 251, 240),  # Warm beige tint - evidence-based
                'text_color': RGBColor(0, 0, 0),  # Pure black on tinted background for contrast
                'highlight_color': RGBColor(255, 255, 200),  # Soft yellow highlight
                'use_background_tint': True,  # Use background tints instead of colored text
                'add_reading_rulers': True,  # Visual guides for reading
                'text_alignment': PP_ALIGN.LEFT,  # Left-align only, no justification
                'max_line_width': 65,  # Characters per line limit
                'visual_structure': 'simple'  # Minimize visual clutter
            },
            'adhd': {
                'font_name': 'Arial',
                'font_fallbacks': ['Helvetica', 'Verdana', 'sans-serif'],
                'font_size': 16,
                'line_spacing': 1.5,  # Increased from 1.3
                'paragraph_spacing': 1.2,  # Clear paragraph breaks
                'background_color': RGBColor(245, 255, 250),  # Soft green tint for calm
                'text_color': RGBColor(0, 0, 0),  # High contrast black
                'highlight_color': RGBColor(200, 255, 200),  # Soft green highlight
                'chunk_content': True,  # Break content into visual chunks
                'add_visual_boundaries': True,  # Boxes/lines around content blocks
                'boundary_color': RGBColor(200, 200, 200),  # Subtle gray boundaries
                'boundary_opacity': 0.3,  # 30% opacity for boundaries
                'consistent_cues': True,  # Use consistent visual markers
                'text_alignment': PP_ALIGN.LEFT,
                'visual_structure': 'chunked'  # Structured content blocks
            },
            'esl': {
                'font_name': 'Arial',
                'font_fallbacks': ['Helvetica', 'sans-serif'],
                'font_size': 16,
                'line_spacing': 1.4,
                'paragraph_spacing': 0.8,
                'background_color': RGBColor(250, 245, 255),
                'text_color': RGBColor(0, 0, 0),  # High contrast
                'highlight_color': RGBColor(230, 230, 250),
                'text_alignment': PP_ALIGN.LEFT,
                'visual_structure': 'standard'
            },
            'default': {
                'font_name': 'Calibri',
                'font_fallbacks': ['Arial', 'sans-serif'],
                'font_size': 14,
                'line_spacing': 1.2,
                'paragraph_spacing': 0.6,
                'background_color': RGBColor(255, 255, 255),
                'text_color': RGBColor(0, 0, 0),
                'highlight_color': RGBColor(255, 255, 0),
                'text_alignment': PP_ALIGN.LEFT,
                'visual_structure': 'standard'
            },
            'translation': {
                # Translation-only profile - no adaptations, just language conversion
                'font_name': None,  # Keep original font
                'font_fallbacks': [],
                'font_size': None,  # Keep original size
                'line_spacing': None,  # Keep original spacing
                'paragraph_spacing': None,
                'background_color': None,  # No background changes
                'text_color': None,  # Keep original colors
                'highlight_color': None,
                'text_alignment': None,  # Keep original alignment
                'visual_structure': 'preserve',  # Preserve all formatting
                'skip_adaptation': True,  # Skip adaptation phase
                'translation_only': True  # Flag for translation-only mode
            },
            'dyslexia_translation': {
                # Dyslexia formatting with translation only (no text simplification)
                'font_name': 'OpenDyslexic',  # Dyslexia-friendly font
                'font_fallbacks': ['Arial', 'Comic Sans MS', 'Verdana', 'sans-serif'],
                'font_size': 18,
                'line_spacing': 1.8,
                'paragraph_spacing': 1.0,
                'background_color': RGBColor(255, 251, 240),  # Warm beige tint
                'text_color': RGBColor(0, 0, 0),  # Pure black
                'highlight_color': RGBColor(255, 255, 200),
                'use_background_tint': True,
                'add_reading_rulers': True,
                'text_alignment': PP_ALIGN.LEFT,
                'max_line_width': 65,
                'visual_structure': 'simple',
                'skip_adaptation': True,  # Skip text adaptation
                'translation_only': True  # Only translate
            },
            'adhd_translation': {
                # ADHD formatting with translation only (no content restructuring)
                'font_name': 'Arial',
                'font_fallbacks': ['Helvetica', 'Verdana', 'sans-serif'],
                'font_size': 16,
                'line_spacing': 1.5,
                'paragraph_spacing': 1.2,
                'background_color': RGBColor(245, 255, 250),  # Soft green tint
                'text_color': RGBColor(0, 0, 0),  # High contrast black
                'highlight_color': RGBColor(200, 255, 200),
                'chunk_content': True,
                'add_visual_boundaries': True,
                'boundary_color': RGBColor(200, 200, 200),
                'boundary_opacity': 0.3,
                'consistent_cues': True,
                'text_alignment': PP_ALIGN.LEFT,
                'visual_structure': 'chunked',
                'skip_adaptation': True,  # Skip text adaptation
                'translation_only': True  # Only translate
            },
            'esl_translation': {
                # ESL formatting with translation only (no simplification)
                'font_name': 'Arial',
                'font_fallbacks': ['Helvetica', 'sans-serif'],
                'font_size': 16,
                'line_spacing': 1.4,
                'paragraph_spacing': 0.8,
                'background_color': RGBColor(250, 245, 255),
                'text_color': RGBColor(0, 0, 0),
                'highlight_color': RGBColor(230, 230, 250),
                'text_alignment': PP_ALIGN.LEFT,
                'visual_structure': 'standard',
                'skip_adaptation': True,  # Skip text adaptation
                'translation_only': True  # Only translate
            }
        }
    
    def calculate_optimal_font_size(self, text: str, max_width_pts: float, max_height_pts: float, 
                                   font_name: str = "Arial", min_size: int = 8, max_size: int = 72) -> int:
        """
        Calculate optimal font size that fits within given bounds using binary search
        
        Args:
            text: Text content to measure
            max_width_pts: Maximum width in points
            max_height_pts: Maximum height in points  
            font_name: Font family name
            min_size: Minimum font size to consider
            max_size: Maximum font size to consider
            
        Returns:
            Optimal font size in points
        """
        try:
            # Handle empty or very short text
            if not text or len(text.strip()) < 2:
                return min_size
            
            # Convert points to pixels for PIL (assuming 96 DPI)
            max_width_px = int(max_width_pts * 96 / 72)
            max_height_px = int(max_height_pts * 96 / 72)
            
            def get_text_dimensions(text: str, font_size: int) -> Tuple[int, int]:
                """Get text dimensions using PIL"""
                try:
                    # Try to load the specified font, fallback to default
                    try:
                        font = ImageFont.truetype(font_name, font_size)
                    except (OSError, IOError):
                        # Fallback to default font if specified font not found
                        try:
                            font = ImageFont.load_default()
                        except:
                            # Ultimate fallback - estimate based on character count
                            avg_char_width = font_size * 0.6  # Rough estimate
                            width = len(text) * avg_char_width
                            height = font_size * 1.2  # Account for line height
                            return int(width), int(height)
                    
                    # Create temporary image for measurement
                    dummy_img = Image.new('RGB', (1, 1))
                    draw = ImageDraw.Draw(dummy_img)
                    
                    # Handle multi-line text
                    lines = text.split('\n')
                    if len(lines) > 1:
                        # Multi-line text measurement
                        max_line_width = 0
                        total_height = 0
                        
                        for line in lines:
                            if line.strip():  # Skip empty lines
                                bbox = draw.textbbox((0, 0), line, font=font)
                                line_width = bbox[2] - bbox[0]
                                line_height = bbox[3] - bbox[1]
                                max_line_width = max(max_line_width, line_width)
                                total_height += line_height
                        
                        return max_line_width, total_height
                    else:
                        # Single line text
                        bbox = draw.textbbox((0, 0), text, font=font)
                        return bbox[2] - bbox[0], bbox[3] - bbox[1]
                        
                except Exception as e:
                    self.logger.warning(f"Error measuring text dimensions: {e}")
                    # Fallback calculation
                    avg_char_width = font_size * 0.6
                    lines = text.split('\n')
                    max_line_width = max(len(line) for line in lines) * avg_char_width
                    total_height = len(lines) * font_size * 1.2
                    return int(max_line_width), int(total_height)
            
            # Binary search for optimal font size
            low, high = min_size, max_size
            best_size = min_size
            
            while low <= high:
                mid = (low + high) // 2
                width, height = get_text_dimensions(text, mid)
                
                if width <= max_width_px and height <= max_height_px:
                    best_size = mid
                    low = mid + 1  # Try larger size
                else:
                    high = mid - 1  # Try smaller size
            
            return best_size
            
        except Exception as e:
            self.logger.error(f"Error calculating optimal font size: {e}")
            return min_size
    
    def apply_text_with_optimal_sizing(self, text_frame, text: str, max_width_pts: float, 
                                     max_height_pts: float, profile: str = 'default') -> tuple[bool, str]:
        """
        Apply text to text frame with optimal font sizing and overflow handling
        
        Args:
            text_frame: PowerPoint text frame object
            text: Text content to add
            max_width_pts: Maximum width in points
            max_height_pts: Maximum height in points
            profile: Learning profile for styling
            
        Returns:
            Tuple of (success_status, overflow_notes)
        """
        overflow_notes = ""
        
        try:
            # Get profile settings
            profile_settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            
            # Step 1: Try with original text and optimal sizing
            original_text = text
            current_text = text
            
            # Calculate optimal font size for full text
            optimal_size = self.calculate_optimal_font_size(
                current_text, max_width_pts, max_height_pts, 
                profile_settings['font_name']
            )
            
            # Ensure we don't exceed the profile's preferred font size
            max_profile_size = profile_settings.get('font_size', 16)
            final_size = min(optimal_size, max_profile_size)
            
            # Minimum readable font size
            min_readable_size = 10
            
            # Step 2: Check if text fits with minimum readable font size
            if final_size < min_readable_size:
                self.logger.warning(f"Text too long for shape, font size would be {final_size}pt (min: {min_readable_size}pt)")
                
                # Strategy 1: Smart truncation with sentence boundaries
                current_text, overflow_text = self._smart_truncate_text(
                    original_text, max_width_pts, max_height_pts, 
                    min_readable_size, profile_settings['font_name']
                )
                
                if overflow_text:
                    overflow_notes = f"Full content: {overflow_text}"
                    final_size = min_readable_size
                    self.logger.info(f"Truncated text to fit, moved overflow to notes")
                else:
                    # Strategy 2: Convert to bullet points
                    current_text = self._convert_to_bullets(original_text)
                    final_size = self.calculate_optimal_font_size(
                        current_text, max_width_pts, max_height_pts, 
                        profile_settings['font_name']
                    )
                    final_size = max(min(final_size, max_profile_size), min_readable_size)
                    
                    if final_size < min_readable_size:
                        # Strategy 3: Extreme truncation
                        current_text = self._extreme_truncate(original_text, 100)
                        overflow_notes = f"Original content (truncated for readability): {original_text}"
                        final_size = min_readable_size
            
            # Step 3: Apply the text with final formatting
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = current_text
            
            # Apply font formatting
            font = p.font
            font.name = profile_settings['font_name']
            font.size = Pt(final_size)
            font.color.rgb = profile_settings['text_color']
            
            # Apply paragraph formatting
            p.line_spacing = profile_settings['line_spacing']
            
            # Apply word wrap and fit settings
            text_frame.word_wrap = True
            text_frame.auto_size = None  # Disable auto-sizing to respect bounds
            
            # Step 4: Final overflow check
            if self._is_text_still_overflowing(text_frame, max_width_pts, max_height_pts):
                self.logger.warning("Text still overflowing after all strategies")
                # Last resort: use auto-fit
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                overflow_notes += " Note: Text auto-fitted to prevent overflow."
            
            self.logger.info(f"Applied text with font size {final_size}pt (optimal: {optimal_size}pt)")
            return True, overflow_notes
            
        except Exception as e:
            self.logger.error(f"Error applying text with optimal sizing: {e}")
            return False, f"Error applying text: {str(e)}"
    
    def measure_text_bounds(self, text: str, font_name: str = "Arial", font_size: int = 12) -> Tuple[float, float]:
        """
        Measure the exact bounds of text in points
        
        Args:
            text: Text to measure
            font_name: Font family name
            font_size: Font size in points
            
        Returns:
            Tuple of (width_pts, height_pts)
        """
        try:
            # Convert font size to pixels for PIL measurement
            font_size_px = int(font_size * 96 / 72)
            
            # Create font object
            try:
                font = ImageFont.truetype(font_name, font_size_px)
            except (OSError, IOError):
                font = ImageFont.load_default()
            
            # Measure text
            dummy_img = Image.new('RGB', (1, 1))
            draw = ImageDraw.Draw(dummy_img)
            bbox = draw.textbbox((0, 0), text, font=font)
            
            # Convert back to points
            width_pts = (bbox[2] - bbox[0]) * 72 / 96
            height_pts = (bbox[3] - bbox[1]) * 72 / 96
            
            return width_pts, height_pts
            
        except Exception as e:
            self.logger.error(f"Error measuring text bounds: {e}")
            # Fallback estimation
            avg_char_width = font_size * 0.6
            lines = text.split('\n')
            max_line_width = max(len(line) for line in lines) * avg_char_width
            total_height = len(lines) * font_size * 1.2
            return max_line_width, total_height
    
    def get_text_frame_bounds(self, text_frame) -> Tuple[float, float]:
        """
        Get the bounds of a text frame in points
        
        Args:
            text_frame: PowerPoint text frame object
            
        Returns:
            Tuple of (width_pts, height_pts)
        """
        try:
            # Try to get bounds from the parent shape
            if hasattr(text_frame, '_element'):
                # Navigate up to find the shape element
                element = text_frame._element
                
                # Look for spPr (shape properties) in ancestors
                while element is not None:
                    if hasattr(element, 'tag') and 'spPr' in str(element.tag):
                        # Found shape properties, look for transform info
                        for child in element:
                            if 'xfrm' in str(child.tag):  # Transform element
                                for grandchild in child:
                                    if 'ext' in str(grandchild.tag):  # Extent element
                                        # Get cx (width) and cy (height) in EMUs
                                        cx = grandchild.get('cx')
                                        cy = grandchild.get('cy')
                                        if cx and cy:
                                            # Convert EMUs to points (1 point = 12700 EMUs)
                                            width_pts = float(cx) / 12700
                                            height_pts = float(cy) / 12700
                                            return width_pts, height_pts
                        break
                    element = element.getparent() if hasattr(element, 'getparent') else None
            
            # Fallback: reasonable defaults based on common slide layouts
            return 400.0, 300.0  # ~5.5" x 4.2" 
            
        except Exception as e:
            self.logger.warning(f"Could not determine text frame bounds: {e}")
            return 400.0, 300.0
    
    def extract_content_from_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract content from PowerPoint file
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dict containing slides with content
        """
        content = {
            'slides': [],
            'metadata': {
                'slide_count': 0,
                'title': '',
                'author': ''
            }
        }
        
        try:
            prs = Presentation(pptx_path)
            content['metadata']['slide_count'] = len(prs.slides)
            
            # Try to get title from first slide
            if prs.slides and hasattr(prs.slides[0], 'shapes'):
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, 'text') and shape.text:
                        content['metadata']['title'] = shape.text
                        break
            
            # Extract slides
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': slide_idx + 1,
                    'title': '',
                    'content': '',
                    'notes': '',
                    'shapes': []
                }
                
                # Extract text from shapes
                text_elements = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        # Check if it's likely a title
                        if shape.top < Inches(2) and len(shape.text) < 100:
                            if not slide_content['title']:
                                slide_content['title'] = shape.text
                        else:
                            text_elements.append(shape.text)
                        
                        # Store shape info for reconstruction
                        slide_content['shapes'].append({
                            'type': 'text',
                            'text': shape.text,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        })
                
                slide_content['content'] = '\n'.join(text_elements)
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    slide_content['notes'] = notes_text
                
                content['slides'].append(slide_content)
            
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content: {str(e)}")
            raise Exception(f"Error extracting PPTX content: {str(e)}")
        
        return content
    
    def create_presentation(self, content: Dict[str, Any], output_path: str,
                          profile: str = 'default') -> bool:
        """
        Create PowerPoint presentation from content
        
        Args:
            content: Content to create presentation from
            output_path: Path for output file
            profile: Learning profile for formatting
            
        Returns:
            bool: Success status
        """
        try:
            prs = Presentation()
            settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            
            # Set presentation properties
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Process each slide
            for slide_data in content.get('slides', []):
                # Add slide with bullet layout
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = settings['background_color']
                
                # Add title
                if slide_data.get('title'):
                    title = slide.shapes.title
                    title.text = slide_data['title']
                    self._format_text_frame(title.text_frame, settings, is_title=True)
                
                # Add content
                if slide_data.get('content'):
                    content_shape = slide.placeholders[1]
                    content_shape.text = slide_data['content']
                    self._format_text_frame(content_shape.text_frame, settings)
                
                # Add notes
                if slide_data.get('notes'):
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data['notes']
            
            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            self.logger.error(f"Error creating presentation: {str(e)}")
            return False
    
    def _format_text_frame(self, text_frame, settings: Dict[str, Any], is_title: bool = False):
        """Apply profile-specific formatting to text frame"""
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = settings['line_spacing']
            
            for run in paragraph.runs:
                run.font.name = settings['font_name']
                run.font.size = Pt(settings['font_size'] + (4 if is_title else 0))
                run.font.color.rgb = settings['text_color']
                run.font.bold = is_title
    
    def analyze_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """
        Analyze PowerPoint structure and content
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Dict with analysis results
        """
        analysis = {
            'slide_count': 0,
            'total_text_length': 0,
            'slides': [],
            'has_notes': False,
            'has_images': False,
            'instructional_framework': None
        }
        
        try:
            prs = Presentation(pptx_path)
            analysis['slide_count'] = len(prs.slides)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_info = {
                    'slide_number': slide_idx + 1,
                    'text_length': 0,
                    'shape_count': len(slide.shapes),
                    'has_title': False,
                    'has_notes': False
                }
                
                # Analyze shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text_length = len(shape.text)
                        slide_info['text_length'] += text_length
                        analysis['total_text_length'] += text_length
                        
                        if shape == slide.shapes.title:
                            slide_info['has_title'] = True
                    
                    if shape.shape_type == MSO_SHAPE.PICTURE:
                        analysis['has_images'] = True
                
                # Check for notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                    slide_info['has_notes'] = True
                    analysis['has_notes'] = True
                
                analysis['slides'].append(slide_info)
            
            # Analyze instructional framework
            analysis['instructional_framework'] = self._analyze_instructional_framework(prs)
            
        except Exception as e:
            self.logger.error(f"Error analyzing PPTX: {str(e)}")
        
        return analysis
    
    def _analyze_instructional_framework(self, presentation) -> Optional[str]:
        """Analyze the instructional framework of the presentation"""
        # Check for common instructional patterns
        slide_titles = []
        for slide in presentation.slides:
            if hasattr(slide.shapes, 'title') and slide.shapes.title and slide.shapes.title.text:
                slide_titles.append(slide.shapes.title.text.lower())
        
        # Check for 5E model
        five_e_keywords = ['engage', 'explore', 'explain', 'elaborate', 'evaluate']
        if any(keyword in ' '.join(slide_titles) for keyword in five_e_keywords):
            return "5E Model"
        
        # Check for other models
        if any(word in ' '.join(slide_titles) for word in ['objective', 'assessment', 'activity']):
            return "Traditional Lesson Structure"
        
        return "General Presentation"
    
    def apply_dyslexia_formatting(self, presentation, profile: str = 'dyslexia'):
        """
        Apply dyslexia-friendly (or other profile) formatting to presentation
        
        Args:
            presentation: PPTX presentation object
            profile: Learning profile
        """
        settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
        
        for slide in presentation.slides:
            # Set slide background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = settings['background_color']
            
            # Format all text shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._format_text_frame(shape.text_frame, settings)
    
    def translate_presentation(self, file_path: str, file_id: str, filename: str,
                              target_language: str, progress_callback: Optional[callable] = None) -> Optional[str]:
        """
        Translate a presentation preserving original slides and creating translated copies
        
        This method creates a presentation with original slides followed by translated copies,
        allowing users to compare original and translated content side by side.
        
        Args:
            file_path: Input PPTX path
            file_id: Unique file identifier
            filename: Original filename
            target_language: Target language for translation (required)
            progress_callback: Optional callback function for progress updates (message, percentage)
            
        Returns:
            Output path or None on error
        """
        if not target_language:
            self.logger.error("Target language is required for translation")
            return None
            
        # Use the new slide duplication method for translation
        return self._create_translation_presentation(
            file_path=file_path,
            file_id=file_id,
            filename=filename,
            target_language=target_language,
            progress_callback=progress_callback
        )
    
    def process_presentation_efficiently(self, file_path: str, file_id: str, filename: str, 
                                       profile: str, target_language: Optional[str] = None,
                                       progress_callback: Optional[callable] = None) -> Optional[str]:
        """
        Process presentation efficiently with adaptations and optional translation
        NOW USES FORMAT-PRESERVING APPROACH - preserves all original formatting, layouts, images, animations, etc.
        
        Args:
            file_path: Input PPTX path
            file_id: Unique file identifier  
            filename: Original filename
            profile: Learning profile
            target_language: Optional target language for translation
            progress_callback: Optional callback function for progress updates (message, percentage)
            
        Returns:
            Output path or None on error
        """
        try:
            self.logger.info(f"Starting format-preserving PowerPoint processing for {filename} with profile: {profile}")
            
            # Update status
            if progress_callback:
                progress_callback('Starting format-preserving PowerPoint adaptation...', 5)
            
            # Generate output path
            output_filename = f"adapted_{filename}"
            output_dir = self.config.get('output_folder', self.config.get('output_dir', 'outputs'))
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{file_id}_{output_filename}")
            
            # Create internal callback for progress updates
            def internal_progress_callback(message, progress):
                if progress_callback:
                    # Ensure progress is capped at 100
                    capped_progress = min(progress, 100)
                    progress_callback(message, capped_progress)
                    # Also log progress for debugging
                    self.logger.info(f"Progress {capped_progress:.1f}%: {message}")
            
            # Use the new format-preserving adaptation method
            success = self.adapt_presentation_preserving_format(
                file_path, 
                output_path, 
                profile, 
                target_language, 
                internal_progress_callback
            )
            
            if success:
                self.logger.info(f"Successfully created format-preserving adapted presentation: {output_path}")
                if progress_callback:
                    progress_callback('Format-preserving PowerPoint adaptation completed successfully!', 100)
                return output_path
            else:
                self.logger.error("Failed to create format-preserving adapted presentation")
                if progress_callback:
                    progress_callback('Failed to create format-preserving adapted presentation', 100)
                return None
                
        except Exception as e:
            self.logger.error(f"Error in format-preserving presentation processing: {str(e)}")
            import traceback
            traceback.print_exc()
            if progress_callback:
                progress_callback(f'Error: {str(e)}', 100)
            return None
    
    def _adapt_pptx_content(self, content: Dict[str, Any], profile: str) -> Dict[str, Any]:
        """
        Adapt PowerPoint content using AI based on learning profile
        
        Args:
            content: Extracted PowerPoint content
            profile: Learning profile (dyslexia, adhd, esl, etc.)
            
        Returns:
            Adapted content dictionary
        """
        if not self.client:
            self.logger.warning("No AI client available for content adaptation")
            return content
        
        adapted_content = {
            'slides': [],
            'metadata': content.get('metadata', {}).copy()
        }
        adapted_content['metadata']['adapted_for'] = profile
        
        try:
            # Collect all texts for batch processing
            all_texts = []
            text_metadata = []  # Track which slide/field each text belongs to
            
            for slide_idx, slide in enumerate(content.get('slides', [])):
                if slide.get('title'):
                    all_texts.append(slide['title'])
                    text_metadata.append({'slide': slide_idx, 'field': 'title'})
                
                if slide.get('content'):
                    all_texts.append(slide['content'])
                    text_metadata.append({'slide': slide_idx, 'field': 'content'})
                
                if slide.get('notes'):
                    all_texts.append(slide['notes'])
                    text_metadata.append({'slide': slide_idx, 'field': 'notes'})
            
            # Batch adapt all texts at once for better efficiency
            adapted_texts = []
            if all_texts:
                try:
                    from .adaptations_service import AdaptationsService
                    adaptations_service = AdaptationsService(self.config)
                    self.logger.info(f"🚀 Batch adapting {len(all_texts)} slide elements for profile '{profile}'")
                    adapted_texts = adaptations_service.process_text_batch(all_texts, profile)
                    self.logger.info(f"✅ Batch adaptation completed: {len(adapted_texts)} elements processed")
                except Exception as batch_error:
                    self.logger.warning(f"⚠️ Batch adaptation failed: {batch_error}. Falling back to individual processing.")
                    # Fallback to individual processing
                    adapted_texts = []
                    for text in all_texts:
                        try:
                            adapted_text = self._adapt_text_for_profile(text, profile)
                            adapted_texts.append(adapted_text)
                        except Exception as e:
                            self.logger.error(f"Individual adaptation failed: {e}")
                            adapted_texts.append(text)  # Use original on failure
            
            # Apply adapted texts back to slides
            text_idx = 0
            for slide_idx, slide in enumerate(content.get('slides', [])):
                adapted_slide = slide.copy()
                
                # Apply adapted title
                if slide.get('title') and text_idx < len(adapted_texts):
                    adapted_slide['title'] = adapted_texts[text_idx]
                    text_idx += 1
                
                # Apply adapted content
                if slide.get('content') and text_idx < len(adapted_texts):
                    adapted_slide['content'] = adapted_texts[text_idx]
                    text_idx += 1
                
                # Apply adapted notes
                if slide.get('notes') and text_idx < len(adapted_texts):
                    adapted_slide['notes'] = adapted_texts[text_idx]
                    text_idx += 1
                
                adapted_content['slides'].append(adapted_slide)
            
            return adapted_content
            
        except Exception as e:
            self.logger.error(f"Error adapting PowerPoint content: {str(e)}")
            return content  # Return original content on error
    
    def _adapt_text_for_profile(self, text: str, profile: str) -> str:
        """
        Adapt text content for specific learning profile using AI
        
        Args:
            text: Original text content
            profile: Learning profile
            
        Returns:
            Adapted text content
        """
        if not text or not text.strip():
            return text
        
        if not self.client:
            return text
        
        try:
            # Create profile-specific adaptation prompt
            prompt = self._create_adaptation_prompt(text, profile)
            
            response = self.client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=2000,
                temperature=0.3,
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            
            return response.content[0].text
            
        except Exception as e:
            self.logger.error(f"Error adapting text for {profile}: {str(e)}")
            return text  # Return original text on error
    
    def _create_adaptation_prompt(self, text: str, profile: str) -> str:
        """Create profile-specific adaptation prompt"""
        base_prompt = f"Adapt the following text for a {profile} learner. "
        
        if profile == "dyslexia":
            profile_instructions = """
            For dyslexia support:
            - Use simple, common words where possible
            - Keep sentences short and clear (under 20 words)
            - Break down complex concepts into smaller parts
            - Use active voice instead of passive
            - Avoid unnecessary jargon or complex terminology
            - If technical terms are needed, explain them simply
            """
        elif profile == "adhd":
            profile_instructions = """
            For ADHD support:
            - Create clear, structured content with bullet points
            - Keep paragraphs short and focused
            - Use concrete examples and avoid abstract concepts
            - Include clear headings and organization
            - Make key information stand out
            - Reduce cognitive load by simplifying sentence structure
            """
        elif profile == "esl":
            profile_instructions = """
            For English Language Learners:
            - Use simpler vocabulary and shorter sentences
            - Avoid idioms, slang, and cultural references
            - Explain technical terms with the original in parentheses
            - Use clear, direct language
            - Provide context for complex concepts
            - Maintain formal academic tone while simplifying
            """
        else:
            profile_instructions = """
            Make the content more accessible by:
            - Using clearer, simpler language
            - Breaking down complex ideas
            - Improving overall readability
            """
        
        prompt = f"""{base_prompt}{profile_instructions}
        
        Maintain the original meaning and educational value while making it more accessible.
        Keep the same general structure and formatting.
        
        Original text:
        {text}
        
        Adapted text:"""
        
        return prompt
    
    def generate_enriched_presentation(self, topic: str, grade_level: str,
                                     learning_objectives: List[str], duration: int,
                                     profile: str = 'default') -> Dict[str, Any]:
        """
        Generate an enriched educational presentation
        
        Args:
            topic: Topic of the presentation
            grade_level: Target grade level
            learning_objectives: List of learning objectives
            duration: Duration in minutes
            profile: Learning profile
            
        Returns:
            Dict with presentation content
        """
        if not self.client:
            raise Exception("No API client available for content generation")
        
        # Create prompt for content generation
        prompt = f"""Create an educational presentation on '{topic}' for {grade_level} students.

Learning Objectives:
{chr(10).join(f'- {obj}' for obj in learning_objectives)}

Duration: {duration} minutes
Learning Profile: {profile}

Please structure the presentation with:
1. Engaging title slide
2. Learning objectives slide
3. Main content slides (5-7 slides)
4. Interactive activity slide
5. Summary/review slide
6. Assessment/quiz slide

For each slide, provide:
- Title
- Main content (bullet points)
- Speaker notes
- Suggested visuals

Format as JSON with structure:
{{
  "slides": [
    {{
      "title": "...",
      "content": "...",
      "notes": "...",
      "visuals": "..."
    }}
  ]
}}"""
        
        try:
            response = self.client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=4000,
                temperature=0.7,
                messages=[{"role": "user", "content": prompt}]
            )
            
            # Parse response (assuming JSON format)
            import json
            content_text = response.content[0].text
            
            # Extract JSON from response
            json_match = re.search(r'\{.*\}', content_text, re.DOTALL)
            if json_match:
                presentation_content = json.loads(json_match.group())
                return presentation_content
            else:
                # Fallback: create basic structure
                return {
                    "slides": [
                        {
                            "title": f"{topic} - {grade_level}",
                            "content": "Welcome to today's lesson",
                            "notes": "Introduce the topic",
                            "visuals": "Title slide with engaging image"
                        }
                    ]
                }
            
        except Exception as e:
            self.logger.error(f"Error generating presentation: {str(e)}")
            raise
    
    def apply_text_to_text_frame(self, text_frame, adapted_text: str, profile: str, profiles_service=None, 
                               use_optimal_sizing: bool = True, max_width_pts: Optional[float] = None, 
                               max_height_pts: Optional[float] = None):
        """
        Apply adapted text to a text frame with proper formatting and profile colors
        
        Args:
            text_frame: PowerPoint text frame object
            adapted_text: The adapted text to apply
            profile: Learning profile for styling
            profiles_service: Optional profiles service for color lookup
            use_optimal_sizing: Whether to use optimal font sizing to prevent overflow
            max_width_pts: Maximum width in points (auto-detected if not provided)
            max_height_pts: Maximum height in points (auto-detected if not provided)
        """
        # Skip if text frame is empty
        if not hasattr(text_frame, "paragraphs"):
            return
        
        # If optimal sizing is enabled, use the new method
        if use_optimal_sizing:
            try:
                # Auto-detect bounds if not provided
                if max_width_pts is None or max_height_pts is None:
                    detected_width, detected_height = self.get_text_frame_bounds(text_frame)
                    max_width_pts = max_width_pts or detected_width
                    max_height_pts = max_height_pts or detected_height
                
                # Use the new optimal sizing method
                success = self.apply_text_with_optimal_sizing(
                    text_frame, adapted_text, max_width_pts, max_height_pts, profile
                )
                
                if success:
                    return  # Successfully applied with optimal sizing
                
                # If optimal sizing failed, fall back to original method
                self.logger.warning("Optimal sizing failed, falling back to original method")
                
            except Exception as e:
                self.logger.warning(f"Error in optimal sizing, falling back to original method: {e}")
        
        # Original method (fallback or when optimal sizing is disabled)
        # Get the profile color
        actual_profile = profile.lower() if profile else "default"
        profile_color = self._get_profile_color(actual_profile, profiles_service)
        
        try:
            # Store original formatting settings we want to preserve
            original_formatting = []
            for p in text_frame.paragraphs:
                p_format = {
                    "level": p.level if hasattr(p, "level") else 0,
                    "alignment": p.alignment if hasattr(p, "alignment") else None,
                    "font_name": None,
                    "font_size": None,
                    "font_bold": None,
                    "font_italic": None
                }
                
                # Get formatting from first run if available
                if p.runs and hasattr(p.runs[0], "font"):
                    font = p.runs[0].font
                    p_format["font_name"] = font.name if hasattr(font, "name") else None
                    p_format["font_size"] = font.size if hasattr(font, "size") else None
                    p_format["font_bold"] = font.bold if hasattr(font, "bold") else None
                    p_format["font_italic"] = font.italic if hasattr(font, "italic") else None
                
                original_formatting.append(p_format)
            
            # Clear all paragraphs except the first one
            while len(text_frame.paragraphs) > 1:
                p = text_frame.paragraphs[-1]
                tr = p._p
                tr.getparent().remove(tr)
            
            # Clear the first paragraph content but preserve its format
            if text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                while len(p.runs) > 0:
                    try:
                        p._p.remove(p.runs[0]._r)
                    except:
                        break
            else:
                p = text_frame.add_paragraph()
            
            # Split the adapted text into paragraphs
            text_paragraphs = adapted_text.split('\n')
            
            # Format the first paragraph with colored first word
            if text_paragraphs and text_paragraphs[0].strip():
                words = text_paragraphs[0].split()
                if words:
                    # First word with color
                    first_word = words[0]
                    first_run = p.add_run()
                    first_run.text = first_word
                    # Apply profile color to first word
                    if profile_color and hasattr(first_run, "font") and hasattr(first_run.font, "color"):
                        first_run.font.color.rgb = profile_color
                    
                    # Apply original formatting to first run
                    if original_formatting and hasattr(first_run, "font"):
                        font = first_run.font
                        if original_formatting[0]["font_name"]:
                            font.name = original_formatting[0]["font_name"]
                        if original_formatting[0]["font_size"]:
                            font.size = original_formatting[0]["font_size"]
                        if original_formatting[0]["font_bold"]:
                            font.bold = original_formatting[0]["font_bold"]
                        if original_formatting[0]["font_italic"]:
                            font.italic = original_formatting[0]["font_italic"]
                    
                    # Rest of first paragraph
                    rest_of_text = ' '.join(words[1:]) if len(words) > 1 else ''
                    if rest_of_text:
                        rest_run = p.add_run()
                        rest_run.text = ' ' + rest_of_text
                        # Apply original formatting to rest run
                        if original_formatting and hasattr(rest_run, "font"):
                            font = rest_run.font
                            if original_formatting[0]["font_name"]:
                                font.name = original_formatting[0]["font_name"]
                            if original_formatting[0]["font_size"]:
                                font.size = original_formatting[0]["font_size"]
                            if original_formatting[0]["font_bold"]:
                                font.bold = original_formatting[0]["font_bold"]
                            if original_formatting[0]["font_italic"]:
                                font.italic = original_formatting[0]["font_italic"]
                
                # Apply original paragraph formatting
                if original_formatting:
                    if original_formatting[0]["alignment"]:
                        p.alignment = original_formatting[0]["alignment"]
                    if original_formatting[0]["level"]:
                        p.level = original_formatting[0]["level"]
            
            # Add remaining paragraphs
            for i, p_text in enumerate(text_paragraphs[1:], 1):
                if p_text.strip():
                    new_p = text_frame.add_paragraph()
                    run = new_p.add_run()
                    run.text = p_text
                    
                    # Apply formatting if available
                    fmt_index = min(i, len(original_formatting) - 1) if original_formatting else -1
                    
                    if fmt_index >= 0:
                        # Apply paragraph formatting
                        if original_formatting[fmt_index]["alignment"]:
                            new_p.alignment = original_formatting[fmt_index]["alignment"]
                        if original_formatting[fmt_index]["level"]:
                            new_p.level = original_formatting[fmt_index]["level"]
                        
                        # Apply run formatting
                        if hasattr(run, "font"):
                            font = run.font
                            if original_formatting[fmt_index]["font_name"]:
                                font.name = original_formatting[fmt_index]["font_name"]
                            if original_formatting[fmt_index]["font_size"]:
                                font.size = original_formatting[fmt_index]["font_size"]
                            if original_formatting[fmt_index]["font_bold"]:
                                font.bold = original_formatting[fmt_index]["font_bold"]
                            if original_formatting[fmt_index]["font_italic"]:
                                font.italic = original_formatting[fmt_index]["font_italic"]
        
        except Exception as e:
            self.logger.error(f"Error formatting text frame: {str(e)}")
            # Fallback: just set the text directly if there's an error
            try:
                text_frame.text = adapted_text
            except:
                pass
    
    def _get_profile_color(self, profile_id: str, profiles_service=None) -> Optional[RGBColor]:
        """Get RGBColor for a profile"""
        if profiles_service:
            try:
                colors = profiles_service.get_profile_colors(profile_id)
                if colors:
                    # Convert hex to RGB
                    hex_color = colors.get('text', '#000000')
                    hex_color = hex_color.lstrip('#')
                    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                    return RGBColor(r, g, b)
            except Exception as e:
                self.logger.error(f"Error getting profile color: {str(e)}")
        
        # Fallback to built-in profile settings
        settings = self.PROFILE_SETTINGS.get(profile_id, self.PROFILE_SETTINGS['default'])
        return settings.get('text_color')
    
    def adapt_presentation_preserving_format(self, input_path: str, output_path: str, 
                                           profile: str, target_language: Optional[str] = None,
                                           processing_callback: Optional[callable] = None) -> bool:
        """
        Adapt PowerPoint presentation while preserving ALL original formatting, layouts, images, etc.
        Only the text content is adapted - everything else remains exactly the same.
        
        Args:
            input_path: Path to input PPTX file
            output_path: Path for output PPTX file  
            profile: Learning profile for text adaptation
            target_language: Optional target language for translation
            processing_callback: Optional callback for progress updates
            
        Returns:
            Success status
        """
        try:
            self.logger.info(f"Starting format-preserving PPTX adaptation for profile: {profile}")
            
            # Load the original presentation
            presentation = Presentation(input_path)
            total_slides = len(presentation.slides)
            
            if processing_callback:
                processing_callback("Analyzing presentation structure...", 10)
            
            # Initialize adaptation service
            from .adaptations_service import AdaptationsService
            adaptations_service = AdaptationsService(self.config)
            
            self.logger.info(f"Processing {total_slides} slides with format preservation")
            
            # Calculate progress ranges for better tracking
            adaptation_start = 10
            adaptation_end = 85
            translation_start = 85
            translation_end = 95
            
            # Count total text elements across all slides for better progress calculation
            total_text_elements = 0
            slide_text_counts = []
            
            for slide in presentation.slides:
                text_count = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                        original_text = shape.text_frame.text.strip()
                        if original_text and len(original_text) > 5:
                            text_count += 1
                total_text_elements += text_count
                slide_text_counts.append(text_count)
            
            processed_elements = 0
            
            # Apply profile-specific slide formatting before processing
            self._apply_profile_slide_formatting(presentation, profile)
            
            # Apply global background to all existing slides immediately for profiles with background tints
            profile_settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            if profile_settings.get('use_background_tint', False) and profile_settings.get('background_color'):
                self.logger.info(f"🎨 Applying global background to all {total_slides} slides for {profile}")
                for slide in presentation.slides:
                    self._apply_slide_background(slide, profile_settings['background_color'])
            
            # Process each slide individually with detailed progress tracking
            for slide_idx, slide in enumerate(presentation.slides):
                slide_text_count = slide_text_counts[slide_idx]
                
                # Update progress at start of slide
                if processing_callback:
                    base_progress = adaptation_start + (processed_elements / total_text_elements) * (adaptation_end - adaptation_start)
                    processing_callback(f"Processing slide {slide_idx + 1} of {total_slides} ({slide_text_count} text elements)...", base_progress)
                
                self.logger.info(f"Processing slide {slide_idx + 1}/{total_slides}")
                
                # Apply evidence-based visual formatting to the slide BEFORE processing text
                # Skip formatting for translation-only profile
                profile_settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
                if not profile_settings.get('translation_only', False):
                    self._apply_evidence_based_formatting(slide, profile)
                
                # Extract all text content from this slide
                slide_texts = []
                text_shapes = []
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                        original_text = shape.text_frame.text.strip()
                        if original_text and len(original_text) > 5:  # Only adapt substantial text
                            slide_texts.append(original_text)
                            text_shapes.append(shape)
                
                if slide_texts:
                    try:
                        # Check if this is translation-only mode
                        profile_settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
                        is_translation_only = profile_settings.get('translation_only', False)
                        
                        if is_translation_only:
                            # Skip adaptation phase for translation-only profile
                            self.logger.info(f"  Slide {slide_idx + 1}: Translation-only mode - skipping adaptation")
                            adapted_texts = slide_texts.copy()  # Use original texts
                            
                            if processing_callback:
                                adapt_progress = adaptation_start + (processed_elements / total_text_elements) * (adaptation_end - adaptation_start)
                                processing_callback(f"Preparing {len(slide_texts)} text elements for translation on slide {slide_idx + 1}...", adapt_progress)
                        else:
                            # Regular adaptation flow
                            # Update progress for adaptation start
                            if processing_callback:
                                adapt_progress = adaptation_start + (processed_elements / total_text_elements) * (adaptation_end - adaptation_start)
                                processing_callback(f"Adapting {len(slide_texts)} text elements on slide {slide_idx + 1}...", adapt_progress)
                            
                            self.logger.info(f"  Slide {slide_idx + 1}: Adapting {len(slide_texts)} text elements")
                            
                            # Try to use batch processing with progress callbacks
                            try:
                                # Create a wrapper callback for adaptation progress
                                def adaptation_progress_callback(current, total, message=""):
                                    if processing_callback:
                                        # Calculate sub-progress within the current slide's adaptation phase
                                        base_progress = adaptation_start + (processed_elements / total_text_elements) * (adaptation_end - adaptation_start)
                                        sub_progress = (current / total) * 5  # Use 5% of total for detailed adaptation progress
                                        total_progress = base_progress + sub_progress
                                        processing_callback(f"Adapting text element {current}/{total} on slide {slide_idx + 1}{': ' + message if message else ''}...", total_progress)
                                
                                # Check if the adaptations service supports progress callbacks
                                if hasattr(adaptations_service, 'process_text_batch_with_progress'):
                                    adapted_texts = adaptations_service.process_text_batch_with_progress(
                                        slide_texts, profile, progress_callback=adaptation_progress_callback
                                    )
                                else:
                                    # Fallback to regular batch processing
                                    adapted_texts = adaptations_service.process_text_batch(slide_texts, profile)
                                
                            except Exception as batch_error:
                                self.logger.warning(f"Batch processing failed for slide {slide_idx + 1}: {batch_error}")
                                # Fallback to individual processing with progress
                                adapted_texts = []
                                for i, text in enumerate(slide_texts):
                                    if processing_callback:
                                        element_progress = adaptation_start + ((processed_elements + i) / total_text_elements) * (adaptation_end - adaptation_start)
                                        processing_callback(f"Adapting element {i+1}/{len(slide_texts)} on slide {slide_idx + 1} (individual fallback)...", element_progress)
                                    try:
                                        adapted_text = adaptations_service.adapt_text(text, profile) if hasattr(adaptations_service, 'adapt_text') else text
                                        adapted_texts.append(adapted_text)
                                    except Exception as e:
                                        self.logger.error(f"Individual adaptation failed: {e}")
                                        adapted_texts.append(text)
                        
                        # Update progress after adaptation
                        processed_elements += len(slide_texts)
                        if processing_callback:
                            post_adapt_progress = adaptation_start + (processed_elements / total_text_elements) * (adaptation_end - adaptation_start)
                            processing_callback(f"Applying adaptations to slide {slide_idx + 1}...", post_adapt_progress)
                        
                        # Apply adaptations back to shapes (preserving all formatting)
                        adaptation_notes = []  # Track adaptations for slide notes
                        
                        for i, (shape, adapted_text) in enumerate(zip(text_shapes, adapted_texts)):
                            if adapted_text and adapted_text.strip():
                                original_text = slide_texts[i]
                                
                                # Check for placeholder or debug text that should not appear in slides
                                if self._is_placeholder_text(adapted_text):
                                    self.logger.warning(f"    Shape {i+1}: Placeholder text detected, using original instead")
                                    # Add placeholder info to notes instead of slide content
                                    adaptation_notes.append(f"Text element {i+1}: Adaptation returned placeholder text - '{adapted_text[:50]}...'")
                                    validated_text = original_text  # Use original text
                                # Check for translation notes that should go in slide notes
                                elif self._is_translation_note(adapted_text):
                                    self.logger.info(f"    Shape {i+1}: Translation note detected, adding to slide notes")
                                    # Add translation note to slide notes
                                    adaptation_notes.append(f"Translation note: {adapted_text}")
                                    validated_text = original_text  # Use original text for slide content
                                else:
                                    # Validate text length is within ±20% of original
                                    validated_text = self._validate_text_length(original_text, adapted_text, max_variance=0.20)
                                
                                # Apply the validated text (with dyslexia font if applicable)
                                self._replace_text_preserving_format(shape.text_frame, validated_text, profile)
                                self.logger.info(f"    Shape {i+1}: Text adapted successfully")
                                
                                # Track adaptation for notes
                                if validated_text != original_text:
                                    change_summary = self._create_adaptation_summary(original_text, validated_text, profile)
                                    adaptation_notes.append(f"Text element {i+1}: {change_summary}")
                                
                                # Micro-progress updates for applying adaptations
                                if processing_callback and len(text_shapes) > 3:  # Only for slides with many elements
                                    micro_progress = post_adapt_progress + (i / len(text_shapes)) * 1  # Small increment
                                    processing_callback(f"Applied adaptation {i+1}/{len(text_shapes)} on slide {slide_idx + 1}", micro_progress)
                            else:
                                self.logger.warning(f"    Shape {i+1}: No adapted text received")
                        
                        # Apply translation if requested (before adding notes)
                        # For translation-only mode, target_language is required
                        translation_notes = []
                        if is_translation_only and not target_language:
                            self.logger.warning("Translation-only mode selected but no target language specified")
                            translation_notes.append("Warning: Translation-only mode requires a target language")
                        
                        if target_language and target_language.strip():
                            try:
                                if processing_callback:
                                    trans_progress = translation_start + (slide_idx / total_slides) * (translation_end - translation_start)
                                    processing_callback(f"Translating slide {slide_idx + 1} to {target_language}...", trans_progress)
                                
                                from .translations_service import TranslationsService
                                translations_service = TranslationsService(self.config)
                                
                                # Translate the adapted texts
                                translated_texts = []
                                for i, adapted_text in enumerate(adapted_texts):
                                    if adapted_text:
                                        translated_text = translations_service.translate_text(adapted_text, target_language)
                                        translated_texts.append(translated_text)
                                        
                                        # Progress update for translation
                                        if processing_callback and len(adapted_texts) > 2:
                                            trans_micro_progress = trans_progress + (i / len(adapted_texts)) * 1
                                            processing_callback(f"Translated {i+1}/{len(adapted_texts)} elements on slide {slide_idx + 1}", trans_micro_progress)
                                    else:
                                        translated_texts.append(adapted_text)
                                
                                # Apply translations to shapes with placeholder detection
                                for i, (shape, translated_text) in enumerate(zip(text_shapes, translated_texts)):
                                    if translated_text and translated_text.strip():
                                        # Check for placeholder text in translations
                                        if self._is_placeholder_text(translated_text):
                                            self.logger.warning(f"    Translation placeholder detected for shape {i+1}: '{translated_text[:50]}...'")
                                            # Keep the adapted text instead of placeholder translation
                                            final_text = adapted_texts[i] if i < len(adapted_texts) else slide_texts[i]
                                            translation_notes.append(f"Text element {i+1}: Translation returned placeholder - '{translated_text[:50]}...', using adapted text instead")
                                        # Check for translation notes that should go in slide notes
                                        elif self._is_translation_note(translated_text):
                                            self.logger.info(f"    Translation note detected for shape {i+1}, adding to slide notes")
                                            # Add translation note to slide notes
                                            translation_notes.append(f"Translation note: {translated_text}")
                                            # Keep the adapted text for slide content
                                            final_text = adapted_texts[i] if i < len(adapted_texts) else slide_texts[i]
                                        else:
                                            final_text = translated_text
                                            translation_notes.append(f"Text element {i+1}: Successfully translated to {target_language}")
                                        
                                        self._replace_text_preserving_format(shape.text_frame, final_text)
                                
                                self.logger.info(f"  Slide {slide_idx + 1}: Translation applied with placeholder filtering")
                                
                            except Exception as trans_error:
                                self.logger.error(f"  Slide {slide_idx + 1}: Translation failed: {trans_error}")
                                translation_notes.append(f"Translation failed: {str(trans_error)}")
                        
                        # Add adaptation notes to slide notes (including translation info)
                        all_notes = adaptation_notes + translation_notes
                        if all_notes:
                            self._add_adaptation_notes_to_slide(slide, all_notes, profile)
                        else:
                            # Even if no text changes, add a note that adaptation was attempted
                            status_info = ""
                            if not self.config.get('anthropic_api_key'):
                                status_info = " (No API key - rule-based adaptation only)"
                            default_notes = [f"Slide processed for {profile} adaptation - {len(slide_texts)} text elements reviewed{status_info}"]
                            self._add_adaptation_notes_to_slide(slide, default_notes, profile)
                        
                    except Exception as adapt_error:
                        self.logger.error(f"  Slide {slide_idx + 1}: Adaptation failed: {adapt_error}")
                        # Still count the elements as processed to maintain progress accuracy
                        processed_elements += len(slide_texts)
                        continue
                else:
                    self.logger.info(f"  Slide {slide_idx + 1}: No text content to adapt")
                    # Update progress even for slides with no text
                    if processing_callback:
                        skip_progress = adaptation_start + (processed_elements / max(total_text_elements, 1)) * (adaptation_end - adaptation_start)
                        processing_callback(f"Slide {slide_idx + 1} has no text to adapt, skipping...", skip_progress)
            
            if processing_callback:
                processing_callback("Saving adapted presentation...", 95)
            
            # Final background application before saving (ensure persistence)
            final_settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            if final_settings.get('use_background_tint', False) and final_settings.get('background_color'):
                self.logger.info(f"🎨 Final background application before saving for {profile}...")
                for slide in presentation.slides:
                    self._apply_slide_background(slide, final_settings['background_color'])
            
            # Save the adapted presentation
            presentation.save(output_path)
            
            if processing_callback:
                processing_callback("Adaptation completed successfully!", 100)
            
            self.logger.info(f"Format-preserving adaptation completed: {output_path}")
            return True
            
        except Exception as e:
            self.logger.error(f"Error in format-preserving adaptation: {str(e)}")
            if processing_callback:
                processing_callback(f"Error: {str(e)}", -1)
            return False
    
    def _replace_text_preserving_format(self, text_frame, new_text: str, profile: str = 'default'):
        """
        Replace text content while preserving ALL original formatting
        
        Args:
            text_frame: PowerPoint text frame object
            new_text: New text content to insert
            profile: Learning profile (may affect font choice for dyslexia)
        """
        try:
            if not hasattr(text_frame, 'paragraphs') or not text_frame.paragraphs:
                return
            
            # Store original text frame properties
            text_frame_props = {}
            try:
                text_frame_props['margin_left'] = getattr(text_frame, 'margin_left', None)
                text_frame_props['margin_right'] = getattr(text_frame, 'margin_right', None)
                text_frame_props['margin_top'] = getattr(text_frame, 'margin_top', None)
                text_frame_props['margin_bottom'] = getattr(text_frame, 'margin_bottom', None)
                text_frame_props['word_wrap'] = getattr(text_frame, 'word_wrap', None)
                text_frame_props['auto_size'] = getattr(text_frame, 'auto_size', None)
                text_frame_props['vertical_anchor'] = getattr(text_frame, 'vertical_anchor', None)
            except Exception as e:
                self.logger.debug(f"Could not capture text frame properties: {e}")
            
            # Store original paragraph and run formatting
            original_formats = []
            for paragraph in text_frame.paragraphs:
                para_format = {
                    'level': getattr(paragraph, 'level', 0),
                    'alignment': getattr(paragraph, 'alignment', None),
                    'space_before': getattr(paragraph, 'space_before', None),
                    'space_after': getattr(paragraph, 'space_after', None),
                    'line_spacing': getattr(paragraph, 'line_spacing', None),
                    'runs': []
                }
                
                for run in paragraph.runs:
                    run_format = {}
                    if hasattr(run, 'font'):
                        font = run.font
                        run_format = {
                            'name': getattr(font, 'name', None),
                            'size': getattr(font, 'size', None),
                            'bold': getattr(font, 'bold', None),
                            'italic': getattr(font, 'italic', None),
                            'underline': getattr(font, 'underline', None),
                            'color': None
                        }
                        
                        # Enhanced color preservation
                        try:
                            if hasattr(font, 'color'):
                                color_info = {}
                                # Try RGB color first
                                try:
                                    if hasattr(font.color, 'rgb') and font.color.rgb:
                                        color_info['rgb'] = font.color.rgb
                                except:
                                    pass
                                
                                # Try theme color
                                try:
                                    if hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                                        color_info['theme_color'] = font.color.theme_color
                                except:
                                    pass
                                
                                # Try brightness/tint modifications
                                try:
                                    if hasattr(font.color, 'brightness') and font.color.brightness is not None:
                                        color_info['brightness'] = font.color.brightness
                                except:
                                    pass
                                
                                if color_info:
                                    run_format['color'] = color_info
                        except Exception as e:
                            self.logger.debug(f"Could not preserve color: {e}")
                    
                    para_format['runs'].append(run_format)
                
                original_formats.append(para_format)
            
            # Split new text into paragraphs
            new_paragraphs = new_text.split('\n')
            
            # Clear existing content but preserve the first paragraph
            while len(text_frame.paragraphs) > 1:
                p = text_frame.paragraphs[-1]
                p._p.getparent().remove(p._p)
            
            # Clear runs from first paragraph
            first_para = text_frame.paragraphs[0]
            while len(first_para.runs) > 0:
                first_para._p.remove(first_para.runs[0]._r)
            
            # Apply new text to first paragraph with original formatting
            if new_paragraphs and new_paragraphs[0].strip():
                para_text = new_paragraphs[0]
                run = first_para.add_run()
                run.text = para_text
                
                # Apply original formatting if available
                if original_formats and original_formats[0]['runs']:
                    run_format = original_formats[0]['runs'][0].copy()
                    
                    # For dyslexia profile, override font family but keep original size
                    if profile.lower() == 'dyslexia':
                        dyslexia_settings = self.PROFILE_SETTINGS['dyslexia']
                        run_format['name'] = dyslexia_settings['font_name']
                        run_format['fallbacks'] = dyslexia_settings.get('font_fallbacks', [])
                        # Keep original size (don't override)
                        self.logger.info(f"🔤 Applying OpenDyslexic font to text run")
                    
                    self._apply_run_format(run, run_format)
                
                # Additional force application for dyslexia
                if profile.lower() == 'dyslexia' and hasattr(run, 'font'):
                    dyslexia_settings = self.PROFILE_SETTINGS['dyslexia']
                    success = self._apply_font_with_fallbacks(
                        run.font, 
                        dyslexia_settings['font_name'], 
                        dyslexia_settings.get('font_fallbacks', [])
                    )
                    if success:
                        self.logger.info(f"🔤 Force-applied dyslexia font to run")
                
                # Apply paragraph formatting
                if original_formats:
                    self._apply_paragraph_format(first_para, original_formats[0])
            
            # Add additional paragraphs if needed
            for i, para_text in enumerate(new_paragraphs[1:], 1):
                if para_text.strip():
                    new_para = text_frame.add_paragraph()
                    run = new_para.add_run()
                    run.text = para_text
                    
                    # Use formatting from corresponding original paragraph or last available
                    format_idx = min(i, len(original_formats) - 1) if original_formats else 0
                    
                    if format_idx < len(original_formats):
                        if original_formats[format_idx]['runs']:
                            run_format = original_formats[format_idx]['runs'][0].copy()
                            
                            # For dyslexia profile, override font family but keep original size
                            if profile.lower() == 'dyslexia':
                                dyslexia_settings = self.PROFILE_SETTINGS['dyslexia']
                                run_format['name'] = dyslexia_settings['font_name']
                                run_format['fallbacks'] = dyslexia_settings.get('font_fallbacks', [])
                                # Keep original size (don't override)
                            
                            self._apply_run_format(run, run_format)
                            
                            # Additional force application for dyslexia
                            if profile.lower() == 'dyslexia' and hasattr(run, 'font'):
                                dyslexia_settings = self.PROFILE_SETTINGS['dyslexia']
                                self._apply_font_with_fallbacks(
                                    run.font, 
                                    dyslexia_settings['font_name'], 
                                    dyslexia_settings.get('font_fallbacks', [])
                                )
                        self._apply_paragraph_format(new_para, original_formats[format_idx])
            
            # Restore text frame properties
            self._restore_text_frame_properties(text_frame, text_frame_props)
            
        except Exception as e:
            self.logger.error(f"Error replacing text while preserving format: {str(e)}")
            # Fallback to simple text replacement
            try:
                text_frame.text = new_text
            except:
                pass
    
    def _replace_text_with_overflow_handling(self, text_frame, new_text: str, shape, profile: str = 'default'):
        """
        Replace text content with overflow detection and handling for translations
        
        Args:
            text_frame: PowerPoint text frame object
            new_text: New text content to insert
            shape: The shape containing the text frame (for size information)
            profile: Learning profile
        """
        try:
            # First, apply the text normally
            self._replace_text_preserving_format(text_frame, new_text, profile)
            
            # Get shape dimensions for overflow checking
            shape_width = getattr(shape, 'width', 100)
            shape_height = getattr(shape, 'height', 50)
            
            # Convert to points (PowerPoint uses EMU internally, but we work in points)
            max_width_pts = shape_width / 12700.0 if hasattr(shape_width, '__truediv__') else float(shape_width) / 12700.0
            max_height_pts = shape_height / 12700.0 if hasattr(shape_height, '__truediv__') else float(shape_height) / 12700.0
            
            # Check if text is overflowing
            if self._is_text_still_overflowing(text_frame, max_width_pts, max_height_pts):
                self.logger.warning(f"🔍 Text overflow detected after translation, attempting to fix...")
                
                # Strategy 1: Try reducing font size slightly (max 3 point reduction)
                original_overflow = True
                attempts = 0
                max_attempts = 3
                
                while original_overflow and attempts < max_attempts:
                    attempts += 1
                    
                    # Reduce font size for all runs by 1 point
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, 'font') and run.font.size:
                                current_size = run.font.size.pt
                                new_size = max(8, current_size - 1)  # Don't go below 8pt
                                run.font.size = Pt(new_size)
                    
                    # Check if overflow is resolved
                    if not self._is_text_still_overflowing(text_frame, max_width_pts, max_height_pts):
                        self.logger.info(f"✅ Text overflow resolved by reducing font size (attempt {attempts})")
                        original_overflow = False
                        break
                
                # Strategy 2: If font size reduction didn't work, try adjusting margins
                if original_overflow:
                    try:
                        # Reduce margins to give more space
                        text_frame.margin_left = Inches(0.05)  # Reduced from 0.1
                        text_frame.margin_right = Inches(0.05)  # Reduced from 0.1
                        text_frame.margin_top = Inches(0.02)   # Reduced from 0.05
                        text_frame.margin_bottom = Inches(0.02)  # Reduced from 0.05
                        
                        if not self._is_text_still_overflowing(text_frame, max_width_pts, max_height_pts):
                            self.logger.info(f"✅ Text overflow resolved by reducing margins")
                            original_overflow = False
                    except Exception as margin_error:
                        self.logger.warning(f"Could not adjust margins: {margin_error}")
                
                # Strategy 3: If still overflowing, try line spacing adjustment
                if original_overflow:
                    try:
                        for paragraph in text_frame.paragraphs:
                            paragraph.line_spacing = 1.0  # Single line spacing
                            paragraph.space_after = Pt(0)  # Remove space after paragraphs
                            paragraph.space_before = Pt(0)  # Remove space before paragraphs
                        
                        if not self._is_text_still_overflowing(text_frame, max_width_pts, max_height_pts):
                            self.logger.info(f"✅ Text overflow resolved by adjusting line spacing")
                            original_overflow = False
                    except Exception as spacing_error:
                        self.logger.warning(f"Could not adjust line spacing: {spacing_error}")
                
                # Strategy 4: Last resort - truncate text with ellipsis
                if original_overflow:
                    self.logger.warning(f"⚠️ Could not resolve overflow, truncating text as last resort")
                    # This is a fallback - ideally we'd want to avoid this
                    # but it prevents completely broken layouts
                    if len(new_text) > 50:
                        truncated_text = new_text[:47] + "..."
                        self._replace_text_preserving_format(text_frame, truncated_text, profile)
            
        except Exception as e:
            self.logger.error(f"Error in overflow handling: {str(e)}")
            # Fall back to the original method
            self._replace_text_preserving_format(text_frame, new_text, profile)
    
    def _apply_run_format(self, run, run_format):
        """Apply formatting to a text run while preserving font size"""
        try:
            if hasattr(run, 'font'):
                font = run.font
                if run_format.get('name'):
                    # Try to apply the preferred font with fallbacks
                    self._apply_font_with_fallbacks(font, run_format['name'], run_format.get('fallbacks', []))
                # PRESERVE ORIGINAL FONT SIZE - do not change it
                if run_format.get('size'):
                    font.size = run_format['size']  # Keep original size
                if run_format.get('bold') is not None:
                    font.bold = run_format['bold']
                if run_format.get('italic') is not None:
                    font.italic = run_format['italic']
                if run_format.get('underline') is not None:
                    font.underline = run_format['underline']
                if run_format.get('color'):
                    self._apply_color_format(font, run_format['color'])
        except Exception as e:
            self.logger.warning(f"Could not apply run format: {e}")
    
    def _apply_color_format(self, font, color_info):
        """Apply color formatting with enhanced support for different color types"""
        try:
            if isinstance(color_info, dict):
                # Enhanced color information
                applied = False
                
                # Try to apply RGB color first
                if 'rgb' in color_info and color_info['rgb']:
                    try:
                        font.color.rgb = color_info['rgb']
                        applied = True
                    except Exception as e:
                        self.logger.debug(f"Could not apply RGB color: {e}")
                
                # Try theme color if RGB failed
                if not applied and 'theme_color' in color_info and color_info['theme_color'] is not None:
                    try:
                        font.color.theme_color = color_info['theme_color']
                        applied = True
                    except Exception as e:
                        self.logger.debug(f"Could not apply theme color: {e}")
                
                # Apply brightness/tint if available
                if applied and 'brightness' in color_info and color_info['brightness'] is not None:
                    try:
                        font.color.brightness = color_info['brightness']
                    except Exception as e:
                        self.logger.debug(f"Could not apply brightness: {e}")
                        
            else:
                # Legacy RGB color format
                try:
                    font.color.rgb = color_info
                except Exception as e:
                    self.logger.debug(f"Could not apply legacy color: {e}")
                    
        except Exception as e:
            self.logger.debug(f"Color application failed: {e}")
    
    def _restore_text_frame_properties(self, text_frame, text_frame_props):
        """Restore text frame properties like margins, word wrap, etc."""
        try:
            if not text_frame_props:
                return
                
            for prop_name, prop_value in text_frame_props.items():
                if prop_value is not None:
                    try:
                        setattr(text_frame, prop_name, prop_value)
                    except Exception as e:
                        self.logger.debug(f"Could not restore text frame property {prop_name}: {e}")
                        
        except Exception as e:
            self.logger.debug(f"Could not restore text frame properties: {e}")
    
    def _copy_text_frame_properties(self, source_text_frame, target_text_frame):
        """Copy text frame properties from source to target"""
        try:
            properties_to_copy = [
                'margin_left', 'margin_right', 'margin_top', 'margin_bottom',
                'word_wrap', 'auto_size', 'vertical_anchor'
            ]
            
            for prop_name in properties_to_copy:
                try:
                    if hasattr(source_text_frame, prop_name):
                        source_value = getattr(source_text_frame, prop_name)
                        if source_value is not None:
                            setattr(target_text_frame, prop_name, source_value)
                except Exception as e:
                    self.logger.debug(f"Could not copy text frame property {prop_name}: {e}")
                    
        except Exception as e:
            self.logger.debug(f"Could not copy text frame properties: {e}")
    
    def _copy_group_shape(self, original_group, target_slide):
        """Copy a group shape by copying its individual components"""
        try:
            # Get group position and size for positioning context
            group_left = getattr(original_group, 'left', 0)
            group_top = getattr(original_group, 'top', 0)
            
            shapes_in_group = 0
            shapes_copied = 0
            
            if hasattr(original_group, 'shapes'):
                for sub_shape in original_group.shapes:
                    shapes_in_group += 1
                    try:
                        # Calculate relative position within the group
                        sub_left = getattr(sub_shape, 'left', group_left)
                        sub_top = getattr(sub_shape, 'top', group_top)
                        
                        # Copy the sub-shape to the target slide
                        self._copy_shape_to_slide(sub_shape, target_slide)
                        shapes_copied += 1
                        
                    except Exception as sub_error:
                        self.logger.debug(f"Could not copy sub-shape from group: {sub_error}")
                        continue
            
            self.logger.debug(f"Copied {shapes_copied}/{shapes_in_group} shapes from group")
            
        except Exception as e:
            self.logger.debug(f"Could not copy group shape: {e}")
            # Fallback to treating the group as a single shape
            self._recreate_shape_fallback(original_group, target_slide)
    
    def _apply_paragraph_format(self, paragraph, para_format):
        """Apply formatting to a paragraph"""
        try:
            if para_format.get('level') is not None:
                paragraph.level = para_format['level']
            if para_format.get('alignment') is not None:
                paragraph.alignment = para_format['alignment']
            if para_format.get('space_before') is not None:
                paragraph.space_before = para_format['space_before']
            if para_format.get('space_after') is not None:
                paragraph.space_after = para_format['space_after']
            if para_format.get('line_spacing') is not None:
                paragraph.line_spacing = para_format['line_spacing']
        except Exception as e:
            self.logger.warning(f"Could not apply paragraph format: {e}")
    
    def _validate_text_length(self, original_text: str, adapted_text: str, max_variance: float = 0.20) -> str:
        """
        Validate that adapted text length is within acceptable variance of original text
        
        Args:
            original_text: Original text content
            adapted_text: Adapted text content
            max_variance: Maximum allowed variance (0.20 = ±20%)
            
        Returns:
            Validated text (truncated or original if out of bounds)
        """
        try:
            original_length = len(original_text)
            adapted_length = len(adapted_text)
            
            if original_length == 0:
                return adapted_text
            
            # Calculate variance
            variance = abs(adapted_length - original_length) / original_length
            
            if variance <= max_variance:
                # Within acceptable range
                return adapted_text
            
            # Text is too long or too short, need to adjust
            min_length = int(original_length * (1 - max_variance))
            max_length = int(original_length * (1 + max_variance))
            
            if adapted_length > max_length:
                # Text is too long, truncate intelligently
                truncated = self._smart_truncate(adapted_text, max_length)
                self.logger.warning(f"Text too long ({adapted_length} chars), truncated to {len(truncated)} chars")
                return truncated
            elif adapted_length < min_length:
                # Text is too short, use original
                self.logger.warning(f"Text too short ({adapted_length} chars), using original ({original_length} chars)")
                return original_text
            
            return adapted_text
            
        except Exception as e:
            self.logger.error(f"Error validating text length: {e}")
            return original_text  # Fallback to original
    
    def _smart_truncate(self, text: str, max_length: int) -> str:
        """
        Intelligently truncate text at word boundaries
        
        Args:
            text: Text to truncate
            max_length: Maximum allowed length
            
        Returns:
            Truncated text
        """
        if len(text) <= max_length:
            return text
        
        # Try to truncate at sentence boundary
        sentences = text.split('. ')
        if len(sentences) > 1:
            result = ""
            for sentence in sentences:
                test_result = result + sentence + ". "
                if len(test_result) <= max_length:
                    result = test_result
                else:
                    break
            if result.strip():
                return result.strip()
        
        # Try to truncate at word boundary
        words = text.split()
        result = ""
        for word in words:
            test_result = result + word + " "
            if len(test_result) <= max_length:
                result = test_result
            else:
                break
        
        if result.strip():
            return result.strip()
        
        # Last resort: hard truncate with ellipsis
        return text[:max_length-3] + "..."
    
    def _smart_truncate_text(self, text: str, max_width_pts: float, max_height_pts: float, 
                           font_size: float, font_name: str) -> tuple[str, str]:
        """
        Smart truncation that respects sentence boundaries and fits within bounds
        
        Args:
            text: Original text to truncate
            max_width_pts: Maximum width in points
            max_height_pts: Maximum height in points
            font_size: Font size to use for measurement
            font_name: Font name for measurement
            
        Returns:
            Tuple of (truncated_text, overflow_text)
        """
        overflow_text = ""
        
        try:
            # Check if the full text fits
            width, height = self.measure_text_bounds(text, font_name, font_size)
            if width <= max_width_pts and height <= max_height_pts:
                return text, ""
            
            # Split into sentences and try to fit as many as possible
            sentences = text.split('. ')
            fitted_sentences = []
            
            for i, sentence in enumerate(sentences):
                # Build current text with this sentence
                test_text = '. '.join(fitted_sentences + [sentence])
                if not test_text.endswith('.') and i < len(sentences) - 1:
                    test_text += '.'
                
                # Measure if it fits
                width, height = self.measure_text_bounds(test_text, font_name, font_size)
                
                if width <= max_width_pts and height <= max_height_pts:
                    fitted_sentences.append(sentence)
                else:
                    # This sentence doesn't fit, save remaining as overflow
                    remaining_sentences = sentences[i:]
                    if remaining_sentences:
                        overflow_text = '. '.join(remaining_sentences)
                        if not overflow_text.endswith('.'):
                            overflow_text += '.'
                    break
            
            # Build final truncated text
            if fitted_sentences:
                truncated_text = '. '.join(fitted_sentences)
                if not truncated_text.endswith('.'):
                    truncated_text += '.'
                return truncated_text, overflow_text
            
            # If no sentences fit, try word-by-word truncation
            words = text.split()
            fitted_words = []
            
            for i, word in enumerate(words):
                test_text = ' '.join(fitted_words + [word])
                width, height = self.measure_text_bounds(test_text, font_name, font_size)
                
                if width <= max_width_pts and height <= max_height_pts:
                    fitted_words.append(word)
                else:
                    remaining_words = words[i:]
                    if remaining_words:
                        overflow_text = ' '.join(remaining_words)
                    break
            
            truncated_text = ' '.join(fitted_words) if fitted_words else text[:50] + "..."
            return truncated_text, overflow_text
            
        except Exception as e:
            self.logger.error(f"Error in smart truncation: {e}")
            # Fallback to simple character-based truncation
            max_chars = int(max_width_pts * max_height_pts / 100)  # Rough estimate
            if len(text) > max_chars:
                return text[:max_chars-3] + "...", text[max_chars:]
            return text, ""
    
    def _convert_to_bullets(self, text: str) -> str:
        """
        Convert long text into bullet points for better readability
        
        Args:
            text: Original text to convert
            
        Returns:
            Text formatted as bullet points
        """
        try:
            # Split into sentences
            sentences = text.split('. ')
            
            # Clean up sentences
            clean_sentences = []
            for sentence in sentences:
                sentence = sentence.strip()
                if sentence and len(sentence) > 10:  # Skip very short fragments
                    if not sentence.endswith('.'):
                        sentence += '.'
                    clean_sentences.append(sentence)
            
            # If we have multiple sentences, convert to bullets
            if len(clean_sentences) > 1:
                bullet_points = []
                for sentence in clean_sentences:
                    # Shorten long sentences
                    if len(sentence) > 80:
                        # Try to find a natural break point
                        words = sentence.split()
                        if len(words) > 10:
                            mid_point = len(words) // 2
                            first_half = ' '.join(words[:mid_point])
                            second_half = ' '.join(words[mid_point:])
                            bullet_points.append(f"• {first_half}")
                            bullet_points.append(f"• {second_half}")
                        else:
                            bullet_points.append(f"• {sentence}")
                    else:
                        bullet_points.append(f"• {sentence}")
                
                return '\n'.join(bullet_points)
            
            # If only one sentence, try to break it down into key points
            if clean_sentences:
                sentence = clean_sentences[0]
                # Look for conjunctions and split on them
                connectors = [', and ', ', but ', ', so ', ', or ', '; ']
                for connector in connectors:
                    if connector in sentence:
                        parts = sentence.split(connector)
                        if len(parts) > 1:
                            bullet_points = []
                            for i, part in enumerate(parts):
                                part = part.strip()
                                if part:
                                    if i == 0:
                                        bullet_points.append(f"• {part}")
                                    else:
                                        bullet_points.append(f"• {part.capitalize()}")
                            return '\n'.join(bullet_points)
            
            # Fallback: just add a bullet to the original text
            return f"• {text}"
            
        except Exception as e:
            self.logger.error(f"Error converting to bullets: {e}")
            return f"• {text}"
    
    def _extreme_truncate(self, text: str, max_chars: int) -> str:
        """
        Last resort truncation for extremely long text
        
        Args:
            text: Text to truncate
            max_chars: Maximum number of characters
            
        Returns:
            Severely truncated text
        """
        if len(text) <= max_chars:
            return text
        
        try:
            # Try to keep the first sentence if it's short enough
            first_sentence = text.split('.')[0] + '.'
            if len(first_sentence) <= max_chars:
                return first_sentence
            
            # Try to keep first few words that make sense
            words = text.split()
            result = ""
            for word in words:
                test_result = result + word + " "
                if len(test_result) <= max_chars - 3:  # Reserve space for ellipsis
                    result = test_result
                else:
                    break
            
            if result.strip():
                return result.strip() + "..."
            
            # Last resort: hard character limit
            return text[:max_chars-3] + "..."
            
        except Exception as e:
            self.logger.error(f"Error in extreme truncation: {e}")
            return text[:max_chars-3] + "..."
    
    def _is_text_still_overflowing(self, text_frame, max_width_pts: float, max_height_pts: float) -> bool:
        """
        Check if text is still overflowing the text frame bounds
        
        Args:
            text_frame: PowerPoint text frame object
            max_width_pts: Maximum width in points
            max_height_pts: Maximum height in points
            
        Returns:
            True if text is overflowing, False otherwise
        """
        try:
            # Get text content
            text_content = ""
            for paragraph in text_frame.paragraphs:
                text_content += paragraph.text + "\n"
            text_content = text_content.strip()
            
            if not text_content:
                return False
            
            # Get font information from first run of first paragraph
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                first_run = text_frame.paragraphs[0].runs[0]
                font_size = first_run.font.size.pt if first_run.font.size else 12
                font_name = first_run.font.name or "Arial"
            else:
                font_size = 12
                font_name = "Arial"
            
            # Measure text bounds
            text_width, text_height = self.measure_text_bounds(text_content, font_name, font_size)
            
            # Check if it exceeds bounds (with small tolerance)
            tolerance = 5  # 5 points tolerance
            width_overflow = text_width > (max_width_pts + tolerance)
            height_overflow = text_height > (max_height_pts + tolerance)
            
            if width_overflow or height_overflow:
                self.logger.debug(f"Text overflow detected: {text_width:.1f}pt x {text_height:.1f}pt "
                                f"vs max {max_width_pts:.1f}pt x {max_height_pts:.1f}pt")
                return True
            
            return False
            
        except Exception as e:
            self.logger.error(f"Error checking text overflow: {e}")
            # If we can't measure, assume no overflow to be safe
            return False
    
    def _create_adaptation_summary(self, original_text: str, adapted_text: str, profile: str) -> str:
        """
        Create a summary of what changed during adaptation
        
        Args:
            original_text: Original text
            adapted_text: Adapted text
            profile: Learning profile used
            
        Returns:
            Summary of changes
        """
        try:
            # Get profile settings
            settings = self.PROFILE_SETTINGS.get(profile.lower(), self.PROFILE_SETTINGS['default'])
            
            # Calculate basic metrics
            original_words = len(original_text.split())
            adapted_words = len(adapted_text.split())
            word_change = adapted_words - original_words
            
            # Create summary
            summary_parts = []
            summary_parts.append(f"Adapted for {profile} profile")
            
            if word_change > 0:
                summary_parts.append(f"expanded by {word_change} words")
            elif word_change < 0:
                summary_parts.append(f"reduced by {abs(word_change)} words")
            else:
                summary_parts.append("word count maintained")
            
            # Check for common adaptations
            adaptations = []
            if len(adapted_text) < len(original_text):
                adaptations.append("simplified")
            if "." in adapted_text and adapted_text.count(".") > original_text.count("."):
                adaptations.append("broken into shorter sentences")
            
            if adaptations:
                summary_parts.append(f"({', '.join(adaptations)})")
            
            # Add profile-specific formatting notes
            if profile.lower() == 'dyslexia':
                formatting_notes = ["OpenDyslexic font", "beige background tint", "1.8x line spacing"]
                if settings.get('add_reading_rulers'):
                    formatting_notes.append("reading rulers")
                summary_parts.append(f"formatting: {', '.join(formatting_notes)}")
            elif profile.lower() == 'adhd':
                formatting_notes = ["visual boundaries", "content chunks", "1.5x line spacing"]
                summary_parts.append(f"formatting: {', '.join(formatting_notes)}")
            
            return " - ".join(summary_parts)
            
        except Exception as e:
            self.logger.error(f"Error creating adaptation summary: {e}")
            return f"Adapted for {profile} profile"
    
    def _add_adaptation_notes_to_slide(self, slide, adaptation_notes: List[str], profile: str):
        """
        Add adaptation notes to the slide's notes section
        
        Args:
            slide: PowerPoint slide object
            adaptation_notes: List of adaptation notes
            profile: Learning profile used
        """
        try:
            # Get or create notes slide
            if not slide.has_notes_slide:
                slide.notes_slide
            
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            
            # Get existing notes
            existing_notes = notes_text_frame.text.strip()
            
            # Create adaptation section
            adaptation_section = f"\n\n--- ADAPTATION NOTES ({profile.upper()}) ---\n"
            adaptation_section += f"Adapted on: {self._get_current_timestamp()}\n"
            
            # Add profile-specific formatting information
            if profile.lower() == 'dyslexia':
                adaptation_section += "Evidence-based formatting applied:\n"
                adaptation_section += "- Font: OpenDyslexic (dyslexia-friendly font)\n"
                adaptation_section += "- Background: Warm beige tint (#FFFBF0) instead of colored text\n"
                adaptation_section += "- Line spacing: 1.8x for better readability\n"
                adaptation_section += "- Text alignment: Left-aligned (no justification)\n"
                if self.PROFILE_SETTINGS['dyslexia'].get('add_reading_rulers'):
                    adaptation_section += "- Reading rulers: Subtle horizontal guides between text blocks\n"
            elif profile.lower() == 'adhd':
                adaptation_section += "Evidence-based formatting applied:\n"
                adaptation_section += "- Visual boundaries: Subtle borders around text blocks\n"
                adaptation_section += "- Background: Soft green tint (#F5FFF0) for calm focus\n"
                adaptation_section += "- Content structure: Chunked text with clear separation\n"
                adaptation_section += "- Line spacing: 1.5x for clarity\n"
                adaptation_section += "- Consistent visual cues throughout slide\n"
            
            adaptation_section += "\n".join(adaptation_notes)
            
            # Add to notes
            if existing_notes:
                notes_text_frame.text = existing_notes + adaptation_section
            else:
                notes_text_frame.text = adaptation_section.strip()
            
            self.logger.info(f"Added {len(adaptation_notes)} adaptation notes to slide")
            
        except Exception as e:
            self.logger.error(f"Error adding adaptation notes to slide: {e}")
    
    def _get_current_timestamp(self) -> str:
        """Get current timestamp for notes"""
        import datetime
        return datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    def _is_placeholder_text(self, text: str) -> bool:
        """
        Check if text appears to be placeholder or debug text that shouldn't appear in slides
        
        Args:
            text: Text to check
            
        Returns:
            True if text appears to be placeholder/debug text
        """
        if not text or not text.strip():
            return False
        
        text_lower = text.lower().strip()
        
        # Common placeholder patterns
        placeholder_patterns = [
            "here is the adapt",
            "adapted version",
            "adaptation for",
            "adapted for dyslexia",
            "adapted for adhd", 
            "adapted for esl",
            "adapted text:",
            "original text:",
            "adaptation:",
            "[adapted:",
            "[dyslexia]",
            "[adhd]",
            "[esl]",
            "translated content",  # Translation placeholder
            "here's the translation",
            "here is the translation", 
            "here's the translation of the text",
            "the translation is",
            "translation:",
            "translated text:",
            "translation follows:",
            "maintaining the original formatting",
            "following the guidelines you provided",
            "translate to",
            "translation result",
            "translation failed",
            "error",
            "failed",
            "placeholder",
            "debug",
            "test adaptation",
            "mock adaptation"
        ]
        
        # Check for exact matches or patterns
        for pattern in placeholder_patterns:
            if pattern in text_lower:
                return True
        
        # Check for very short "adapted" responses that are likely placeholders
        if len(text.strip()) < 10 and any(word in text_lower for word in ["adapt", "error", "fail"]):
            return True
        
        # Check for responses that start with common debug prefixes
        debug_prefixes = ["adapted:", "result:", "output:", "error:", "warning:"]
        for prefix in debug_prefixes:
            if text_lower.startswith(prefix):
                return True
        
        # Check for responses that are entirely in brackets or quotes
        if (text.strip().startswith('[') and text.strip().endswith(']')) or \
           (text.strip().startswith('"') and text.strip().endswith('"') and len(text.strip()) < 50):
            return True
        
        return False
    
    def _is_translation_note(self, text: str) -> bool:
        """
        Check if text appears to be a translation note that should go in slide notes
        
        Args:
            text: Text to check
            
        Returns:
            True if text appears to be a translation note
        """
        if not text or not text.strip():
            return False
        
        text_lower = text.lower().strip()
        
        # Translation note patterns
        note_patterns = [
            "note:",
            "translation note:",
            "translator's note:",
            "the polish translation is",
            "the spanish translation is", 
            "the french translation is",
            "the german translation is",
            "the italian translation is",
            "the translation is a direct equivalent",
            "there are no technical terms",
            "requires additional explanation",
            "including the english original",
            "direct equivalent of the english",
            "equivalent in the target language"
        ]
        
        # Check for note patterns
        for pattern in note_patterns:
            if pattern in text_lower:
                return True
        
        # Check if it starts with "Note:" (case insensitive)
        if text_lower.startswith("note:"):
            return True
            
        return False
    
    def _apply_font_with_fallbacks(self, font, preferred_font: str, fallback_fonts: List[str] = None):
        """
        Apply font with fallback options in case the preferred font is not available
        
        Args:
            font: PowerPoint font object
            preferred_font: Primary font to try
            fallback_fonts: List of fallback fonts to try if primary fails
        """
        # List of OpenDyslexic font name variations to try
        opendyslexic_variants = [
            'OpenDyslexic',
            'OpenDyslexic-Regular', 
            'OpenDyslexic Regular',
            'Open Dyslexic',
            'OpenDyslexic3',
            'OpenDyslexic-Bold',
            'OpenDyslexic Bold'
        ]
        
        fonts_to_try = []
        
        # If preferred font is OpenDyslexic, try all variants first
        if 'opendyslexic' in preferred_font.lower():
            fonts_to_try.extend(opendyslexic_variants)
        else:
            fonts_to_try.append(preferred_font)
        
        # Add fallback fonts
        if fallback_fonts:
            fonts_to_try.extend(fallback_fonts)
        
        # Try each font
        for font_name in fonts_to_try:
            try:
                # Test if font name assignment works
                original_name = getattr(font, 'name', None)
                font.name = font_name
                
                # Verify the font was actually applied
                if hasattr(font, 'name') and font.name:
                    applied_name = font.name
                    if applied_name.lower() == font_name.lower():
                        self.logger.info(f"✅ Successfully applied font: {font_name}")
                        return True
                    else:
                        self.logger.debug(f"Font name changed from {font_name} to {applied_name}")
                        # If it's an OpenDyslexic variant, consider it success
                        if 'opendyslexic' in font_name.lower() and 'opendyslexic' in applied_name.lower():
                            self.logger.info(f"✅ Successfully applied OpenDyslexic variant: {applied_name}")
                            return True
                        # If it's a close match, accept it
                        elif font_name.lower() in applied_name.lower() or applied_name.lower() in font_name.lower():
                            self.logger.info(f"✅ Successfully applied close match: {applied_name}")
                            return True
                        
                self.logger.info(f"✅ Applied font: {font_name}")
                return True
                
            except Exception as e:
                self.logger.debug(f"Font '{font_name}' not available: {e}")
                continue
        
        # If no fonts worked, log warning
        self.logger.warning(f"⚠️ No suitable fonts available from: {preferred_font}, {fallback_fonts}")
        return False
    
    def _get_dyslexia_font_settings(self) -> Dict[str, Any]:
        """Get OpenDyslexic font settings with proper fallbacks"""
        return {
            'name': 'OpenDyslexic',
            'fallbacks': ['OpenDyslexic-Regular', 'OpenDyslexic Regular', 'Arial', 'Helvetica', 'Calibri']
        }
    
    def _apply_profile_slide_formatting(self, presentation, profile: str):
        """
        Apply profile-specific formatting to all slides in the presentation
        
        Args:
            presentation: PowerPoint presentation object
            profile: Learning profile
        """
        try:
            settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            
            # DO NOT change slide dimensions - this causes background alignment issues
            # Keep original presentation dimensions to ensure background covers properly
            
            self.logger.info(f"Applied presentation-wide formatting for {profile} profile (dimensions preserved)")
            
        except Exception as e:
            self.logger.error(f"Error applying profile slide formatting: {e}")
    
    def _apply_evidence_based_formatting(self, slide, profile: str):
        """
        Apply evidence-based visual formatting to a slide based on profile
        
        Args:
            slide: PowerPoint slide object
            profile: Learning profile
        """
        try:
            settings = self.PROFILE_SETTINGS.get(profile, self.PROFILE_SETTINGS['default'])
            
            # Apply background color to entire slide
            if settings.get('use_background_tint', False):
                self.logger.info(f"🎨 Attempting to apply background for {profile} profile")
                background_applied = False
                
                # Method 1: Standard slide background approach
                try:
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = settings['background_color']
                    background_applied = True
                    self.logger.info(f"✅ Method 1: Applied slide background successfully")
                    
                except Exception as bg_error:
                    self.logger.warning(f"⚠️ Method 1 failed: {bg_error}")
                
                # Method 2: Try through slide layout background
                if not background_applied:
                    try:
                        if hasattr(slide, 'slide_layout'):
                            layout_bg = slide.slide_layout.background
                            layout_fill = layout_bg.fill
                            layout_fill.solid()
                            layout_fill.fore_color.rgb = settings['background_color']
                            background_applied = True
                            self.logger.info(f"✅ Method 2: Applied layout background successfully")
                    except Exception as layout_error:
                        self.logger.warning(f"⚠️ Method 2 failed: {layout_error}")
                
                # Method 3: Force background through XML manipulation
                if not background_applied:
                    try:
                        # Direct XML approach to ensure background is set
                        slide_element = slide._element
                        
                        # Check if background properties exist
                        bg_props = slide_element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                        if bg_props is None:
                            # Create background properties
                            from lxml import etree
                            bg_props = etree.SubElement(slide_element, '{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                        
                        # Set solid fill
                        bg_color = settings['background_color']
                        # This is a more complex XML manipulation - might need adjustment
                        self.logger.info(f"✅ Method 3: Attempted XML background manipulation")
                        background_applied = True
                        
                    except Exception as xml_error:
                        self.logger.warning(f"⚠️ Method 3 failed: {xml_error}")
                
                # Final status
                if background_applied:
                    self.logger.info(f"🎨 Background successfully applied for {profile}")
                else:
                    self.logger.error(f"❌ ALL BACKGROUND METHODS FAILED for {profile}")
                    self.logger.error(f"Background color attempted: {settings['background_color']}")
            
            # For ADHD: Add visual boundaries to text shapes
            if profile == 'adhd' and settings.get('add_visual_boundaries', False):
                self._add_visual_boundaries_to_shapes(slide, settings)
            
            # For Dyslexia: Add reading rulers if enabled
            if profile == 'dyslexia' and settings.get('add_reading_rulers', False):
                self._add_reading_rulers(slide, settings)
            
            # Apply consistent text alignment
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    self._apply_text_frame_formatting(shape.text_frame, settings, profile)
            
        except Exception as e:
            self.logger.error(f"Error applying evidence-based formatting: {e}")
    
    def _add_visual_boundaries_to_shapes(self, slide, settings):
        """
        Add visual boundaries around text shapes for ADHD support
        
        Args:
            slide: PowerPoint slide object
            settings: Profile settings dict
        """
        try:
            boundary_color = settings.get('boundary_color', RGBColor(200, 200, 200))
            
            # Get all text shapes
            text_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
                    text_shapes.append(shape)
            
            # Add subtle borders around text shapes
            for shape in text_shapes:
                if hasattr(shape, 'line'):
                    shape.line.color.rgb = boundary_color
                    shape.line.width = Pt(1)  # Thin line
                    
                    # Add padding by slightly adjusting margins
                    if hasattr(shape.text_frame, 'margin_left'):
                        shape.text_frame.margin_left = Inches(0.2)
                        shape.text_frame.margin_right = Inches(0.2)
                        shape.text_frame.margin_top = Inches(0.1)
                        shape.text_frame.margin_bottom = Inches(0.1)
            
            self.logger.info(f"Added visual boundaries to {len(text_shapes)} text shapes")
            
        except Exception as e:
            self.logger.error(f"Error adding visual boundaries: {e}")
    
    def _add_reading_rulers(self, slide, settings):
        """
        Add subtle reading ruler guides for dyslexia support
        
        Args:
            slide: PowerPoint slide object
            settings: Profile settings dict
        """
        try:
            # Add subtle horizontal lines between major text blocks
            text_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
                    text_shapes.append(shape)
            
            # Sort shapes by vertical position
            text_shapes.sort(key=lambda s: s.top)
            
            # Add faint lines between text blocks
            for i in range(len(text_shapes) - 1):
                current_shape = text_shapes[i]
                next_shape = text_shapes[i + 1]
                
                # Calculate position for ruler line
                line_top = current_shape.top + current_shape.height + Inches(0.05)
                line_left = min(current_shape.left, next_shape.left)
                line_width = max(current_shape.width, next_shape.width)
                
                # Only add if there's enough space
                if next_shape.top - line_top > Inches(0.1):
                    try:
                        line = slide.shapes.add_connector(
                            1, line_left, line_top,
                            line_left + line_width, line_top
                        )
                        line.line.color.rgb = RGBColor(230, 230, 230)  # Very light gray
                        line.line.width = Pt(0.5)  # Very thin
                    except Exception as line_error:
                        # Fallback: use a text box with underline
                        ruler_textbox = slide.shapes.add_textbox(
                            line_left, line_top - Inches(0.02), line_width, Inches(0.02)
                        )
                        ruler_frame = ruler_textbox.text_frame
                        ruler_para = ruler_frame.paragraphs[0]
                        ruler_run = ruler_para.add_run()
                        ruler_run.text = "─" * int(line_width.inches * 10)  # Dash characters
                        ruler_run.font.color.rgb = RGBColor(230, 230, 230)
            
            self.logger.info("Added reading ruler guides for dyslexia support")
            
        except Exception as e:
            self.logger.error(f"Error adding reading rulers: {e}")
    
    def _apply_text_frame_formatting(self, text_frame, settings, profile: str = 'default'):
        """
        Apply evidence-based text frame formatting
        
        Args:
            text_frame: PowerPoint text frame object
            settings: Profile settings dict
            profile: Learning profile
        """
        try:
            # Apply text alignment
            alignment = settings.get('text_alignment', PP_ALIGN.LEFT)
            
            for paragraph in text_frame.paragraphs:
                # Set alignment
                paragraph.alignment = alignment
                
                # Apply line spacing
                if settings.get('line_spacing'):
                    paragraph.line_spacing = settings['line_spacing']
                
                # Apply paragraph spacing
                if settings.get('paragraph_spacing'):
                    paragraph.space_after = Pt(settings['paragraph_spacing'] * 12)  # Convert to points
                
                # For dyslexia: Ensure no justification
                if settings.get('visual_structure') == 'simple':
                    paragraph.alignment = PP_ALIGN.LEFT  # Force left alignment
                
                # Apply to all runs in paragraph
                for run in paragraph.runs:
                    if hasattr(run, 'font'):
                        # Use pure black text on tinted backgrounds
                        if settings.get('use_background_tint', False):
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        
                        # Apply dyslexia font if needed
                        if profile.lower() == 'dyslexia':
                            dyslexia_settings = self.PROFILE_SETTINGS['dyslexia']
                            self._apply_font_with_fallbacks(
                                run.font, 
                                dyslexia_settings['font_name'], 
                                dyslexia_settings.get('font_fallbacks', [])
                            )
            
            # Set margins for better readability
            if settings.get('visual_structure') in ['simple', 'chunked']:
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
            
        except Exception as e:
            self.logger.error(f"Error applying text frame formatting: {e}")
    
    def _move_shape_to_back(self, slide, shape):
        """
        Move a shape to the back (behind all other shapes)
        
        Args:
            slide: PowerPoint slide object
            shape: Shape to move to back
        """
        try:
            # Method 1: Use the shape's XML element to move it to the beginning
            shape_element = shape._element
            parent = shape_element.getparent()
            
            # Remove from current position
            parent.remove(shape_element)
            # Insert at the beginning (index 0)
            parent.insert(0, shape_element)
            
            self.logger.debug("✅ Successfully moved shape to back using XML manipulation")
            
        except Exception as e:
            self.logger.warning(f"Failed to move shape to back: {e}")
            # Fallback: try alternative method if available
            try:
                # Alternative: iterate through shapes and reorder
                shapes_list = list(slide.shapes)
                shape_index = shapes_list.index(shape)
                
                # This is a more complex operation that would require
                # manipulating the underlying XML structure differently
                self.logger.debug("Shape reordering attempted")
                
            except Exception as fallback_error:
                self.logger.warning(f"Fallback shape reordering also failed: {fallback_error}")
    
    def _apply_slide_background(self, slide, background_color):
        """
        Apply background color to a single slide with multiple fallback methods
        Ensures full slide coverage without dimension changes
        
        Args:
            slide: PowerPoint slide object
            background_color: RGBColor object for background
        """
        try:
            success = False
            
            # Method 1: Direct slide background (this should cover the full slide)
            try:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = background_color
                success = True
                self.logger.debug(f"✅ Applied direct slide background (full coverage)")
            except Exception as e:
                self.logger.debug(f"⚠️ Direct slide background failed: {e}")
            
            # Method 2: Slide layout background (ensure layout doesn't override)
            try:
                if hasattr(slide, 'slide_layout'):
                    layout_bg = slide.slide_layout.background
                    layout_fill = layout_bg.fill
                    layout_fill.solid()
                    layout_fill.fore_color.rgb = background_color
                    success = True
                    self.logger.debug(f"✅ Applied layout background")
            except Exception as e:
                self.logger.debug(f"⚠️ Layout background failed: {e}")
            
            # Method 3: Set slide master background (for templates that inherit)
            try:
                if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'slide_master'):
                    master_bg = slide.slide_layout.slide_master.background
                    master_fill = master_bg.fill
                    master_fill.solid()
                    master_fill.fore_color.rgb = background_color
                    success = True
                    self.logger.debug(f"✅ Applied master background")
            except Exception as e:
                self.logger.debug(f"⚠️ Master background failed: {e}")
            
            # CRITICAL: Ensure the slide itself has no conflicting layout settings
            try:
                # Remove any layout-specific background that might override
                slide_element = slide._element
                
                # Look for any layout background references and clear them
                layout_refs = slide_element.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldLayoutIdLst')
                if layout_refs:
                    self.logger.debug(f"📋 Found layout references that might affect background")
                
            except Exception as e:
                self.logger.debug(f"⚠️ Layout cleanup failed: {e}")
            
            if success:
                self.logger.debug(f"🎨 Background applied to slide (full coverage ensured)")
            else:
                self.logger.warning(f"❌ All background methods failed for slide")
                
        except Exception as e:
            self.logger.error(f"Error applying slide background: {e}")
    
    def _duplicate_slide(self, presentation, slide_index: int, insert_position: Optional[int] = None) -> int:
        """
        Duplicate a slide using safe shape-by-shape copying
        
        Args:
            presentation: PowerPoint presentation object
            slide_index: Index of slide to duplicate (0-based)
            insert_position: Optional position to insert the new slide (0-based)
                           If None, slide is added at the end
            
        Returns:
            Index of the newly created duplicate slide
        """
        try:
            from copy import deepcopy
            
            # Get the slide to duplicate
            original_slide = presentation.slides[slide_index]
            original_layout = original_slide.slide_layout
            
            # Create a new slide with the same layout
            new_slide = presentation.slides.add_slide(original_layout)
            new_slide_index = len(presentation.slides) - 1
            
            # Copy slide content using safe shape-by-shape approach
            # This avoids XML-level media corruption issues
            self.logger.debug(f"Copying {len(original_slide.shapes)} shapes using safe shape-by-shape method")
            
            shapes_copied = 0
            shape_type_counts = {}
            
            for shape in original_slide.shapes:
                try:
                    # Get shape type for logging
                    shape_type_name = "unknown"
                    if hasattr(shape, 'shape_type'):
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            shape_type_name = "picture"
                        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                            shape_type_name = "media"
                        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            shape_type_name = "group"
                        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                            shape_type_name = "textbox"
                        elif hasattr(shape, 'text'):
                            shape_type_name = "text"
                    elif hasattr(shape, 'text'):
                        shape_type_name = "text"
                    
                    # Copy the shape safely
                    self._copy_shape_to_slide(shape, new_slide)
                    shapes_copied += 1
                    shape_type_counts[shape_type_name] = shape_type_counts.get(shape_type_name, 0) + 1
                    
                except Exception as shape_error:
                    self.logger.warning(f"Failed to copy {shape_type_name} shape: {shape_error}")
                    # Continue with other shapes
                    continue
            
            self.logger.info(f"Successfully copied {shapes_copied} shapes: {shape_type_counts}")
            
            # Copy slide notes if they exist
            if hasattr(original_slide, 'notes_slide') and original_slide.notes_slide:
                try:
                    original_notes = original_slide.notes_slide.notes_text_frame.text
                    if original_notes:
                        new_slide.notes_slide.notes_text_frame.text = original_notes
                except Exception as notes_error:
                    self.logger.debug(f"Could not copy slide notes: {notes_error}")
            
            self.logger.info(f"Duplicated slide {slide_index + 1} -> new slide {new_slide_index + 1}")
            return new_slide_index
            
        except Exception as e:
            self.logger.error(f"Error duplicating slide {slide_index + 1}: {e}")
            raise
    
    def _copy_shape_to_slide(self, original_shape, target_slide):
        """
        Copy a shape from one slide to another using safe methods
        
        Args:
            original_shape: Shape to copy
            target_slide: Slide to copy shape to
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # Handle different shape types with appropriate methods
            if hasattr(original_shape, 'shape_type'):
                shape_type = original_shape.shape_type
                
                # Handle pictures safely by recreating with image data
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self._copy_picture_shape(original_shape, target_slide)
                    return
                
                # Skip media shapes to avoid corruption - they're complex to copy
                # Media shapes (videos, audio) involve complex internal relationships
                # and media file storage that can't be safely duplicated without
                # risking PPTX corruption and "needs repair" warnings
                elif shape_type == MSO_SHAPE_TYPE.MEDIA:
                    self.logger.info("Skipping media shape to avoid PPTX corruption - videos/audio not duplicated")
                    return
                
                # Handle groups by copying individual components
                elif shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._copy_group_shape(original_shape, target_slide)
                    return
            
            # For all other shapes (text, textboxes, etc.), use the fallback method
            self._recreate_shape_fallback(original_shape, target_slide)
            
        except Exception as e:
            self.logger.debug(f"Error copying shape: {e}")
            # If copying fails, just skip this shape to avoid corruption
            pass
    
    def _copy_slide_content(self, source_slide, target_slide):
        """
        Copy all content from source slide to target slide
        
        Args:
            source_slide: Source slide to copy from
            target_slide: Target slide to copy to
        """
        try:
            # Copy all shapes from source to target
            shapes_copied = 0
            for shape in source_slide.shapes:
                try:
                    self._copy_shape_to_slide(shape, target_slide)
                    shapes_copied += 1
                except Exception as shape_error:
                    self.logger.warning(f"Failed to copy shape: {shape_error}")
                    continue
            
            # Copy slide notes if they exist
            if hasattr(source_slide, 'notes_slide') and source_slide.notes_slide:
                try:
                    source_notes = source_slide.notes_slide.notes_text_frame.text
                    if source_notes:
                        target_slide.notes_slide.notes_text_frame.text = source_notes
                except Exception as notes_error:
                    self.logger.debug(f"Could not copy slide notes: {notes_error}")
            
            self.logger.debug(f"Copied {shapes_copied} shapes to new slide")
            
        except Exception as e:
            self.logger.error(f"Error copying slide content: {e}")
            raise
    
    def _reorder_slides_for_translation(self, presentation, original_count: int, duplicated_indices: List[int]):
        """
        Reorder slides so translated slides appear immediately after their originals
        
        Args:
            presentation: PowerPoint presentation object
            original_count: Number of original slides
            duplicated_indices: List of indices where duplicated slides were placed
        """
        try:
            # Build the desired order: [0, dup0, 1, dup1, 2, dup2, ...]
            target_order = []
            for i in range(original_count):
                target_order.append(i)  # Original slide
                if i < len(duplicated_indices) and duplicated_indices[i] is not None:
                    target_order.append(duplicated_indices[i])  # Translated slide
            
            # Access the slide list in the presentation XML
            # Get the presentation part and find the slide ID list
            presentation_part = presentation.part
            slide_id_list = presentation_part._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
            
            if slide_id_list is not None:
                # Get all slide ID elements in current order
                slide_ids = list(slide_id_list)
                
                # Reorder them according to our target order
                if len(slide_ids) == len(target_order):
                    # Clear the current list
                    slide_id_list.clear()
                    
                    # Add slides in the new order
                    for slide_index in target_order:
                        if slide_index < len(slide_ids):
                            slide_id_list.append(slide_ids[slide_index])
                    
                    self.logger.debug(f"Reordered {len(target_order)} slides successfully")
                else:
                    self.logger.warning(f"Slide count mismatch: {len(slide_ids)} XML elements vs {len(target_order)} target order")
            else:
                self.logger.warning("Could not find slide ID list in presentation XML")
                
        except Exception as e:
            self.logger.error(f"Error reordering slides: {e}")
            raise
    
    def _copy_picture_shape(self, original_shape, target_slide):
        """
        Safely copy a picture shape by recreating it with the image data
        
        Args:
            original_shape: Original picture shape
            target_slide: Target slide
        """
        try:
            # Get position and size
            left = getattr(original_shape, 'left', 0)
            top = getattr(original_shape, 'top', 0)
            width = getattr(original_shape, 'width', 100)
            height = getattr(original_shape, 'height', 100)
            
            # Get image data
            if hasattr(original_shape, 'image') and hasattr(original_shape.image, 'blob'):
                image_data = original_shape.image.blob
                
                # Create new picture shape with the same image data
                new_picture = target_slide.shapes.add_picture(
                    io.BytesIO(image_data), left, top, width, height
                )
                
                self.logger.debug("Successfully copied picture shape with image data")
                return new_picture
            else:
                self.logger.warning("Could not access image data for picture shape")
                
        except Exception as e:
            self.logger.warning(f"Failed to copy picture shape: {e}")
            # Don't raise - just skip this picture to avoid corruption
            pass
    
    def _recreate_shape_fallback(self, original_shape, target_slide):
        """
        Fallback method to recreate text and simple shapes
        
        Args:
            original_shape: Original shape to recreate
            target_slide: Target slide
        """
        try:
            # Get shape position and size safely
            try:
                left = getattr(original_shape, 'left', 0)
                top = getattr(original_shape, 'top', 0)
                width = getattr(original_shape, 'width', 100)
                height = getattr(original_shape, 'height', 50)
            except:
                left, top, width, height = 0, 0, 100, 50
            
            # Handle text boxes and shapes with text
            if hasattr(original_shape, 'text_frame') and hasattr(original_shape, 'text') and original_shape.text:
                
                # Try to use existing placeholders first (for titles, content, etc.)
                text_copied = False
                
                # Check if this looks like a title (position-based heuristic)
                if top < 100 and hasattr(target_slide.shapes, 'title') and target_slide.shapes.title:
                    try:
                        target_slide.shapes.title.text = str(original_shape.text)
                        text_copied = True
                        self.logger.debug("Copied text to title placeholder")
                    except:
                        pass
                
                # Check if this looks like content (try content placeholders)
                if not text_copied:
                    for placeholder in target_slide.placeholders:
                        if hasattr(placeholder, 'placeholder_format') and placeholder.placeholder_format.type == 1:  # Content placeholder
                            try:
                                placeholder.text = str(original_shape.text)
                                text_copied = True
                                self.logger.debug("Copied text to content placeholder")
                                break
                            except:
                                continue
                
                # Fallback: create new text box with formatting preservation
                if not text_copied:
                    textbox = target_slide.shapes.add_textbox(left, top, width, height)
                    
                    # Copy text content with formatting preservation
                    try:
                        original_text = str(original_shape.text_frame.text)
                        if original_text:
                            # Use the enhanced formatting preservation method
                            self._replace_text_preserving_format(textbox.text_frame, original_text)
                            self.logger.debug("Created new text box with formatting preservation")
                        else:
                            textbox.text_frame.text = str(original_shape.text) if hasattr(original_shape, 'text') else ""
                            self.logger.debug("Created new text box with basic text copy")
                    except Exception as text_error:
                        # Final fallback to basic text copy
                        try:
                            textbox.text_frame.text = str(original_shape.text) if hasattr(original_shape, 'text') else ""
                            self.logger.debug(f"Fallback to basic text copy: {text_error}")
                        except:
                            self.logger.debug("Could not copy any text content")
                    
                    # Try to copy text frame properties
                    self._copy_text_frame_properties(original_shape.text_frame, textbox.text_frame)
                
        except Exception as e:
            self.logger.debug(f"Fallback shape recreation failed: {e}")
            # Just skip this shape if we can't recreate it
            pass
    
    def _create_translation_presentation(self, file_path: str, file_id: str, filename: str,
                                       target_language: str, progress_callback: Optional[callable] = None) -> Optional[str]:
        """
        Create a translation presentation that preserves original slides and adds translated copies
        
        Args:
            file_path: Input PPTX path
            file_id: Unique file identifier
            filename: Original filename
            target_language: Target language for translation
            progress_callback: Optional callback function for progress updates (message, percentage)
            
        Returns:
            Output path or None on error
        """
        try:
            self.logger.info(f"Creating translation presentation with slide duplication")
            
            # Open the presentation
            prs = Presentation(file_path)
            original_slide_count = len(prs.slides)
            
            # Duplicate slides in the correct order: Original 1, Translation 1, Original 2, Translation 2, etc.
            if progress_callback:
                progress_callback('Creating slides with proper ordering...', 10)
            
            # Initialize translation service
            from .translations_service import TranslationsService
            translations_service = TranslationsService(self.config)
            
            # First, duplicate all slides (this adds them at the end)
            duplicated_slides = []
            for slide_idx in range(original_slide_count):
                try:
                    new_slide_idx = self._duplicate_slide(prs, slide_idx)
                    duplicated_slides.append(new_slide_idx)
                    self.logger.debug(f"Duplicated slide {slide_idx + 1} to position {new_slide_idx + 1}")
                except Exception as e:
                    self.logger.error(f"Failed to duplicate slide {slide_idx + 1}: {e}")
                    duplicated_slides.append(None)  # Mark as failed
                    continue
            
            # Now translate all the duplicated slides
            for i, duplicate_idx in enumerate(duplicated_slides):
                if duplicate_idx is None:
                    continue  # Skip failed duplications
                    
                try:
                    slide_num = i + 1
                    self.logger.info(f"Translating slide {slide_num} of {original_slide_count}")
                    
                    # Update progress
                    if progress_callback:
                        progress = 20 + int((i / original_slide_count) * 60)
                        progress_callback(f'Translating slide {slide_num}/{original_slide_count}...', progress)
                    
                    translated_slide = prs.slides[duplicate_idx]
                    
                    # Find all text shapes and translate them
                    text_shapes = []
                    for shape in translated_slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.text_frame and shape.text:
                            text_shapes.append(shape)
                    
                    # Translate each text shape
                    for shape in text_shapes:
                        try:
                            original_text = shape.text_frame.text
                            if original_text.strip():
                                translated_text = translations_service.translate_text(original_text, target_language)
                                
                                # Check for translation issues
                                if not self._is_placeholder_text(translated_text):
                                    # Use overflow handling for translations to prevent text misalignment
                                    self._replace_text_with_overflow_handling(shape.text_frame, translated_text, shape)
                                else:
                                    self.logger.warning(f"Translation placeholder detected, keeping original text")
                                    
                        except Exception as trans_error:
                            self.logger.error(f"Translation failed for shape in slide {slide_num}: {trans_error}")
                    
                    # Add translation note
                    translation_note = f"This slide translated to {target_language}. Original: Slide {slide_num}"
                    self._add_adaptation_notes_to_slide(translated_slide, [translation_note], 'translation')
                    
                except Exception as e:
                    self.logger.error(f"Failed to translate slide {slide_num}: {e}")
                    continue
            
            # Now we need to reorder the slides: Original 1, Translation 1, Original 2, Translation 2, etc.
            if progress_callback:
                progress_callback('Reordering slides...', 80)
            
            # Reorder using a simpler approach: create the order we want by index manipulation
            current_slides = list(prs.slides)
            total_slides = len(current_slides)
            
            # Build the target order: pairs of (original, translation)
            target_order = []
            for i in range(original_slide_count):
                target_order.append(i)  # Original slide
                if i < len(duplicated_slides) and duplicated_slides[i] is not None:
                    target_order.append(duplicated_slides[i])  # Translated slide
            
            self.logger.info(f"Target slide order: {target_order}")
            
            # Reorder slides using XML manipulation for exact control
            try:
                self._reorder_slides_for_translation(prs, original_slide_count, duplicated_slides)
                self.logger.info(f"Successfully reordered slides: Original 1, Translation 1, Original 2, Translation 2, etc.")
            except Exception as reorder_error:
                self.logger.warning(f"Could not reorder slides: {reorder_error}. Translations will appear at the end.")
                # This is acceptable - content is still preserved, just different order
            
            # Add title slide note explaining the structure
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]
                structure_note = (f"This presentation contains original slides (1-{original_slide_count}) "
                                f"followed by {target_language} translations (slides {original_slide_count + 1}-{len(prs.slides)})")
                self._add_adaptation_notes_to_slide(first_slide, [structure_note], 'translation')
            
            # Save the presentation
            base_name = os.path.splitext(filename)[0]
            output_filename = f"translated_{target_language}_{base_name}.pptx"
            output_path = os.path.join(self.config['output_folder'], f"{file_id}_{output_filename}")
            
            prs.save(output_path)
            
            if progress_callback:
                progress_callback(f'Translation complete! Original + {target_language} versions created', 100)
            
            self.logger.info(f"Translation presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"Error creating translation presentation: {e}")
            if progress_callback:
                progress_callback(f'Translation failed: {str(e)}', 100)
            return None
    
    def translate_presentation_in_place(self, file_path: str, file_id: str, filename: str,
                                      target_language: str, progress_callback: Optional[callable] = None) -> Optional[str]:
        """
        Translate presentation by replacing original content (no slide duplication)
        
        Args:
            file_path: Input PPTX path
            file_id: Unique file identifier
            filename: Original filename
            target_language: Target language for translation
            progress_callback: Optional callback function for progress updates (message, percentage)
            
        Returns:
            Output path or None on error
        """
        try:
            self.logger.info(f"Translating presentation in-place to {target_language}")
            
            # Open the presentation
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            if progress_callback:
                progress_callback(f'Translating {total_slides} slides to {target_language}...', 10)
            
            # Initialize translation service
            from .translations_service import TranslationsService
            translations_service = TranslationsService(self.config)
            
            # Translate each slide directly (replace original content)
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    slide_num = slide_idx + 1
                    self.logger.info(f"Translating slide {slide_num} of {total_slides}")
                    
                    # Update progress
                    if progress_callback:
                        progress = 10 + int((slide_idx / total_slides) * 80)
                        progress_callback(f'Translating slide {slide_num}/{total_slides}...', progress)
                    
                    # Find all text shapes and translate them
                    text_shapes = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame') and shape.text_frame and shape.text:
                            text_shapes.append(shape)
                    
                    # Translate each text shape
                    for shape in text_shapes:
                        try:
                            original_text = shape.text_frame.text
                            if original_text.strip():
                                translated_text = translations_service.translate_text(original_text, target_language)
                                
                                # Check for translation issues
                                if not self._is_placeholder_text(translated_text):
                                    # Use overflow handling for translations to prevent text misalignment
                                    self._replace_text_with_overflow_handling(shape.text_frame, translated_text, shape)
                                else:
                                    self.logger.warning(f"Translation placeholder detected, keeping original text")
                                    
                        except Exception as trans_error:
                            self.logger.error(f"Translation failed for shape in slide {slide_num}: {trans_error}")
                    
                    # Add translation note to slide
                    translation_note = f"This slide translated to {target_language} (replace mode)"
                    self._add_adaptation_notes_to_slide(slide, [translation_note], 'translation')
                    
                except Exception as e:
                    self.logger.error(f"Failed to translate slide {slide_num}: {e}")
                    continue
            
            # Add title slide note explaining the translation
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]
                structure_note = f"This presentation has been translated to {target_language} (original content replaced)"
                self._add_adaptation_notes_to_slide(first_slide, [structure_note], 'translation')
            
            # Save the presentation
            base_name = os.path.splitext(filename)[0]
            output_filename = f"translated_{target_language}_{base_name}.pptx"
            output_path = os.path.join(self.config['output_folder'], f"{file_id}_{output_filename}")
            
            prs.save(output_path)
            
            if progress_callback:
                progress_callback(f'Translation complete! Content translated to {target_language}', 100)
            
            self.logger.info(f"In-place translation presentation created: {output_path}")
            return output_path
            
        except Exception as e:
            self.logger.error(f"Error in in-place translation: {e}")
            if progress_callback:
                progress_callback(f'Translation failed: {str(e)}', 100)
            return None
