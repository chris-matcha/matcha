"""
Formats Service

Handles extraction and creation of different file formats (PDF, PowerPoint).
"""
import os
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_JUSTIFY
from .base_service import BaseService
from .pdf_visual_handler import PDFVisualHandler


class FormatsService(BaseService):
    """Service for handling different file formats"""
    
    def _initialize(self):
        """Initialize format handlers"""
        self.pdf_handler = PDFHandler()
        self.pdf_visual_handler = PDFVisualHandler()
        self.pptx_handler = PowerPointHandler()
    
    def extract_content(self, file_path: str, file_type: str, 
                       include_formatting: bool = False) -> Dict[str, Any]:
        """
        Extract content from a file
        
        Args:
            file_path: Path to the file
            file_type: Type of file ('pdf' or 'pptx')
            include_formatting: Include formatting details (PDF only)
            
        Returns:
            Dict containing extracted content
        """
        if file_type == 'pdf':
            if include_formatting:
                return self.pdf_visual_handler.extract_content_with_formatting(file_path)
            else:
                return self.pdf_handler.extract_content(file_path)
        elif file_type == 'pptx':
            return self.pptx_handler.extract_content(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def create_file(self, content: Dict[str, Any], output_path: str, file_type: str, 
                   profile: str = 'default', preserve_visuals: bool = False,
                   original_path: Optional[str] = None) -> bool:
        """
        Create a file from content
        
        Args:
            content: Content to write
            output_path: Path for output file
            file_type: Type of file to create
            profile: Learning profile for formatting
            preserve_visuals: Preserve visual layout (PDF only)
            original_path: Path to original file (required for visual preservation)
            
        Returns:
            bool: Success status
        """
        if file_type == 'pdf':
            if preserve_visuals and original_path:
                # Try visual preservation methods in order
                print(f"Attempting visual preservation for {os.path.basename(output_path)}")
                success = self.pdf_visual_handler.create_visual_preserved_pdf(
                    original_path, content, output_path, profile
                )
                
                if success:
                    print("✅ Primary visual preservation method succeeded")
                    return True
                
                print("⚠️ Primary method failed, trying overlay method...")
                # Fallback to overlay method
                success = self.pdf_visual_handler.create_visual_preserved_with_overlay(
                    original_path, content, output_path, profile
                )
                
                if success:
                    print("✅ Overlay method succeeded")
                    return True
                
                print("⚠️ Overlay method failed, trying simple preservation...")
                # Final fallback to simple visual preservation
                success = self.pdf_visual_handler.create_simple_visual_preserved(
                    original_path, output_path, profile
                )
                
                if success:
                    print("⚠️ Using simple preservation (original PDF copy)")
                    return True
                
                # Ultimate fallback: create non-visual PDF
                print("❌ All visual preservation methods failed, creating non-visual PDF")
                success = self.pdf_handler.create_file(content, output_path, profile)
                
                return success
            else:
                # Non-visual PDF creation
                return self.pdf_handler.create_file(content, output_path, profile)
        elif file_type == 'pptx':
            return self.pptx_handler.create_file(content, output_path, profile)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")


class PDFHandler:
    """Handler for PDF files"""
    
    # Profile-specific formatting settings
    PROFILE_STYLES = {
        'dyslexia': {
            'fontName': 'Helvetica',
            'fontSize': 14,
            'leading': 20,
            'textColor': HexColor('#000080'),
            'backColor': HexColor('#FFFEF5'),
            'spaceAfter': 18
        },
        'adhd': {
            'fontName': 'Helvetica',
            'fontSize': 13,
            'leading': 18,
            'textColor': HexColor('#004D00'),
            'backColor': HexColor('#F0FFF0'),
            'spaceAfter': 16
        },
        'esl': {
            'fontName': 'Helvetica',
            'fontSize': 13,
            'leading': 17,
            'textColor': HexColor('#4B0082'),
            'backColor': HexColor('#FAF5FF'),
            'spaceAfter': 15
        },
        'default': {
            'fontName': 'Helvetica',
            'fontSize': 12,
            'leading': 14,
            'textColor': HexColor('#000000'),
            'backColor': HexColor('#FFFFFF'),
            'spaceAfter': 12
        }
    }
    
    def extract_content(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PDF"""
        content = {
            'pages': [],
            'metadata': {}
        }
        
        try:
            doc = fitz.open(file_path)
            
            # Extract metadata
            content['metadata'] = {
                'page_count': doc.page_count,
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', '')
            }
            
            # Extract pages
            for page_num, page in enumerate(doc):
                page_content = {
                    'page_number': page_num + 1,
                    'text': page.get_text(),
                    'images': [],
                    'tables': []
                }
                
                # Extract images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    page_content['images'].append({
                        'index': img_index,
                        'width': img[2],
                        'height': img[3]
                    })
                
                content['pages'].append(page_content)
            
            doc.close()
            
        except Exception as e:
            raise Exception(f"Error extracting PDF content: {str(e)}")
        
        return content
    
    def create_file(self, content: Dict[str, Any], output_path: str, profile: str) -> bool:
        """Create PDF from content with profile-specific formatting"""
        try:
            # Get profile style settings
            profile_style = self.PROFILE_STYLES.get(profile, self.PROFILE_STYLES['default'])
            
            # Create document with background color
            doc = SimpleDocTemplate(
                output_path, 
                pagesize=A4,
                leftMargin=inch,
                rightMargin=inch,
                topMargin=inch,
                bottomMargin=inch
            )
            
            # Create custom style based on profile
            styles = getSampleStyleSheet()
            profile_para_style = ParagraphStyle(
                'ProfileStyle',
                parent=styles['Normal'],
                fontName=profile_style['fontName'],
                fontSize=profile_style['fontSize'],
                leading=profile_style['leading'],
                textColor=profile_style['textColor'],
                alignment=TA_JUSTIFY,
                spaceAfter=profile_style['spaceAfter']
            )
            
            story = []
            
            # Add content with profile formatting
            for page in content.get('pages', []):
                text = page.get('text', '')
                if text:
                    # Split into paragraphs
                    paragraphs = text.split('\n\n')
                    for para_text in paragraphs:
                        if para_text.strip():
                            para = Paragraph(para_text.strip(), profile_para_style)
                            story.append(para)
                            story.append(Spacer(1, profile_style['spaceAfter']))
                    
                    # Add page break between pages
                    if page != content['pages'][-1]:
                        story.append(PageBreak())
            
            # Build document
            doc.build(story)
            return True
            
        except Exception as e:
            print(f"Error creating PDF: {str(e)}")
            return False


class PowerPointHandler:
    """Handler for PowerPoint files"""
    
    def extract_content(self, file_path: str) -> Dict[str, Any]:
        """Extract content from PowerPoint"""
        content = {
            'slides': [],
            'metadata': {}
        }
        
        try:
            prs = Presentation(file_path)
            
            # Extract metadata
            content['metadata'] = {
                'slide_count': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
            
            # Extract slides
            for slide_num, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': slide_num + 1,
                    'title': '',
                    'content': '',
                    'notes': '',
                    'shapes': []
                }
                
                # Extract title
                if slide.shapes.title:
                    slide_content['title'] = slide.shapes.title.text
                
                # Extract text from shapes
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        if shape != slide.shapes.title:
                            text_parts.append(shape.text)
                
                slide_content['content'] = '\n'.join(text_parts)
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_content['notes'] = slide.notes_slide.notes_text_frame.text
                
                content['slides'].append(slide_content)
            
        except Exception as e:
            raise Exception(f"Error extracting PowerPoint content: {str(e)}")
        
        return content
    
    def create_file(self, content: Dict[str, Any], output_path: str, profile: str) -> bool:
        """Create PowerPoint from content"""
        try:
            prs = Presentation()
            
            # Add slides
            for slide_data in content.get('slides', []):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                
                # Add title
                if slide_data.get('title'):
                    slide.shapes.title.text = slide_data['title']
                
                # Add content
                if slide_data.get('content'):
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = slide_data['content']
                
                # Add notes
                if slide_data.get('notes'):
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data['notes']
            
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Error creating PowerPoint: {str(e)}")
            return False