"""
Document Conversion Service

Handles conversions between different document formats
"""
import os
import subprocess
import platform
import tempfile
import shutil
from typing import Optional, Dict, Any, List, Tuple
from .base_service import BaseService
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from PIL import Image
import fitz  # PyMuPDF


class ConversionService(BaseService):
    """Service for converting between document formats"""
    
    def _initialize(self):
        """Initialize conversion service resources"""
        self.system = platform.system()
        self.libreoffice_path = self._find_libreoffice()
        
        # Quality settings
        self.quality_settings = {
            'high': {
                'dpi': 300,
                'image_quality': 95,
                'compress': False
            },
            'medium': {
                'dpi': 150,
                'image_quality': 85,
                'compress': True
            },
            'low': {
                'dpi': 96,
                'image_quality': 70,
                'compress': True
            }
        }
    
    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice installation"""
        possible_paths = [
            '/usr/bin/libreoffice',
            '/usr/local/bin/libreoffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
            'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe'
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                return path
        
        # Try to find in PATH
        try:
            result = subprocess.run(['which', 'libreoffice'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout.strip()
        except:
            pass
        
        return None
    
    def convert_pptx_to_pdf(self, pptx_path: str, output_path: str,
                           quality: str = 'high', method: Optional[str] = None) -> bool:
        """
        Convert PowerPoint to PDF with quality options
        
        Args:
            pptx_path: Path to input PPTX file
            output_path: Path for output PDF file
            quality: 'high', 'medium', 'low' - affects file size and fidelity
            method: Force specific method ('libreoffice', 'windows', 'pymupdf', 'reportlab')
            
        Returns:
            bool: Success status
        """
        # Store quality setting for methods to use
        self.current_quality = quality if quality in self.quality_settings else 'high'
        
        if method:
            method_map = {
                'libreoffice': self._convert_pptx_to_pdf_libreoffice,
                'windows': self._convert_pptx_to_pdf_windows,
                'pymupdf': self._convert_pptx_to_pdf_pymupdf,
                'reportlab': self._convert_pptx_to_pdf_reportlab
            }
            if method in method_map:
                try:
                    return method_map[method](pptx_path, output_path)
                except Exception as e:
                    self.logger.error(f"Method {method} failed: {str(e)}")
                    return False
            else:
                self.logger.warning(f"Unknown method: {method}")
        
        # Use default priority order
        methods = [
            self._convert_pptx_to_pdf_libreoffice,
            self._convert_pptx_to_pdf_windows,
            self._convert_pptx_to_pdf_pymupdf,
            self._convert_pptx_to_pdf_reportlab
        ]
        
        for method in methods:
            try:
                if method(pptx_path, output_path):
                    self.logger.info(f"Successfully converted using {method.__name__}")
                    return True
            except Exception as e:
                self.logger.warning(f"Method {method.__name__} failed: {str(e)}")
                continue
        
        self.logger.error("All conversion methods failed")
        return False
    
    def _convert_pptx_to_pdf_libreoffice(self, pptx_path: str, output_path: str) -> bool:
        """Convert using LibreOffice"""
        if not self.libreoffice_path:
            raise Exception("LibreOffice not found")
        
        # Create temp directory for output
        temp_dir = tempfile.mkdtemp()
        
        try:
            # Convert to PDF in temp directory
            cmd = [
                self.libreoffice_path,
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir,
                pptx_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
            
            # Find the output file
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            temp_pdf = os.path.join(temp_dir, f"{base_name}.pdf")
            
            if not os.path.exists(temp_pdf):
                raise Exception("Output PDF not found")
            
            # Move to final location
            shutil.move(temp_pdf, output_path)
            return True
            
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
    
    def _convert_pptx_to_pdf_windows(self, pptx_path: str, output_path: str) -> bool:
        """Convert using Windows COM if available"""
        if self.system != 'Windows':
            raise Exception("Not on Windows")
        
        try:
            import comtypes.client
            
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            # Open presentation
            presentation = powerpoint.Presentations.Open(pptx_path)
            
            # Save as PDF
            presentation.SaveAs(output_path, 32)  # 32 = ppSaveAsPDF
            presentation.Close()
            powerpoint.Quit()
            
            return True
            
        except Exception as e:
            raise Exception(f"Windows COM conversion failed: {str(e)}")
    
    def _convert_pptx_to_pdf_pymupdf(self, pptx_path: str, output_path: str) -> bool:
        """Enhanced PyMuPDF conversion with better slide rendering"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.enum.dml import MSO_FILL_TYPE
            import fitz
            from PIL import Image, ImageDraw, ImageFont
            import io
            import sys
            
            prs = Presentation(pptx_path)
            doc = fitz.open()
            
            # Standard slide dimensions (16:9)
            slide_width = 1920
            slide_height = 1080
            
            # Get system fonts
            def get_font(name, size, bold=False, italic=False):
                """Get font with platform-specific paths"""
                font_names = []
                
                if bold and italic:
                    font_names = [f"{name}bi.ttf", f"{name}-BoldItalic.ttf", f"{name} Bold Italic.ttf"]
                elif bold:
                    font_names = [f"{name}bd.ttf", f"{name}-Bold.ttf", f"{name} Bold.ttf", f"{name}b.ttf"]
                elif italic:
                    font_names = [f"{name}i.ttf", f"{name}-Italic.ttf", f"{name} Italic.ttf"]
                else:
                    font_names = [f"{name}.ttf", f"{name}-Regular.ttf", f"{name} Regular.ttf"]
                
                # Platform-specific font directories
                font_dirs = []
                if sys.platform == "win32":
                    font_dirs = ["C:/Windows/Fonts/"]
                elif sys.platform == "darwin":
                    font_dirs = ["/System/Library/Fonts/", "/Library/Fonts/", "~/Library/Fonts/"]
                else:  # Linux
                    font_dirs = ["/usr/share/fonts/truetype/", "/usr/local/share/fonts/", "~/.fonts/"]
                
                # Try to find font
                for font_dir in font_dirs:
                    for font_name in font_names:
                        try:
                            font_path = os.path.expanduser(os.path.join(font_dir, font_name))
                            return ImageFont.truetype(font_path, size)
                        except:
                            continue
                
                # Try common fonts
                for common_font in ["Arial", "Helvetica", "Liberation Sans", "DejaVu Sans"]:
                    for font_dir in font_dirs:
                        for font_name in [f"{common_font}.ttf", f"{common_font}-Regular.ttf"]:
                            try:
                                font_path = os.path.expanduser(os.path.join(font_dir, font_name))
                                return ImageFont.truetype(font_path, size)
                            except:
                                continue
                
                # Fallback to default
                try:
                    return ImageFont.load_default()
                except:
                    return None
            
            for slide_idx, slide in enumerate(prs.slides):
                # Create image for slide with white background
                img = Image.new('RGB', (slide_width, slide_height), 'white')
                draw = ImageDraw.Draw(img)
                
                # Handle slide background
                try:
                    if hasattr(slide, 'background') and slide.background:
                        if hasattr(slide.background, 'fill'):
                            fill = slide.background.fill
                            if fill.type == MSO_FILL_TYPE.SOLID:
                                if hasattr(fill.fore_color, 'rgb') and fill.fore_color.rgb:
                                    rgb = fill.fore_color.rgb
                                    bg_color = (rgb.red, rgb.green, rgb.blue) if hasattr(rgb, 'red') else rgb
                                    img = Image.new('RGB', (slide_width, slide_height), bg_color)
                                    draw = ImageDraw.Draw(img)
                except Exception as e:
                    self.logger.debug(f"Could not process slide background: {str(e)}")
                
                # Process shapes in order (back to front)
                for shape in slide.shapes:
                    try:
                        # Calculate shape position and size
                        if hasattr(shape, 'left') and hasattr(shape, 'top'):
                            x = int((shape.left / prs.slide_width) * slide_width)
                            y = int((shape.top / prs.slide_height) * slide_height)
                            w = int((shape.width / prs.slide_width) * slide_width) if hasattr(shape, 'width') else slide_width - x
                            h = int((shape.height / prs.slide_height) * slide_height) if hasattr(shape, 'height') else slide_height - y
                        else:
                            x, y, w, h = 100, 100, slide_width - 200, slide_height - 200
                        
                        # Handle different shape types
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # Handle images
                            try:
                                image_bytes = shape.image.blob
                                shape_img = Image.open(io.BytesIO(image_bytes))
                                
                                # Maintain aspect ratio
                                img_ratio = shape_img.width / shape_img.height
                                shape_ratio = w / h
                                
                                if img_ratio > shape_ratio:
                                    # Image is wider
                                    new_w = w
                                    new_h = int(w / img_ratio)
                                else:
                                    # Image is taller
                                    new_h = h
                                    new_w = int(h * img_ratio)
                                
                                # Center image in shape bounds
                                x_offset = x + (w - new_w) // 2
                                y_offset = y + (h - new_h) // 2
                                
                                # Resize and paste
                                shape_img = shape_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                                
                                # Handle transparency
                                if shape_img.mode == 'RGBA':
                                    img.paste(shape_img, (x_offset, y_offset), shape_img)
                                else:
                                    img.paste(shape_img, (x_offset, y_offset))
                                    
                            except Exception as e:
                                self.logger.warning(f"Could not render image on slide {slide_idx + 1}: {str(e)}")
                        
                        elif hasattr(shape, 'text_frame') and shape.text_frame:
                            # Handle text shapes
                            text_frame = shape.text_frame
                            
                            # Start position for text
                            text_y = y + 10  # Add padding
                            
                            for paragraph in text_frame.paragraphs:
                                if paragraph.text.strip():
                                    # Determine font properties
                                    font_size = 24  # Default
                                    font_name = "Arial"
                                    bold = False
                                    italic = False
                                    color = 'black'
                                    
                                    # Get font properties from first run
                                    if paragraph.runs:
                                        run = paragraph.runs[0]
                                        if hasattr(run.font, 'size') and run.font.size:
                                            font_size = int(run.font.size.pt * 1.5) if hasattr(run.font.size, 'pt') else 24
                                        if hasattr(run.font, 'name') and run.font.name:
                                            font_name = run.font.name
                                        if hasattr(run.font, 'bold'):
                                            bold = run.font.bold or False
                                        if hasattr(run.font, 'italic'):
                                            italic = run.font.italic or False
                                        if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                            rgb = run.font.color.rgb
                                            color = (rgb.red, rgb.green, rgb.blue) if hasattr(rgb, 'red') else (0, 0, 0)
                                    
                                    # Handle title shapes with larger font
                                    if shape == slide.shapes.title:
                                        font_size = int(font_size * 1.5)
                                    
                                    # Get appropriate font
                                    font = get_font(font_name, font_size, bold, italic)
                                    
                                    # Draw text with word wrapping
                                    words = paragraph.text.split()
                                    lines = []
                                    current_line = []
                                    
                                    for word in words:
                                        current_line.append(word)
                                        test_line = ' '.join(current_line)
                                        
                                        # Simple width check (approximate)
                                        if len(test_line) * font_size * 0.6 > w - 20:
                                            if len(current_line) > 1:
                                                current_line.pop()
                                                lines.append(' '.join(current_line))
                                                current_line = [word]
                                    
                                    if current_line:
                                        lines.append(' '.join(current_line))
                                    
                                    # Draw each line
                                    for line in lines:
                                        if text_y + font_size > y + h:
                                            break  # Don't exceed shape bounds
                                        
                                        draw.text((x + 10, text_y), line, fill=color, font=font)
                                        text_y += int(font_size * 1.2)
                                    
                                    # Add paragraph spacing
                                    text_y += int(font_size * 0.5)
                        
                        # Handle tables
                        elif hasattr(shape, 'table'):
                            table = shape.table
                            
                            # Simple table rendering
                            cell_height = h // table.rows.count if table.rows.count > 0 else 30
                            cell_width = w // table.columns.count if table.columns.count > 0 else 100
                            
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    cell_x = x + col_idx * cell_width
                                    cell_y = y + row_idx * cell_height
                                    
                                    # Draw cell border
                                    draw.rectangle([cell_x, cell_y, cell_x + cell_width, cell_y + cell_height], 
                                                 outline='gray', width=1)
                                    
                                    # Draw cell text
                                    if cell.text.strip():
                                        font = get_font("Arial", 16)
                                        # Simple text centering
                                        text_x = cell_x + 5
                                        text_y_pos = cell_y + (cell_height - 20) // 2
                                        draw.text((text_x, text_y_pos), cell.text[:50], fill='black', font=font)
                        
                        # Handle basic shapes (rectangles, circles, etc.)
                        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                            # Draw a simple rectangle as placeholder
                            if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL_TYPE.SOLID:
                                if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                                    rgb = shape.fill.fore_color.rgb
                                    fill_color = (rgb.red, rgb.green, rgb.blue) if hasattr(rgb, 'red') else (200, 200, 200)
                                    draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline='gray')
                            
                    except Exception as e:
                        self.logger.debug(f"Could not process shape on slide {slide_idx + 1}: {str(e)}")
                
                # Convert PIL image to PDF page
                img_bytes = io.BytesIO()
                
                # Get quality settings
                quality = getattr(self, 'current_quality', 'high')
                settings = self.quality_settings.get(quality, self.quality_settings['high'])
                
                # Save with quality settings
                img.save(img_bytes, 
                        format='PNG', 
                        dpi=(settings['dpi'], settings['dpi']), 
                        optimize=settings['compress'],
                        quality=settings['image_quality'])
                img_bytes.seek(0)
                
                # Create PDF page from image
                img_pdf = fitz.open("png", img_bytes.getvalue())
                pdf_page = doc.new_page(width=612, height=344)  # 16:9 aspect ratio
                pdf_page.show_pdf_page(pdf_page.rect, img_pdf, 0)
                img_pdf.close()
                
                # Add slide number
                slide_num_text = f"{slide_idx + 1}"
                pdf_page.insert_text(
                    (pdf_page.rect.width - 30, pdf_page.rect.height - 10),
                    slide_num_text,
                    fontsize=10,
                    color=(0.5, 0.5, 0.5)
                )
            
            # Add metadata
            doc.set_metadata({
                'title': os.path.splitext(os.path.basename(pptx_path))[0],
                'producer': 'Matcha Document Conversion Service',
                'creator': 'PyMuPDF Enhanced Converter v2',
                'subject': 'Converted from PowerPoint presentation',
                'author': 'Matcha'
            })
            
            # Save PDF with optimization
            doc.save(output_path, 
                    garbage=4,  # Maximum garbage collection
                    deflate=True,  # Compress streams
                    clean=True,  # Clean up redundant objects
                    ascii=False,  # Allow binary format for smaller size
                    expand=255,  # Maximum expansion of images
                    linear=True,  # Optimize for web
                    pretty=False)  # Don't pretty-print for smaller size
            
            doc.close()
            self.logger.info(f"Successfully converted {pptx_path} to PDF")
            return True
            
        except Exception as e:
            self.logger.error(f"Enhanced PyMuPDF conversion failed: {str(e)}")
            raise Exception(f"Enhanced PyMuPDF conversion failed: {str(e)}")
    
    def _convert_pptx_to_pdf_reportlab(self, pptx_path: str, output_path: str) -> bool:
        """Convert using ReportLab (fallback method)"""
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            
            # Extract content from PPTX
            prs = Presentation(pptx_path)
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Title style
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30
            )
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                # Add slide number
                story.append(Paragraph(f"Slide {slide_idx + 1}", styles['Normal']))
                story.append(Spacer(1, 0.2 * inch))
                
                # Add title
                if hasattr(slide.shapes, 'title') and slide.shapes.title and slide.shapes.title.text:
                    story.append(Paragraph(slide.shapes.title.text, title_style))
                
                # Add content
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                        # Split into paragraphs
                        for para in shape.text.split('\n'):
                            if para.strip():
                                story.append(Paragraph(para, styles['Normal']))
                                story.append(Spacer(1, 0.1 * inch))
                
                # Add page break except for last slide
                if slide_idx < len(prs.slides) - 1:
                    story.append(PageBreak())
            
            # Build PDF
            doc.build(story)
            return True
            
        except Exception as e:
            raise Exception(f"ReportLab conversion failed: {str(e)}")
    
    def convert_pdf_to_pptx(self, pdf_path: str, output_path: str) -> bool:
        """
        Convert PDF to PowerPoint presentation
        
        Args:
            pdf_path: Path to input PDF file
            output_path: Path for output PPTX file
            
        Returns:
            bool: Success status
        """
        try:
            # Open PDF
            pdf_doc = fitz.open(pdf_path)
            
            # Create presentation
            prs = Presentation()
            
            # Process each page
            for page_num in range(pdf_doc.page_count):
                page = pdf_doc[page_num]
                
                # Add slide
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Extract text
                text = page.get_text()
                
                if text.strip():
                    # Add text to slide
                    from pptx.util import Inches, Pt
                    left = top = Inches(0.5)
                    width = Inches(9)
                    height = Inches(6.5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text
                    
                    # Basic formatting
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(12)
            
            # Save presentation
            prs.save(output_path)
            pdf_doc.close()
            return True
            
        except Exception as e:
            self.logger.error(f"Error converting PDF to PPTX: {str(e)}")
            return False
    
    def convert_to_format(self, input_path: str, output_path: str,
                         input_format: str, output_format: str) -> bool:
        """
        General conversion method
        
        Args:
            input_path: Path to input file
            output_path: Path for output file
            input_format: Input format (pdf, pptx, etc.)
            output_format: Output format (pdf, pptx, etc.)
            
        Returns:
            bool: Success status
        """
        conversion_map = {
            ('pptx', 'pdf'): self.convert_pptx_to_pdf,
            ('pdf', 'pptx'): self.convert_pdf_to_pptx,
        }
        
        key = (input_format.lower(), output_format.lower())
        if key in conversion_map:
            return conversion_map[key](input_path, output_path)
        else:
            self.logger.error(f"Unsupported conversion: {input_format} to {output_format}")
            return False
    
    def batch_convert(self, conversions: List[Dict[str, str]], 
                     max_workers: int = 4) -> Dict[str, Dict[str, Any]]:
        """
        Batch convert multiple files
        
        Args:
            conversions: List of dicts with 'input', 'output', 'from_format', 'to_format'
            max_workers: Maximum parallel conversions
            
        Returns:
            Dict mapping input paths to results
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import time
        
        results = {}
        
        def convert_single(conversion: Dict[str, str]) -> Tuple[str, Dict[str, Any]]:
            input_path = conversion['input']
            start_time = time.time()
            
            try:
                # Add optional quality parameter
                kwargs = {}
                if 'quality' in conversion:
                    kwargs['quality'] = conversion['quality']
                if 'method' in conversion:
                    kwargs['method'] = conversion['method']
                
                # Handle format-specific conversions
                from_fmt = conversion['from_format'].lower()
                to_fmt = conversion['to_format'].lower()
                
                if from_fmt == 'pptx' and to_fmt == 'pdf':
                    success = self.convert_pptx_to_pdf(
                        input_path,
                        conversion['output'],
                        **kwargs
                    )
                else:
                    success = self.convert_to_format(
                        input_path,
                        conversion['output'],
                        from_fmt,
                        to_fmt
                    )
                
                return input_path, {
                    'success': success,
                    'output': conversion['output'],
                    'duration': time.time() - start_time,
                    'error': None
                }
            except Exception as e:
                return input_path, {
                    'success': False,
                    'output': None,
                    'duration': time.time() - start_time,
                    'error': str(e)
                }
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(convert_single, conv): conv 
                for conv in conversions
            }
            
            for future in as_completed(futures):
                input_path, result = future.result()
                results[input_path] = result
                
                # Log progress
                self.logger.info(f"Completed conversion of {input_path}: {result['success']}")
        
        return results
    
    def detect_format(self, file_path: str) -> Optional[str]:
        """
        Detect file format using both extension and magic bytes
        
        Returns:
            Format string or None if unknown
        """
        if not os.path.exists(file_path):
            return None
        
        # Check extension
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            # Check magic bytes
            with open(file_path, 'rb') as f:
                header = f.read(16)
            
            # PPTX signature (ZIP with specific structure)
            if header.startswith(b'PK\x03\x04') or header.startswith(b'PK\x05\x06'):
                # Further check for PPTX structure
                try:
                    from zipfile import ZipFile
                    with ZipFile(file_path, 'r') as z:
                        if '[Content_Types].xml' in z.namelist():
                            return 'pptx'
                except:
                    pass
            
            # PDF signature
            if header.startswith(b'%PDF'):
                return 'pdf'
            
            # Old PPT format
            if header.startswith(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1'):
                return 'ppt'
            
        except Exception as e:
            self.logger.debug(f"Error reading file header: {str(e)}")
        
        # Fallback to extension
        extension_map = {
            '.pptx': 'pptx',
            '.pdf': 'pdf',
            '.ppt': 'ppt',
            '.doc': 'doc',
            '.docx': 'docx',
            '.xls': 'xls',
            '.xlsx': 'xlsx'
        }
        
        return extension_map.get(ext)
    
    def validate_conversion(self, input_path: str, output_path: str,
                           input_format: str, output_format: str) -> Dict[str, Any]:
        """
        Validate conversion quality and completeness
        
        Returns:
            Validation report with warnings and errors
        """
        report = {
            'valid': True,
            'warnings': [],
            'errors': [],
            'stats': {}
        }
        
        # Check input exists
        if not os.path.exists(input_path):
            report['valid'] = False
            report['errors'].append("Input file not found")
            return report
        
        # Check output exists
        if not os.path.exists(output_path):
            report['valid'] = False
            report['errors'].append("Output file not created")
            return report
        
        # Check file sizes
        input_size = os.path.getsize(input_path)
        output_size = os.path.getsize(output_path)
        
        report['stats']['input_size'] = input_size
        report['stats']['output_size'] = output_size
        report['stats']['size_ratio'] = output_size / input_size if input_size > 0 else 0
        
        # Warn about suspicious size ratios
        if report['stats']['size_ratio'] > 10:
            report['warnings'].append(f"Output file is {report['stats']['size_ratio']:.1f}x larger than input")
        elif report['stats']['size_ratio'] < 0.1:
            report['warnings'].append(f"Output file is {report['stats']['size_ratio']:.1f}x smaller than input")
        
        # Validate based on formats
        if input_format == 'pptx' and output_format == 'pdf':
            try:
                # Check page count matches slide count
                from pptx import Presentation
                prs = Presentation(input_path)
                slide_count = len(prs.slides)
                
                pdf_doc = fitz.open(output_path)
                page_count = pdf_doc.page_count
                pdf_doc.close()
                
                report['stats']['source_slides'] = slide_count
                report['stats']['output_pages'] = page_count
                
                if slide_count != page_count:
                    report['warnings'].append(
                        f"Slide count ({slide_count}) != Page count ({page_count})"
                    )
                
                # Check if PDF is readable
                try:
                    pdf_doc = fitz.open(output_path)
                    for page_num in range(min(3, pdf_doc.page_count)):  # Check first 3 pages
                        page = pdf_doc[page_num]
                        text = page.get_text()
                        if not text.strip():
                            report['warnings'].append(f"Page {page_num + 1} appears to have no text")
                    pdf_doc.close()
                except Exception as e:
                    report['errors'].append(f"PDF validation error: {str(e)}")
                    report['valid'] = False
                    
            except Exception as e:
                report['errors'].append(f"Validation error: {str(e)}")
        
        elif input_format == 'pdf' and output_format == 'pptx':
            try:
                # Check if PPTX is valid
                from pptx import Presentation
                prs = Presentation(output_path)
                slide_count = len(prs.slides)
                
                pdf_doc = fitz.open(input_path)
                page_count = pdf_doc.page_count
                pdf_doc.close()
                
                report['stats']['source_pages'] = page_count
                report['stats']['output_slides'] = slide_count
                
                if page_count != slide_count:
                    report['warnings'].append(
                        f"Page count ({page_count}) != Slide count ({slide_count})"
                    )
                    
            except Exception as e:
                report['errors'].append(f"Validation error: {str(e)}")
        
        return report