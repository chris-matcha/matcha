"""
Test script for evidence-based PPTX adaptations
Tests the new formatting features for dyslexia and ADHD profiles
"""

import os
import sys
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from services.pptx_service import PowerPointService
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def create_test_presentation():
    """Create a test presentation with various content types"""
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Evidence-Based Adaptation Test"
    subtitle.text = "Testing dyslexia and ADHD formatting features"
    
    # Slide 2: Content with bullets
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = bullet_slide.shapes.title
    content = bullet_slide.placeholders[1]
    title.text = "Key Concepts"
    content.text = """• Complex scientific terminology requires careful adaptation
• Multiple interconnected concepts need clear structure
• Visual formatting aids comprehension significantly
• Evidence-based approaches improve accessibility"""
    
    # Slide 3: Dense paragraph
    content_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = content_slide.shapes.title
    title.text = "Detailed Explanation"
    
    # Add text box with dense content
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    textbox = content_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = """The process of photosynthesis involves the conversion of light energy into chemical energy. This complex biochemical reaction occurs in chloroplasts and requires specific wavelengths of light, carbon dioxide, and water molecules. The light-dependent reactions produce ATP and NADPH, while the Calvin cycle uses these products to synthesize glucose. Understanding these interconnected processes is essential for comprehending plant biology and ecosystem function."""
    
    # Save test presentation
    test_path = "test_evidence_based_original.pptx"
    prs.save(test_path)
    return test_path

def test_dyslexia_adaptation():
    """Test dyslexia profile adaptations with evidence-based formatting"""
    print("\n=== Testing Dyslexia Adaptations ===")
    
    # Create service instance
    config = {
        'anthropic_api_key': os.getenv('ANTHROPIC_API_KEY'),
        'output_dir': 'test_outputs'
    }
    service = PowerPointService(config)
    
    # Create test presentation
    input_path = create_test_presentation()
    output_path = "test_outputs/dyslexia_adapted.pptx"
    
    # Create output directory
    os.makedirs("test_outputs", exist_ok=True)
    
    # Adapt with dyslexia profile
    success = service.adapt_presentation_preserving_format(
        input_path, output_path, 'dyslexia'
    )
    
    if success:
        print("✅ Dyslexia adaptation completed successfully")
        
        # Verify formatting
        prs = Presentation(output_path)
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\nSlide {slide_idx + 1}:")
            
            # Check background tint
            if hasattr(slide.background.fill, 'fore_color'):
                bg_color = slide.background.fill.fore_color.rgb
                expected_bg = RGBColor(255, 251, 240)  # Beige tint
                if bg_color == expected_bg:
                    print("  ✅ Background tint applied correctly")
                else:
                    print(f"  ❌ Background color mismatch: {bg_color} vs {expected_bg}")
            
            # Check text formatting
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Check alignment
                        if paragraph.alignment == 0:  # PP_ALIGN.LEFT
                            print("  ✅ Text is left-aligned")
                        
                        # Check line spacing
                        if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing:
                            if paragraph.line_spacing >= 1.8:
                                print(f"  ✅ Line spacing: {paragraph.line_spacing}")
            
            # Check for notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if "Evidence-based formatting applied" in notes_text:
                    print("  ✅ Evidence-based formatting documented in notes")
    else:
        print("❌ Dyslexia adaptation failed")
    
    # Cleanup
    if os.path.exists(input_path):
        os.remove(input_path)

def test_adhd_adaptation():
    """Test ADHD profile adaptations with evidence-based formatting"""
    print("\n=== Testing ADHD Adaptations ===")
    
    # Create service instance
    config = {
        'anthropic_api_key': os.getenv('ANTHROPIC_API_KEY'),
        'output_dir': 'test_outputs'
    }
    service = PowerPointService(config)
    
    # Create test presentation
    input_path = create_test_presentation()
    output_path = "test_outputs/adhd_adapted.pptx"
    
    # Adapt with ADHD profile
    success = service.adapt_presentation_preserving_format(
        input_path, output_path, 'adhd'
    )
    
    if success:
        print("✅ ADHD adaptation completed successfully")
        
        # Verify formatting
        prs = Presentation(output_path)
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\nSlide {slide_idx + 1}:")
            
            # Check for visual boundaries
            boundary_count = 0
            for shape in slide.shapes:
                if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                    if shape.line.color.rgb == RGBColor(200, 200, 200):
                        boundary_count += 1
            
            if boundary_count > 0:
                print(f"  ✅ Visual boundaries added: {boundary_count} shapes")
            
            # Check text margins (padding)
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if hasattr(shape.text_frame, 'margin_left'):
                        margin = shape.text_frame.margin_left
                        if margin >= Inches(0.2):
                            print(f"  ✅ Text padding applied: {margin}")
    else:
        print("❌ ADHD adaptation failed")
    
    # Cleanup
    if os.path.exists(input_path):
        os.remove(input_path)

def main():
    """Run all evidence-based adaptation tests"""
    print("Testing Evidence-Based PPTX Adaptations")
    print("=" * 50)
    
    # Test dyslexia adaptations
    test_dyslexia_adaptation()
    
    # Test ADHD adaptations
    test_adhd_adaptation()
    
    print("\n" + "=" * 50)
    print("Testing complete!")
    print("\nCheck the 'test_outputs' directory for adapted presentations.")

if __name__ == "__main__":
    main()