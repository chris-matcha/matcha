"""
Detailed test for background coverage issues
"""

import os
import sys
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from services.pptx_service import PowerPointService
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def create_complex_test_presentation():
    """Create a test presentation with complex layouts to test background coverage"""
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Background Coverage Test"
    subtitle.text = "Testing complete slide background coverage"
    
    # Slide 2: Complex layout with multiple shapes
    complex_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title = complex_slide.shapes.title
    title.text = "Complex Layout Test"
    
    # Add multiple text boxes in different positions
    positions = [
        (Inches(0.5), Inches(1.5), Inches(3), Inches(1)),  # Top left
        (Inches(6), Inches(1.5), Inches(3), Inches(1)),    # Top right
        (Inches(0.5), Inches(4), Inches(3), Inches(1)),    # Bottom left
        (Inches(6), Inches(4), Inches(3), Inches(1)),      # Bottom right
        (Inches(3), Inches(2.5), Inches(3), Inches(1))     # Center
    ]
    
    texts = [
        "Top left corner text",
        "Top right corner text", 
        "Bottom left corner text",
        "Bottom right corner text",
        "Center text"
    ]
    
    for (left, top, width, height), text in zip(positions, texts):
        textbox = complex_slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = text
    
    # Slide 3: Full content slide
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = content_slide.shapes.title
    content = content_slide.placeholders[1]
    title.text = "Full Content Test"
    content.text = """This slide has both title and content areas filled.
The background should cover the entire slide area including:
• Behind the title
• Behind the content
• Behind all margins and empty areas
• Edge to edge coverage"""
    
    # Save test presentation
    test_path = "test_background_coverage_detailed.pptx"
    prs.save(test_path)
    return test_path

def analyze_background_coverage(presentation_path):
    """Analyze the background coverage in detail"""
    print(f"\n🔍 Analyzing background coverage in: {presentation_path}")
    
    prs = Presentation(presentation_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 Slide {slide_idx + 1} Analysis:")
        
        # Count all shapes
        total_shapes = len(slide.shapes)
        background_shapes = 0
        text_shapes = 0
        other_shapes = 0
        
        # Analyze each shape
        background_colors_found = []
        shape_details = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                'index': shape_idx,
                'type': str(type(shape)),
                'has_fill': hasattr(shape, 'fill'),
                'has_text': hasattr(shape, 'text_frame') and shape.text_frame,
                'position': f"({getattr(shape, 'left', 'N/A')}, {getattr(shape, 'top', 'N/A')})",
                'size': f"{getattr(shape, 'width', 'N/A')} x {getattr(shape, 'height', 'N/A')}"
            }
            
            # Check if this could be a background shape
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                try:
                    fill_color = shape.fill.fore_color.rgb
                    shape_info['fill_color'] = f"RGB({fill_color.r}, {fill_color.g}, {fill_color.b})"
                    background_colors_found.append(fill_color)
                    
                    # Check if it's our dyslexia background color
                    expected_bg = RGBColor(255, 251, 240)
                    if fill_color == expected_bg:
                        background_shapes += 1
                        shape_info['is_background'] = True
                        print(f"  🎨 Background shape found at index {shape_idx}: {shape_info['size']} at {shape_info['position']}")
                    else:
                        shape_info['is_background'] = False
                except Exception as color_error:
                    shape_info['fill_color'] = f"Error: {color_error}"
                    shape_info['is_background'] = False
            else:
                shape_info['fill_color'] = "No fill"
                shape_info['is_background'] = False
            
            # Check if it's a text shape
            if shape_info['has_text']:
                text_shapes += 1
                if shape.text_frame.text.strip():
                    shape_info['text_content'] = shape.text_frame.text[:30] + "..."
                else:
                    shape_info['text_content'] = "[Empty]"
            else:
                other_shapes += 1
            
            shape_details.append(shape_info)
        
        # Print summary
        print(f"  📊 Shape Summary:")
        print(f"    Total shapes: {total_shapes}")
        print(f"    Background shapes: {background_shapes}")
        print(f"    Text shapes: {text_shapes}")
        print(f"    Other shapes: {other_shapes}")
        
        # Print detailed shape info
        print(f"  📋 Detailed Shape Analysis:")
        for shape_info in shape_details:
            marker = "🎨" if shape_info['is_background'] else ("📝" if shape_info['has_text'] else "🔹")
            print(f"    {marker} Shape {shape_info['index']}: {shape_info['size']} | {shape_info['fill_color']} | {shape_info.get('text_content', 'No text')}")
        
        # Check slide background
        print(f"  🖼️ Slide Background Check:")
        try:
            if hasattr(slide.background.fill, 'fore_color'):
                slide_bg_color = slide.background.fill.fore_color.rgb
                expected_bg = RGBColor(255, 251, 240)
                if slide_bg_color == expected_bg:
                    print(f"    ✅ Slide background correctly set to dyslexia beige")
                else:
                    print(f"    ❌ Slide background is RGB({slide_bg_color.r}, {slide_bg_color.g}, {slide_bg_color.b}), expected RGB(255, 251, 240)")
            else:
                print(f"    ⚠️ Could not read slide background color")
        except Exception as slide_bg_error:
            print(f"    ❌ Error reading slide background: {slide_bg_error}")
        
        # Overall assessment
        if background_shapes > 0:
            print(f"  ✅ Background coverage: {background_shapes} background shape(s) found")
        else:
            print(f"  ❌ Background coverage: No background shapes found!")

def test_background_coverage():
    """Test background coverage with detailed analysis"""
    print("🧪 Testing Background Coverage with Detailed Analysis")
    print("=" * 60)
    
    # Create service instance
    config = {
        'anthropic_api_key': os.getenv('ANTHROPIC_API_KEY'),
        'output_dir': 'test_outputs'
    }
    service = PowerPointService(config)
    
    # Create test presentation
    input_path = create_complex_test_presentation()
    output_path = "test_outputs/background_coverage_detailed_test.pptx"
    
    # Create output directory
    os.makedirs("test_outputs", exist_ok=True)
    
    print("\n🔄 Running dyslexia adaptation...")
    
    # Adapt with dyslexia profile
    success = service.adapt_presentation_preserving_format(
        input_path, output_path, 'dyslexia'
    )
    
    if success:
        print("✅ Adaptation completed successfully")
        
        # Analyze the original for comparison
        print("\n📊 ORIGINAL PRESENTATION ANALYSIS:")
        analyze_background_coverage(input_path)
        
        # Analyze the adapted version
        print("\n📊 ADAPTED PRESENTATION ANALYSIS:")
        analyze_background_coverage(output_path)
        
    else:
        print("❌ Adaptation failed")
    
    # Cleanup
    if os.path.exists(input_path):
        os.remove(input_path)
    
    print(f"\n💡 Open {output_path} to manually verify background coverage")
    print("   Check that the warm beige background covers the entire slide area")

if __name__ == "__main__":
    test_background_coverage()