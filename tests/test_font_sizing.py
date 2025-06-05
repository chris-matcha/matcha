#!/usr/bin/env python3
"""
Test script for the new font sizing functionality in PowerPointService
"""
import sys
import os
sys.path.append('.')

from services.pptx_service import PowerPointService
from pptx import Presentation
from pptx.util import Inches, Pt


def test_font_calculation():
    """Test the font size calculation functionality"""
    print("🔧 Testing Font Size Calculation Enhancement")
    print("=" * 50)
    
    # Initialize service
    config = {
        'output_folder': './test_outputs',
        'upload_folder': './uploads'
    }
    pptx_service = PowerPointService(config)
    
    # Test 1: Basic font size calculation
    print("\n📏 Test 1: Basic Font Size Calculation")
    test_text = "This is a sample text for testing optimal font sizing"
    optimal_size = pptx_service.calculate_optimal_font_size(
        text=test_text,
        max_width_pts=400,
        max_height_pts=100,
        font_name="Arial"
    )
    print(f"✓ Optimal font size for '{test_text[:30]}...': {optimal_size}pt")
    
    # Test 2: Multi-line text
    print("\n📏 Test 2: Multi-line Text Calculation")
    multiline_text = """This is a longer text that spans multiple lines.
It should be properly measured to ensure it fits within the given bounds.
The font size should be calculated to prevent overflow."""
    
    optimal_size_multiline = pptx_service.calculate_optimal_font_size(
        text=multiline_text,
        max_width_pts=300,
        max_height_pts=150,
        font_name="Arial"
    )
    print(f"✓ Optimal font size for multi-line text: {optimal_size_multiline}pt")
    
    # Test 3: Text bounds measurement
    print("\n📏 Test 3: Text Bounds Measurement")
    width, height = pptx_service.measure_text_bounds(
        text="Sample text", 
        font_name="Arial", 
        font_size=12
    )
    print(f"✓ Text bounds for 'Sample text' at 12pt: {width:.1f}pt x {height:.1f}pt")
    
    # Test 4: Create a test presentation with optimal sizing
    print("\n🎯 Test 4: Creating Test Presentation with Optimal Sizing")
    try:
        # Create a new presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
        
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = "Font Sizing Test"
        
        # Add content with long text that would normally overflow
        long_text = """This is a very long piece of text that would normally cause overflow issues in PowerPoint slides. 
The new optimal font sizing feature should automatically calculate the appropriate font size to ensure this text fits 
within the available space while maintaining readability. This demonstrates the enhanced text handling capabilities 
that prevent common layout issues in educational content adaptation."""
        
        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Content placeholder
                content_placeholder = shape
                break
        
        if content_placeholder:
            # Test the enhanced apply_text_to_text_frame method
            pptx_service.apply_text_to_text_frame(
                text_frame=content_placeholder.text_frame,
                adapted_text=long_text,
                profile='dyslexia',
                use_optimal_sizing=True
            )
            print("✓ Applied long text with optimal sizing for dyslexia profile")
        
        # Save test presentation
        os.makedirs('test_outputs', exist_ok=True)
        test_file = 'test_outputs/font_sizing_test.pptx'
        prs.save(test_file)
        print(f"✓ Test presentation saved: {test_file}")
        
    except Exception as e:
        print(f"❌ Error creating test presentation: {e}")
    
    print("\n🎉 Font sizing enhancement testing completed!")
    print("\n📋 Summary of Enhancements:")
    print("   • Binary search algorithm for optimal font sizing")
    print("   • Multi-line text support with accurate measurement")
    print("   • Automatic text frame bounds detection")
    print("   • Profile-specific font preferences integration")
    print("   • Fallback handling for font loading errors")
    print("   • Backward compatibility with existing code")


if __name__ == "__main__":
    test_font_calculation()