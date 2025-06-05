"""
Test script for OpenDyslexic font and full background color fixes
"""

import os
import sys
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from services.pptx_service import PowerPointService
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def create_test_presentation():
    """Create a test presentation for font and background testing"""
    prs = Presentation()
    
    # Slide 1: Title slide with multiple text elements
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "OpenDyslexic Font Test"
    subtitle.text = "Testing font application and background coverage"
    
    # Slide 2: Content with different text formatting
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = content_slide.shapes.title
    content = content_slide.placeholders[1]
    title.text = "Font Testing Content"
    content.text = """This text should appear in OpenDyslexic font when adapted for dyslexia profile.
The background should be a warm beige color covering the entire slide.
Check that both title and content text use the dyslexia-friendly font."""
    
    # Slide 3: Manual text box to test different shape types
    manual_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = manual_slide.shapes.title
    title.text = "Manual Text Box Test"
    
    # Add manual text box
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    textbox = manual_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "This text is in a manually created text box. It should also be converted to OpenDyslexic font with proper background coverage."
    
    # Save test presentation
    test_path = "test_font_background_original.pptx"
    prs.save(test_path)
    return test_path

def test_dyslexia_font_and_background():
    """Test OpenDyslexic font application and full background coverage"""
    print("\n=== Testing OpenDyslexic Font and Background Coverage ===")
    
    # Create service instance
    config = {
        'anthropic_api_key': os.getenv('ANTHROPIC_API_KEY'),
        'output_dir': 'test_outputs'
    }
    service = PowerPointService(config)
    
    # Create test presentation
    input_path = create_test_presentation()
    output_path = "test_outputs/dyslexia_font_background_test.pptx"
    
    # Create output directory
    os.makedirs("test_outputs", exist_ok=True)
    
    # Adapt with dyslexia profile
    print("🔄 Starting dyslexia adaptation...")
    success = service.adapt_presentation_preserving_format(
        input_path, output_path, 'dyslexia'
    )
    
    if success:
        print("✅ Dyslexia adaptation completed successfully")
        
        # Verify the adaptations
        print("\n📊 Analyzing adapted presentation...")
        prs = Presentation(output_path)
        
        font_check_results = []
        background_check_results = []
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n🔍 Slide {slide_idx + 1} Analysis:")
            
            # Check background coverage
            background_found = False
            
            # Method 1: Check slide background
            try:
                if hasattr(slide.background.fill, 'fore_color'):
                    bg_color = slide.background.fill.fore_color.rgb
                    expected_bg = RGBColor(255, 251, 240)  # Beige tint
                    if bg_color == expected_bg:
                        print("  ✅ Slide background tint applied correctly")
                        background_found = True
                    else:
                        print(f"  ⚠️ Slide background color: {bg_color} (expected: {expected_bg})")
            except Exception as e:
                print(f"  ⚠️ Could not read slide background: {e}")
            
            # Method 2: Check for background rectangle (fallback method)
            if not background_found:
                rectangle_backgrounds = 0
                for shape in slide.shapes:
                    if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                        try:
                            shape_color = shape.fill.fore_color.rgb
                            expected_bg = RGBColor(255, 251, 240)
                            if shape_color == expected_bg:
                                rectangle_backgrounds += 1
                        except:
                            pass
                
                if rectangle_backgrounds > 0:
                    print(f"  ✅ Background rectangle found: {rectangle_backgrounds} background shapes")
                    background_found = True
            
            if not background_found:
                print("  ❌ No background tint found!")
            
            background_check_results.append(background_found)
            
            # Check font application
            fonts_found = []
            text_elements_checked = 0
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, 'font') and hasattr(run.font, 'name'):
                                text_elements_checked += 1
                                font_name = run.font.name
                                fonts_found.append(font_name)
                                
                                # Check for OpenDyslexic variants
                                if font_name and 'opendyslexic' in font_name.lower():
                                    print(f"  ✅ OpenDyslexic font found: '{font_name}'")
                                elif font_name and font_name.lower() in ['arial', 'comic sans ms', 'verdana']:
                                    print(f"  ⚠️ Fallback font applied: '{font_name}' (OpenDyslexic may not be installed)")
                                else:
                                    print(f"  ❌ Unexpected font: '{font_name}'")
            
            print(f"  📝 Text elements checked: {text_elements_checked}")
            print(f"  🔤 Fonts found: {set(fonts_found)}")
            
            # Consider it a success if we find OpenDyslexic OR dyslexia-friendly fallbacks
            dyslexia_friendly_fonts = any(
                font and ('opendyslexic' in font.lower() or 
                         font.lower() in ['arial', 'comic sans ms', 'verdana'])
                for font in fonts_found
            )
            font_check_results.append(dyslexia_friendly_fonts)
        
        # Overall results
        print(f"\n📋 Overall Results:")
        print(f"  Background coverage: {sum(background_check_results)}/{len(background_check_results)} slides")
        print(f"  Font application: {sum(font_check_results)}/{len(font_check_results)} slides")
        
        if all(background_check_results):
            print("  ✅ Background tint successfully applied to all slides")
        else:
            print("  ❌ Some slides missing background tint")
        
        if all(font_check_results):
            print("  ✅ Dyslexia-friendly fonts applied to all slides")
        else:
            print("  ❌ Some slides missing dyslexia-friendly fonts")
        
        # Check notes for documentation
        notes_documented = 0
        for slide in prs.slides:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if "OpenDyslexic" in notes_text or "dyslexia" in notes_text.lower():
                    notes_documented += 1
        
        print(f"  📝 Documentation in notes: {notes_documented}/{len(prs.slides)} slides")
        
    else:
        print("❌ Dyslexia adaptation failed")
    
    # Cleanup
    if os.path.exists(input_path):
        os.remove(input_path)
    
    return success

def main():
    """Run font and background fixes test"""
    print("Testing OpenDyslexic Font and Background Fixes")
    print("=" * 60)
    
    success = test_dyslexia_font_and_background()
    
    print("\n" + "=" * 60)
    if success:
        print("✅ Testing completed! Check the test_outputs/dyslexia_font_background_test.pptx file")
        print("\n💡 Tips for verifying manually:")
        print("   1. Open the generated PPTX file")
        print("   2. Check if text appears in OpenDyslexic font (or Arial/Comic Sans as fallback)")
        print("   3. Verify the background has a warm beige color across entire slides")
        print("   4. Check slide notes for adaptation documentation")
    else:
        print("❌ Testing failed")

if __name__ == "__main__":
    main()