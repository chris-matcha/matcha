# Grade Level (lower is better)
import os
from flask import Flask, request, send_file, render_template_string, redirect, url_for, jsonify
import uuid
from pptx import Presentation
from pptx.dml.color import RGBColor
import anthropic
import re
import io
import base64
import threading
import json

# Handle matplotlib import gracefully
try:
    import matplotlib
    matplotlib.use('Agg')  # Use non-interactive backend
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
except ImportError as e:
    print(f"Warning: matplotlib not available - chart generation disabled: {e}")
    MATPLOTLIB_AVAILABLE = False
    # Create dummy plt object
    class DummyPlt:
        @staticmethod
        def figure(*args, **kwargs): pass
        @staticmethod
        def bar(*args, **kwargs): pass
        @staticmethod
        def barh(*args, **kwargs): pass
        @staticmethod
        def axhline(*args, **kwargs): pass
        @staticmethod
        def xlabel(*args, **kwargs): pass
        @staticmethod
        def ylabel(*args, **kwargs): pass
        @staticmethod
        def title(*args, **kwargs): pass
        @staticmethod
        def legend(*args, **kwargs): pass
        @staticmethod
        def grid(*args, **kwargs): pass
        @staticmethod
        def text(*args, **kwargs): pass
        @staticmethod
        def ylim(*args, **kwargs): pass
        @staticmethod
        def savefig(*args, **kwargs): pass
        @staticmethod
        def close(*args, **kwargs): pass
        @staticmethod
        def axis(*args, **kwargs): pass
    
    plt = DummyPlt()
from api_utils import ApiUtils, API_CHECK_SUCCESS_TEMPLATE, API_CHECK_ERROR_TEMPLATE
from migrate_pdf_functions import PDFMigrationHelper


# Global dictionaries to store status information
global_analysis_status = {}
processing_tasks = {}

# Define the upload folder path
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')

# Set up Flask
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER  # Use the absolute path
app.config['OUTPUT_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'outputs')
app.config['REPORT_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'reports')
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB max upload

# Create directories if they don't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
os.makedirs(app.config['REPORT_FOLDER'], exist_ok=True)


# Set up API client (not exposed in UI)
import os
from dotenv import load_dotenv
load_dotenv()

api_key = os.getenv('ANTHROPIC_API_KEY')
if not api_key:
    raise ValueError("ANTHROPIC_API_KEY environment variable is required")

api_utils = ApiUtils(api_key)
client = anthropic.Anthropic(api_key=api_utils.api_key)

# Initialize PDF Migration Helper for new service-based PDF processing
pdf_migration_helper = PDFMigrationHelper({
    'output_folder': app.config['OUTPUT_FOLDER'],
    'upload_folder': app.config['UPLOAD_FOLDER'],
    'anthropic_api_key': api_utils.api_key
})

# Initialize enhanced services
try:
    from services import (
        PDFService, PowerPointService, ConversionService, 
        EducationalContentService, LearningProfilesService, UploadService,
        DownloadsService, FileStoreService, AdaptationsService, TranslationsService,
        AssessmentsService
    )
except ImportError as e:
    print(f"Error importing services: {e}")
    print("Some features may not be available. Please check your services module.")
    # Create dummy classes to prevent crashes
    class DummyService:
        def __init__(self, config): pass
        def __getattr__(self, name): return lambda *args, **kwargs: None
    
    PDFService = PowerPointService = ConversionService = DummyService
    EducationalContentService = LearningProfilesService = UploadService = DummyService
    DownloadsService = FileStoreService = AdaptationsService = DummyService
    TranslationsService = AssessmentsService = DummyService

service_config = {
    'output_folder': app.config['OUTPUT_FOLDER'],
    'upload_folder': app.config['UPLOAD_FOLDER'],
    'anthropic_api_key': api_utils.api_key,
    'upload_dir': app.config['UPLOAD_FOLDER'],
    'output_dir': app.config['OUTPUT_FOLDER'],
    'temp_dir': os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp')
}

pdf_service = PDFService(service_config)
pptx_service = PowerPointService(service_config)
conversion_service = ConversionService(service_config)
educational_service = EducationalContentService(service_config)
profiles_service = LearningProfilesService(service_config)
upload_service = UploadService(service_config)
downloads_service = DownloadsService(service_config)
filestore_service = FileStoreService(service_config)
adaptations_service = AdaptationsService(service_config)
translations_service = TranslationsService(service_config)
assessments_service = AssessmentsService(service_config)

# Helper function to generate output file path
def get_output_file_path(file_id, filename):
    """Generate output file path using consistent naming"""
    return os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")

# Helper function to register a file that was created externally
def register_output_file(file_id, filename, file_path):
    """Register a file with FileStoreService after it's been created"""
    # For now, just log it - in future could track in database
    if os.path.exists(file_path):
        file_info = filestore_service.get_file_info(file_path)
        print(f"Registered output file: {filename} for {file_id} - Size: {file_info.get('size', 0)} bytes")
        return True
    return False

# Helper function to find output files for downloads
def find_output_file(file_id, filename):
    """Find output file using FileStoreService and fallback patterns"""
    # First try FileStoreService
    file_path = filestore_service.get_file_path(file_id, filename, 'output')
    if file_path:
        return file_path
    
    # Try with adapted prefix
    if not filename.startswith('adapted_'):
        file_path = filestore_service.get_file_path(file_id, f"adapted_{filename}", 'output')
        if file_path:
            return file_path
    
    # Fallback to direct file system check for legacy files
    output_folder = app.config['OUTPUT_FOLDER']
    possible_paths = [
        os.path.join(output_folder, f"{file_id}_{filename}"),
        os.path.join(output_folder, filename),
        os.path.join(output_folder, f"adapted_{filename}" if not filename.startswith('adapted_') else filename),
        os.path.join(output_folder, f"{file_id}_adapted_{filename}" if not filename.startswith('adapted_') else f"{file_id}_{filename}")
    ]
    
    for path in possible_paths:
        if os.path.exists(path):
            return path
    
    return None

# Get profile colors from service and convert to RGBColor objects
def get_profile_color(profile_id):
    """Get RGBColor for a profile"""
    colors = profiles_service.get_profile_colors(profile_id)
    if not colors:
        return None
    # Convert hex to RGB
    hex_color = colors.get('text', '#000000')
    hex_color = hex_color.lstrip('#')
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return RGBColor(r, g, b)

def get_profile_names():
    """Get profile display names from service"""
    profiles = {}
    for profile in profiles_service.list_profiles():
        profiles[profile['id']] = profile['name']
    return profiles

def get_readability_thresholds(profile_id):
    """Get readability thresholds from service"""
    thresholds = profiles_service.get_thresholds(profile_id)
    # Map service threshold names to app's expected format
    return {
        'flesch_reading_ease': thresholds.get('flesch_ease', 60),
        'flesch_kincaid_grade': thresholds.get('grade_level', 8),
        'smog_index': thresholds.get('smog_index', 8),
        'sentence_length': 15,  # Default, not in service
        'complex_word_percent': thresholds.get('complex_word_threshold', 10)
    }

# Cache for storing adaptation results

# HTML Templates
from html_templates import (DOWNLOAD_TEMPLATE_UNIVERSAL, INDEX_TEMPLATE, SCAFFOLDING_TEMPLATE, FRAMEWORK_TEMPLATE,
                           ASSESSMENT_TEMPLATE_SIMPLIFIED, ASSESSMENT_TEMPLATE, 
                           PROCESSING_TEMPLATE_WITH_PROGRESS, DOWNLOAD_TEMPLATE_WITH_TRANSLATION, ERROR_TEMPLATE)

@app.route('/')
def index():
    """Render the home page with a form"""
    return render_template_string(INDEX_TEMPLATE)

@app.route('/analysis_status/<file_id>', methods=['GET'])
def check_analysis_status(file_id):
    """Check the status of structural analysis"""
    if file_id in global_analysis_status:
        return jsonify(global_analysis_status[file_id])
    else:
        return jsonify({'status': 'not_found'})

# Clean Upload Route - Replace the entire broken upload section with this

@app.route('/upload', methods=['POST'])
def upload():
    """Handle file upload and direct to appropriate process - supports both PPTX and PDF"""
    try:
        # Get the uploaded file (check common field names)
        file = None
        if 'pptx' in request.files and request.files['pptx'].filename:
            file = request.files['pptx']
        elif 'pdf' in request.files and request.files['pdf'].filename:
            file = request.files['pdf']
        elif 'file' in request.files and request.files['file'].filename:
            file = request.files['file']
        
        # Validate file presence
        if not file or file.filename == '':
            return render_template_string(ERROR_TEMPLATE, 
                message="No file selected. Please choose a PowerPoint (.pptx) or PDF file."), 400
        
        # Get form parameters
        profile = request.form.get('profile')
        action = request.form.get('action', 'assess')
        target_language = request.form.get('target_language', '')
        translation_mode = request.form.get('translation_mode', 'copy')  # Default to copy mode
        export_format = request.form.get('export_format', 'pdf')
        
        # Debug form parameters
        print(f"DEBUG: Form parameters - profile: {profile}, action: {action}, target_language: '{target_language}', translation_mode: '{translation_mode}', export_format: {export_format}")
        
        # Validate required parameters
        if not profile:
            return render_template_string(ERROR_TEMPLATE, 
                message="Please select a learning profile."), 400
        
        # Process upload through service
        metadata = {
            'profile': profile,
            'action': action,
            'target_language': target_language,
            'translation_mode': translation_mode,
            'export_format': export_format
        }
        
        upload_result = upload_service.process_upload(file, metadata)
        
        if not upload_result['success']:
            return render_template_string(ERROR_TEMPLATE, 
                message=upload_result.get('error', 'File upload failed')), 400
        
        # Extract data from upload result
        file_id = upload_result['file_id']
        filename = upload_result['filename']
        file_path = upload_result['file_path']
        file_ext = f".{upload_result['file_type']}"
        
        # Initialize processing task
        processing_tasks[file_id] = {
            'status': 'upload', 
            'filename': filename, 
            'profile': profile,
            'file_type': file_ext,
            'metadata': metadata
        }
        
        # Route based on action
        if action == 'assess':
            # Assessment route
            if target_language:
                return redirect(url_for('analyze_scaffolding', file_id=file_id) + f'?profile={profile}&target_language={target_language}')
            else:
                return redirect(url_for('analyze_scaffolding', file_id=file_id) + f'?profile={profile}')
        
        else:  # action == 'adapt'
            # Start adaptation process
            processing_tasks[file_id].update({
                'status': 'processing', 
                'message': 'Starting adaptation...',
                'export_format': export_format,
                'progress': {'total': 0, 'processed': 0, 'percentage': 0}
            })
            
            # Choose processing method based on file type
            if file_ext == '.pdf':
                # Use new service-based PDF processing
                print(f"Processing PDF with service-based system: {filename}")
                thread_target = process_pdf_with_services
                thread_args = (file_path, file_id, filename, profile, export_format, target_language)
            else:  # .pptx
                print(f"Processing PowerPoint: {filename}")
                # Create wrapper function to pass processing_tasks
                def pptx_thread_target():
                    try:
                        # Check if this is a translation request (any profile with target language)
                        has_target_language = (target_language and target_language.strip())
                        
                        # Use translation method based on mode
                        if has_target_language:
                            # Determine which translation method to use based on mode
                            if translation_mode == 'replace':
                                print(f"Using replace mode translation to {target_language}")
                                processing_tasks[file_id]['message'] = f'Translating content to {target_language} (replace mode)...'
                                processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 5, 'percentage': 5}
                                
                                # Use in-place translation (replaces original content)
                                result = pptx_service.translate_presentation_in_place(
                                    file_path, file_id, filename, target_language, processing_tasks[file_id]
                                )
                            else:  # copy mode (default)
                                print(f"Using copy mode translation to {target_language}")
                                processing_tasks[file_id]['message'] = f'Creating translated copies for {target_language}...'
                                processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 5, 'percentage': 5}
                                
                                # Use slide duplication method (keeps originals + adds translations)
                                result = pptx_service.translate_presentation(
                                    file_path, file_id, filename, target_language, processing_tasks[file_id]
                                )
                        else:
                            # Use the standard adaptation method (with optional translation)
                            print(f"Using standard adaptation method for profile: {profile}")
                            result = pptx_service.process_presentation_efficiently(
                                file_path, file_id, filename, profile, target_language, processing_tasks
                            )
                        
                        if result:
                            processing_tasks[file_id]['adapted_path'] = result
                            processing_tasks[file_id]['status'] = 'completed'
                            print(f"DEBUG: Processing completed successfully for {file_id}")
                        else:
                            processing_tasks[file_id]['status'] = 'error'
                            processing_tasks[file_id]['message'] = 'Processing returned no result'
                            print(f"DEBUG: Processing failed - no result returned for {file_id}")
                            
                    except Exception as e:
                        print(f"ERROR in PowerPoint processing: {str(e)}")
                        import traceback
                        traceback.print_exc()
                        processing_tasks[file_id]['status'] = 'error'
                        processing_tasks[file_id]['message'] = f'Processing error: {str(e)}'
                
                thread_target = pptx_thread_target
                thread_args = ()
            
            # Start background processing
            thread = threading.Thread(
                target=thread_target,
                args=thread_args,
                name=f"process-{file_id}"
            )
            thread.daemon = True
            thread.start()
            
            # Profile display names
            profile_names = get_profile_names()
            
            # Return processing page
            return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                file_id=file_id, 
                filename=f"adapted_{filename}",
                profile=profile,
                profile_name=profile_names.get(profile, profile)
            )
            
    except Exception as e:
        print(f"Error in upload route: {str(e)}")
        return render_template_string(ERROR_TEMPLATE, 
            message=f"Upload failed: {str(e)}"), 500

# Add this helper function to gracefully handle missing PDF libraries


def analyze_instructional_framework(pptx_path):
    """Analyze instructional framework using educational service"""
    try:
        return educational_service.analyze_instructional_framework(pptx_path)
    except Exception as e:
        print(f"Error in framework analysis: {str(e)}")
        return {"error": f"Analysis failed: {str(e)}"}

def process_pdf_with_services(file_path, file_id, filename, profile, export_format='pdf', target_language=None):
    """Service-based PDF processing function using the new architecture"""
    print(f"Starting service-based PDF processing for {filename} with profile: {profile}")
    
    try:
        processing_tasks[file_id]['message'] = 'Initializing PDF processing services...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 5, 'percentage': 5}
        
        # Use the migration helper to process the PDF
        processing_tasks[file_id]['message'] = 'Processing PDF with visual preservation...'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 20, 'percentage': 20}
        
        # Process with services - this handles adaptation and visual preservation
        result = pdf_migration_helper.process_with_pdf_template_system(
            file_path, 
            profile, 
            direct_adapt=True,
            preserve_visuals=True,
            translate=(target_language and target_language.strip() != ''),
            target_language=target_language
        )
        
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 60, 'percentage': 60}
        
        if not result['success']:
            raise Exception(f"PDF processing failed: {result.get('error', 'Unknown error')}")
        
        # Update processing status
        processing_tasks[file_id]['message'] = 'PDF adaptation completed successfully!'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 90, 'percentage': 90}
        
        # Store the output path for downloads
        output_path = result['output_path']
        processing_tasks[file_id]['output_path'] = output_path
        processing_tasks[file_id]['adapted_content'] = result.get('adapted_content', {})
        
        # Handle translation if requested and available
        if target_language and target_language.strip():
            processing_tasks[file_id]['message'] = f'Creating translation in {target_language}...'
            processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 95, 'percentage': 95}
            
            # Check if translation was successful
            if 'translated_output_path' in result:
                processing_tasks[file_id]['translated_output_path'] = result['translated_output_path']
                processing_tasks[file_id]['translated_language'] = result['translated_language']
                processing_tasks[file_id]['translation_languages'] = [target_language]
                print(f"✓ Translation completed: {result['translated_output_path']}")
            else:
                print(f"⚠ Translation was requested but failed for {target_language}")
                processing_tasks[file_id]['translation_error'] = f'Translation to {target_language} failed'
        
        # Mark as completed
        processing_tasks[file_id]['status'] = 'completed'
        processing_tasks[file_id]['message'] = 'Processing completed successfully!'
        processing_tasks[file_id]['progress'] = {'total': 100, 'processed': 100, 'percentage': 100}
        
        print(f"✓ Service-based PDF processing completed successfully for {filename}")
        return True
        
    except Exception as e:
        print(f"✗ Error in service-based PDF processing: {str(e)}")
        import traceback
        traceback.print_exc()
        
        # Update task with error
        processing_tasks[file_id]['status'] = 'error'
        processing_tasks[file_id]['message'] = f'Error: {str(e)}'
        processing_tasks[file_id]['error'] = str(e)
        
        return False

# Original PDF processing function (kept for fallback)
def update_processing_status(file_id, message, percentage):
    """Helper function to update processing status"""
    if percentage < 0:  # Error state
        processing_tasks[file_id] = {
            'status': 'error',
            'message': message
        }
    else:
        if file_id not in processing_tasks:
            processing_tasks[file_id] = {}
        processing_tasks[file_id]['message'] = message
        processing_tasks[file_id]['progress'] = {
            'total': 100,
            'processed': percentage,
            'percentage': percentage
        }
        if percentage >= 100:
            processing_tasks[file_id]['status'] = 'complete'

def process_with_pdf_template_system(file_path, file_id, filename, profile, export_format='pdf', target_language=None):
    """Main PDF processing function - delegates to pdf_service"""
    # Use the pdf_service's comprehensive processing method
    output_path = pdf_service.process_with_template_system(
        file_path=file_path,
        file_id=file_id,
        filename=filename,
        profile=profile,
        export_format=export_format,
        target_language=target_language,
        processing_callback=update_processing_status,
        output_path_callback=get_output_file_path
    )
    
    if output_path:
        # Handle successful processing
        status_data = {
            'status': 'complete',
            'message': f'PDF content successfully adapted to {export_format.upper()}',
            'adapted_path': output_path,
            'export_format': export_format,
            'profile': profile,
            'progress': {
                'total': 100,
                'processed': 100,
                'percentage': 100
            }
        }
        
        # Check for translation
        if target_language:
            # Check if translated file exists
            base_name = os.path.splitext(filename)[0]
            if base_name.startswith('adapted_'):
                base_name = base_name[8:]
            translated_filename = f"translated_{target_language}_{base_name}.pdf"
            translated_path = os.path.join(os.path.dirname(output_path), f"{file_id}_{translated_filename}")
            
            if os.path.exists(translated_path):
                status_data['has_translation'] = True
                status_data['translated_path'] = translated_path
                status_data['target_language'] = target_language
                status_data['translated_language'] = target_language.title()
                status_data['translated_filename'] = translated_filename
                status_data['message'] += f' and translated to {target_language}'
        
        processing_tasks[file_id] = status_data
        register_output_file(file_id, os.path.basename(output_path), output_path)
        
    return output_path

# Fix the upload route's PDF handling



def create_pptx_from_pdf_content(adapted_content, output_path, profile):
    """Create PowerPoint from adapted PDF content"""
    try:
        prs = Presentation()
        
        for page in adapted_content['pages']:
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = f"Page {page['page_number']}"
                # Apply profile formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    pptx_service.apply_text_to_text_frame(slide.shapes.title.text_frame, 
                                            slide.shapes.title.text, profile, profiles_service)
            
            # Add content
            if page['text']:
                # Find content placeholder
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Content placeholder
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    pptx_service.apply_text_to_text_frame(content_placeholder.text_frame, 
                                           page['text'], profile, profiles_service)
            
            # Note: Images would need additional handling to be added to PPTX
            # This is a simplified version focusing on text content
        
        prs.save(output_path)
        return True
        
    except Exception as e:
        print(f"Error creating PPTX from PDF content: {e}")
        return False


@app.route('/convert_to_pdf/<file_id>')
def convert_to_pdf(file_id):
    """Convert an uploaded PPTX to PDF template"""
    try:
        if file_id not in processing_tasks:
            return jsonify({'error': 'File not found'}), 404
        
        filename = processing_tasks[file_id].get('filename', '')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found on server'}), 404
        
        # Convert to PDF
        pdf_filename = f"{os.path.splitext(filename)[0]}_template.pdf"
        pdf_path = get_output_file_path(file_id, pdf_filename)
        
        success = conversion_service.convert_pptx_to_pdf_template(file_path, pdf_path)
        
        if success:
            processing_tasks[file_id]['pdf_template'] = pdf_path
            return jsonify({
                'status': 'success',
                'pdf_path': pdf_path,
                'message': 'PDF template created successfully'
            })
        else:
            return jsonify({'error': 'Failed to convert to PDF'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/check_api')
def check_api():
    """API connection check endpoint"""
    result = api_utils.check_connection()
    
    if result["status"] == "connected":
        # Return success page
        return render_template_string(API_CHECK_SUCCESS_TEMPLATE, result=result)
    else:
        # Return error page
        return render_template_string(API_CHECK_ERROR_TEMPLATE, result=result)

# Replace the download route in app.py with this fixed version

@app.route('/download/<file_id>/<filename>')
def download(file_id, filename):
    """Universal download page for both PDF and PPTX adaptations - uses downloads_service"""
    try:
        # Get download page info from service
        download_info = downloads_service.get_download_page_info(file_id, filename, processing_tasks)
        
        # Get profile display names
        profile_names = get_profile_names()
        profile_name = profile_names.get(download_info['profile'], download_info['profile'])
        
        # Use the universal template
        template = DOWNLOAD_TEMPLATE_UNIVERSAL
        
        return render_template_string(
            template,
            file_id=download_info['file_id'],
            filename=download_info.get('original_filename', download_info['filename']),
            profile=download_info['profile'],
            profile_name=profile_name,
            original_format=download_info['original_format'],
            export_format=download_info['export_format'],
            pdf_available=download_info['pdf_available'],
            pptx_available=download_info['pptx_available'],
            pdf_filename=download_info['pdf_filename'],
            pptx_filename=download_info['pptx_filename'],
            template_available=download_info['template_available'],
            has_translation=download_info['has_translation'],
            translated_filename=download_info['translated_filename'],
            translated_language=download_info['translated_language'],
            processing_time=None,
            file_count=download_info['file_count']
        )
    
    except Exception as e:
        print(f"Error in download route: {e}")
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, 
                                     message=f"Error loading download page: {str(e)}"), 500

@app.route('/download_file/<file_id>/<filename>')
def download_file(file_id, filename):
    """Actually download the file with proper filename handling - uses downloads_service"""
    try:
        print(f"Download request: file_id={file_id}, filename={filename}")
        
        # Use downloads service to find and prepare the file
        result = downloads_service.get_file_for_download(file_id, filename, processing_tasks)
        
        if result:
            file_path, clean_filename = result
            print(f"Found file at: {file_path}, downloading as: {clean_filename}")
            return send_file(file_path, as_attachment=True, download_name=clean_filename)
        
        # No file found
        return render_template_string(ERROR_TEMPLATE, 
            message=f"File not found: {filename}"), 404
        
    except Exception as e:
        print(f"Error downloading file: {e}")
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE,
            message=f"Error downloading file: {str(e)}"), 500

@app.route('/debug/files/<file_id>')
def debug_files(file_id):
    """Debug route to check file status and locations"""
    try:
        # Get task info
        task_info = processing_tasks.get(file_id, {})
        
        # List all files in output folder
        output_files = []
        try:
            output_files = os.listdir(app.config['OUTPUT_FOLDER'])
        except:
            output_files = ["Error reading output folder"]
        
        # Find files matching this file_id
        matching_files = [f for f in output_files if file_id in f]
        
        # Debug info
        debug_info = {
            'file_id': file_id,
            'task_info': task_info,
            'output_folder': app.config['OUTPUT_FOLDER'],
            'all_output_files_count': len(output_files),
            'matching_files': matching_files,
            'processing_tasks_keys': list(processing_tasks.keys())
        }
        
        return jsonify(debug_info)
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/analyze/framework/<file_id>')
def analyze_framework(file_id):
    """Analyze a presentation's instructional framework"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing_framework'
        processing_tasks[file_id]['message'] = 'Analyzing instructional framework...'
        
        # Analyze framework
        framework_data = analyze_instructional_framework(file_path)
        
        # Store the framework data
        processing_tasks[file_id]['framework'] = framework_data
        processing_tasks[file_id]['status'] = 'complete'
        
        # Create a template to display the results
        return render_template_string(
            FRAMEWORK_TEMPLATE,
            framework=framework_data,
            file_id=file_id
        )
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing framework: {str(e)}"
        }), 500


@app.route('/analyze/scaffolding/<file_id>')
def analyze_scaffolding(file_id):
    """Analyze a presentation for learning scaffolding elements"""
    try:
        # Get query parameters
        profile = request.args.get('profile', '')
        target_language = request.args.get('target_language', '')
        
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        # Debug output
        print(f"Looking for file at: {file_path}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Check file type - now supports both PDF and PowerPoint
        file_type = processing_tasks[file_id].get('file_type', '').lower()
        
        # Both PDF and PowerPoint files are now supported for analysis
        if file_type not in ['.pdf', '.pptx']:
            return jsonify({
                "status": "error",
                "message": "Scaffolding analysis is available for PowerPoint presentations (.pptx) and PDF files (.pdf). Please upload a supported file type."
            }), 400
        
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing'
        file_type_name = "PDF" if file_type == '.pdf' else "PowerPoint"
        processing_tasks[file_id]['message'] = f'Analyzing {file_type_name} for learning scaffolding elements...'
        
        # Store profile information for potential use
        if profile:
            processing_tasks[file_id]['profile'] = profile
        if target_language:
            processing_tasks[file_id]['target_language'] = target_language
        
        # Extract scaffolding elements with profile assessment
        scaffolding_data = educational_service.extract_learning_scaffolding(file_path, profile)
        if isinstance(scaffolding_data, dict):
            if 'error' in scaffolding_data:
                error_msg = scaffolding_data['error']
                
                # Check for authentication errors specifically
                if 'authentication_error' in str(error_msg) or 'invalid x-api-key' in str(error_msg):
                    # Store a user-friendly error message
                    processing_tasks[file_id]['status'] = 'error'
                    processing_tasks[file_id]['message'] = 'API authentication failed. Please check the API key configuration.'
                    return jsonify({
                        "status": "error",
                        "message": "API authentication failed. Please check the API key configuration."
                    }), 500
                
                # For other errors, create a meaningful scaffolding response
                scaffolding_data = {
                    'slides': [],
                    'scaffolding_elements': {
                        'learning_objectives': [],
                        'key_concepts': [],
                        'examples': [],
                        'practice_activities': [],
                        'assessment_items': [],
                        'review_elements': []
                    },
                    'scaffolding_score': 0,
                    'scaffolding_analysis': f'Analysis failed due to: {error_msg}'
                }
                print(f"DEBUG: Created fallback scaffolding data due to error")
            else:
                elements = scaffolding_data.get('scaffolding_elements', {})
                print(f"DEBUG: Found {len(elements.get('assessment_items', []))} assessment items")
        
        # Store the scaffolding data in the processing task
        processing_tasks[file_id]['scaffolding'] = scaffolding_data
        processing_tasks[file_id]['status'] = 'complete'
        processing_tasks[file_id]['message'] = 'Scaffolding analysis complete'
        
        print(f"DEBUG: Redirecting to view_presentation_scaffolding for {file_id}")
        # Redirect to the scaffolding results page
        return redirect(url_for('view_presentation_scaffolding', file_id=file_id))        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        if file_id in processing_tasks:
            processing_tasks[file_id]['status'] = 'error'
            processing_tasks[file_id]['message'] = f"Error analyzing scaffolding: {str(e)}"
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing scaffolding: {str(e)}"
        }), 500

@app.route('/view_scaffolding/<file_id>')
def view_presentation_scaffolding(file_id):
    """Show the scaffolding analysis results"""
    print(f"DEBUG: view_scaffolding called for file_id: {file_id}")
    print(f"DEBUG: processing_tasks keys: {list(processing_tasks.keys())}")
    
    if file_id not in processing_tasks:
        print(f"DEBUG: file_id {file_id} not found in processing_tasks")
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Presentation not found. Please upload again.")
                                     
    print(f"DEBUG: processing_tasks[{file_id}] keys: {list(processing_tasks[file_id].keys())}")
    
    if 'scaffolding' not in processing_tasks[file_id]:
        print(f"DEBUG: 'scaffolding' key not found in processing_tasks[{file_id}]")
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Scaffolding analysis not found. Please analyze the presentation first.")
    
    scaffolding_data = processing_tasks[file_id]['scaffolding']
    print(f"DEBUG: scaffolding_data type: {type(scaffolding_data)}")
    
    if isinstance(scaffolding_data, dict):
        elements = scaffolding_data.get('scaffolding_elements', {})
        print(f"DEBUG: About to render template with {len(elements.get('assessment_items', []))} assessment items")
    
    # Extract scaffolding elements for template compatibility
    elements = scaffolding_data.get('scaffolding_elements', {}) if isinstance(scaffolding_data, dict) else {}
    debug_info = scaffolding_data.get('debug_info', {}) if isinstance(scaffolding_data, dict) else {}
    profile_assessment = scaffolding_data.get('profile_assessment', {}) if isinstance(scaffolding_data, dict) else {}
    
    # Get file type for proper adaptation routing
    file_type = processing_tasks[file_id].get('file_type', '').lower()
    
    return render_template_string(
        SCAFFOLDING_TEMPLATE,
        scaffolding=elements,  # Pass the elements directly, not the whole structure
        slides=scaffolding_data.get('slides', []) if isinstance(scaffolding_data, dict) else [],
        debug_info=debug_info,  # Pass debug information
        profile_assessment=profile_assessment,  # Pass profile assessment
        file_id=file_id,  # Pass file_id for adaptation links
        file_type=file_type  # Pass file type for conditional logic
    )

@app.route('/assess_content/<file_id>')
def assess_content(file_id):
    """Assess content quality and provide recommendations"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "File not found. Please upload again."
            }), 404
        
        # Get file information
        filename = processing_tasks[file_id].get('filename', '')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "File not found on server"
            }), 404
        
        # Get profile from query parameters
        profile = request.args.get('profile', 'default')
        
        # Extract content based on file type
        file_type = processing_tasks[file_id].get('file_type', '').lower()
        
        if file_type == '.pdf':
            content = pdf_service.extract_content_from_pdf(file_path)
        elif file_type == '.pptx':
            content = pptx_service.extract_content_from_pptx(file_path)
        else:
            return jsonify({
                "status": "error", 
                "message": "Unsupported file type for assessment"
            }), 400
        
        # Perform assessment
        assessment_result = assessments_service.assess_content(content, profile)
        
        # Store assessment results
        processing_tasks[file_id]['assessment'] = assessment_result
        
        return jsonify({
            "status": "success",
            "assessment": assessment_result
        })
        
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error assessing content: {str(e)}"
        }), 500

@app.route('/content_recommendations/<file_id>')
def content_recommendations(file_id):
    """Get content improvement recommendations"""
    try:
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "File not found"
            }), 404
        
        # Check if assessment exists
        if 'assessment' not in processing_tasks[file_id]:
            return jsonify({
                "status": "error",
                "message": "No assessment found. Please assess content first."
            }), 400
        
        assessment = processing_tasks[file_id]['assessment']
        recommendations = assessment.get('recommendations', [])
        
        return jsonify({
            "status": "success",
            "recommendations": recommendations,
            "readability_metrics": assessment.get('readability_metrics', {}),
            "profile_suitability": assessment.get('profile_suitability', {})
        })
        
    except Exception as e:
        return jsonify({
            "status": "error", 
            "message": f"Error getting recommendations: {str(e)}"
        }), 500

@app.route('/assess_readability', methods=['POST'])
def assess_readability():
    """Assess readability of provided text content"""
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({
                "status": "error",
                "message": "No text content provided"
            }), 400
        
        text = data['text']
        profile = data.get('profile', 'default')
        
        # Create content structure for assessment
        content = {
            'pages': [{'text': text, 'page_number': 1}],
            'metadata': {'title': 'Text Assessment'}
        }
        
        # Perform assessment
        assessment_result = assessments_service.assess_content(content, profile)
        
        return jsonify({
            "status": "success",
            "assessment": assessment_result
        })
        
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error assessing readability: {str(e)}"
        }), 500

@app.route('/advanced_pdf_processing/<file_id>')
def advanced_pdf_processing(file_id):
    """Process PDF with advanced visual preservation features"""
    try:
        # Validate file exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "File not found"
            }), 404
        
        # Get parameters
        profile = request.args.get('profile', 'default')
        use_gradients = request.args.get('gradients', 'false').lower() == 'true'
        optimize_layout = request.args.get('optimize', 'false').lower() == 'true'
        add_reading_guides = request.args.get('guides', 'false').lower() == 'true'
        
        # Get file information
        filename = processing_tasks[file_id].get('filename', '')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "File not found on server"
            }), 404
        
        # Extract content
        content = pdf_service.extract_content_from_pdf(file_path)
        
        # Set up advanced options
        options = {
            'use_gradients': use_gradients,
            'optimize_layout': optimize_layout,
            'add_reading_guides': add_reading_guides,
            'preserve_images': True
        }
        
        # Generate output path
        output_filename = f"advanced_{os.path.splitext(filename)[0]}.pdf"
        output_path = get_output_file_path(file_id, output_filename)
        
        # Process with advanced features
        success = pdf_service.create_advanced_pdf(
            file_path, content, output_path, profile, options
        )
        
        if success:
            # Calculate quality metrics
            quality_metrics = pdf_service.calculate_quality_metrics(file_path, output_path)
            
            return jsonify({
                "status": "success",
                "output_path": output_path,
                "filename": output_filename,
                "quality_metrics": quality_metrics,
                "options_used": options
            })
        else:
            return jsonify({
                "status": "error",
                "message": "Advanced processing failed"
            }), 500
            
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error in advanced processing: {str(e)}"
        }), 500

@app.route('/batch_pdf_processing', methods=['POST'])
def batch_pdf_processing():
    """Process multiple PDFs in batch with progress tracking"""
    try:
        data = request.get_json()
        if not data or 'files' not in data:
            return jsonify({
                "status": "error",
                "message": "No files provided for batch processing"
            }), 400
        
        files = data['files']
        profile = data.get('profile', 'default')
        options = data.get('options', {})
        
        # Generate unique batch ID
        batch_id = str(uuid.uuid4())
        
        # Process files
        results = pdf_service.batch_process_pdfs(files, profile, None, options)
        
        return jsonify({
            "status": "success",
            "batch_id": batch_id,
            "results": results
        })
        
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error in batch processing: {str(e)}"
        }), 500

@app.route('/test_adaptation/<profile>')
def test_adaptation(profile):
    """Test if adaptation service is working properly"""
    try:
        # Test the adaptation service
        test_result = adaptations_service.test_adaptation(profile)
        
        return jsonify({
            "status": "success" if test_result['success'] else "error",
            "test_result": test_result,
            "cache_stats": adaptations_service.get_cache_stats()
        })
        
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error testing adaptation: {str(e)}"
        }), 500

@app.route('/pdf_accessibility/<file_id>')
def pdf_accessibility(file_id):
    """Optimize PDF for screen reader accessibility"""
    try:
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "File not found"
            }), 404
        
        # Get file information
        filename = processing_tasks[file_id].get('filename', '')
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "File not found on server"
            }), 404
        
        # Generate accessible output path
        output_filename = f"accessible_{os.path.splitext(filename)[0]}.pdf"
        output_path = get_output_file_path(file_id, output_filename)
        
        # Optimize for accessibility
        success = pdf_service.optimize_for_accessibility(file_path, output_path)
        
        if success:
            return jsonify({
                "status": "success",
                "output_path": output_path,
                "filename": output_filename,
                "message": "PDF optimized for screen reader accessibility"
            })
        else:
            return jsonify({
                "status": "error",
                "message": "Accessibility optimization failed"
            }), 500
            
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": f"Error optimizing for accessibility: {str(e)}"
        }), 500

def update_task_status(file_id, status, message=None, progress=None, **kwargs):
    """Helper function to update task status consistently"""
    if file_id not in processing_tasks:
        processing_tasks[file_id] = {}
    
    processing_tasks[file_id]['status'] = status
    
    if message is not None:
        processing_tasks[file_id]['message'] = message
    
    if progress is not None:
        processing_tasks[file_id]['progress'] = progress
    
    # Add any additional fields
    for key, value in kwargs.items():
        processing_tasks[file_id][key] = value
    
    # Debug logging
    print(f"Status update for {file_id}: {status} - {message}")

def update_task_progress(file_id, processed, total, message=None):
    """Helper function to update task progress"""
    percentage = min(100, int((processed / max(1, total)) * 100))
    
    progress = {
        'total': total,
        'processed': processed,
        'percentage': percentage
    }
    
    if message:
        update_task_status(file_id, 'processing', message, progress)
    else:
        if file_id in processing_tasks:
            processing_tasks[file_id]['progress'] = progress



## Text Processing Functions
# Batch processing has been moved to AdaptationsService


@app.route('/status/<file_id>')
def status(file_id):
    """Check the status of a processing task"""
    try:
        if file_id not in processing_tasks:
            return jsonify({
                'status': 'not_found',
                'message': 'Task not found'
            }), 404
        
        task = processing_tasks[file_id]
        
        # Return the task status as JSON
        response_data = {
            'status': task.get('status', 'unknown'),
            'message': task.get('message', ''),
            'progress': task.get('progress', {
                'total': 0,
                'processed': 0,
                'percentage': 0
            })
        }
        
        # Add additional fields if they exist
        if 'error' in task:
            response_data['error'] = task['error']
        
        if 'has_translation' in task:
            response_data['has_translation'] = task['has_translation']
            
        if 'export_format' in task:
            response_data['export_format'] = task['export_format']
        
        # Debug logging
        print(f"Status check for {file_id}: {response_data}")
        
        return jsonify(response_data)
        
    except Exception as e:
        print(f"Error in status route: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': f'Error checking status: {str(e)}'
        }), 500


# PDF capabilities are now handled by pdf_service


def calculate_simple_readability(text):
    """Calculate simplified readability metrics to determine if adaptation is needed"""
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = max(1, len(sentences))  # Avoid division by zero
    
    # Calculate average sentence length
    avg_sentence_length = word_count / sentence_count
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter
        vowels = "aeiouy"
        word = word.lower()
        count = 0
        if word and word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word and word.endswith("e"):
            count -= 1
        return max(1, count)  # Always at least 1 syllable
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = (len(complex_words) / max(1, word_count)) * 100
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / max(1, word_count))
        flesch = max(0, min(100, flesch))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / max(1, word_count)) - 15.59
        fk_grade = max(0, fk_grade)
    else:
        fk_grade = 0
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

def collect_text_elements_from_slide(slide, slide_index=0):
    """Collect text elements from a slide and their references"""
    elements = []
    
    # Process shapes with text
    for shape_index, shape in enumerate(slide.shapes):
        if hasattr(shape, "text_frame") and shape.text:
            elements.append({
                'text': shape.text,
                'reference': {
                    'type': 'shape', 
                    'object': shape,
                    'slide_num': slide_index,
                    'shape_index': shape_index
                }
            })
        
        # Tables
        if hasattr(shape, "table"):
            for row_index, row in enumerate(shape.table.rows):
                for cell_index, cell in enumerate(row.cells):
                    if hasattr(cell, "text_frame") and cell.text:
                        elements.append({
                            'text': cell.text,
                            'reference': {
                                'type': 'cell', 
                                'object': cell,
                                'slide_num': slide_index,
                                'shape_index': shape_index,
                                'row_index': row_index,
                                'cell_index': cell_index
                            }
                        })
        
        # Group shapes
        if hasattr(shape, "shapes"):
            for subshape_index, subshape in enumerate(shape.shapes):
                if hasattr(subshape, "text_frame") and subshape.text:
                    elements.append({
                        'text': subshape.text,
                        'reference': {
                            'type': 'subshape', 
                            'object': subshape,
                            'slide_num': slide_index,
                            'shape_index': shape_index,
                            'subshape_index': subshape_index
                        }
                    })
    
    return elements

def apply_adapted_text(reference, adapted_text, profile):
    """Apply adapted text to the appropriate element with formatting"""
    element_type = reference['type']
    obj = reference['object']
    
    if element_type == 'shape' or element_type == 'subshape' or element_type == 'cell':
        if hasattr(obj, "text_frame"):
            pptx_service.apply_text_to_text_frame(obj.text_frame, adapted_text, profile, profiles_service)
            
            # For dyslexia, apply additional formatting
            if profile == "dyslexia" and hasattr(obj, "text_frame"):
                pptx_service.apply_dyslexia_formatting(obj.text_frame)

# Function to apply text to a text frame with color formatting


@app.route('/error')
def error():
    """Show error page"""
    message = request.args.get('message', 'An unknown error occurred')
    return render_template_string(ERROR_TEMPLATE, message=message)

def calculate_readability_metrics(text):
    """Calculate readability metrics for the text"""
    # Simplified calculation of readability metrics
    # In a production environment, use a proper textstat library
    
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = len(sentences)
    
    # Calculate average sentence length
    avg_sentence_length = round(word_count / max(1, sentence_count), 1)
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter - not perfect but adequate for demo
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = round((len(complex_words) / max(1, word_count)) * 100, 1)
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / word_count)
        flesch = max(0, min(100, round(flesch)))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / word_count) - 15.59
        fk_grade = max(0, round(fk_grade, 1))
    else:
        fk_grade = 0
    
    # Calculate SMOG Index (simplified)
    if sentence_count >= 30:
        smog = 1.043 * ((len(complex_words) * (30 / sentence_count)) ** 0.5) + 3.1291
    else:
        smog = 1.043 * ((len(complex_words) * (30 / max(1, sentence_count))) ** 0.5) + 3.1291
    smog = max(0, round(smog, 1))
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "smog_index": smog,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

def find_complex_words(text):
    """Find and count complex words (3+ syllables)"""
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    complex_words = {}
    
    # Count syllables function (simplified)
    def count_syllables(word):
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    for word in words:
        syllables = count_syllables(word)
        if syllables >= 3:
            if word in complex_words:
                complex_words[word] += 1
            else:
                complex_words[word] = 1
    
    # Sort by frequency (most frequent first)
    sorted_words = sorted(complex_words.items(), key=lambda x: x[1], reverse=True)
    
    # Return the top 10 most frequent complex words
    return dict(sorted_words[:10])


def generate_complexity_chart(slide_texts, profile):
    """Generate a chart showing content complexity across slides"""
    # Check if matplotlib is available
    if not MATPLOTLIB_AVAILABLE:
        return generate_placeholder_chart("Chart generation unavailable - matplotlib not installed")
    
    # Create data for visualization
    slide_numbers = []
    complexity_scores = []
    threshold_values = []
    
    # Get threshold for the selected profile
    threshold = get_readability_thresholds(profile).get('flesch_kincaid_grade', 8)
    
    # Calculate complexity for each slide
    for slide in slide_texts:
        # Skip slides with very little text
        if len(slide['text'].strip()) < 20:
            continue
            
        metrics = calculate_simple_readability(slide['text'])
        slide_numbers.append(slide['slide_number'])
        complexity_scores.append(metrics['flesch_kincaid_grade'])
        threshold_values.append(threshold)
    
    # If no valid slides, return placeholder
    if not slide_numbers:
        return generate_placeholder_chart("No substantial text content found in slides")
    
    try:
        # Create the chart
        plt.figure(figsize=(8, 4))
        plt.bar(slide_numbers, complexity_scores, color='#3498db', alpha=0.7)
        plt.axhline(y=threshold, color='#e74c3c', linestyle='-', label=f'Target ({threshold})')
        
        plt.xlabel('Slide Number')
        plt.ylabel('Reading Grade Level')
        plt.title('Content Complexity by Slide')
        plt.legend()
        plt.grid(True, alpha=0.3)
        
        # Add value labels on bars
        for i, v in enumerate(complexity_scores):
            plt.text(slide_numbers[i], v + 0.3, str(round(v, 1)), ha='center')
        
        # Ensure y-axis starts at 0
        plt.ylim(bottom=0)
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
    
    except Exception as e:
        print(f"Error generating complexity chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_complex_words_chart(complex_words):
    """Generate a bar chart showing most frequent complex words"""
    # Check if matplotlib is available
    if not MATPLOTLIB_AVAILABLE:
        return generate_placeholder_chart("Chart generation unavailable - matplotlib not installed")
    
    try:
        # If no complex words found
        if not complex_words:
            return generate_placeholder_chart("No complex words found")
            
        # Sort and limit to top 10 words
        words = list(complex_words.keys())[:10]
        frequencies = list(complex_words.values())[:10]
        
        # Create the chart
        plt.figure(figsize=(8, 4))
        bars = plt.barh(words, frequencies, color='#2ecc71', alpha=0.7)
        
        plt.xlabel('Frequency')
        plt.ylabel('Words')
        plt.title('Most Frequent Complex Words')
        plt.grid(True, axis='x', alpha=0.3)
        
        # Add value labels
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 0.3, bar.get_y() + bar.get_height()/2, str(int(width)), va='center')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating complex words chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_placeholder_chart(message="No data available"):
    """Generate a placeholder chart with a message"""
    # If matplotlib is not available, return a simple fallback
    if not MATPLOTLIB_AVAILABLE:
        # Return a minimal base64 encoded 1x1 transparent PNG
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
    
    try:
        plt.figure(figsize=(8, 4))
        
        # Create an empty chart with a message
        plt.text(0.5, 0.5, message, ha='center', va='center', fontsize=14)
        plt.axis('off')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating placeholder chart: {e}")
        # Last resort - return empty image data
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="


def adapt_text_with_matcha(text, profile):
    """Adapt text using the adaptations service"""
    return adaptations_service.adapt_text(text, profile)

def translate_text(text, target_language):
    """Translate text using the translations service"""
    print(f"DEBUG: translate_text called with target_language: '{target_language}', text length: {len(text) if text else 0}")
    
    try:
        translated_text = translations_service.translate_text(text, target_language)
        print(f"Successfully translated text of length {len(text) if text else 0} to {target_language}")
        return translated_text
    except Exception as e:
        print(f"Error translating text: {str(e)}")
        return text

# Error handling and user feedback routes
@app.route('/adapt/<file_id>/<profile>')
def adapt(file_id, profile):
    print(f"DEBUG: Adaptation requested with profile: {profile}")
    """Adapt the PowerPoint based on assessment with API check"""
    try:
        # Check API connection first
        api_status = api_utils.check_connection()
        if api_status["status"] != "connected":
            # API connection failed - show error
            return render_template_string(ERROR_TEMPLATE, 
                                         message=f"Claude API connection error: {api_status['message']}"), 400
        
        # Validate if file_id exists in our system
        if file_id not in processing_tasks:
            return render_template_string(ERROR_TEMPLATE, 
                                         message="Presentation not found. Please upload again."), 404
        
        # Set initial processing status with empty progress
        processing_tasks[file_id] = {
            'status': 'processing', 
            'message': 'Starting adaptation...',
            'progress': {
                'total': 0,
                'processed': 0,
                'percentage': 0
            }
        }
        
        # Get the original filename from the processing tasks
        filename = processing_tasks[file_id].get('filename', 'presentation.pptx')
        
        # IMPORTANT FIX: Use the correct file path pattern
        # Look for files in the format "file_id_filename.pptx"
        file_path = None
        for file in os.listdir(app.config['UPLOAD_FOLDER']):
            if file.startswith(f"{file_id}_"):
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], file)
                break
        
        # If file not found, check alternate patterns
        if not file_path:
            # Try with exact file_id
            potential_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}.pptx")
            if os.path.exists(potential_path):
                file_path = potential_path
        
        # Double check the file exists
        if not file_path or not os.path.exists(file_path):
            # Debug info
            print(f"File not found. Looking for file with ID: {file_id}")
            print(f"Files in upload folder: {os.listdir(app.config['UPLOAD_FOLDER'])}")
            
            processing_tasks[file_id] = {'status': 'error', 'message': 'Presentation file not found on server'}
            return redirect(url_for('error', message="File not found on server"))
        
        # Validate profile type
        valid_profiles = ['dyslexia', 'esl', 'adhd', 'visual', 'cognitive']
        if profile not in valid_profiles:
            processing_tasks[file_id] = {'status': 'error', 'message': f'Invalid profile: {profile}'}
            return redirect(url_for('error', message=f"Invalid profile: {profile}"))
        
        # Start processing in a background thread to avoid blocking
        def process_task():
            try:
                # Get target language from processing tasks metadata
                target_language = processing_tasks[file_id].get('metadata', {}).get('target_language')
                
                # Check file type and use appropriate service
                file_type = processing_tasks[file_id].get('file_type', '').lower()
                print(f"DEBUG: Adaptation - file_type detected: '{file_type}' for file_id: {file_id}")
                print(f"DEBUG: Adaptation - processing_tasks[{file_id}]: {processing_tasks[file_id]}")
                
                if file_type == '.pdf':
                    # Use PDF service for adaptation
                    result = pdf_service.process_with_template_system(
                        file_path, file_id, filename, profile, 'pdf', target_language,
                        lambda fid, msg, prog: update_task_progress(fid, prog, 100, msg),
                        lambda fid, name: os.path.join(app.config['OUTPUT_FOLDER'], f"{fid}_{name}")
                    )
                else:
                    # Use PowerPoint service for .pptx files
                    result = pptx_service.process_presentation_efficiently(
                        file_path, file_id, filename, profile, target_language, processing_tasks
                    )
                
                if result:
                    processing_tasks[file_id]['adapted_path'] = result
            except Exception as e:
                import traceback
                traceback.print_exc()
                processing_tasks[file_id] = {'status': 'error', 'message': str(e)}
        
        # Start the processing thread with a name containing the file_id
        processing_thread = threading.Thread(
            target=process_task,
            name=f"process-{file_id}"  # This helps identify which file_id is being processed
        )
        processing_thread.daemon = True
        processing_thread.start()
        
        # Map profile to display name for the processing template
        profile_names = get_profile_names()
        
        # Determine which template to use (with or without progress)
        template = PROCESSING_TEMPLATE_WITH_PROGRESS if 'PROCESSING_TEMPLATE_WITH_PROGRESS' in globals() else PROCESSING_TEMPLATE_WITH_PROGRESS
        
        # Return the processing template with progress tracking
        return render_template_string(
            template,
            file_id=file_id,
            filename=f"adapted_{filename}",
            profile=profile,
            profile_name=profile_names.get(profile, profile)
        )
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        error_message = f"Error starting adaptation: {str(e)}"
        return redirect(url_for('error', message=error_message))

# Route to handle the status check for long-running tasks    
@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Handle generation of a new presentation"""
    try:
        # Get parameters from form
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Generate a unique ID for this presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        filename = f"{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_new_presentation, 
            args=(file_id, filename, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = get_profile_names()
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                                   file_id=file_id, 
                                   filename=filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

@app.route('/enrich_and_generate', methods=['POST'])
def enrich_and_generate():
    """Handle generating an enhanced presentation based on an existing adapted file"""
    try:
        # Get parameters from form
        original_file_id = request.form.get('original_file_id')
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Find the original adapted file
        adapted_file_path = None
        
        for filename in os.listdir(app.config['OUTPUT_FOLDER']):
            if filename.startswith(f"{original_file_id}_"):
                adapted_file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
                break
        
        if not adapted_file_path:
            return render_template_string(ERROR_TEMPLATE, message="Original adapted file not found"), 404
        
        # Generate a unique ID for this new presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename for the new presentation
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        new_filename = f"enriched-{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_enriched_presentation, 
            args=(file_id, new_filename, adapted_file_path, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = get_profile_names()
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE_WITH_PROGRESS, 
                                   file_id=file_id, 
                                   filename=new_filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

def generate_new_presentation(file_id, filename, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate a new presentation based on the specified parameters"""
    try:
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Step 1: Plan the content structure using Claude
        lesson_plan = educational_service.generate_lesson_plan(
            topic, grade_level, slide_count,
            learning_objectives=[f"Understand {topic}", f"Apply {topic} concepts"],
            profile=profile
        )
        
        # Step 2: Create slides based on the plan
        for slide_data in lesson_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    pptx_service.apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile, profiles_service)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            pptx_service.apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content
                    adapted_content = adapt_text_with_matcha(content, profile)
                    text_frame = content_placeholder.text_frame
                    pptx_service.apply_text_to_text_frame(text_frame, adapted_content, profile, profiles_service)
        
        # Save the presentation
        output_path = get_output_file_path(file_id, filename)
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

def process_presentation_with_pdf(file_path, file_id, filename, profile, target_language=None):
    """Modified presentation processing function that also generates PDF"""
    # First do the normal PowerPoint adaptation
    output_path = pptx_service.process_presentation_efficiently(file_path, file_id, filename, profile, target_language, processing_tasks)
    
    # Also generate PDF version
    if output_path:
        pdf_filename = f"adapted_{os.path.splitext(filename)[0]}.pdf"
        pdf_path = get_output_file_path(file_id, pdf_filename)
        
        try:
            # Convert the adapted PowerPoint to PDF
            pdf_success = conversion_service.convert_pptx_to_pdf(output_path, pdf_path)
            
            # Update processing task with PDF info
            if pdf_success and os.path.exists(pdf_path):
                processing_tasks[file_id]['has_pdf'] = True
                processing_tasks[file_id]['pdf_filename'] = pdf_filename
            else:
                processing_tasks[file_id]['has_pdf'] = False
                
        except Exception as e:
            print(f"Error generating PDF: {str(e)}")
            processing_tasks[file_id]['has_pdf'] = False
    
    return output_path

def generate_enriched_presentation(file_id, filename, adapted_file_path, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate an enriched presentation based on an adapted file"""
    try:
        # Extract content from the adapted file
        adapted_content = pptx_service.extract_content_from_pptx(adapted_file_path)
        
        # Generate an enriched presentation plan
        enriched_plan = educational_service.generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Create slides based on the plan
        for slide_data in enriched_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    pptx_service.apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile, profiles_service)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            pptx_service.apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content - already adapted content
                    text_frame = content_placeholder.text_frame
                    pptx_service.apply_text_to_text_frame(text_frame, content, profile, profiles_service)
        
        # Save the presentation
        output_path = get_output_file_path(file_id, filename)
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

# HTML templates module
# Create a separate file html_templates.py with all the HTML templates

if __name__ == "__main__":
    # Run the Flask app
    app.run(debug=True, host='0.0.0.0', port=5000)