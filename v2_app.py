import os
from dotenv import load_dotenv
from flask import Flask, request, send_file, render_template_string, redirect, url_for, jsonify
from werkzeug.utils import secure_filename
import uuid

# Load environment variables
load_dotenv()
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import anthropic
import re
import io
import base64
import matplotlib.pyplot as plt
import numpy as np
import threading
import hashlib
import time
from tenacity import retry, stop_after_attempt, wait_exponential
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import json


# Global dictionary to store analysis status
global_analysis_status = {}
processing_tasks = {}

# Define the upload folder path
UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')

# Create the directory if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


# Set up Flask
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'
app.config['REPORT_FOLDER'] = 'reports'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB max upload

# Create directories if they don't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
os.makedirs(app.config['REPORT_FOLDER'], exist_ok=True)

# Set up API client (not exposed in UI)
api_key = os.getenv('ANTHROPIC_API_KEY')
if not api_key:
    raise ValueError("ANTHROPIC_API_KEY not found in environment variables")

client = anthropic.Anthropic(api_key=api_key)

# Define colors for different learning profiles
PROFILE_COLORS = {
    "dyslexia": RGBColor(0, 102, 204),  # Blue
    "adhd": RGBColor(46, 139, 87),      # Sea Green
    "esl": RGBColor(148, 0, 211)        # Purple
}

# Readability thresholds for recommendations
READABILITY_THRESHOLDS = {
    "dyslexia": {
        "flesch_reading_ease": 70,  # Higher is easier
        "flesch_kincaid_grade": 6,  # Lower grade level is easier
        "smog_index": 8,            # Lower is easier
        "sentence_length": 12,      # Average words per sentence
        "complex_word_percent": 15  # Percentage of complex words
    },
    "adhd": {
        "flesch_reading_ease": 60,
        "flesch_kincaid_grade": 8,
        "smog_index": 10,
        "sentence_length": 15,
        "complex_word_percent": 20
    },
    "esl": {
        "flesch_reading_ease": 80,
        "flesch_kincaid_grade": 5,
        "smog_index": 7,
        "sentence_length": 10,
        "complex_word_percent": 10
    }
}

# 1. Caching Implementation
class AdaptationCache:
    """Cache for storing text adaptation results"""
    def __init__(self, max_size=1000):
        self.cache = {}
        self.max_size = max_size
        self.access_times = {}
    
    def get_key(self, text, profile):
        """Generate a cache key based on text content and profile"""
        # Use hash for efficiency, but truncate text to avoid excessive computation
        text_to_hash = text[:1000] if len(text) > 1000 else text
        text_hash = hashlib.md5(text_to_hash.encode()).hexdigest()
        return f"{text_hash}-{profile}"
    
    def get(self, text, profile):
        """Get adaptation from cache if available"""
        key = self.get_key(text, profile)
        if key in self.cache:
            # Update access time
            self.access_times[key] = time.time()
            return self.cache[key]
        return None
    
    def set(self, text, profile, adapted_text):
        """Store adaptation in cache"""
        key = self.get_key(text, profile)
        
        # Check if we need to evict entries because cache is full
        if len(self.cache) >= self.max_size:
            # Evict least recently used item
            oldest_key = min(self.access_times, key=self.access_times.get)
            del self.cache[oldest_key]
            del self.access_times[oldest_key]
        
        # Store new item
        self.cache[key] = adapted_text
        self.access_times[key] = time.time()

# Initialize the cache globally
adaptation_cache = AdaptationCache()

# 2. Token-Efficient Prompting
def create_efficient_prompt(text, profile):
    """Create a token-efficient prompt for the specified profile"""
    # First line contains essential instructions, rest is just the content
    if profile == "dyslexia":
        instructions = "Adapt for dyslexia: short sentences (max 15 words), simple words, active voice, bullet points where appropriate. Keep meaning intact."
    elif profile == "adhd":
        instructions = "Adapt for ADHD: clear structure, short chunks, bullet points, highlight key info, remove unnecessary details. Keep meaning intact."
    else:  # ESL
        instructions = "Adapt for English learners: simpler words (original in parentheses), short sentences, explain idioms, consistent terms. Keep meaning intact."
    
    # Minimal format with clear sections
    prompt = f"{instructions}\n\nOriginal:\n{text}\n\nAdapted:"
    return prompt

# Main page template with PPTX upload only
INDEX_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>Matcha PowerPoint Adaptor</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .container { max-width: 800px; margin: 0 auto; }
        form { margin: 20px 0; }
        .form-group { margin-bottom: 15px; }
        label { display: block; margin-bottom: 5px; }
        select, input[type="file"] { width: 100%; padding: 8px; }
        button { background: #4CAF50; color: white; border: none; padding: 10px 15px; cursor: pointer; margin-right: 10px; }
        .secondary-btn { background: #3498db; }
        .color-samples { margin-top: 20px; }
        .color-sample { display: inline-block; width: 20px; height: 20px; margin-right: 5px; border: 1px solid #ccc; }
        .dyslexia-color { background-color: #0066cc; }
        .adhd-color { background-color: #2e8b57; }
        .esl-color { background-color: #9400d3; }
        .profile-info { margin-top: 20px; padding: 15px; background-color: #f5f5f5; border-radius: 5px; }
        .profile-info h3 { margin-top: 0; }
        .toggle-section { cursor: pointer; color: #2c3e50; }
        .toggle-section:hover { text-decoration: underline; }
        .hidden { display: none; }
        .logo { font-size: 28px; color: #4CAF50; margin-bottom: 10px; font-weight: bold; }
        .info-box { background-color: #e8f4f8; padding: 15px; border-left: 5px solid #4CAF50; margin: 20px 0; }
        .buttons { display: flex; margin-top: 20px; }
    </style>
    <script>
        function toggleSection(id) {
            var section = document.getElementById(id);
            if (section.classList.contains('hidden')) {
                section.classList.remove('hidden');
            } else {
                section.classList.add('hidden');
            }
        }
        
        // Function to show/hide language dropdown
        function updateLanguageOptions() {
            var profile = document.getElementById('profile').value;
            var languageDiv = document.getElementById('language_div');
            if (profile === 'esl') {
                languageDiv.style.display = 'block';
            } else {
                languageDiv.style.display = 'none';
            }
        }
        
        // Initialize when page loads
        document.addEventListener('DOMContentLoaded', function() {
            updateLanguageOptions();
            document.getElementById('profile').addEventListener('change', updateLanguageOptions);
        });
    </script>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        <h1>PowerPoint Adaptor</h1>
        <p>Upload a PowerPoint file to adapt its content for different learning needs.</p>
        
        <div class="info-box">
            <strong>Two Ways to Use Matcha:</strong> 
            <ul>
                <li><strong>Assess First:</strong> Get a detailed analysis of your presentation's readability</li>
                <li><strong>Direct Adaptation:</strong> Skip assessment and directly adapt your presentation</li>
            </ul>
        </div>
        
        <form action="/upload" method="POST" enctype="multipart/form-data">
            <div class="form-group">
                <label for="pptx">Select PowerPoint file:</label>
                <input type="file" id="pptx" name="pptx" accept=".pptx" required>
            </div>
            
            <div class="form-group">
                <label for="profile">Select learning profile:</label>
                <select id="profile" name="profile" required>
                    <option value="dyslexia">Dyslexia Support</option>
                    <option value="adhd">ADHD Support</option>
                    <option value="esl">English Language Learner</option>
                </select>
            </div>
            
            <div id="language_div" class="form-group" style="display:none;">
                <label for="target_language">Target Language (for ESL):</label>
                <select id="target_language" name="target_language">
                    <option value="">English only (no translation)</option>
                    <option value="spanish">Spanish</option>
                    <option value="french">French</option>
                    <option value="german">German</option>
                    <option value="italian">Italian</option>
                    <option value="portuguese">Portuguese</option>
                    <option value="chinese">Chinese</option>
                    <option value="japanese">Japanese</option>
                    <option value="arabic">Arabic</option>
                    <option value="hindi">Hindi</option>
                    <option value="russian">Russian</option>
                </select>
            </div>
            
            <div class="buttons">
                <button type="submit" name="action" value="assess">Assess Content</button>
                <button type="submit" name="action" value="adapt" class="secondary-btn">Direct Adaptation</button>
            </div>
        </form>
        <div id="generation_options" class="hidden">
    <div class="form-group">
        <label for="lesson_topic">Lesson Topic:</label>
        <input type="text" id="lesson_topic" name="lesson_topic" placeholder="e.g., Photosynthesis, World War II">
    </div>

    <div class="form-group">
        <label for="subject_area">Subject Area:</label>
        <select id="subject_area" name="subject_area" required>
            <option value="english">English</option>
            <option value="maths">Mathematics</option>
            <option value="science">Science (General)</option>
            <option value="biology">Biology</option>
            <option value="chemistry">Chemistry</option>
            <option value="physics">Physics</option>
            <option value="history">History</option>
            <option value="geography">Geography</option>
            <option value="religious_studies">Religious Studies</option>
            <option value="modern_languages">Modern Foreign Languages</option>
            <option value="computing">Computing</option>
            <option value="design_technology">Design & Technology</option>
            <option value="art_design">Art & Design</option>
            <option value="music">Music</option>
            <option value="physical_education">Physical Education</option>
            <option value="citizenship">Citizenship</option>
            <option value="pshe">PSHE (Personal, Social, Health & Economic)</option>
            <option value="business">Business Studies</option>
            <option value="sociology">Sociology</option>
            <option value="psychology">Psychology</option>
            <option value="economics">Economics</option>
            <option value="media_studies">Media Studies</option>
        </select>
    </div>
    
    <div class="form-group">
        <label for="grade_level">Education Level:</label>
        <select id="grade_level" name="grade_level">
            <option value="ks3-year7-8">Key Stage 3 (Year 7-8)</option>
            <option value="ks3-year9">Key Stage 3 (Year 9)</option>
            <option value="ks4-gcse">Key Stage 4 (GCSE, Year 10-11)</option>
            <option value="ks5-alevel">Key Stage 5 (A-Level, Year 12-13)</option>
            <option value="btec">BTEC</option>
        </select>
    </div>
    
    <div class="form-group">
        <label for="slide_count">Number of Slides:</label>
        <input type="number" id="slide_count" name="slide_count" min="5" max="20" value="10">
    </div>
    
    <div class="form-group checkbox-container">
        <input type="checkbox" id="include_images" name="include_images" checked>
        <label for="include_images">Include Generated Images</label>
    </div>
    
    <div class="form-group">
        <label for="extra_notes">Additional Instructions (optional):</label>
        <textarea id="extra_notes" name="extra_notes" rows="3" placeholder="Any specific content or structure requests"></textarea>
    </div>
</div>
        <div class="color-samples">
            <p><strong>Color Key:</strong> First word in adapted text is highlighted</p>
            <div class="color-sample dyslexia-color"></div> Dyslexia Support
            <div class="color-sample adhd-color"></div> ADHD Support
            <div class="color-sample esl-color"></div> English Language Learner
        </div>
        
        <h3 class="toggle-section" onclick="toggleSection('features-info')">Learn About Adaptation Features ▼</h3>
        <div id="features-info" class="profile-info hidden">
            <h4>Dyslexia Support Features</h4>
            <ul>
                <li>Simplified vocabulary and shorter sentences</li>
                <li>Left-aligned text with increased spacing</li>
                <li>Sans-serif fonts for better readability</li>
                <li>Key information highlighted</li>
                <li>Research-based text adaptations</li>
            </ul>
            
            <h4>ADHD Support Features</h4>
            <ul>
                <li>Structured content with clear organization</li>
                <li>Bullet points for easier focus</li>
                <li>Broken-down paragraphs into manageable chunks</li>
                <li>Highlighted important information</li>
            </ul>
            
            <h4>ESL Support Features</h4>
            <ul>
                <li>Simplified vocabulary with original terms in parentheses</li>
                <li>Shorter sentences for easier comprehension</li>
                <li>Clear explanations of complex terms</li>
                <li>Consistent terminology throughout</li>
                <li>Optional translation to support bilingual learning</li>
            </ul>
        </div>
    </div>
</body>
</html>
"""

SCAFFOLDING_TEMPLATE = """<!DOCTYPE html>
<html>
<head>
    <title>Learning Scaffolding Analysis</title>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1">
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; line-height: 1.6; }
        .container { max-width: 900px; margin: 0 auto; }
        h1, h2, h3 { color: #333; }
        .card { border: 1px solid #ddd; border-radius: 4px; padding: 15px; margin-bottom: 20px; }
        .section { margin-bottom: 30px; }
        .slide-purpose { font-size: 14px; color: #666; }
        .key-concept { margin-bottom: 10px; }
        .term { font-weight: bold; }
        .example { background-color: #f9f9f9; padding: 10px; border-radius: 4px; margin-bottom: 10px; }
    </style>
</head>
<body>
    <div class="container">
        <h1>Learning Scaffolding Analysis</h1>
        
        <div class="section">
            <h2>Learning Objectives</h2>
            <ul>
                {% for obj in scaffolding.learning_objectives %}
                <li>{{ obj }}</li>
                {% else %}
                <li>No learning objectives detected</li>
                {% endfor %}
            </ul>
        </div>
        
        <div class="section">
            <h2>Key Concepts</h2>
            {% for concept in scaffolding.key_concepts %}
            <div class="key-concept">
                <span class="term">{{ concept.term }}:</span> {{ concept.definition }}
            </div>
            {% else %}
            <p>No key concepts detected</p>
            {% endfor %}
        </div>
        
        <div class="section">
            <h2>Examples</h2>
            {% for example in scaffolding.examples %}
            <div class="example">
                <strong>Slide {{ example.slide_number }}:</strong> {{ example.content }}
            </div>
            {% else %}
            <p>No examples detected</p>
            {% endfor %}
        </div>
        
        <div class="section">
            <h2>Practice Activities</h2>
            <ul>
                {% for activity in scaffolding.practice_activities %}
                <li>{{ activity.content }} (Slide {{ activity.slide_number }})</li>
                {% else %}
                <li>No practice activities detected</li>
                {% endfor %}
            </ul>
        </div>
        
        <div class="section">
            <h2>Assessment Items</h2>
            <ul>
                {% for item in scaffolding.assessment_items %}
                <li>{{ item.content }} (Slide {{ item.slide_number }})</li>
                {% else %}
                <li>No assessment items detected</li>
                {% endfor %}
            </ul>
        </div>
        
        <div class="section">
            <h2>Review Elements</h2>
            <ul>
                {% for item in scaffolding.review_elements %}
                <li>{{ item.content }} (Slide {{ item.slide_number }})</li>
                {% else %}
                <li>No review elements detected</li>
                {% endfor %}
            </ul>
        </div>
        
        <div class="section">
            <h2>Slide Purposes</h2>
            {% for slide in slides %}
            <div class="card">
                <h3>Slide {{ slide.slide_number }}</h3>
                <p class="slide-purpose">Purpose: {{ slide.purpose }}</p>
            </div>
            {% endfor %}
        </div>
    </div>
</body>
</html>"""

FRAMEWORK_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>Instructional Framework Analysis</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; line-height: 1.6; }
        .container { max-width: 900px; margin: 0 auto; }
        h1, h2, h3 { color: #333; }
        .framework-section { margin-bottom: 30px; background-color: #f8f9fa; padding: 20px; border-radius: 5px; }
        .score-container { display: flex; align-items: center; margin: 20px 0; }
        .score-circle {
            width: 120px;
            height: 120px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 36px;
            font-weight: bold;
            color: white;
            margin-right: 20px;
        }
        .high-score { background-color: #27ae60; }
        .medium-score { background-color: #f39c12; }
        .low-score { background-color: #e74c3c; }
        .analysis-list { margin: 15px 0; }
        .analysis-list li { margin-bottom: 8px; }
        .phase-visualization {
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin: 20px 0;
        }
        .phase-block {
            padding: 15px;
            border-radius: 5px;
            min-width: 150px;
            text-align: center;
        }
        .present { background-color: #d4edda; border: 1px solid #c3e6cb; }
        .missing { background-color: #f8d7da; border: 1px solid #f5c6cb; }
        .phase-name { font-weight: bold; margin-bottom: 5px; }
        .slide-container { margin-top: 40px; }
        .slide-row {
            display: flex;
            margin-bottom: 15px;
            padding: 15px;
            border: 1px solid #dee2e6;
            border-radius: 5px;
        }
        .slide-number {
            font-weight: bold;
            width: 80px;
        }
        .slide-phase {
            width: 150px;
        }
        .slide-details {
            flex: 1;
        }
        .strong { color: #27ae60; }
        .adequate { color: #f39c12; }
        .weak { color: #e74c3c; }
        .btn {
            display: inline-block;
            background: #4CAF50;
            color: white;
            text-decoration: none;
            padding: 10px 15px;
            margin: 10px 5px;
            border: none;
            cursor: pointer;
            border-radius: 4px;
        }
        .btn-blue { background-color: #3498db; }
    </style>
</head>
<body>
    <div class="container">
        <h1>Instructional Framework Analysis</h1>
        
        <div class="framework-section">
            <h2>{{ framework.framework.identified_framework }}</h2>
            
            <div class="score-container">
                <div class="score-circle {% if framework.framework.framework_alignment_score >= 80 %}high-score{% elif framework.framework.framework_alignment_score >= 60 %}medium-score{% else %}low-score{% endif %}">
                    {{ framework.framework.framework_alignment_score }}
                </div>
                <div>
                    <h3>Framework Alignment Score</h3>
                    <p>{{ framework.framework.balance_analysis }}</p>
                </div>
            </div>
            
            <h3>Strengths</h3>
            <ul class="analysis-list">
                {% for strength in framework.framework.strengths %}
                <li>{{ strength }}</li>
                {% endfor %}
            </ul>
            
            <h3>Areas for Improvement</h3>
            <ul class="analysis-list">
                {% for weakness in framework.framework.weaknesses %}
                <li>{{ weakness }}</li>
                {% endfor %}
            </ul>
            
            <h3>Recommendation</h3>
            <p>{{ framework.framework.recommendations }}</p>
            
            <h3>Instructional Phases</h3>
            <div class="phase-visualization">
                {% if framework.framework.identified_framework.startswith('Recall') %}
                    {% set phases = ['Recall', 'Teacher Demo', 'Guided Practice', 'Independent Work', 'Review'] %}
                {% elif framework.framework.identified_framework.startswith('5E') %}
                    {% set phases = ['Engage', 'Explore', 'Explain', 'Elaborate', 'Evaluate'] %}
                {% else %}
                    {% set phases = [] %}
                {% endif %}
                
                {% for phase in phases %}
                    {% set phase_display = phase %}
                    {% set phase_match = phase %}
                    
                    {% if phase == 'Teacher Demo' %}
                        {% set phase_match = 'I do' %}
                    {% elif phase == 'Guided Practice' %}
                        {% set phase_match = 'We do' %}
                    {% elif phase == 'Independent Work' %}
                        {% set phase_match = 'You do' %}
                    {% endif %}
                    
                    {% set is_missing = phase_match in framework.framework.missing_phases %}
                    <div class="phase-block {{ 'missing' if is_missing else 'present' }}">
                        <div class="phase-name">{{ phase_display }}</div>
                        {% if is_missing %}
                            <div>Missing Phase</div>
                        {% else %}
                            <div>
                                {% set phase_slides = [] %}
                                {% for slide in framework.slides %}
                                    {% if slide.framework_phase == phase_match %}
                                        {% set __ = phase_slides.append(slide.slide_number) %}
                                    {% endif %}
                                {% endfor %}
                                
                                Slides: {{ phase_slides|join(', ') if phase_slides else 'None' }}
                            </div>
                        {% endif %}
                    </div>
                {% endfor %}
            </div>
        </div>
        
        <div class="slide-container">
            <h2>Slide Analysis</h2>
            
            {% for slide in framework.slides %}
            <div class="slide-row">
                <div class="slide-number">Slide {{ slide.slide_number }}</div>
                <div class="slide-phase">{{ slide.framework_phase }}</div>
                <div class="slide-details">
                    <div class="{{ slide.effectiveness }}">Effectiveness: {{ slide.effectiveness|capitalize }}</div>
                    {% if slide.elements %}
                    <ul>
                        {% for element in slide.elements %}
                        <li>{{ element }}</li>
                        {% endfor %}
                    </ul>
                    {% endif %}
                </div>
            </div>
            {% endfor %}
        </div>
        
        <div style="text-align:center; margin-top:30px;">
            <a href="/adapt/{{ file_id }}/dyslexia" class="btn">Adapt for Dyslexia</a>
            <a href="/adapt/{{ file_id }}/adhd" class="btn">Adapt for ADHD</a>
            <a href="/adapt/{{ file_id }}/esl" class="btn">Adapt for ESL</a>
            <a href="/" class="btn btn-blue">Back to Home</a>
        </div>
    </div>
</body>
</html>
"""

ASSESSMENT_TEMPLATE_SIMPLIFIED = """
<!DOCTYPE html>
<html>
<head>
    <title>Content Assessment</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .container { max-width: 800px; margin: 0 auto; }
        .assessment-header { text-align: center; margin-bottom: 30px; }
        .score-container { 
            display: flex; 
            justify-content: center; 
            align-items: center; 
            margin: 20px 0;
        }
        .score-circle {
            width: 120px;
            height: 120px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 36px;
            font-weight: bold;
            color: white;
            margin-right: 20px;
        }
        .high-need { background-color: #e74c3c; }      /* Red */
        .medium-need { background-color: #f39c12; }    /* Orange */
        .low-need { background-color: #27ae60; }       /* Green */
        .recommendation {
            flex: 1;
            padding: 15px;
            background-color: #f5f5f5;
            border-radius: 5px;
        }
        .metrics-table {
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }
        .metrics-table th, .metrics-table td {
            padding: 8px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }
        .metrics-table th {
            background-color: #f5f5f5;
        }
        .bad { color: #e74c3c; }
        .okay { color: #f39c12; }
        .good { color: #27ae60; }
        .btn {
            display: inline-block;
            background: #4CAF50;
            color: white;
            text-decoration: none;
            padding: 10px 15px;
            margin: 10px 5px;
            border: none;
            cursor: pointer;
            border-radius: 4px;
        }
        .btn-blue { background-color: #3498db; }
        .action-container { text-align: center; margin-top: 30px; }
        .chart-container { margin: 30px 0; text-align: center; }
        .color-sample { display: inline-block; width: 20px; height: 20px; margin-right: 5px; vertical-align: middle; border: 1px solid #ccc; }
        .profile-color {
            background-color: {% if profile == "dyslexia" %}#0066cc{% elif profile == "adhd" %}#2e8b57{% else %}#9400d3{% endif %};
        }
        .logo { font-size: 24px; color: #4CAF50; margin-bottom: 10px; font-weight: bold; text-align: center; }
    </style>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        
        <div class="assessment-header">
            <h1>Content Assessment Results</h1>
            <p>
                Profile: 
                <span>
                    <div class="color-sample profile-color"></div>
                    {% if profile == "dyslexia" %}Dyslexia Support
                    {% elif profile == "adhd" %}ADHD Support
                    {% else %}English Language Learner{% endif %}
                </span>
            </p>
        </div>
        
        <div class="score-container">
            <div class="score-circle {% if adaptation_score < 60 %}high-need{% elif adaptation_score < 80 %}medium-need{% else %}low-need{% endif %}">
                {{ adaptation_score }}
            </div>
            <div class="recommendation">
                <h3>Recommendation:</h3>
                <p>{{ recommendation }}</p>
            </div>
        </div>
        
        <h2>Readability Analysis</h2>
        <table class="metrics-table">
            <tr>
                <th>Metric</th>
                <th>Current Value</th>
                <th>Target for {{ profile_name }}</th>
                <th>Status</th>
            </tr>
            <tr>
                <td>Flesch Reading Ease</td>
                <td>{{ metrics.flesch_reading_ease }}</td>
                <td>{{ thresholds.flesch_reading_ease }}+</td>
                <td class="{% if metrics.flesch_reading_ease >= thresholds.flesch_reading_ease %}good{% elif metrics.flesch_reading_ease >= thresholds.flesch_reading_ease - 15 %}okay{% else %}bad{% endif %}">
                    {% if metrics.flesch_reading_ease >= thresholds.flesch_reading_ease %}Good{% elif metrics.flesch_reading_ease >= thresholds.flesch_reading_ease - 15 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                </td>
            </tr>
            <tr>
                <td>Grade Level</td>
                <td>{{ metrics.flesch_kincaid_grade }}</td>
                <td>{{ thresholds.flesch_kincaid_grade }} or lower</td>
                <td class="{% if metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade %}good{% elif metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade + 3 %}okay{% else %}bad{% endif %}">
                    {% if metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade %}Good{% elif metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade + 3 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                </td>
            </tr>
            <tr>
                <td>SMOG Index</td>
                <td>{{ metrics.smog_index }}</td>
                <td>{{ thresholds.smog_index }} or lower</td>
                <td class="{% if metrics.smog_index <= thresholds.smog_index %}good{% elif metrics.smog_index <= thresholds.smog_index + 2 %}okay{% else %}bad{% endif %}">
                    {% if metrics.smog_index <= thresholds.smog_index %}Good{% elif metrics.smog_index <= thresholds.smog_index + 2 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                </td>
            </tr>
            <tr>
                <td>Avg. Sentence Length</td>
                <td>{{ metrics.sentence_length }} words</td>
                <td>{{ thresholds.sentence_length }} words or fewer</td>
                <td class="{% if metrics.sentence_length <= thresholds.sentence_length %}good{% elif metrics.sentence_length <= thresholds.sentence_length + 5 %}okay{% else %}bad{% endif %}">
                    {% if metrics.sentence_length <= thresholds.sentence_length %}Good{% elif metrics.sentence_length <= thresholds.sentence_length + 5 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                </td>
            </tr>
            <tr>
                <td>Complex Word %</td>
                <td>{{ metrics.complex_word_percent }}%</td>
                <td>{{ thresholds.complex_word_percent }}% or lower</td>
                <td class="{% if metrics.complex_word_percent <= thresholds.complex_word_percent %}good{% elif metrics.complex_word_percent <= thresholds.complex_word_percent + 10 %}okay{% else %}bad{% endif %}">
                    {% if metrics.complex_word_percent <= thresholds.complex_word_percent %}Good{% elif metrics.complex_word_percent <= thresholds.complex_word_percent + 10 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                </td>
            </tr>
        </table>
        
        <div class="chart-container">
            <h3>Content Complexity Analysis</h3>
            <img src="{{ complexity_chart }}" alt="Content Complexity Chart" style="max-width:100%;">
        </div>
        
        <div class="chart-container">
            <h3>Most Frequent Complex Words</h3>
            <img src="{{ complex_words_chart }}" alt="Complex Words Chart" style="max-width:100%;">
        </div>
        
<div class="action-container">
    <a href="/adapt/{{ file_id }}/{{ profile }}" class="btn">Adapt PowerPoint Now</a>
    <a href="/analyze/framework/{{ file_id }}" class="btn btn-blue">Analyze Teaching Framework</a>
    <a href="/analyze/scaffolding/{{ file_id }}" class="btn btn-blue">Analyze Learning Scaffolding</a>
    <a href="/" class="btn btn-blue">Start Over</a>
</div>
    </div>
</body>
</html>
"""

ASSESSMENT_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>Content Assessment</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .container { max-width: 800px; margin: 0 auto; }
        .assessment-header { text-align: center; margin-bottom: 30px; }
        .score-container { 
            display: flex; 
            justify-content: center; 
            align-items: center; 
            margin: 20px 0;
        }
        .score-circle {
            width: 120px;
            height: 120px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 36px;
            font-weight: bold;
            color: white;
            margin-right: 20px;
        }
        .high-need { background-color: #e74c3c; }      /* Red */
        .medium-need { background-color: #f39c12; }    /* Orange */
        .low-need { background-color: #27ae60; }       /* Green */
        .recommendation {
            flex: 1;
            padding: 15px;
            background-color: #f5f5f5;
            border-radius: 5px;
        }
        .metrics-table {
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }
        .metrics-table th, .metrics-table td {
            padding: 8px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }
        .metrics-table th {
            background-color: #f5f5f5;
        }
        .bad { color: #e74c3c; }
        .okay { color: #f39c12; }
        .good { color: #27ae60; }
        .btn {
            display: inline-block;
            background: #4CAF50;
            color: white;
            text-decoration: none;
            padding: 10px 15px;
            margin: 10px 5px;
            border: none;
            cursor: pointer;
            border-radius: 4px;
        }
        .btn-blue { background-color: #3498db; }
        .action-container { text-align: center; margin-top: 30px; }
        .chart-container { margin: 30px 0; text-align: center; }
        .color-sample { display: inline-block; width: 20px; height: 20px; margin-right: 5px; vertical-align: middle; border: 1px solid #ccc; }
        .profile-color {
            background-color: {% if profile == "dyslexia" %}#0066cc{% elif profile == "adhd" %}#2e8b57{% else %}#9400d3{% endif %};
        }
        .logo { font-size: 24px; color: #4CAF50; margin-bottom: 10px; font-weight: bold; text-align: center; }
        
        /* New styles for scaffolding analysis */
        .scaffolding-container {
            margin: 30px 0;
            padding: 20px;
            background-color: #f8f9fa;
            border-radius: 5px;
            border-left: 5px solid #4CAF50;
        }
        .framework-analysis {
            margin-bottom: 30px;
        }
        .framework-score {
            display: flex;
            align-items: center;
            margin: 15px 0;
        }
        .score-label {
            margin-left: 15px;
            font-weight: bold;
        }
        .framework-visualization {
            display: flex;
            flex-wrap: wrap;
            gap: 10px;
            margin: 20px 0;
        }
        .phase-block {
            padding: 10px;
            border-radius: 5px;
            min-width: 120px;
            text-align: center;
        }
        .present {
            background-color: #d4edda;
            border: 1px solid #c3e6cb;
        }
        .missing {
            background-color: #f8d7da;
            border: 1px solid #f5c6cb;
        }
        .phase-name {
            font-weight: bold;
            margin-bottom: 5px;
        }
        .phase-slides {
            font-size: 0.9em;
            color: #666;
        }
        .phase-missing {
            font-size: 0.9em;
            color: #721c24;
        }
        .learning-elements {
            margin-top: 20px;
        }
        .objective-list {
            padding-left: 20px;
        }
        .objective-list li {
            margin-bottom: 5px;
        }
        .tab-container {
            margin: 30px 0;
        }
        .tab {
            overflow: hidden;
            border: 1px solid #ccc;
            background-color: #f1f1f1;
            border-radius: 5px 5px 0 0;
        }
        .tab button {
            background-color: inherit;
            float: left;
            border: none;
            outline: none;
            cursor: pointer;
            padding: 10px 16px;
            transition: 0.3s;
            font-size: 16px;
        }
        .tab button:hover {
            background-color: #ddd;
        }
        .tab button.active {
            background-color: #4CAF50;
            color: white;
        }
        .tabcontent {
            display: none;
            padding: 20px;
            border: 1px solid #ccc;
            border-top: none;
            border-radius: 0 0 5px 5px;
        }
        .tabcontent.show {
            display: block;
        }
    </style>
    <script>
        function openTab(evt, tabName) {
            var i, tabcontent, tablinks;
            tabcontent = document.getElementsByClassName("tabcontent");
            for (i = 0; i < tabcontent.length; i++) {
                tabcontent[i].style.display = "none";
            }
            tablinks = document.getElementsByClassName("tablinks");
            for (i = 0; i < tablinks.length; i++) {
                tablinks[i].className = tablinks[i].className.replace(" active", "");
            }
            document.getElementById(tabName).style.display = "block";
            evt.currentTarget.className += " active";
        }
        
        // Initialize first tab as active when page loads
        document.addEventListener('DOMContentLoaded', function() {
            document.getElementsByClassName('tablinks')[0].click();
        });
    </script>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        
        <div class="assessment-header">
            <h1>Content Assessment Results</h1>
            <p>
                Profile: 
                <span>
                    <div class="color-sample profile-color"></div>
                    {% if profile == "dyslexia" %}Dyslexia Support
                    {% elif profile == "adhd" %}ADHD Support
                    {% else %}English Language Learner{% endif %}
                </span>
            </p>
        </div>
        
        <div class="score-container">
            <div class="score-circle {% if adaptation_score < 60 %}high-need{% elif adaptation_score < 80 %}medium-need{% else %}low-need{% endif %}">
                {{ adaptation_score }}
            </div>
            <div class="recommendation">
                <h3>Recommendation:</h3>
                <p>{{ recommendation }}</p>
            </div>
        </div>
        
        <div class="tab-container">
            <div class="tab">
                <button class="tablinks" onclick="openTab(event, 'Readability')">Readability Analysis</button>
                <button class="tablinks" onclick="openTab(event, 'Structure')">Learning Structure</button>
            </div>
            
            <div id="Readability" class="tabcontent">
                <h2>Readability Analysis</h2>
                <table class="metrics-table">
                    <tr>
                        <th>Metric</th>
                        <th>Current Value</th>
                        <th>Target for {{ profile_name }}</th>
                        <th>Status</th>
                    </tr>
                    <tr>
                        <td>Flesch Reading Ease</td>
                        <td>{{ metrics.flesch_reading_ease }}</td>
                        <td>{{ thresholds.flesch_reading_ease }}+</td>
                        <td class="{% if metrics.flesch_reading_ease >= thresholds.flesch_reading_ease %}good{% elif metrics.flesch_reading_ease >= thresholds.flesch_reading_ease - 15 %}okay{% else %}bad{% endif %}">
                            {% if metrics.flesch_reading_ease >= thresholds.flesch_reading_ease %}Good{% elif metrics.flesch_reading_ease >= thresholds.flesch_reading_ease - 15 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                        </td>
                    </tr>
                    <tr>
                        <td>Grade Level</td>
                        <td>{{ metrics.flesch_kincaid_grade }}</td>
                        <td>{{ thresholds.flesch_kincaid_grade }} or lower</td>
                        <td class="{% if metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade %}good{% elif metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade + 3 %}okay{% else %}bad{% endif %}">
                            {% if metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade %}Good{% elif metrics.flesch_kincaid_grade <= thresholds.flesch_kincaid_grade + 3 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                        </td>
                    </tr>
                    <tr>
                        <td>SMOG Index</td>
                        <td>{{ metrics.smog_index }}</td>
                        <td>{{ thresholds.smog_index }} or lower</td>
                        <td class="{% if metrics.smog_index <= thresholds.smog_index %}good{% elif metrics.smog_index <= thresholds.smog_index + 2 %}okay{% else %}bad{% endif %}">
                            {% if metrics.smog_index <= thresholds.smog_index %}Good{% elif metrics.smog_index <= thresholds.smog_index + 2 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                        </td>
                    </tr>
                    <tr>
                        <td>Avg. Sentence Length</td>
                        <td>{{ metrics.sentence_length }} words</td>
                        <td>{{ thresholds.sentence_length }} words or fewer</td>
                        <td class="{% if metrics.sentence_length <= thresholds.sentence_length %}good{% elif metrics.sentence_length <= thresholds.sentence_length + 5 %}okay{% else %}bad{% endif %}">
                            {% if metrics.sentence_length <= thresholds.sentence_length %}Good{% elif metrics.sentence_length <= thresholds.sentence_length + 5 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                        </td>
                    </tr>
                    <tr>
                        <td>Complex Word %</td>
                        <td>{{ metrics.complex_word_percent }}%</td>
                        <td>{{ thresholds.complex_word_percent }}% or lower</td>
                        <td class="{% if metrics.complex_word_percent <= thresholds.complex_word_percent %}good{% elif metrics.complex_word_percent <= thresholds.complex_word_percent + 10 %}okay{% else %}bad{% endif %}">
                            {% if metrics.complex_word_percent <= thresholds.complex_word_percent %}Good{% elif metrics.complex_word_percent <= thresholds.complex_word_percent + 10 %}Needs Improvement{% else %}Needs Significant Improvement{% endif %}
                        </td>
                    </tr>
                </table>
                
                <div class="chart-container">
                    <h3>Content Complexity Analysis</h3>
                    <img src="{{ complexity_chart }}" alt="Content Complexity Chart" style="max-width:100%;">
                </div>
                
                <div class="chart-container">
                    <h3>Most Frequent Complex Words</h3>
                    <img src="{{ complex_words_chart }}" alt="Complex Words Chart" style="max-width:100%;">
                </div>
            </div>
            
            <div id="Structure" class="tabcontent">
                <div class="scaffolding-container">
                    <h2>Learning Structure Analysis</h2>
                    
                    <!-- Instructional Framework Section -->
                    <div class="framework-analysis">
                        <h3>Instructional Approach: {{ framework.identified_framework }}</h3>
                        
                        <div class="score-container">
                            <div class="score-circle {% if framework.framework_alignment_score >= 80 %}low-need{% elif framework.framework_alignment_score >= 60 %}medium-need{% else %}high-need{% endif %}">
                                {{ framework.framework_alignment_score }}
                            </div>
                            <div class="recommendation">
                                <h3>Framework Alignment:</h3>
                                <p>{{ framework.recommendations }}</p>
                            </div>
                        </div>
                        
                        <h4>Lesson Structure</h4>
                        <div class="framework-visualization">
                            {% set phases = ['Recall', 'Check', 'I do', 'We do', 'You do', 'Recap'] if framework.identified_framework == 'Recall, Check, I do, We do, You do, Recap' else ['Engage', 'Explore', 'Explain', 'Elaborate', 'Evaluate'] if framework.identified_framework == '5E Model' else [] %}
                            
                            {% for phase in phases %}
                                {% set has_phase = phase not in framework.missing_phases %}
                                <div class="phase-block {{ 'present' if has_phase else 'missing' }}">
                                    <div class="phase-name">{{ phase }}</div>
                                    
                                    {% if has_phase %}
                                        {% set phase_slides = [] %}
                                        {% for slide in framework.slides if slide.framework_phase == phase %}
                                            {% do phase_slides.append(slide.slide_number) %}
                                        {% endfor %}
                                        
                                        <div class="phase-slides">Slides: {{ phase_slides|join(', ') }}</div>
                                    {% else %}
                                        <div class="phase-missing">Missing</div>
                                    {% endif %}
                                </div>
                            {% endfor %}
                        </div>
                    </div>
                <script>
    // Only run this if analysis is in progress
    {% if analysis_in_progress %}
    // Check analysis status every 5 seconds
    function checkAnalysisStatus() {
        fetch('/analysis_status/{{ file_id }}')
            .then(response => response.json())
            .then(data => {
                if (data.status === 'complete') {
                    // Refresh the page to show complete analysis
                    window.location.reload();
                }
                else if (data.status === 'error') {
                    // Show error message
                    document.getElementById('analysis-error').style.display = 'block';
                    document.getElementById('analysis-error-message').textContent = 
                        "Error analyzing presentation structure: " + data.error;
                    clearInterval(statusInterval);
                }
            })
            .catch(error => console.error('Error checking status:', error));
    }
    
    // Set interval to check status
    var statusInterval = setInterval(checkAnalysisStatus, 5000);
    
    // Initial check after 2 seconds
    setTimeout(checkAnalysisStatus, 2000);
    {% endif %}
</script>

<!-- Add this in your Structure tab content -->
<div id="analysis-error" style="display: none; margin: 20px 0; padding: 15px; background-color: #f8d7da; border-left: 5px solid #e74c3c;">
    <p id="analysis-error-message">Error analyzing structure</p>
</div>

<!-- Show loading indicator if analysis is in progress -->
{% if analysis_in_progress %}
<div class="loading-container" style="text-align: center; margin: 30px 0;">
    <div class="loader" style="border: 8px solid #f3f3f3; border-top: 8px solid #4CAF50; border-radius: 50%; width: 60px; height: 60px; animation: spin 2s linear infinite; margin: 20px auto;"></div>
    <p>Analyzing presentation structure... This may take a minute.</p>
</div>
{% endif %}    
                    <!-- Learning Elements Section -->
                    <div class="learning-elements">
                        <h3>Learning Objectives</h3>
                        <ul class="objective-list">
                            {% for objective in scaffolding.scaffolding.learning_objectives %}
                            <li>{{ objective }}</li>
                            {% endfor %}
                        </ul>
                        
                        <h3>Key Concepts</h3>
                        <table class="metrics-table">
                            <tr>
                                <th>Term</th>
                                <th>Definition</th>
                            </tr>
                            {% for concept in scaffolding.scaffolding.key_concepts %}
                            <tr>
                                <td>{{ concept.term }}</td>
                                <td>{{ concept.definition }}</td>
                            </tr>
                            {% endfor %}
                        </table>
                        
                        {% if scaffolding.scaffolding.examples %}
                        <h3>Examples</h3>
                        <ul class="objective-list">
                            {% for example in scaffolding.scaffolding.examples %}
                            <li>Slide {{ example.slide_number }}: {{ example.content|truncate(100) }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                        
                        {% if scaffolding.scaffolding.assessment_items %}
                        <h3>Assessment Items</h3>
                        <ul class="objective-list">
                            {% for item in scaffolding.scaffolding.assessment_items %}
                            <li>{{ item.question }}</li>
                            {% endfor %}
                        </ul>
                        {% endif %}
                    </div>
                </div>
            </div>
        </div>
        
        <div class="action-container">
            <a href="/adapt/{{ file_id }}/{{ profile }}" class="btn">Adapt PowerPoint Now</a>
            <a href="/" class="btn btn-blue">Start Over</a>
        </div>
    </div>
</body>
</html>
"""

PROCESSING_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>Processing Your Presentation</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; text-align: center; }
        .container { max-width: 800px; margin: 0 auto; }
        .loader { 
            border: 16px solid #f3f3f3;
            border-top: 16px solid #4CAF50;
            border-radius: 50%;
            width: 120px;
            height: 120px;
            animation: spin 2s linear infinite;
            margin: 40px auto;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        .info { margin: 30px 0; line-height: 1.6; }
        .logo { font-size: 24px; color: #4CAF50; margin-bottom: 20px; font-weight: bold; }
    </style>
    <script>
        function checkStatus() {
    fetch('/status/{{ file_id }}')
        .then(response => {
            if (!response.ok) {
                throw new Error('Response not OK');
            }
            return response.json();
        })
        .then(data => {
            if (data.status === 'complete') {
                window.location.href = '/download/{{ file_id }}/{{ filename }}';
            }
            else if (data.status === 'error') {
                window.location.href = '/error?message=' + encodeURIComponent(data.message || "Unknown error occurred");
            }
            // Otherwise continue checking
        })
        .catch(error => {
            console.error('Error checking status:', error);
            // Don't redirect on network errors, just continue checking
        });
}
        // Check status every few seconds
        function checkStatus() {
            fetch('/status/{{ file_id }}')
                .then(response => response.json())
                .then(data => {
                    if (data.status === 'complete') {
                        window.location.href = '/download/{{ file_id }}/{{ filename }}';
                    }
                    else if (data.status === 'error') {
                        window.location.href = '/error?message=' + encodeURIComponent(data.message);
                    }
                })
                .catch(error => console.error('Error checking status:', error));
        }
        
        // Check every 5 seconds
        setInterval(checkStatus, 5000);
        
        // Initial check after 2 seconds
        setTimeout(checkStatus, 2000);
    </script>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        <h1>Processing Your Presentation</h1>
        
        <div class="loader"></div>
        
        <div class="info">
            <p>Matcha is adapting your PowerPoint for {{ profile_name }}...</p>
            <p>This process may take a few minutes depending on the size of your presentation.</p>
            <p>The page will automatically refresh when your presentation is ready.</p>
        </div>
    </div>
</body>
</html>
"""

DOWNLOAD_TEMPLATE_WITH_TRANSLATION = """
<!DOCTYPE html>
<html>
<head>
    <title>Download Adapted File</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .container { max-width: 800px; margin: 0 auto; text-align: center; }
        .btn { background: #4CAF50; color: white; text-decoration: none; padding: 10px 15px; display: inline-block; margin-top: 20px; border: none; cursor: pointer; border-radius: 4px; }
        .home-btn { background: #3498db; margin-left: 10px; }
        .generate-btn { background: #9b59b6; margin-left: 10px; }
        .translated-btn { background: #e74c3c; margin-top: 10px; }
        .color-sample { display: inline-block; width: 20px; height: 20px; margin-right: 5px; vertical-align: middle; border: 1px solid #ccc; }
        .profile-color {
            background-color: {% if profile == "dyslexia" %}#0066cc{% elif profile == "adhd" %}#2e8b57{% else %}#9400d3{% endif %};
        }
        .adaptation-info { text-align: left; margin: 20px auto; max-width: 600px; padding: 15px; background-color: #f5f5f5; border-radius: 5px; }
        .success-message { margin: 20px 0; padding: 15px; background-color: #d4edda; border-left: 5px solid #28a745; text-align: left; }
        .logo { font-size: 24px; color: #4CAF50; margin-bottom: 10px; font-weight: bold; }
        .translation-info { margin-top: 20px; padding: 15px; background-color: #fff3cd; border-left: 5px solid #ffc107; text-align: left; display: {% if has_translation %}block{% else %}none{% endif %}; }
        .generation-form { display: none; margin-top: 20px; text-align: left; padding: 20px; background-color: #f0f4f8; border-radius: 5px; }
        .form-group { margin-bottom: 15px; }
        .form-group label { display: block; margin-bottom: 5px; }
        select, input[type="text"], input[type="number"], textarea { width: 100%; padding: 8px; }
        .checkbox-container { display: flex; align-items: center; }
        .checkbox-container input[type="checkbox"] { width: auto; margin-right: 8px; }
    </style>
    <script>
        function toggleGenerationForm() {
            var form = document.getElementById('generation-form');
            if (form.style.display === 'block') {
                form.style.display = 'none';
            } else {
                form.style.display = 'block';
            }
        }
    </script>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        <h1>Adaptation Complete!</h1>
        
        <div class="success-message">
            <p>Your PowerPoint has been successfully adapted for
               <span>
                   <div class="color-sample profile-color"></div>
                   {% if profile == "dyslexia" %}Dyslexia Support
                   {% elif profile == "adhd" %}ADHD Support
                   {% else %}English Language Learners{% endif %}
               </span>
            </p>
        </div>
        
        <div class="adaptation-info">
            {% if profile == "dyslexia" %}
            <h3>Dyslexia Support Adaptations</h3>
            <p>Your presentation has been adapted following research-based guidelines for dyslexia support:</p>
            <ul>
                <li>Text has been simplified with shorter sentences and clearer vocabulary</li>
                <li>Formatting has been adjusted for better readability</li>
                <li>The first word of each adapted paragraph is highlighted in blue</li>
                <li>Content is left-aligned with increased spacing where possible</li>
            </ul>
            {% elif profile == "adhd" %}
            <h3>ADHD Support Adaptations</h3>
            <p>Your presentation has been adapted with ADHD-friendly features:</p>
            <ul>
                <li>Content is structured with bullet points and clear organization</li>
                <li>Long paragraphs have been broken into manageable chunks</li>
                <li>The first word of each adapted paragraph is highlighted in green</li>
                <li>Information is presented in a more focused way</li>
            </ul>
            {% else %}
            <h3>English Language Learner Adaptations</h3>
            <p>Your presentation has been adapted for English language learners:</p>
            <ul>
                <li>Complex vocabulary has been simplified with original terms in parentheses</li>
                <li>Long sentences have been divided for easier comprehension</li>
                <li>The first word of each adapted paragraph is highlighted in purple</li>
                <li>Idiomatic expressions are explained or replaced</li>
            </ul>
            {% endif %}
        </div>
        
        {% if has_translation %}
        <div class="translation-info">
            <h3>{{ translated_language }} Version Available</h3>
            <p>A translated version of your presentation is also available in {{ translated_language }}:</p>
            <ul>
                <li>Contains a direct translation of the original presentation content to {{ translated_language }}</li>
                <li>Maintains the same slide structure as the English version for easy reference</li>
                <li>Technical terms are preserved in their original form with translations</li>
                <li>Perfect for bilingual classrooms or students transitioning between languages</li>
            </ul>
            <p><strong>Pro tip:</strong> Use both presentations side by side for bilingual learning support.</p>
        </div>
        {% endif %}
        
        <a href="/download_file/{{ file_id }}/{{ filename }}" class="btn">Download Adapted PowerPoint (English)</a>
        {% if has_translation %}
        <a href="/download_file/{{ file_id }}/{{ translated_filename }}" class="btn translated-btn">Download {{ translated_language }} Version</a>
        {% endif %}
        <a href="/" class="btn home-btn">Back to Home</a>
        <button onclick="toggleGenerationForm()" class="btn generate-btn">Enrich & Generate New</button>
        
        <div id="generation-form" class="generation-form">
            <h3>Create an Enhanced Presentation</h3>
            <p>Generate a brand new presentation based on the content in your adapted file, with additional learning resources and engaging visuals.</p>
            
            <form action="/enrich_and_generate" method="POST">
                <input type="hidden" name="original_file_id" value="{{ file_id }}">
                <input type="hidden" name="profile" value="{{ profile }}">
                
                <div class="form-group">
                    <label for="lesson_topic">Presentation Topic:</label>
                    <input type="text" id="lesson_topic" name="lesson_topic" placeholder="e.g., Photosynthesis, World War II" required>
                </div>
                
                <div class="form-group">
                    <label for="subject_area">Subject Area:</label>
                    <select id="subject_area" name="subject_area" required>
                        <option value="english">English</option>
                        <option value="maths">Mathematics</option>
                        <option value="science">Science (General)</option>
                        <option value="biology">Biology</option>
                        <option value="chemistry">Chemistry</option>
                        <option value="physics">Physics</option>
                        <option value="history">History</option>
                        <option value="geography">Geography</option>
                        <option value="religious_studies">Religious Studies</option>
                        <option value="modern_languages">Modern Foreign Languages</option>
                        <option value="computing">Computing</option>
                        <option value="design_technology">Design & Technology</option>
                        <option value="art_design">Art & Design</option>
                        <option value="music">Music</option>
                        <option value="physical_education">Physical Education</option>
                        <option value="citizenship">Citizenship</option>
                        <option value="pshe">PSHE (Personal, Social, Health & Economic)</option>
                        <option value="business">Business Studies</option>
                        <option value="sociology">Sociology</option>
                        <option value="psychology">Psychology</option>
                        <option value="economics">Economics</option>
                        <option value="media_studies">Media Studies</option>
                    </select>
                </div>
                <div class="framework-analysis">
    <h3>Instructional Framework: {{ framework.identified_framework }}</h3>
    
    <div class="framework-score">
        <div class="score-circle {{ 'high' if framework.framework_alignment_score >= 80 else 'medium' if framework.framework_alignment_score >= 60 else 'low' }}">
            {{ framework.framework_alignment_score }}
        </div>
        <div class="score-label">Framework Alignment</div>
    </div>
    
    <div class="framework-visualization">
        {% set phases = ['Recall', 'Check', 'I do', 'We do', 'You do', 'Recap'] if framework.identified_framework == 'Recall, Check, I do, We do, You do, Recap' else ['Engage', 'Explore', 'Explain', 'Elaborate', 'Evaluate'] if framework.identified_framework == '5E Model' else [] %}
        
        {% for phase in phases %}
            {% set has_phase = phase not in framework.missing_phases %}
            <div class="phase-block {{ 'present' if has_phase else 'missing' }}">
                <div class="phase-name">{{ phase }}</div>
                
                {% if has_phase %}
                    {% set phase_slides = [] %}
                    {% for slide in framework.slides if slide.framework_phase == phase %}
                        {% do phase_slides.append(slide.slide_number) %}
                    {% endfor %}
                    
                    <div class="phase-slides">Slides: {{ phase_slides|join(', ') }}</div>
                {% else %}
                    <div class="phase-missing">Missing</div>
                {% endif %}
            </div>
        {% endfor %}
    </div>
    
    <div class="framework-recommendations">
        <h4>Recommendations</h4>
        <p>{{ framework.recommendations }}</p>
    </div>
</div>
                <div class="form-group">
                    <label for="grade_level">Education Level:</label>
                    <select id="grade_level" name="grade_level" required>
                        <option value="ks3-year7-8">Key Stage 3 (Year 7-8)</option>
                        <option value="ks3-year9">Key Stage 3 (Year 9)</option>
                        <option value="ks4-gcse">Key Stage 4 (GCSE, Year 10-11)</option>
                        <option value="ks5-alevel">Key Stage 5 (A-Level, Year 12-13)</option>
                        <option value="btec">BTEC</option>
                    </select>
                </div>
                
                <div class="form-group">
                    <label for="slide_count">Number of Slides:</label>
                    <input type="number" id="slide_count" name="slide_count" min="5" max="20" value="10" required>
                </div>
                
                <div class="form-group checkbox-container">
                    <input type="checkbox" id="include_images" name="include_images" checked>
                    <label for="include_images">Include Generated Images</label>
                </div>
                
                <div class="form-group">
                    <label for="extra_notes">Additional Instructions (optional):</label>
                    <textarea id="extra_notes" name="extra_notes" rows="3" placeholder="Any specific content or structure requests"></textarea>
                </div>
                
                <button type="submit" class="btn">Generate Enhanced Presentation</button>
            </form>
        </div>
    </div>
</body>
</html>
"""

ERROR_TEMPLATE = """
<!DOCTYPE html>
<html>
<head>
    <title>Error</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .container { max-width: 800px; margin: 0 auto; text-align: center; }
        .error { color: #e74c3c; }
        .error-box { background-color: #f8d7da; padding: 15px; text-align: left; border-left: 5px solid #e74c3c; margin: 20px 0; }
        .btn { background: #3498db; color: white; text-decoration: none; padding: 10px 15px; display: inline-block; margin-top: 20px; }
        .logo { font-size: 24px; color: #4CAF50; margin-bottom: 10px; font-weight: bold; }
    </style>
</head>
<body>
    <div class="container">
        <div class="logo">Matcha</div>
        <h1 class="error">Error</h1>
        
        <div class="error-box">
            <p>{{ message }}</p>
        </div>
        
        <a href="/" class="btn">Back to Home</a>
    </div>
</body>
</html>
"""

@app.route('/')
def index():
    """Render the home page with a form"""
    return render_template_string(INDEX_TEMPLATE)

@app.route('/analysis_status/<file_id>', methods=['GET'])
def check_analysis_status(file_id):
    """Check the status of structural analysis"""
    if file_id in global_analysis_status:
        return jsonify(global_analysis_status[file_id])
    else:
        return jsonify({'status': 'not_found'})

@app.route('/upload', methods=['POST'])
def upload():
    """Handle file upload and direct to appropriate process"""
    # Check if file is in the request
    if 'pptx' not in request.files:
        return render_template_string(ERROR_TEMPLATE, message="No file part in the request"), 400
        
    file = request.files['pptx']
    profile = request.form.get('profile')
    action = request.form.get('action', 'assess')
    target_language = request.form.get('target_language', '')
    
    # Check if file was selected
    if file.filename == '':
        return render_template_string(ERROR_TEMPLATE, message="No file selected"), 400
        
    # Save the file
    filename = secure_filename(file.filename)
    file_id = str(uuid.uuid4())
    processing_tasks[file_id] = {'status': 'upload', 'filename': filename}
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
    file.save(file_path)
    
    if action == 'assess':
        # Redirect to assessment
        if profile == 'esl' and target_language:
            return redirect(url_for('assess', file_id=file_id, profile=profile, target_language=target_language))
        else:
            return redirect(url_for('assess', file_id=file_id, profile=profile))
    else:
        # Start adaptation process
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Start processing in background thread
        thread = threading.Thread(
            target=process_presentation_efficiently, 
            args=(file_path, file_id, filename, profile, target_language)
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE, 
                                   file_id=file_id, 
                                   filename=f"adapted_{filename}",
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))

@app.route('/scaffolding/<file_id>')
def show_scaffolding(file_id):
    """Show the scaffolding analysis results"""
    if file_id not in processing_tasks or 'scaffolding' not in processing_tasks[file_id]:
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Scaffolding analysis not found. Please analyze the presentation first.")
    
    scaffolding_data = processing_tasks[file_id]['scaffolding']
    
    return render_template_string(
        SCAFFOLDING_TEMPLATE,
        scaffolding=scaffolding_data['scaffolding'],
        slides=scaffolding_data['slides']
    )

@app.route('/status/<file_id>')
def status(file_id):
    """Check the status of presentation processing"""
    if file_id in processing_tasks:
        return jsonify({
            "status": processing_tasks[file_id]['status'], 
            "message": processing_tasks[file_id].get('message', '')
        })
    else:
        return jsonify({
            "status": "error", 
            "message": "Processing task not found"
        }), 404

def analyze_instructional_framework(pptx_path):
    """Use Claude to analyze the instructional framework of a presentation"""
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Extract all slide content with metadata
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = ""
        title = ""
        
        # Get slide title
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title.text
            
        # Extract all text from the slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text:
                slide_content += shape.text + "\n"
        
        slides_data.append({
            "slide_number": i+1,
            "title": title,
            "content": slide_content
        })
    
    # Prepare prompt for Claude
    slide_descriptions = "\n\n".join([
        f"SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
        for s in slides_data
    ])
    
    # Use triple quotes and raw strings to avoid formatting issues with JSON examples
    prompt = f"""
    You are an expert in instructional design. Analyze this PowerPoint presentation and identify its instructional framework and teaching pattern.
    
    SLIDES:
    {slide_descriptions}
    
    Focus on identifying these specific instructional patterns:
    
    1. Recall/I do/We do/You do/Review framework:
       - Recall: Activating prior knowledge at beginning
       - I do: Teacher demonstration, direct instruction
       - We do: Guided practice where teacher and students work together
       - You do: Independent practice by students
       - Review: Summarizing learning at end
    
    2. 5E Instructional Model:
       - Engage: Hook students' interest
       - Explore: Students explore concepts through activities
       - Explain: Concepts are explained clearly
       - Elaborate: Extend understanding to new situations
       - Evaluate: Assessment of understanding
    
    3. Other common frameworks (specify which one)
    
    For each slide, determine:
    1. Which specific phase of the instructional framework it represents
    2. How effectively it fulfills that role (strong, adequate, weak)
    3. What specific elements make it part of that phase
    
    Then provide:
    1. Framework alignment score (0-100)
    2. Missing or underrepresented phases
    3. Specific recommendations to strengthen the instructional flow
    4. Analysis of balance between direct instruction, guided practice, and independent work
    
    Format your analysis as detailed JSON with this structure (don't include this example in your output):
    
    {{
      "framework": {{
        "identified_framework": "Recall, I do, We do, You do, Review",
        "framework_alignment_score": 75,
        "strengths": ["Clear direct instruction", "Strong recall elements"],
        "weaknesses": ["Limited guided practice", "Insufficient review"],
        "missing_phases": ["We do"],
        "recommendations": "Add more guided practice activities between slides 15-20",
        "balance_analysis": "This presentation is heavily weighted toward direct instruction (65%) with limited guided practice (15%) and independent practice (20%)"
      }},
      "slides": [
        {{
          "slide_number": 1,
          "framework_phase": "Recall",
          "effectiveness": "strong",
          "elements": ["Activates prior knowledge", "Questions prompt thinking"]
        }}
      ]
    }}
    """
    
    try:
        # Call Claude API
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        
        # Extract and parse JSON response
        content = response.content[0].text.strip()
        json_match = re.search(r'(\{.*\})', content, re.DOTALL)
        
        if json_match:
            try:
                framework_data = json.loads(json_match.group(0))
                return framework_data
            except json.JSONDecodeError:
                print("Error parsing JSON from LLM response")
                return {"error": "Could not parse JSON response"}
        
        return {"error": "Could not extract JSON from response"}
        
    except Exception as e:
        print(f"Error in framework analysis: {str(e)}")
        return {"error": f"Analysis failed: {str(e)}"}

@app.route('/analyze/framework/<file_id>')
def analyze_framework(file_id):
    """Analyze a presentation's instructional framework"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing_framework'
        processing_tasks[file_id]['message'] = 'Analyzing instructional framework...'
        
        # Analyze framework
        framework_data = analyze_instructional_framework(file_path)
        
        # Store the framework data
        processing_tasks[file_id]['framework'] = framework_data
        processing_tasks[file_id]['status'] = 'complete'
        
        # Create a template to display the results
        return render_template_string(
            FRAMEWORK_TEMPLATE,
            framework=framework_data,
            file_id=file_id
        )
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing framework: {str(e)}"
        }), 500

def extract_learning_scaffolding_with_llm(pptx_path):
    """Use Claude to identify learning scaffolding elements in a presentation"""
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Extract all slide content with metadata
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = ""
        title = ""
        
        # Get slide title
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title.text
            
        # Extract all text from the slide
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text:
                slide_content += shape.text + "\n"
        
        slides_data.append({
            "slide_number": i+1,
            "title": title,
            "content": slide_content
        })
    
    # Prepare prompt for Claude
    slide_descriptions = "\n\n".join([
        f"SLIDE {s['slide_number']}: {s['title']}\n{s['content']}" 
        for s in slides_data
    ])
    
    prompt = f"""
    Analyze the following PowerPoint presentation slides and identify the learning scaffolding elements.
    
    SLIDES:
    {slide_descriptions}
    
    For each slide, determine its instructional purpose (e.g., objectives, key concept, example, practice, assessment, review).
    
    Then extract the following scaffolding elements from the entire presentation:
    1. Learning Objectives
    2. Key Concepts and their definitions
    3. Examples
    4. Practice Activities
    5. Assessment Items
    6. Review Elements
    
    Format your response as JSON with the following structure:
    {{
      "slides": [
        {{
          "slide_number": 1,
          "purpose": "title",
          "elements": []
        }},
        // other slides...
      ],
      "scaffolding": {{
        "learning_objectives": ["objective 1", "objective 2"],
        "key_concepts": [
          {{ "term": "concept 1", "definition": "definition 1" }}
        ],
        "examples": [
          {{ "slide_number": 5, "content": "example content" }}
        ],
        "practice_activities": [],
        "assessment_items": [],
        "review_elements": []
      }}
    }}
    """
    
    try:
        # Call Claude API
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        
        # Extract and parse JSON response
        content = response.content[0].text.strip()
        json_match = re.search(r'(\{.*\})', content, re.DOTALL)
        
        if json_match:
            try:
                scaffolding_data = json.loads(json_match.group(0))
                return scaffolding_data
            except json.JSONDecodeError:
                print("Error parsing JSON from LLM response")
        
        # Fallback to simplified extraction if JSON parsing fails
        return {
            "slides": [{"slide_number": s["slide_number"], "purpose": "content"} for s in slides_data],
            "scaffolding": {
                "learning_objectives": [],
                "key_concepts": [],
                "examples": [],
                "practice_activities": [],
                "assessment_items": [],
                "review_elements": []
            }
        }
    except Exception as e:
        print(f"Error in scaffolding extraction: {str(e)}")
        # Return empty scaffolding structure as fallback
        return {
            "slides": [],
            "scaffolding": {
                "learning_objectives": [],
                "key_concepts": [],
                "examples": [],
                "practice_activities": [],
                "assessment_items": [],
                "review_elements": []
            }
        }

@app.route('/analyze/scaffolding/<file_id>')
def analyze_scaffolding(file_id):
    """Analyze a presentation for learning scaffolding elements"""
    try:
        # Validate if file_id exists
        if file_id not in processing_tasks:
            return jsonify({
                "status": "error",
                "message": "Presentation not found. Please upload again."
            }), 404
        
        # Get original filename
        filename = processing_tasks[file_id].get('filename', '')
        
        # Use the same naming pattern as in your upload function
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{file_id}_{filename}")
        
        # Debug output
        print(f"Looking for file at: {file_path}")
        
        if not os.path.exists(file_path):
            return jsonify({
                "status": "error",
                "message": "Presentation file not found on server"
            }), 404
            
        # Update processing status
        processing_tasks[file_id]['status'] = 'analyzing'
        processing_tasks[file_id]['message'] = 'Extracting learning scaffolding elements...'
        
        # Extract scaffolding elements
        scaffolding_data = extract_learning_scaffolding_with_llm(file_path)
        
        # Store the scaffolding data in the processing task
        processing_tasks[file_id]['scaffolding'] = scaffolding_data
        processing_tasks[file_id]['status'] = 'complete'
        processing_tasks[file_id]['message'] = 'Scaffolding analysis complete'
        
        # Redirect to the scaffolding results page
        return redirect(url_for('view_presentation_scaffolding', file_id=file_id))        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        if file_id in processing_tasks:
            processing_tasks[file_id]['status'] = 'error'
            processing_tasks[file_id]['message'] = f"Error analyzing scaffolding: {str(e)}"
        
        return jsonify({
            "status": "error",
            "message": f"Error analyzing scaffolding: {str(e)}"
        }), 500

@app.route('/view_scaffolding/<file_id>')
def view_presentation_scaffolding(file_id):
    """Show the scaffolding analysis results"""
    if file_id not in processing_tasks:
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Presentation not found. Please upload again.")
                                     
    if 'scaffolding' not in processing_tasks[file_id]:
        return render_template_string(ERROR_TEMPLATE, 
                                     message="Scaffolding analysis not found. Please analyze the presentation first.")
    
    scaffolding_data = processing_tasks[file_id]['scaffolding']
    
    return render_template_string(
        SCAFFOLDING_TEMPLATE,
        scaffolding=scaffolding_data,
        slides=scaffolding_data.get('slides', [])
    )

def identify_instructional_framework(pptx_path):
    """Use Claude to identify if presentation follows specific instructional frameworks"""
    # Load presentation
    prs = Presentation(pptx_path)
    
    # Extract slide sequence with titles and brief content summary
    slide_sequence = []
    
    for i, slide in enumerate(prs.slides):
        title = ""
        if hasattr(slide.shapes, "title") and slide.shapes.title:
            title = slide.shapes.title.text
            
        # Get content summary
        content_summary = ""
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text:
                # Just get first 200 chars for summary
                content_summary += shape.text[:200] + "... "
        
        slide_sequence.append({
            "slide_number": i+1,
            "title": title,
            "content_summary": content_summary
        })
    
    # Convert to text for the prompt
    sequence_text = "\n\n".join([
        f"SLIDE {s['slide_number']}: {s['title']}\n{s['content_summary']}" 
        for s in slide_sequence
    ])
    
    prompt = f"""
    Analyze this PowerPoint presentation sequence to identify if it follows a specific instructional framework.
    
    PRESENTATION SEQUENCE:
    {sequence_text}
    
    Specifically determine if this presentation follows any of these instructional frameworks:
    
    1. "Recall, Check, I do, We do, You do, Recap" - A common classroom framework where:
       - Recall: Students recall prior knowledge
       - Check: Teacher checks understanding
       - I do: Teacher demonstrates or models
       - We do: Guided practice together
       - You do: Independent practice
       - Recap: Summary and reflection
    
    2. "5E Model" (Engage, Explore, Explain, Elaborate, Evaluate)
    
    3. Direct Instruction (Objectives, Input, Modeling, Check, Guided Practice, Independent Practice)
    
    4. Problem-Based Learning (Problem Presentation, Investigation, Solution Development, Reflection)
    
    5. Other recognized framework (please specify)
    
    For each slide, indicate which phase of the instructional framework it represents.
    Then provide an overall analysis of how well the presentation aligns with the identified framework.
    
    Format your response as JSON with this structure:
    {{
      "identified_framework": "framework name",
      "framework_alignment_score": 85,
      "slides": [
        {{
          "slide_number": 1,
          "framework_phase": "Recall",
          "confidence": "high"
        }},
        // remaining slides...
      ],
      "missing_phases": ["phase name"],
      "recommendations": "recommendations for improving framework alignment"
    }}
    """
    
    try:
        # Call Claude API
        response = call_claude_api_with_retry(prompt, max_tokens=2000)
        
        # Extract JSON from response
        content = response.content[0].text.strip()
        json_match = re.search(r'(\{.*\})', content, re.DOTALL)
        
        if json_match:
            try:
                framework_data = json.loads(json_match.group(0))
                return framework_data
            except json.JSONDecodeError:
                print("Error parsing JSON from LLM response")
        
        # Fallback response if JSON parsing fails
        return {
            "identified_framework": "unknown",
            "framework_alignment_score": 0,
            "slides": [],
            "missing_phases": [],
            "recommendations": "Unable to analyze framework"
        }
    except Exception as e:
        print(f"Error in framework identification: {str(e)}")
        # Return basic framework data as fallback
        return {
            "identified_framework": "unknown",
            "framework_alignment_score": 0,
            "slides": [{"slide_number": i+1, "framework_phase": "unknown"} for i in range(len(prs.slides))],
            "missing_phases": [],
            "recommendations": "Unable to analyze instructional framework due to an error."
        }
    

@app.route('/assess/<file_id>/<profile>')
def assess(file_id, profile):
    try:
        # Find the uploaded file
        filename = None
        for file in os.listdir(app.config['UPLOAD_FOLDER']):
            if file.startswith(f"{file_id}_"):
                filename = file.replace(f"{file_id}_", "")
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], file)
                break
        
        if not filename:
            return render_template_string(ERROR_TEMPLATE, message="File not found"), 404
            
        # Analyze the PowerPoint
        all_text, slide_texts, metrics, complex_words = analyze_pptx(file_path)
        
        # Generate placeholder charts instead of real ones
        complexity_chart = generate_complexity_chart(slide_texts, profile)
        complex_words_chart = generate_complex_words_chart(complex_words)
        
        # Get thresholds for the selected profile
        thresholds = READABILITY_THRESHOLDS.get(profile)
        
        # Calculate an adaptation score
        adaptation_score = calculate_adaptation_score(metrics, thresholds)
        
        # Generate recommendation
        recommendation = generate_recommendation(adaptation_score, profile, metrics, thresholds)
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Initial scaffolding and framework data (empty placeholders)
        scaffolding_analysis = {
            "slides": [],
            "scaffolding": {
                "learning_objectives": ["Learning objectives will appear here"],
                "key_concepts": [{"term": "Key concepts", "definition": "will appear here"}],
                "examples": [],
                "practice_activities": [],
                "assessment_items": [],
                "review_elements": []
            }
        }
        
        framework_analysis = {
            "identified_framework": "Direct Instruction",
            "framework_alignment_score": 70,
            "slides": [{"slide_number": 1, "framework_phase": "Introduction"}],
            "missing_phases": [],
            "recommendations": "Structure analysis in progress. Basic readability analysis is complete."
        }
        
        return render_template_string(ASSESSMENT_TEMPLATE_SIMPLIFIED, 
                                    file_id=file_id,
                                    profile=profile,
                                    profile_name=profile_names.get(profile, profile),
                                    metrics=metrics,
                                    thresholds=thresholds,
                                    adaptation_score=adaptation_score,
                                    recommendation=recommendation,
                                    complexity_chart=complexity_chart,
                                    complex_words_chart=complex_words_chart
                                    )
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error during assessment: {str(e)}"), 500

def extract_content_from_pptx(file_path):
    """Extract all text content from a PowerPoint file"""
    try:
        # Load the presentation
        prs = Presentation(file_path)
        
        # Extract text from all slides
        all_content = ""
        
        for slide_index, slide in enumerate(prs.slides):
            # Add slide number
            all_content += f"Slide {slide_index + 1}:\n"
            
            # Process every shape that might contain text
            for shape in slide.shapes:
                # Process text content in text frames
                if hasattr(shape, "text_frame") and shape.text:
                    all_content += shape.text + "\n\n"
                
                # Check for tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_content = []
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text:
                                row_content.append(cell.text)
                        all_content += " | ".join(row_content) + "\n"
                    all_content += "\n"
                
                # Group shapes
                if hasattr(shape, "shapes"):
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text_frame") and subshape.text:
                            all_content += subshape.text + "\n\n"
            
            all_content += "\n---\n\n"
        
        return all_content
        
    except Exception as e:
        print(f"Error extracting content from PPTX: {str(e)}")
        return "Error extracting content from presentation"

def process_structure_analysis(file_id, file_path, profile):
    """Process the structural analysis in a background thread"""
    try:
        # Store analysis status
        analysis_status = {
            'status': 'processing',
            'scaffolding': None,
            'framework': None,
            'error': None
        }
        
        # Store in global dictionary or database
        global_analysis_status[file_id] = analysis_status
        
        # Extract learning scaffolding with LLM
        scaffolding_analysis = extract_learning_scaffolding_with_llm(file_path)
        analysis_status['scaffolding'] = scaffolding_analysis
        
        # Identify instructional framework
        framework_analysis = identify_instructional_framework(file_path)
        analysis_status['framework'] = framework_analysis
        
        # Mark as complete
        analysis_status['status'] = 'complete'
        
    except Exception as e:
        print(f"Error in structure analysis: {str(e)}")
        analysis_status['status'] = 'error'
        analysis_status['error'] = str(e)

def process_single_batch(texts, profile):
    """Process a single batch of texts using one API call"""
    # For very small batches, process individually
    if len(texts) <= 1:
        return [adapt_text_with_matcha(texts[0], profile)]
    
    # Create a single prompt with multiple texts
    combined_prompt = f"Adapt the following texts for {profile} users. Format your response using exactly '### TEXT N ###' before each adapted text (where N is the text number).\n\n"
    
    for i, text in enumerate(texts):
        combined_prompt += f"### TEXT {i+1} ###\n{text}\n\n"
    
    try:
        # Call API with the combined prompt
        response = call_claude_api_with_retry(combined_prompt)
        
        # Parse the response to extract individual adapted texts
        adapted_texts = []
        content = response.content[0].text
        
        # Better parsing approach
        pattern = r'###\s*TEXT\s*(\d+)\s*###\s*(.*?)(?=###\s*TEXT\s*\d+\s*###|$)'
        matches = re.findall(pattern, content, re.DOTALL)
        
        # Create a dictionary to maintain order
        text_dict = {}
        for match in matches:
            text_num = int(match[0])
            if 1 <= text_num <= len(texts):  # Ensure valid text number
                text_dict[text_num] = match[1].strip()
        
        # If we don't have matches for all texts, fall back to individual processing
        if len(text_dict) != len(texts):
            print(f"Warning: Batch processing response parsing failed. Falling back to individual processing.")
            return [adapt_text_with_matcha(text, profile) for text in texts]
            
        # Convert dictionary to ordered list
        for i in range(1, len(texts) + 1):
            if i in text_dict:
                adapted_texts.append(text_dict[i])
            else:
                # If missing a specific text, adapt it individually
                adapted_texts.append(adapt_text_with_matcha(texts[i-1], profile))
        
        return adapted_texts
        
    except Exception as e:
        print(f"Error in batch processing: {str(e)}")
        # Fall back to individual processing
        return [adapt_text_with_matcha(text, profile) for text in texts]

# 4. Content Analysis for Selective Processing
def needs_adaptation(text, profile):
    """Determine if text needs API-based adaptation"""
    # Skip very short text or placeholder text
    if len(text.strip()) < 15 or "Click to add text" in text:
        return False
    
    # Calculate metrics
    metrics = calculate_simple_readability(text)
    
    # Check against thresholds for the profile
    threshold = READABILITY_THRESHOLDS.get(profile, {})
    
    # If already meeting key thresholds, no need to adapt
    if (metrics['flesch_reading_ease'] >= threshold.get('flesch_reading_ease', 70) and
        metrics['flesch_kincaid_grade'] <= threshold.get('flesch_kincaid_grade', 8) and
        metrics['sentence_length'] <= threshold.get('sentence_length', 15)):
        return False
    
    return True

def calculate_simple_readability(text):
    """Calculate simplified readability metrics to determine if adaptation is needed"""
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = max(1, len(sentences))  # Avoid division by zero
    
    # Calculate average sentence length
    avg_sentence_length = word_count / sentence_count
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter
        vowels = "aeiouy"
        word = word.lower()
        count = 0
        if word and word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word and word.endswith("e"):
            count -= 1
        return max(1, count)  # Always at least 1 syllable
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = (len(complex_words) / max(1, word_count)) * 100
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / max(1, word_count))
        flesch = max(0, min(100, flesch))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / max(1, word_count)) - 15.59
        fk_grade = max(0, fk_grade)
    else:
        fk_grade = 0
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

# Claude API call with retry logic and timeout
@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
def call_claude_api_with_retry(prompt, model="claude-3-5-sonnet-20240620", max_tokens=1024, timeout=30):
    """Call Claude API with retry logic and timeout"""
    try:
        # Create a client with timeout
        client_with_timeout = anthropic.Anthropic(
            api_key=api_key,
            timeout=timeout  # Set timeout in seconds
        )
        
        # Call API with timeout
        response = client_with_timeout.messages.create(
            model=model,
            max_tokens=max_tokens,
            messages=[{"role": "user", "content": prompt}]
        )
        return response
    except Exception as e:
        # Check if it's a timeout error
        if "timeout" in str(e).lower() or "timed out" in str(e).lower():
            print(f"API call timed out after {timeout} seconds. Retrying...")
        raise  # Re-raise for retry handling

def process_text_batch(texts, profile, max_batch_size=5, max_tokens_per_batch=4000):
    """Process multiple text elements in efficient batches"""
    results = []
    current_batch = []
    current_batch_tokens = 0
    
    def estimate_tokens(text):
        """Roughly estimate token count based on character count"""
        return len(text) // 4  # Approximation: ~4 chars per token
    
    # Group texts into efficient batches
    for text in texts:
        text_tokens = estimate_tokens(text)
        
        # If this text would make the batch too large, process the current batch first
        if len(current_batch) >= max_batch_size or (current_batch_tokens + text_tokens > max_tokens_per_batch and current_batch):
            # Process the current batch
            batch_results = process_single_batch(current_batch, profile)
            results.extend(batch_results)
            current_batch = []
            current_batch_tokens = 0
            
            # Small delay between batches to avoid rate limits
            time.sleep(0.5)
        
        # Add this text to the current batch
        current_batch.append(text)
        current_batch_tokens += text_tokens
    
    # Process any remaining texts in the final batch
    if current_batch:
        batch_results = process_single_batch(current_batch, profile)
        results.extend(batch_results)
    
    return results

# Helper function to process a single batch of texts with optimizations
def process_single_batch(texts, profile):
    """Process a single batch of texts with optimizations"""
    # For very small batches, process individually
    if len(texts) <= 1:
        return [adapt_text_with_matcha(texts[0], profile)]
    
    # Debounce identical texts in the batch to avoid redundant processing
    # Create a mapping from unique texts to their positions in the original list
    unique_texts = {}
    for i, text in enumerate(texts):
        if text in unique_texts:
            unique_texts[text].append(i)
        else:
            unique_texts[text] = [i]
    
    # Process only the unique texts
    unique_text_list = list(unique_texts.keys())
    
    # Create a more efficient combined prompt
    combined_prompt = f"Adapt these texts for {profile} users. Format exactly as '### TEXT N ###' before each adapted text.\n\n"
    
    for i, text in enumerate(unique_text_list):
        # Format each text with minimal tokens
        combined_prompt += f"### TEXT {i+1} ###\n{text}\n\n"
    
    try:
        # Use token-efficient prompt with higher max_tokens
        # This reduces the chance of truncation issues
        response = call_claude_api_with_retry(
            combined_prompt, 
            max_tokens=max(4000, sum(len(text) // 2 for text in unique_text_list)),
            timeout=45  # Increased timeout for larger batches
        )
        
        # Parse the response to extract individual adapted texts
        content = response.content[0].text
        
        # Improved regex pattern for more reliable extraction
        pattern = r'###\s*TEXT\s*(\d+)\s*###\s*(.*?)(?=###\s*TEXT\s*\d+\s*###|$)'
        matches = re.findall(pattern, content, re.DOTALL)
        
        # Create a dictionary to maintain order
        unique_results = {}
        for match in matches:
            text_num = int(match[0])
            if 1 <= text_num <= len(unique_text_list):  # Ensure valid text number
                unique_results[text_num] = match[1].strip()
        
        # Check if we have all results
        if len(unique_results) != len(unique_text_list):
            missing_indices = set(range(1, len(unique_text_list) + 1)) - set(unique_results.keys())
            print(f"Warning: Missing indices in batch response: {missing_indices}")
            
            # Process missing texts individually
            for idx in missing_indices:
                if 1 <= idx <= len(unique_text_list):
                    unique_results[idx] = adapt_text_with_matcha(unique_text_list[idx-1], profile)
        
        # Map the unique results back to the original positions
        results = [None] * len(texts)
        for unique_idx, positions in enumerate(unique_texts.values(), 1):
            adapted_text = unique_results.get(unique_idx, "")
            for pos in positions:
                results[pos] = adapted_text
        
        return results
        
    except Exception as e:
        print(f"Error in batch processing: {str(e)}")
        # Fall back to individual processing
        return [adapt_text_with_matcha(text, profile) for text in texts]

# Efficient implementation of process_presentation
def process_presentation_efficiently(file_path, file_id, filename, profile, target_language=None):
    try:
        # Set initial status if not already set
        if file_id not in processing_tasks:
            processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        # Load the presentation
        prs = Presentation(file_path)
        
        # For bilingual ESL, create a copy of the presentation for the translated version
        translated_prs = None
        if profile == "esl" and target_language:
            translated_prs = Presentation(file_path)
        
        # Collect all text elements that need processing
        text_elements = []
        element_references = []
        
        print(f"Analyzing presentation with {len(prs.slides)} slides...")
        
        # First pass: collect all text
        for slide_index, slide in enumerate(prs.slides):
            slide_elements = collect_text_elements_from_slide(slide, slide_index)
            
            for element in slide_elements:
                text = element['text']
                # Only add if it needs adaptation
                if needs_adaptation(text, profile):
                    text_elements.append(text)
                    element_references.append(element['reference'])
        
        print(f"Found {len(text_elements)} text elements requiring adaptation")
        
        # Process collected text in batches
        if text_elements:
            print(f"Processing {len(text_elements)} text elements in batches")
            adapted_texts = process_text_batch(text_elements, profile)
            
            # For bilingual ESL, create translations of the original text
            translated_texts = []
            if profile == "esl" and target_language and translated_prs:
                print(f"Translating {len(text_elements)} text elements to {target_language}")
                # Create translations directly from original text (not from adapted text)
                for original_text in text_elements:
                    translated_text = translate_text(original_text, target_language)
                    translated_texts.append(translated_text)
            
            # Apply adapted texts back to their elements in the English version
            for i, adapted_text in enumerate(adapted_texts):
                reference = element_references[i]
                apply_adapted_text(reference, adapted_text, profile)
                
                # Apply translations to the translated presentation
                if profile == "esl" and target_language and translated_prs and i < len(translated_texts):
                    # Find the corresponding element in the translated presentation
                    slide_num = reference.get('slide_num', 0)
                    shape_index = reference.get('shape_index', 0)
                    
                    if slide_num < len(translated_prs.slides):
                        trans_slide = translated_prs.slides[slide_num]
                        shapes = list(trans_slide.shapes)
                        if shape_index < len(shapes):
                            trans_shape = shapes[shape_index]
                            
                            if reference['type'] == 'shape' and hasattr(trans_shape, "text_frame"):
                                apply_text_to_text_frame(trans_shape.text_frame, translated_texts[i], profile)
                            elif reference['type'] == 'cell':
                                # Handle table cells
                                row_index = reference.get('row_index', 0)
                                cell_index = reference.get('cell_index', 0)
                                if hasattr(trans_shape, "table"):
                                    if row_index < len(trans_shape.table.rows):
                                        row = trans_shape.table.rows[row_index]
                                        if cell_index < len(row.cells):
                                            cell = row.cells[cell_index]
                                            if hasattr(cell, "text_frame"):
                                                apply_text_to_text_frame(cell.text_frame, translated_texts[i], profile)
                            elif reference['type'] == 'subshape':
                                # Handle grouped shapes
                                subshape_index = reference.get('subshape_index', 0)
                                if hasattr(trans_shape, "shapes") and subshape_index < len(trans_shape.shapes):
                                    subshape = trans_shape.shapes[subshape_index]
                                    if hasattr(subshape, "text_frame"):
                                        apply_text_to_text_frame(subshape.text_frame, translated_texts[i], profile)
        
        # Save the adapted presentation
        output_filename = f"adapted_{filename}"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{output_filename}")
        prs.save(output_path)
        
        # Save the translated version if created
        translated_path = None
        translated_filename = None
        lang_code = ""
        if translated_prs and target_language:
            lang_code = target_language[:3]  # First 3 chars as code
            translated_filename = f"translated_{lang_code}_{filename}"
            translated_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{translated_filename}")
            translated_prs.save(translated_path)
        
        # Update status to complete
        processing_tasks[file_id] = {
            'status': 'complete', 
            'message': '',
            'has_translation': bool(translated_path),
            'translated_filename': translated_filename
        }
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}


# Helper function to collect text elements from a slide
def collect_text_elements_from_slide(slide, slide_index=0):
    """Collect text elements from a slide and their references"""
    elements = []
    
    # Process shapes with text
    for shape_index, shape in enumerate(slide.shapes):
        if hasattr(shape, "text_frame") and shape.text:
            elements.append({
                'text': shape.text,
                'reference': {
                    'type': 'shape', 
                    'object': shape,
                    'slide_num': slide_index,
                    'shape_index': shape_index
                }
            })
        
        # Tables
        if hasattr(shape, "table"):
            for row_index, row in enumerate(shape.table.rows):
                for cell_index, cell in enumerate(row.cells):
                    if hasattr(cell, "text_frame") and cell.text:
                        elements.append({
                            'text': cell.text,
                            'reference': {
                                'type': 'cell', 
                                'object': cell,
                                'slide_num': slide_index,
                                'shape_index': shape_index,
                                'row_index': row_index,
                                'cell_index': cell_index
                            }
                        })
        
        # Group shapes
        if hasattr(shape, "shapes"):
            for subshape_index, subshape in enumerate(shape.shapes):
                if hasattr(subshape, "text_frame") and subshape.text:
                    elements.append({
                        'text': subshape.text,
                        'reference': {
                            'type': 'subshape', 
                            'object': subshape,
                            'slide_num': slide_index,
                            'shape_index': shape_index,
                            'subshape_index': subshape_index
                        }
                    })
    
    return elements

# Helper function to apply adapted text back to elements
def apply_adapted_text(reference, adapted_text, profile):
    """Apply adapted text to the appropriate element with formatting"""
    element_type = reference['type']
    obj = reference['object']
    
    if element_type == 'shape' or element_type == 'subshape' or element_type == 'cell':
        if hasattr(obj, "text_frame"):
            apply_text_to_text_frame(obj.text_frame, adapted_text, profile)
            
            # For dyslexia, apply additional formatting
            if profile == "dyslexia" and hasattr(obj, "text_frame"):
                apply_dyslexia_formatting(obj.text_frame)

# Function to apply text to a text frame with color formatting
def apply_text_to_text_frame(text_frame, adapted_text, profile):
    """Apply adapted text to a text frame with proper formatting"""
    # Skip if text frame is empty
    if not hasattr(text_frame, "paragraphs"):
        return
    
    # Get the profile color
    profile_color = PROFILE_COLORS.get(profile)
    
    try:
        # Store original formatting settings we want to preserve
        original_formatting = []
        for p in text_frame.paragraphs:
            p_format = {
                "level": p.level if hasattr(p, "level") else 0,
                "alignment": p.alignment if hasattr(p, "alignment") else None,
                "font_name": None,
                "font_size": None,
                "font_bold": None,
                "font_italic": None
            }
            
            # Get formatting from first run if available
            if p.runs and hasattr(p.runs[0], "font"):
                font = p.runs[0].font
                p_format["font_name"] = font.name if hasattr(font, "name") else None
                p_format["font_size"] = font.size if hasattr(font, "size") else None
                p_format["font_bold"] = font.bold if hasattr(font, "bold") else None
                p_format["font_italic"] = font.italic if hasattr(font, "italic") else None
            
            original_formatting.append(p_format)
        
        # Clear all paragraphs except the first one
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            tr = p._p
            tr.getparent().remove(tr)
        
        # Clear the first paragraph content but preserve its format
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            while len(p.runs) > 0:
                try:
                    p._p.remove(p.runs[0]._r)
                except:
                    break
        else:
            p = text_frame.add_paragraph()
        
        # Split the adapted text into paragraphs
        text_paragraphs = adapted_text.split('\n')
        
        # Format the first paragraph with colored first word
        if text_paragraphs and text_paragraphs[0].strip():
            words = text_paragraphs[0].split()
            if words:
                # First word with color
                first_word = words[0]
                first_run = p.add_run()
                first_run.text = first_word
                # Apply profile color to first word
                if profile_color and hasattr(first_run, "font") and hasattr(first_run.font, "color"):
                    first_run.font.color.rgb = profile_color
                
                # Apply original formatting to first run
                if original_formatting and hasattr(first_run, "font"):
                    font = first_run.font
                    if original_formatting[0]["font_name"]:
                        font.name = original_formatting[0]["font_name"]
                    if original_formatting[0]["font_size"]:
                        font.size = original_formatting[0]["font_size"]
                    if original_formatting[0]["font_bold"]:
                        font.bold = original_formatting[0]["font_bold"]
                    if original_formatting[0]["font_italic"]:
                        font.italic = original_formatting[0]["font_italic"]
                
                # Rest of first paragraph
                rest_of_text = ' '.join(words[1:]) if len(words) > 1 else ''
                if rest_of_text:
                    rest_run = p.add_run()
                    rest_run.text = ' ' + rest_of_text
                    # Apply original formatting to rest run
                    if original_formatting and hasattr(rest_run, "font"):
                        font = rest_run.font
                        if original_formatting[0]["font_name"]:
                            font.name = original_formatting[0]["font_name"]
                        if original_formatting[0]["font_size"]:
                            font.size = original_formatting[0]["font_size"]
                        if original_formatting[0]["font_bold"]:
                            font.bold = original_formatting[0]["font_bold"]
                        if original_formatting[0]["font_italic"]:
                            font.italic = original_formatting[0]["font_italic"]
            
            # Apply original paragraph formatting
            if original_formatting:
                if original_formatting[0]["alignment"]:
                    p.alignment = original_formatting[0]["alignment"]
                if original_formatting[0]["level"]:
                    p.level = original_formatting[0]["level"]
        
        # Add remaining paragraphs
        for i, p_text in enumerate(text_paragraphs[1:], 1):
            if p_text.strip():
                new_p = text_frame.add_paragraph()
                run = new_p.add_run()
                run.text = p_text
                
                # Apply formatting if available
                fmt_index = min(i, len(original_formatting) - 1) if original_formatting else -1
                
                if fmt_index >= 0:
                    # Apply paragraph formatting
                    if original_formatting[fmt_index]["alignment"]:
                        new_p.alignment = original_formatting[fmt_index]["alignment"]
                    if original_formatting[fmt_index]["level"]:
                        new_p.level = original_formatting[fmt_index]["level"]
                    
                    # Apply run formatting
                    if hasattr(run, "font"):
                        font = run.font
                        if original_formatting[fmt_index]["font_name"]:
                            font.name = original_formatting[fmt_index]["font_name"]
                        if original_formatting[fmt_index]["font_size"]:
                            font.size = original_formatting[fmt_index]["font_size"]
                        if original_formatting[fmt_index]["font_bold"]:
                            font.bold = original_formatting[fmt_index]["font_bold"]
                        if original_formatting[fmt_index]["font_italic"]:
                            font.italic = original_formatting[fmt_index]["font_italic"]
    
    except Exception as e:
        print(f"Error formatting text frame: {str(e)}")
        # Fallback: just set the text directly if there's an error
        try:
            text_frame.text = adapted_text
        except:
            pass

@app.route('/download/<file_id>/<filename>')
def download(file_id, filename):
    """Show the download page"""
    # Get the profile from the filename or default to dyslexia
    profile = "dyslexia"
    if "adhd" in filename.lower():
        profile = "adhd"
    elif "esl" in filename.lower():
        profile = "esl"
    
    # Get translation information if available
    task_info = processing_tasks.get(file_id, {})
    has_translation = task_info.get('has_translation', False)
    translated_filename = task_info.get('translated_filename', '')
    
    # Get language name from filename if available
    translated_language = "Translated"
    if translated_filename:
        lang_code = translated_filename.split('_')[1]
        language_map = {
            'spa': 'Spanish', 'fre': 'French', 'ger': 'German', 
            'ita': 'Italian', 'por': 'Portuguese', 'chi': 'Chinese', 
            'jap': 'Japanese', 'ara': 'Arabic', 'hin': 'Hindi', 'rus': 'Russian'
        }
        translated_language = language_map.get(lang_code, 'Translated')
        
    return render_template_string(DOWNLOAD_TEMPLATE_WITH_TRANSLATION, 
                           file_id=file_id, 
                           filename=filename,
                           profile=profile,
                           has_translation=has_translation,
                           translated_filename=translated_filename,
                           translated_language=translated_language)

@app.route('/download_file/<file_id>/<filename>')
def download_file(file_id, filename):
    """Download the actual PowerPoint file"""
    try:
        return send_file(
            os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}"),
            as_attachment=True,
            download_name=filename
        )
    except Exception as e:
        return render_template_string(ERROR_TEMPLATE, message=f"File not found: {str(e)}"), 404

@app.route('/error')
def error():
    """Show error page"""
    message = request.args.get('message', 'An unknown error occurred')
    return render_template_string(ERROR_TEMPLATE, message=message)

def analyze_pptx(file_path):
    """Analyze PowerPoint content for readability"""
    # Load the presentation
    prs = Presentation(file_path)
    
    # Extract all text from the presentation
    all_text = ""
    slide_texts = []  # Text per slide for analysis
    
    for slide_index, slide in enumerate(prs.slides):
        slide_text = ""
        
        # Process every shape that might contain text
        for shape in slide.shapes:
            # Process text content in text frames
            if hasattr(shape, "text_frame") and shape.text:
                all_text += shape.text + "\n\n"
                slide_text += shape.text + "\n\n"
            
            # Check for tables
            if hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, "text_frame") and cell.text:
                            all_text += cell.text + "\n"
                            slide_text += cell.text + "\n"
            
            # Group shapes
            if hasattr(shape, "shapes"):
                for subshape in shape.shapes:
                    if hasattr(subshape, "text_frame") and subshape.text:
                        all_text += subshape.text + "\n\n"
                        slide_text += subshape.text + "\n\n"
        
        slide_texts.append({
            "slide_number": slide_index + 1,
            "text": slide_text
        })
    
    # Calculate readability metrics
    metrics = calculate_readability_metrics(all_text)
    
    # Get complex words
    complex_words = find_complex_words(all_text)
    
    return all_text, slide_texts, metrics, complex_words

def calculate_readability_metrics(text):
    """Calculate readability metrics for the text"""
    # Simplified calculation of readability metrics
    # In a production environment, use a proper textstat library
    
    # Count words, sentences, and syllables
    words = re.findall(r'\b[a-zA-Z]+\b', text.lower())
    sentences = re.split(r'[.!?]+', text)
    sentences = [s for s in sentences if s.strip()]
    
    # Calculate metrics
    word_count = len(words)
    sentence_count = len(sentences)
    
    # Calculate average sentence length
    avg_sentence_length = round(word_count / max(1, sentence_count), 1)
    
    # Count syllables (simplified approach)
    def count_syllables(word):
        # Simple syllable counter - not perfect but adequate for demo
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    syllable_count = sum(count_syllables(word) for word in words)
    
    # Calculate complex words (words with 3+ syllables)
    complex_words = [word for word in words if count_syllables(word) >= 3]
    complex_word_percent = round((len(complex_words) / max(1, word_count)) * 100, 1)
    
    # Calculate Flesch Reading Ease
    if word_count > 0 and sentence_count > 0:
        flesch = 206.835 - 1.015 * (word_count / sentence_count) - 84.6 * (syllable_count / word_count)
        flesch = max(0, min(100, round(flesch)))
    else:
        flesch = 0
    
    # Calculate Flesch-Kincaid Grade Level
    if word_count > 0 and sentence_count > 0:
        fk_grade = 0.39 * (word_count / sentence_count) + 11.8 * (syllable_count / word_count) - 15.59
        fk_grade = max(0, round(fk_grade, 1))
    else:
        fk_grade = 0
    
    # Calculate SMOG Index (simplified)
    if sentence_count >= 30:
        smog = 1.043 * ((len(complex_words) * (30 / sentence_count)) ** 0.5) + 3.1291
    else:
        smog = 1.043 * ((len(complex_words) * (30 / max(1, sentence_count))) ** 0.5) + 3.1291
    smog = max(0, round(smog, 1))
    
    return {
        "flesch_reading_ease": flesch,
        "flesch_kincaid_grade": fk_grade,
        "smog_index": smog,
        "sentence_length": avg_sentence_length,
        "complex_word_percent": complex_word_percent
    }

def generate_lesson_plan(topic, grade_level, slide_count, profile, extra_notes, subject_area):
    """Generate a lesson plan structure using Claude with UK curriculum"""
    # Updated UK education level descriptions
    level_descriptions = {
        "ks3-year7-8": "Key Stage 3 (Year 7-8, ages 11-13)",
        "ks3-year9": "Key Stage 3 (Year 9, age 14)",
        "ks4-gcse": "Key Stage 4 (GCSE, Year 10-11, ages 15-16)",
        "ks5-alevel": "Key Stage 5 (A-Level, Year 12-13, ages 17-18)",
        "btec": "BTEC qualification level"
    }
    
    level_desc = level_descriptions.get(grade_level, "secondary school")
    
    adaptation_needs = {
        "dyslexia": "simplified vocabulary, short sentences, clear structure, minimal text per slide",
        "adhd": "visually engaging, chunked information, clear headings, bulleted lists, minimal distractions",
        "esl": "basic vocabulary with more complex terms in parentheses, visual support, simple sentence structure"
    }
    
    adaptation_desc = adaptation_needs.get(profile, "")
    
    # Subject-specific guidance
    subject_guidance = {
        "english": "Include key literary terms, text analysis techniques, and writing structures appropriate for this level.",
        "maths": "Include clear mathematical notations, step-by-step explanations, and visual representations of concepts.",
        "science": "Include scientific diagrams, experimental procedures, and key terminology with explanations.",
        "biology": "Include biological diagrams, processes, and scientific terminology with explanations.",
        "chemistry": "Include chemical formulas, reaction mechanisms, and key principles with clear explanations.",
        "physics": "Include physics equations, diagrams, and real-world applications of concepts.",
        "history": "Include timelines, primary sources, and analysis of historical events.",
        "geography": "Include maps, diagrams, and case studies relevant to the topic.",
        "religious_studies": "Include key beliefs, practices, and comparative analysis where appropriate.",
        "modern_languages": "Include vocabulary lists, language structures, and cultural context.",
        "computing": "Include code examples, system diagrams, and step-by-step procedures.",
        "design_technology": "Include design specifications, technical drawings, and manufacturing processes.",
        "art_design": "Include visual examples, techniques, and artistic context.",
        "music": "Include notation examples, listening activities, and music theory concepts.",
        "physical_education": "Include diagrams of techniques, rules of activities, and fitness concepts.",
        "citizenship": "Include case studies, legal frameworks, and civic participation examples.",
        "pshe": "Include scenarios, discussion points, and practical guidance.",
        "business": "Include business models, case studies, and key terminology.",
        "sociology": "Include sociological theories, research methods, and social examples.",
        "psychology": "Include psychological theories, studies, and applications.",
        "economics": "Include economic models, graphs, and real-world examples.",
        "media_studies": "Include media analysis frameworks, case studies, and production techniques."
    }
    
    subject_specific_guidance = subject_guidance.get(subject_area, "")
    
    prompt = f"""
    Create a detailed lesson plan for a {slide_count}-slide PowerPoint presentation about "{topic}" for {level_desc} students following the UK National Curriculum for {subject_area}.
    
    This presentation needs to be adapted for students with {profile} and should include: {adaptation_desc}.
    
    Additional subject-specific guidance: {subject_specific_guidance}
    
    Additional instructions from the teacher: {extra_notes}
    
    For each slide, provide:
    1. Slide type (title slide, content slide, image slide, etc.)
    2. Title of the slide
    3. Key content points (in bullet form)
    4. A brief image description that would support the content (for slides that should have images)
    
    Format your response as a structured JSON array where each element is a slide with the following properties:
    - layout_index: (0 for title slide, 1 for content slide with title)
    - title: The slide title
    - content: The main content (bulleted list where appropriate)
    - image_prompt: A description for image generation (if applicable)
    
    Keep the language appropriate for {profile} needs and align with UK curriculum standards.
    """
    
    try:
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        content = response.content[0].text.strip()
        
        # Extract JSON from the response
        json_match = re.search(r'\[\s*{.*}\s*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            lesson_plan = json.loads(json_str)
            return lesson_plan
        else:
            # If no JSON found, try to parse the content in a more forgiving way
            try:
                # Try to find anything that looks like JSON
                json_start = content.find('[')
                json_end = content.rfind(']') + 1
                if json_start >= 0 and json_end > json_start:
                    json_str = content[json_start:json_end]
                    lesson_plan = json.loads(json_str)
                    return lesson_plan
            except json.JSONDecodeError:
                pass
            
            # If all JSON parsing fails, create a simple structure from the text
            slides = []
            sections = content.split('\n\n')
            for i, section in enumerate(sections):
                if i == 0:
                    # First section is the title slide
                    slides.append({
                        "layout_index": 0,
                        "title": topic,
                        "content": section
                    })
                else:
                    # Extract title and content
                    lines = section.split('\n')
                    if lines:
                        title = lines[0].strip('#- ').strip()
                        content = '\n'.join(lines[1:]).strip()
                        slides.append({
                            "layout_index": 1,
                            "title": title,
                            "content": content,
                            "image_prompt": f"Educational illustration about {title} for {level_desc} students studying {subject_area}"
                        })
            
            return slides
    
    except Exception as e:
        print(f"Error generating lesson plan: {str(e)}")
        # Create a basic fallback plan
        return [
            {"layout_index": 0, "title": topic, "content": f"A presentation about {topic} for {level_desc} students studying {subject_area}"},
            {"layout_index": 1, "title": "Introduction", "content": f"• Introduction to {topic}\n• Key concepts\n• Learning objectives", "image_prompt": f"Introduction to {topic} for UK {subject_area} curriculum"},
            {"layout_index": 1, "title": "Main Concepts", "content": f"• Important aspects of {topic}\n• Key information", "image_prompt": f"Main concepts of {topic} for UK {subject_area} curriculum"}
        ]
def generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area):
    """Generate an enriched lesson plan based on the adapted content for UK curriculum"""
    # Updated UK education level descriptions
    level_descriptions = {
        "ks3-year7-8": "Key Stage 3 (Year 7-8, ages 11-13)",
        "ks3-year9": "Key Stage 3 (Year 9, age 14)",
        "ks4-gcse": "Key Stage 4 (GCSE, Year 10-11, ages 15-16)",
        "ks5-alevel": "Key Stage 5 (A-Level, Year 12-13, ages 17-18)",
        "btec": "BTEC qualification level"
    }
    
    level_desc = level_descriptions.get(grade_level, "secondary school")
    
    adaptation_needs = {
        "dyslexia": "simplified vocabulary, short sentences, clear structure, minimal text per slide",
        "adhd": "visually engaging, chunked information, clear headings, bulleted lists, minimal distractions",
        "esl": "basic vocabulary with more complex terms in parentheses, visual support, simple sentence structure"
    }
    
    adaptation_desc = adaptation_needs.get(profile, "")
    
    # Subject-specific guidance
    subject_guidance = {
        "english": "Include key literary terms, text analysis techniques, and writing structures appropriate for this level.",
        "maths": "Include clear mathematical notations, step-by-step explanations, and visual representations of concepts.",
        "science": "Include scientific diagrams, experimental procedures, and key terminology with explanations.",
        "biology": "Include biological diagrams, processes, and scientific terminology with explanations.",
        "chemistry": "Include chemical formulas, reaction mechanisms, and key principles with clear explanations.",
        "physics": "Include physics equations, diagrams, and real-world applications of concepts.",
        "history": "Include timelines, primary sources, and analysis of historical events.",
        "geography": "Include maps, diagrams, and case studies relevant to the topic.",
        "religious_studies": "Include key beliefs, practices, and comparative analysis where appropriate.",
        "modern_languages": "Include vocabulary lists, language structures, and cultural context.",
        "computing": "Include code examples, system diagrams, and step-by-step procedures.",
        "design_technology": "Include design specifications, technical drawings, and manufacturing processes.",
        "art_design": "Include visual examples, techniques, and artistic context.",
        "music": "Include notation examples, listening activities, and music theory concepts.",
        "physical_education": "Include diagrams of techniques, rules of activities, and fitness concepts.",
        "citizenship": "Include case studies, legal frameworks, and civic participation examples.",
        "pshe": "Include scenarios, discussion points, and practical guidance.",
        "business": "Include business models, case studies, and key terminology.",
        "sociology": "Include sociological theories, research methods, and social examples.",
        "psychology": "Include psychological theories, studies, and applications.",
        "economics": "Include economic models, graphs, and real-world examples.",
        "media_studies": "Include media analysis frameworks, case studies, and production techniques."
    }
    
    subject_specific_guidance = subject_guidance.get(subject_area, "")
    
    prompt = f"""
    Create an enriched, enhanced lesson plan for a PowerPoint presentation about "{topic}" for {level_desc} students following the UK National Curriculum for {subject_area}.
    
    This presentation should build upon the existing content below, but enhance and expand it to create a more comprehensive, engaging presentation.
    
    EXISTING CONTENT:
    {adapted_content}
    
    The enriched presentation should:
    1. Maintain the core content from the original
    2. Add engaging visuals and interactive elements
    3. Include additional examples and learning activities
    4. Be adapted for students with {profile} needs: {adaptation_desc}
    5. Align with UK {subject_area} curriculum standards for {level_desc}
    6. {subject_specific_guidance}
    7. {extra_notes}
    
    For each slide, provide:
    1. Slide type (title slide, content slide, image slide, etc.)
    2. Title of the slide
    3. Key content points (in bullet form)
    4. A brief image description that would support the content
    
    Format your response as a structured JSON array where each element is a slide with the following properties:
    - layout_index: (0 for title slide, 1 for content slide with title)
    - title: The slide title
    - content: The main content (bulleted list where appropriate)
    - image_prompt: A description for image generation (if applicable)
    
    Keep the language appropriate for {profile} needs and UK curriculum standards.
    """
    
    try:
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        content = response.content[0].text.strip()
        
        # Extract JSON from the response
        json_match = re.search(r'\[\s*{.*}\s*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            lesson_plan = json.loads(json_str)
            return lesson_plan
        else:
            # If no JSON found, try to parse the content in a more forgiving way
            try:
                # Try to find anything that looks like JSON
                json_start = content.find('[')
                json_end = content.rfind(']') + 1
                if json_start >= 0 and json_end > json_start:
                    json_str = content[json_start:json_end]
                    lesson_plan = json.loads(json_str)
                    return lesson_plan
            except json.JSONDecodeError:
                pass
            
            # If all JSON parsing fails, create a simple structure from the text
            slides = []
            sections = content.split('\n\n')
            for i, section in enumerate(sections):
                if i == 0:
                    # First section is the title slide
                    slides.append({
                        "layout_index": 0,
                        "title": topic,
                        "content": section
                    })
                else:
                    # Extract title and content
                    lines = section.split('\n')
                    if lines:
                        title = lines[0].strip('#- ').strip()
                        content = '\n'.join(lines[1:]).strip()
                        slides.append({
                            "layout_index": 1,
                            "title": title,
                            "content": content,
                            "image_prompt": f"Educational illustration about {title} for {level_desc} students studying {subject_area}"
                        })
            
            return slides
    
    except Exception as e:
        print(f"Error generating enriched lesson plan: {str(e)}")
        # Create a basic fallback plan
        return [
            {"layout_index": 0, "title": topic, "content": f"An enriched presentation about {topic} for {level_desc} students studying {subject_area}"},
            {"layout_index": 1, "title": "Introduction", "content": f"• Introduction to {topic}\n• Key concepts\n• Learning objectives", "image_prompt": f"Introduction to {topic} for UK {subject_area} curriculum"},
            {"layout_index": 1, "title": "Main Concepts", "content": f"• Important aspects of {topic}\n• Key information", "image_prompt": f"Main concepts of {topic} for UK {subject_area} curriculum"}
        ]


def find_complex_words(text):
    """Find and count complex words (3+ syllables)"""
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    complex_words = {}
    
    # Count syllables function (simplified)
    def count_syllables(word):
        count = 0
        vowels = "aeiouy"
        word = word.lower()
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i-1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count = 1
        return count
    
    for word in words:
        syllables = count_syllables(word)
        if syllables >= 3:
            if word in complex_words:
                complex_words[word] += 1
            else:
                complex_words[word] = 1
    
    # Sort by frequency (most frequent first)
    sorted_words = sorted(complex_words.items(), key=lambda x: x[1], reverse=True)
    
    # Return the top 10 most frequent complex words
    return dict(sorted_words[:10])

def calculate_adaptation_score(metrics, thresholds):
    """Calculate an overall adaptation score based on metrics and thresholds"""
    # Initialize score components
    flesch_score = 0
    grade_score = 0
    smog_score = 0
    sentence_score = 0
    complex_score = 0
    
    # Flesch Reading Ease (higher is better)
    flesch_threshold = thresholds['flesch_reading_ease']
    if metrics['flesch_reading_ease'] >= flesch_threshold:
        flesch_score = 100
    elif metrics['flesch_reading_ease'] >= flesch_threshold - 20:
        # Partial score for being close
        flesch_score = 50 + ((metrics['flesch_reading_ease'] - (flesch_threshold - 20)) / 20) * 50
    else:
        flesch_score = max(0, (metrics['flesch_reading_ease'] / (flesch_threshold - 20)) * 50)
    
    # Grade Level (lower is better)
    grade_threshold = thresholds['flesch_kincaid_grade']
    if metrics['flesch_kincaid_grade'] <= grade_threshold:
        grade_score = 100
    elif metrics['flesch_kincaid_grade'] <= grade_threshold + 4:
        # Partial score for being close
        grade_score = 100 - ((metrics['flesch_kincaid_grade'] - grade_threshold) / 4) * 50
    else:
        grade_score = max(0, 50 - ((metrics['flesch_kincaid_grade'] - (grade_threshold + 4)) / 4) * 50)
    
    # SMOG Index (lower is better)
    smog_threshold = thresholds['smog_index']
    if metrics['smog_index'] <= smog_threshold:
        smog_score = 100
    elif metrics['smog_index'] <= smog_threshold + 3:
        # Partial score for being close
        smog_score = 100 - ((metrics['smog_index'] - smog_threshold) / 3) * 50
    else:
        smog_score = max(0, 50 - ((metrics['smog_index'] - (smog_threshold + 3)) / 3) * 50)
    
    # Sentence Length (lower is better)
    sentence_threshold = thresholds['sentence_length']
    if metrics['sentence_length'] <= sentence_threshold:
        sentence_score = 100
    elif metrics['sentence_length'] <= sentence_threshold + 8:
        # Partial score for being close
        sentence_score = 100 - ((metrics['sentence_length'] - sentence_threshold) / 8) * 50
    else:
        sentence_score = max(0, 50 - ((metrics['sentence_length'] - (sentence_threshold + 8)) / 5) * 50)
    
    # Complex Word Percentage (lower is better)
    complex_threshold = thresholds['complex_word_percent']
    if metrics['complex_word_percent'] <= complex_threshold:
        complex_score = 100
    elif metrics['complex_word_percent'] <= complex_threshold + 15:
        # Partial score for being close
        complex_score = 100 - ((metrics['complex_word_percent'] - complex_threshold) / 15) * 50
    else:
        complex_score = max(0, 50 - ((metrics['complex_word_percent'] - (complex_threshold + 15)) / 10) * 50)
    
    # Weighted average (weights can be adjusted based on importance)
    weights = {
        'flesch': 0.3,
        'grade': 0.2,
        'smog': 0.1,
        'sentence': 0.2,
        'complex': 0.2
    }
    
    final_score = (
        flesch_score * weights['flesch'] +
        grade_score * weights['grade'] +
        smog_score * weights['smog'] +
        sentence_score * weights['sentence'] +
        complex_score * weights['complex']
    )
    
    # Round to whole number
    return round(final_score)

# 1. First, let's implement the missing extract_content_from_pptx function:

def extract_content_from_pptx(file_path):
    """Extract all text content from a PowerPoint file"""
    try:
        # Load the presentation
        prs = Presentation(file_path)
        
        # Extract text from all slides
        all_content = ""
        
        for slide_index, slide in enumerate(prs.slides):
            # Add slide number
            all_content += f"Slide {slide_index + 1}:\n"
            
            # Process every shape that might contain text
            for shape in slide.shapes:
                # Process text content in text frames
                if hasattr(shape, "text_frame") and shape.text:
                    all_content += shape.text + "\n\n"
                
                # Check for tables
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_content = []
                        for cell in row.cells:
                            if hasattr(cell, "text_frame") and cell.text:
                                row_content.append(cell.text)
                        all_content += " | ".join(row_content) + "\n"
                    all_content += "\n"
                
                # Group shapes
                if hasattr(shape, "shapes"):
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text_frame") and subshape.text:
                            all_content += subshape.text + "\n\n"
            
            all_content += "\n---\n\n"
        
        return all_content
        
    except Exception as e:
        print(f"Error extracting content from PPTX: {str(e)}")
        return "Error extracting content from presentation"

# 2. Complete the generate_new_presentation function:

def generate_new_presentation(file_id, filename, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate a new presentation based on the specified parameters"""
    try:
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Step 1: Plan the content structure using Claude
        lesson_plan = generate_lesson_plan(topic, grade_level, slide_count, profile, extra_notes, subject_area)
        
        # Step 2: Create slides based on the plan
        for slide_data in lesson_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content
                    adapted_content = adapt_text_with_matcha(content, profile)
                    text_frame = content_placeholder.text_frame
                    apply_text_to_text_frame(text_frame, adapted_content, profile)
        
        # Save the presentation
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

# 3. Complete the generate_enriched_presentation function:

def generate_enriched_presentation(file_id, filename, adapted_file_path, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate an enriched presentation based on an adapted file"""
    try:
        # Extract content from the adapted file
        adapted_content = extract_content_from_pptx(adapted_file_path)
        
        # Generate an enriched presentation plan
        enriched_plan = generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Create slides based on the plan
        for slide_data in enriched_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content - already adapted content
                    text_frame = content_placeholder.text_frame
                    apply_text_to_text_frame(text_frame, content, profile)
        
        # Save the presentation
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

# 4. Remove the incomplete generate_recommendation function and fix the typo in the complex_work check
# Replace the entire function with this corrected version:

def generate_recommendation(score, profile, metrics, thresholds):
    """Generate a recommendation based on the score and metrics"""
    profile_names = {
        "dyslexia": "readers with dyslexia",
        "adhd": "readers with ADHD",
        "esl": "English language learners"
    }
    
    audience = profile_names.get(profile, profile)
    
    if score >= 85:
        return f"Your content is already well-suited for {audience}. Adaptation will provide minimal benefits but may still enhance readability."
    
    elif score >= 70:
        # Medium need for adaptation
        issues = []
        
        if metrics['flesch_reading_ease'] < thresholds['flesch_reading_ease']:
            issues.append("general readability")
            
        if metrics['flesch_kincaid_grade'] > thresholds['flesch_kincaid_grade']:
            issues.append("reading level")
            
        if metrics['sentence_length'] > thresholds['sentence_length']:
            issues.append("sentence length")
            
        if metrics['complex_word_percent'] > thresholds['complex_word_percent']:
            issues.append("vocabulary complexity")
            
        if issues:
            issue_text = ", ".join(issues[:-1])
            if len(issues) > 1:
                issue_text += f", and {issues[-1]}"
            else:
                issue_text = issues[0]
                
            return f"Your content would benefit from moderate adaptation for {audience}, particularly improving {issue_text}."
        else:
            return f"Your content would benefit from moderate adaptation to better suit {audience}."
    
    else:
        # High need for adaptation
        key_issues = []
        
        if metrics['flesch_reading_ease'] < thresholds['flesch_reading_ease'] - 15:
            key_issues.append("significantly improve general readability")
            
        if metrics['flesch_kincaid_grade'] > thresholds['flesch_kincaid_grade'] + 3:
            key_issues.append(f"reduce the reading level (currently at grade {metrics['flesch_kincaid_grade']})")
            
        if metrics['sentence_length'] > thresholds['sentence_length'] + 5:
            key_issues.append(f"shorten sentences (currently averaging {metrics['sentence_length']} words)")
            
        if metrics['complex_word_percent'] > thresholds['complex_word_percent'] + 10:
            key_issues.append(f"simplify vocabulary (currently {metrics['complex_word_percent']}% complex words)")
            
        if key_issues:
            issue_text = ", ".join(key_issues[:-1])
            if len(key_issues) > 1:
                issue_text += f", and {key_issues[-1]}"
            else:
                issue_text = key_issues[0]
                
            return f"Your content requires significant adaptation for {audience}. We recommend you: {issue_text}."
        else:
            return f"Your content requires significant adaptation to better suit the needs of {audience}."

def generate_complexity_chart(slide_texts, profile):
    """Generate a chart showing content complexity across slides"""
    # Create data for visualization
    slide_numbers = []
    complexity_scores = []
    threshold_values = []
    
    # Get threshold for the selected profile
    threshold = READABILITY_THRESHOLDS.get(profile, {}).get('flesch_kincaid_grade', 8)
    
    # Calculate complexity for each slide
    for slide in slide_texts:
        # Skip slides with very little text
        if len(slide['text'].strip()) < 20:
            continue
            
        metrics = calculate_simple_readability(slide['text'])
        slide_numbers.append(slide['slide_number'])
        complexity_scores.append(metrics['flesch_kincaid_grade'])
        threshold_values.append(threshold)
    
    # If no valid slides, return placeholder
    if not slide_numbers:
        return generate_placeholder_chart("No substantial text content found in slides")
    
    try:
        # Create the chart
        plt.figure(figsize=(8, 4))
        plt.bar(slide_numbers, complexity_scores, color='#3498db', alpha=0.7)
        plt.axhline(y=threshold, color='#e74c3c', linestyle='-', label=f'Target ({threshold})')
        
        plt.xlabel('Slide Number')
        plt.ylabel('Reading Grade Level')
        plt.title('Content Complexity by Slide')
        plt.legend()
        plt.grid(True, alpha=0.3)
        
        # Add value labels on bars
        for i, v in enumerate(complexity_scores):
            plt.text(slide_numbers[i], v + 0.3, str(round(v, 1)), ha='center')
        
        # Ensure y-axis starts at 0
        plt.ylim(bottom=0)
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
    
    except Exception as e:
        print(f"Error generating complexity chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_complex_words_chart(complex_words):
    """Generate a bar chart showing most frequent complex words"""
    try:
        # If no complex words found
        if not complex_words:
            return generate_placeholder_chart("No complex words found")
            
        # Sort and limit to top 10 words
        words = list(complex_words.keys())[:10]
        frequencies = list(complex_words.values())[:10]
        
        # Create the chart
        plt.figure(figsize=(8, 4))
        bars = plt.barh(words, frequencies, color='#2ecc71', alpha=0.7)
        
        plt.xlabel('Frequency')
        plt.ylabel('Words')
        plt.title('Most Frequent Complex Words')
        plt.grid(True, axis='x', alpha=0.3)
        
        # Add value labels
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 0.3, bar.get_y() + bar.get_height()/2, str(int(width)), va='center')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating complex words chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_placeholder_chart(message="No data available"):
    """Generate a placeholder chart with a message"""
    try:
        plt.figure(figsize=(8, 4))
        
        # Create an empty chart with a message
        plt.text(0.5, 0.5, message, ha='center', va='center', fontsize=14)
        plt.axis('off')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating placeholder chart: {e}")
        # Last resort - return empty image data
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="

# Update the generate_presentation function:
@app.route('/generate', methods=['POST'])
def generate_presentation():
    """Handle generation of a new presentation"""
    try:
        # Get parameters from form
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')  # Add this line
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Generate a unique ID for this presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        filename = f"{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_new_presentation, 
            args=(file_id, filename, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)  # Add subject_area
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE, 
                                   file_id=file_id, 
                                   filename=filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

# Update the enrich_and_generate function:
@app.route('/enrich_and_generate', methods=['POST'])
def enrich_and_generate():
    """Handle generating an enhanced presentation based on an existing adapted file"""
    try:
        # Get parameters from form
        original_file_id = request.form.get('original_file_id')
        profile = request.form.get('profile')
        lesson_topic = request.form.get('lesson_topic')
        grade_level = request.form.get('grade_level')
        subject_area = request.form.get('subject_area')  # Add this line
        slide_count = int(request.form.get('slide_count', 10))
        include_images = 'include_images' in request.form
        extra_notes = request.form.get('extra_notes', '')
        
        # Find the original adapted file
        adapted_filename = None
        adapted_file_path = None
        
        for filename in os.listdir(app.config['OUTPUT_FOLDER']):
            if filename.startswith(f"{original_file_id}_"):
                adapted_filename = filename.replace(f"{original_file_id}_", "")
                adapted_file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
                break
        
        if not adapted_file_path:
            return render_template_string(ERROR_TEMPLATE, message="Original adapted file not found"), 404
        
        # Generate a unique ID for this new presentation
        file_id = str(uuid.uuid4())
        
        # Set initial status
        processing_tasks[file_id] = {'status': 'processing', 'message': ''}
        
        # Create a sanitized filename for the new presentation
        sanitized_topic = re.sub(r'[^\w\s-]', '', lesson_topic).strip().lower()
        sanitized_topic = re.sub(r'[-\s]+', '-', sanitized_topic)
        new_filename = f"enriched-{sanitized_topic}-presentation.pptx"
        
        # Start processing in background thread
        thread = threading.Thread(
            target=generate_enriched_presentation, 
            args=(file_id, new_filename, adapted_file_path, profile, lesson_topic, grade_level, slide_count, include_images, extra_notes, subject_area)  # Add subject_area
        )
        thread.daemon = True
        thread.start()
        
        # Map profile to display name
        profile_names = {
            "dyslexia": "Dyslexia Support",
            "adhd": "ADHD Support",
            "esl": "English Language Learners"
        }
        
        # Return processing page
        return render_template_string(PROCESSING_TEMPLATE, 
                                   file_id=file_id, 
                                   filename=new_filename,
                                   profile=profile,
                                   profile_name=profile_names.get(profile, profile))
            
    except Exception as e:
        import traceback
        traceback.print_exc()
        return render_template_string(ERROR_TEMPLATE, message=f"Error: {str(e)}"), 500

def generate_new_presentation(file_id, filename, profile, topic, grade_level, slide_count, include_images, extra_notes, subject_area):
    """Generate a new presentation based on the specified parameters"""
    try:
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Step 1: Plan the content structure using Claude
        lesson_plan = generate_lesson_plan(topic, grade_level, slide_count, profile, extra_notes, subject_area)
        
        # Step 2: Create slides based on the plan
        for slide_data in lesson_plan:
            # Determine slide layout
            layout_index = slide_data.get("layout_index", 1)
            
            # Select slide layout
            if layout_index == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
            else:
                # Content slide with title
                slide_layout = prs.slide_layouts[1]
            
            # Add a slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                title = slide_data.get("title", "")
                slide.shapes.title.text = title
                
                # Apply profile-specific formatting to title
                if hasattr(slide.shapes.title, "text_frame"):
                    apply_text_to_text_frame(slide.shapes.title.text_frame, title, profile)
            
            # Set slide content
            content = slide_data.get("content", "")
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            # If content placeholder found, add text
            if content_placeholder and content:
                if isinstance(content, list):
                    # Handle list content
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    
                    for item in content:
                        p = text_frame.add_paragraph()
                        p.text = "• " + item
                        
                        # Apply profile-specific formatting
                        if profile == "dyslexia":
                            apply_dyslexia_formatting(text_frame)
                else:
                    # Handle string content
                    adapted_content = adapt_text_with_matcha(content, profile)
                    text_frame = content_placeholder.text_frame
                    apply_text_to_text_frame(text_frame, adapted_content, profile)
        
        # Save the presentation
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{file_id}_{filename}")
        prs.save(output_path)
        
        # Update status to complete
        processing_tasks[file_id] = {'status': 'complete', 'message': ''}
        
        return output_path
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        
        # Update status to error
        processing_tasks[file_id] = {'status': 'error', 'message': str(e)}

def generate_enriched_lesson_plan(adapted_content, topic, grade_level, profile, extra_notes, subject_area):
    """Generate an enriched lesson plan based on the adapted content for UK curriculum"""
    # Updated UK education level descriptions
    level_descriptions = {
        "ks3-year7-8": "Key Stage 3 (Year 7-8, ages 11-13)",
        "ks3-year9": "Key Stage 3 (Year 9, age 14)",
        "ks4-gcse": "Key Stage 4 (GCSE, Year 10-11, ages 15-16)",
        "ks5-alevel": "Key Stage 5 (A-Level, Year 12-13, ages 17-18)",
        "btec": "BTEC qualification level"
    }
    
    level_desc = level_descriptions.get(grade_level, "secondary school")
    
    adaptation_needs = {
        "dyslexia": "simplified vocabulary, short sentences, clear structure, minimal text per slide",
        "adhd": "visually engaging, chunked information, clear headings, bulleted lists, minimal distractions",
        "esl": "basic vocabulary with more complex terms in parentheses, visual support, simple sentence structure"
    }
    
    adaptation_desc = adaptation_needs.get(profile, "")
    
    # Subject-specific guidance
    subject_guidance = {
        "english": "Include key literary terms, text analysis techniques, and writing structures appropriate for this level.",
        "maths": "Include clear mathematical notations, step-by-step explanations, and visual representations of concepts.",
        "science": "Include scientific diagrams, experimental procedures, and key terminology with explanations.",
        "biology": "Include biological diagrams, processes, and scientific terminology with explanations.",
        "chemistry": "Include chemical formulas, reaction mechanisms, and key principles with clear explanations.",
        "physics": "Include physics equations, diagrams, and real-world applications of concepts.",
        "history": "Include timelines, primary sources, and analysis of historical events.",
        "geography": "Include maps, diagrams, and case studies relevant to the topic.",
        "religious_studies": "Include key beliefs, practices, and comparative analysis where appropriate.",
        "modern_languages": "Include vocabulary lists, language structures, and cultural context.",
        "computing": "Include code examples, system diagrams, and step-by-step procedures.",
        "design_technology": "Include design specifications, technical drawings, and manufacturing processes.",
        "art_design": "Include visual examples, techniques, and artistic context.",
        "music": "Include notation examples, listening activities, and music theory concepts.",
        "physical_education": "Include diagrams of techniques, rules of activities, and fitness concepts.",
        "citizenship": "Include case studies, legal frameworks, and civic participation examples.",
        "pshe": "Include scenarios, discussion points, and practical guidance.",
        "business": "Include business models, case studies, and key terminology.",
        "sociology": "Include sociological theories, research methods, and social examples.",
        "psychology": "Include psychological theories, studies, and applications.",
        "economics": "Include economic models, graphs, and real-world examples.",
        "media_studies": "Include media analysis frameworks, case studies, and production techniques."
    }
    
    subject_specific_guidance = subject_guidance.get(subject_area, "")
    
    prompt = f"""
    Create an enriched, enhanced lesson plan for a PowerPoint presentation about "{topic}" for {level_desc} students following the UK National Curriculum for {subject_area}.
    
    This presentation should build upon the existing content below, but enhance and expand it to create a more comprehensive, engaging presentation.
    
    EXISTING CONTENT:
    {adapted_content}
    
    The enriched presentation should:
    1. Maintain the core content from the original
    2. Add engaging visuals and interactive elements
    3. Include additional examples and learning activities
    4. Be adapted for students with {profile} needs: {adaptation_desc}
    5. Align with UK {subject_area} curriculum standards for {level_desc}
    6. {subject_specific_guidance}
    7. {extra_notes}
    
    For each slide, provide:
    1. Slide type (title slide, content slide, image slide, etc.)
    2. Title of the slide
    3. Key content points (in bullet form)
    4. A brief image description that would support the content
    
    Format your response as a structured JSON array where each element is a slide with the following properties:
    - layout_index: (0 for title slide, 1 for content slide with title)
    - title: The slide title
    - content: The main content (bulleted list where appropriate)
    - image_prompt: A description for image generation (if applicable)
    
    Keep the language appropriate for {profile} needs and UK curriculum standards.
    """
    
    try:
        response = call_claude_api_with_retry(prompt, max_tokens=4000)
        content = response.content[0].text.strip()
        
        # Extract JSON from the response
        json_match = re.search(r'\[\s*{.*}\s*\]', content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            lesson_plan = json.loads(json_str)
            return lesson_plan
        else:
            # If no JSON found, try to parse the content in a more forgiving way
            try:
                # Try to find anything that looks like JSON
                json_start = content.find('[')
                json_end = content.rfind(']') + 1
                if json_start >= 0 and json_end > json_start:
                    json_str = content[json_start:json_end]
                    lesson_plan = json.loads(json_str)
                    return lesson_plan
            except json.JSONDecodeError:
                pass
            
            # If all JSON parsing fails, create a simple structure from the text
            slides = []
            sections = content.split('\n\n')
            for i, section in enumerate(sections):
                if i == 0:
                    # First section is the title slide
                    slides.append({
                        "layout_index": 0,
                        "title": topic,
                        "content": section
                    })
                else:
                    # Extract title and content
                    lines = section.split('\n')
                    if lines:
                        title = lines[0].strip('#- ').strip()
                        content = '\n'.join(lines[1:]).strip()
                        slides.append({
                            "layout_index": 1,
                            "title": title,
                            "content": content,
                            "image_prompt": f"Educational illustration about {title} for {level_desc} students studying {subject_area}"
                        })
            
            return slides
    
    except Exception as e:
        print(f"Error generating enriched lesson plan: {str(e)}")
        # Create a basic fallback plan
        return [
            {"layout_index": 0, "title": topic, "content": f"An enriched presentation about {topic} for {level_desc} students studying {subject_area}"},
            {"layout_index": 1, "title": "Introduction", "content": f"• Introduction to {topic}\n• Key concepts\n• Learning objectives", "image_prompt": f"Introduction to {topic} for UK {subject_area} curriculum"},
            {"layout_index": 1, "title": "Main Concepts", "content": f"• Important aspects of {topic}\n• Key information", "image_prompt": f"Main concepts of {topic} for UK {subject_area} curriculum"}
        ]
    

def generate_complex_words_chart(complex_words):
    """Generate a bar chart showing most frequent complex words"""
    try:
        # If no complex words found
        if not complex_words:
            return generate_placeholder_chart("No complex words found")
            
        # Sort and limit to top 10 words
        words = list(complex_words.keys())[:10]
        frequencies = list(complex_words.values())[:10]
        
        # Create the chart
        plt.figure(figsize=(8, 4))
        bars = plt.barh(words, frequencies, color='#2ecc71', alpha=0.7)
        
        plt.xlabel('Frequency')
        plt.ylabel('Words')
        plt.title('Most Frequent Complex Words')
        plt.grid(True, axis='x', alpha=0.3)
        
        # Add value labels
        for bar in bars:
            width = bar.get_width()
            plt.text(width + 0.3, bar.get_y() + bar.get_height()/2, str(int(width)), va='center')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating complex words chart: {e}")
        return generate_placeholder_chart("Error generating chart")

def generate_placeholder_chart(message="No data available"):
    """Generate a placeholder chart with a message"""
    try:
        plt.figure(figsize=(8, 4))
        
        # Create an empty chart with a message
        plt.text(0.5, 0.5, message, ha='center', va='center', fontsize=14)
        plt.axis('off')
        
        # Convert to base64 for embedding in HTML
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', bbox_inches='tight')
        buffer.seek(0)
        image_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
        plt.close()
        
        return f"data:image/png;base64,{image_base64}"
        
    except Exception as e:
        print(f"Error generating placeholder chart: {e}")
        # Last resort - return empty image data
        return "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="

@app.route('/adapt/<file_id>/<profile>')
def adapt(file_id, profile):
    """Adapt the PowerPoint based on assessment"""
    try:
        # Validate if file_id exists in our system
        if file_id not in processing_tasks:
            return render_template_string(ERROR_TEMPLATE, 
                                         message="Presentation not found. Please upload again."), 404
        
        # Set initial processing status
        processing_tasks[file_id] = {'status': 'processing', 'message': 'Starting adaptation...'}
        
        # Get the original filename from the processing tasks if available
        filename = processing_tasks[file_id].get('filename', 'presentation.pptx')
        
        # Determine the file path based on your storage structure
        file_path = os.path.join(UPLOAD_FOLDER, f"{file_id}.pptx")
        
        if not os.path.exists(file_path):
            processing_tasks[file_id] = {'status': 'error', 'message': 'Presentation file not found on server'}
            return redirect(url_for('error', message="File not found on server"))
        
        # Validate profile type
        valid_profiles = ['dyslexia', 'esl', 'adhd', 'visual', 'cognitive']
        if profile not in valid_profiles:
            processing_tasks[file_id] = {'status': 'error', 'message': f'Invalid profile: {profile}'}
            return redirect(url_for('error', message=f"Invalid profile: {profile}"))
        
        # Start processing in a background thread to avoid blocking
        def process_task():
            try:
                process_presentation_efficiently(file_path, file_id, filename, profile)
            except Exception as e:
                import traceback
                traceback.print_exc()
                processing_tasks[file_id] = {'status': 'error', 'message': str(e)}
        
        # Start the processing thread
        processing_thread = threading.Thread(target=process_task)
        processing_thread.daemon = True
        processing_thread.start()
        
        # Return the processing template
        return render_template_string(
            PROCESSING_TEMPLATE,
            file_id=file_id,
            filename=filename,
            profile=profile
        )
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        error_message = f"Error starting adaptation: {str(e)}"
        return redirect(url_for('error', message=error_message))

def adapt_text_with_matcha(text, profile):
    """Adapt text using Claude for the specified learning profile"""
    # Skip processing very short text
    if len(text.strip()) < 15:
        return text
        
    # Check cache first
    cached_result = adaptation_cache.get(text, profile)
    if cached_result:
        print(f"Cache hit for {profile} adaptation")
        return cached_result
    
    try:
        # Create efficient prompt
        prompt = create_efficient_prompt(text, profile)
        
        # Call Claude API
        response = call_claude_api_with_retry(prompt)
        
        # Extract adapted text from response
        adapted_text = response.content[0].text.strip()
        
        # Store in cache for future use
        adaptation_cache.set(text, profile, adapted_text)
        
        return adapted_text
        
    except Exception as e:
        print(f"Error in text adaptation: {e}")
        # Return original text if adaptation fails
        return text

def translate_text(text, target_language):
    """
    Translate text to the specified target language using an AI API
    
    Args:
        text (str): The original text to translate
        target_language (str): The target language code (e.g., 'fr' for French)
        
    Returns:
        str: The translated text
    """
    if not text or not text.strip():
        return text
        
    try:
        # Create a prompt for the translation
        prompt = f"""Translate the following text to {target_language}. 
Maintain the original formatting, including line breaks and bullet points.
Ensure the translation sounds natural and conversational.

Text to translate:
"{text}"

Translation:"""

        # Call Claude API with the translation prompt
        response = call_claude_api_with_retry(prompt)
        
        # Extract the translated content from the response
        translated_text = response.content[0].text.strip()
        
        # Log success for monitoring
        print(f"Successfully translated text of length {len(text)} to {target_language}")
        
        return translated_text
        
    except Exception as e:
        print(f"Error translating text: {str(e)}")
        # Return original text if translation fails
        return text

def apply_dyslexia_formatting(text_frame):
    """Apply special formatting for dyslexia support"""
    try:
        # Adjust text alignment to left
        for paragraph in text_frame.paragraphs:
            if hasattr(paragraph, "alignment"):
                paragraph.alignment = PP_ALIGN.LEFT
                
            # Increase line spacing if not already set
            # Check if line_spacing is None before comparing
            if hasattr(paragraph, "line_spacing") and paragraph.line_spacing is not None and paragraph.line_spacing < 1.2:
                paragraph.line_spacing = 1.5
            elif hasattr(paragraph, "line_spacing") and paragraph.line_spacing is None:
                # If it's None, just set it without comparison
                paragraph.line_spacing = 1.5
                
            # Apply sans-serif font to all runs if not already a sans-serif font
            for run in paragraph.runs:
                if hasattr(run, "font") and hasattr(run.font, "name"):
                    current_font = run.font.name.lower() if hasattr(run.font.name, "lower") else ""
                    # Only change if not already using a dyslexia-friendly font
                    if not any(font in current_font for font in ["arial", "verdana", "tahoma", "calibri", "helvetica", "open sans"]):
                        run.font.name = "Arial"
                    
                    # Increase font size slightly for better readability if it's too small
                    if hasattr(run.font, "size") and hasattr(run.font.size, "pt"):
                        current_size = run.font.size.pt
                        # Only increase if not already large
                        if current_size < 12:
                            run.font.size = Pt(max(12, current_size * 1.2))
    except Exception as e:
        # Log error but continue processing
        print(f"Error applying dyslexia formatting: {e}")
        pass
    
if __name__ == "__main__":
    # Run the Flask app
    app.run(debug=True, host='0.0.0.0', port=5000)
