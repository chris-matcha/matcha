#!/usr/bin/env python3
"""
Test script for the enhanced text overflow handling in PowerPointService
"""
import sys
import os
sys.path.append('.')

from services.pptx_service import PowerPointService
from pptx import Presentation
from pptx.util import Inches, Pt

def test_text_overflow_handling():
    """Test the complete text overflow handling system"""
    print("🔧 Testing Enhanced Text Overflow Handling")
    print("=" * 55)
    
    # Initialize service
    config = {
        'output_folder': './test_outputs',
        'upload_folder': './uploads'
    }
    pptx_service = PowerPointService(config)
    
    # Test scenarios
    test_scenarios = [
        {
            "name": "Very Long Text",
            "text": """This is an extremely long piece of text that contains multiple sentences and would definitely cause overflow issues in a normal PowerPoint slide. The text is designed to test the smart truncation functionality which should intelligently break the content at sentence boundaries. This sentence is even longer than the previous ones to really push the limits of what can fit in a standard text box. We want to ensure that the overflow handling maintains readability while preserving as much content as possible. The system should also provide useful feedback about what content was moved to overflow notes.""",
            "profile": "dyslexia",
            "max_width": 300,
            "max_height": 150
        },
        {
            "name": "Medium Text for Bullets",
            "text": "This is a moderately long text that should be converted to bullet points for better readability. It contains enough content to benefit from bullet formatting, but not so much that it requires extreme truncation. The system should break this down into digestible bullet points.",
            "profile": "adhd", 
            "max_width": 400,
            "max_height": 200
        },
        {
            "name": "Short Text (No Overflow)",
            "text": "This is a short text that should fit perfectly without any overflow issues.",
            "profile": "esl",
            "max_width": 500,
            "max_height": 100
        },
        {
            "name": "Extreme Overflow",
            "text": "This text is so extremely long that even with all the truncation strategies, it will still need to be dramatically shortened using the extreme truncation method. " * 10,
            "profile": "visual",
            "max_width": 200,
            "max_height": 80
        }
    ]
    
    # Test each scenario
    for i, scenario in enumerate(test_scenarios, 1):
        print(f"\n📝 Test {i}: {scenario['name']}")
        print("-" * 40)
        
        # Test smart truncation
        try:
            truncated, overflow = pptx_service._smart_truncate_text(
                scenario['text'],
                scenario['max_width'],
                scenario['max_height'],
                12,  # font size
                "Arial"
            )
            
            print(f"✓ Original length: {len(scenario['text'])} chars")
            print(f"✓ Truncated length: {len(truncated)} chars")
            print(f"✓ Overflow length: {len(overflow)} chars")
            print(f"✓ Truncated text: {truncated[:100]}{'...' if len(truncated) > 100 else ''}")
            
            if overflow:
                print(f"✓ Overflow text: {overflow[:100]}{'...' if len(overflow) > 100 else ''}")
            else:
                print("✓ No overflow - text fits perfectly")
                
        except Exception as e:
            print(f"❌ Error in smart truncation: {e}")
        
        # Test bullet conversion
        try:
            bullets = pptx_service._convert_to_bullets(scenario['text'])
            print(f"✓ Bullet conversion: {bullets[:150]}{'...' if len(bullets) > 150 else ''}")
        except Exception as e:
            print(f"❌ Error in bullet conversion: {e}")
        
        # Test extreme truncation
        try:
            extreme = pptx_service._extreme_truncate(scenario['text'], 100)
            print(f"✓ Extreme truncation (100 chars): {extreme}")
        except Exception as e:
            print(f"❌ Error in extreme truncation: {e}")
    
    # Integration test with actual PowerPoint
    print(f"\n🎯 Integration Test: Creating PowerPoint with Overflow Handling")
    print("-" * 55)
    
    try:
        # Create test presentation
        prs = Presentation()
        
        for i, scenario in enumerate(test_scenarios, 1):
            # Add a new slide for each scenario
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = f"Test {i}: {scenario['name']}"
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                # Apply text with overflow handling
                success, notes = pptx_service.apply_text_with_optimal_sizing(
                    content_placeholder.text_frame,
                    scenario['text'],
                    scenario['max_width'],
                    scenario['max_height'],
                    scenario['profile']
                )
                
                print(f"✓ Slide {i} processed successfully: {success}")
                if notes:
                    print(f"  Notes: {notes[:100]}{'...' if len(notes) > 100 else ''}")
        
        # Save test presentation
        os.makedirs('test_outputs', exist_ok=True)
        output_file = 'test_outputs/text_overflow_test.pptx'
        prs.save(output_file)
        print(f"✓ Test presentation saved: {output_file}")
        
    except Exception as e:
        print(f"❌ Error in integration test: {e}")
        import traceback
        traceback.print_exc()
    
    # Summary
    print(f"\n🎉 Text Overflow Handling Test Complete!")
    print("\n📋 Features Tested:")
    print("   • Smart truncation with sentence boundaries")
    print("   • Intelligent bullet point conversion")  
    print("   • Extreme truncation for last resort scenarios")
    print("   • Overflow detection and measurement")
    print("   • Integration with PowerPoint text frames")
    print("   • Profile-specific font handling")
    print("   • Multi-strategy overflow resolution")

def test_edge_cases():
    """Test edge cases and error conditions"""
    print(f"\n🔍 Testing Edge Cases")
    print("-" * 30)
    
    config = {
        'output_folder': './test_outputs',
        'upload_folder': './uploads'
    }
    pptx_service = PowerPointService(config)
    
    edge_cases = [
        {"name": "Empty Text", "text": ""},
        {"name": "Single Word", "text": "Word"},
        {"name": "No Periods", "text": "Text without any sentence endings at all"},
        {"name": "Only Periods", "text": "...."},
        {"name": "Very Small Bounds", "text": "Test", "width": 10, "height": 10},
        {"name": "Unicode Text", "text": "Текст на русском языке с ümlauts éñ français 中文"},
    ]
    
    for case in edge_cases:
        print(f"\n🧪 Edge Case: {case['name']}")
        try:
            width = case.get('width', 200)
            height = case.get('height', 100)
            
            truncated, overflow = pptx_service._smart_truncate_text(
                case['text'], width, height, 12, "Arial"
            )
            print(f"   ✓ Handled successfully: '{truncated}' (overflow: {len(overflow)} chars)")
        except Exception as e:
            print(f"   ❌ Error: {e}")

if __name__ == "__main__":
    test_text_overflow_handling()
    test_edge_cases()